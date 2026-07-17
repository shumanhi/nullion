"""Typed tool registry for UI-neutral Nullion capabilities."""

from __future__ import annotations

from collections import Counter
from concurrent.futures import ThreadPoolExecutor
from contextlib import contextmanager
import bz2
import csv
from dataclasses import dataclass, field, replace
from datetime import UTC, datetime
from email.message import EmailMessage
from enum import Enum
from fnmatch import fnmatch
import hashlib
import html
from html import unescape
from html.parser import HTMLParser
from ipaddress import ip_address
from io import BytesIO
import inspect
import json
import logging
import gzip
import lzma
import mimetypes
import os
from pathlib import Path
import re
import shlex
import shutil
import socket
import ssl
import subprocess
import tempfile
import tarfile
import textwrap
import threading
from time import perf_counter
from typing import Callable, Iterable, Mapping, Protocol
import base64
import unicodedata
import urllib.error
import urllib.request
from urllib.parse import parse_qsl, quote, unquote, unquote_to_bytes, urlencode, urljoin, urlparse
from uuid import uuid4
from xml.etree import ElementTree
import zipfile

from nullion.attachment_format_graph import VALID_ATTACHMENT_EXTENSIONS
from nullion.artifacts import (
    ARTIFACT_ROLE_DELIVERABLE,
    ARTIFACT_ROLE_INTERMEDIATE,
    ARTIFACT_ROLE_SOURCE,
    artifact_output_descriptor,
    promote_supporting_asset_artifact_paths,
)
from nullion.approval_context import FLOW_TRIGGER_CONTEXT_KEY, build_trigger_flow_context
from nullion.approvals import (
    ApprovalRequest,
    TERMINAL_DESTRUCTIVE_ACTION_REQUEST_KIND,
    TERMINAL_DESTRUCTIVE_PERMISSION_PREFIX,
    consume_boundary_permit as consume_boundary_permit_record,
    create_approval_request,
    is_boundary_permit_active,
    is_permission_grant_active,
    revoke_permission_grant as revoke_permission_grant_record,
)
from nullion.audit import make_audit_record
from nullion.events import make_event
from nullion.policy import (
    BoundaryFact,
    BoundaryKind,
    BoundaryPolicyRequest,
    BoundaryPolicyRule,
    PolicyDecision,
    PolicyMode,
    SentinelPolicy,
    GLOBAL_PERMISSION_PRINCIPAL,
    OPERATOR_PERMISSION_PRINCIPAL,
    evaluate_boundary_request,
    normalize_outbound_network_selector,
    permission_scope_principal,
)
from nullion.prompt_injection import scan_tool_output
from nullion.redaction import redact_value
from nullion.runtime_store import RuntimeStore
from nullion.tips import MEDIA_PROVIDER_SETUP_TIP, format_setup_tip
from nullion.tool_boundaries import extract_boundary_facts


_WEB_FETCH_MAX_REDIRECTS = 5
_WEB_FETCH_MAX_BODY_BYTES = 2_000_000
_CONNECTOR_JSON_MAX_BODY_BYTES = 1_000_000
_EMAIL_ATTACHMENT_CONNECTOR_MAX_BODY_BYTES = 36 * 1024 * 1024
_FILE_DOWNLOAD_MAX_BYTES = 512 * 1024 * 1024
_BROWSER_IMAGE_COLLECT_MAX_IMAGES = 20
_BROWSER_IMAGE_MAX_BYTES = 6_000_000
_TERMINAL_OUTPUT_ATTACHMENT_THRESHOLD_CHARS = 3200
_TERMINAL_DESTRUCTIVE_PREVIEW_LIMIT = 40
_TERMINAL_DESTRUCTIVE_SCRIPT_COMMANDS = frozenset({
    "node",
    "node.exe",
    "perl",
    "perl.exe",
    "php",
    "php.exe",
    "python",
    "python.exe",
    "python3",
    "python3.exe",
    "ruby",
    "ruby.exe",
})
_TERMINAL_DESTRUCTIVE_SCRIPT_MARKERS = frozenset({
    ".rmtree(",
    ".unlink(",
    "deletefile(",
    "fs.rm(",
    "fs.rmdir(",
    "fs.unlink(",
    "os.remove(",
    "os.rmdir(",
    "os.unlink(",
    "path.unlink(",
    "rmdirsync(",
    "rmsync(",
    "shutil.rmtree(",
    "unlink(",
    "unlinksync(",
})
_TERMINAL_PATH_LITERAL_RE = re.compile(
    r"(?P<quote>['\"])(?P<path>(?:~|/|\.{1,2}/)[^'\"]+)(?P=quote)"
)
_ARCHIVE_READ_ENTRY_LIMIT = 500
_ARCHIVE_EXTRACT_ENTRY_LIMIT = 2_000
_ARCHIVE_MANIFEST_IMAGE_SUFFIXES = frozenset({".bmp", ".gif", ".jpeg", ".jpg", ".png", ".webp"})
_TAR_ARCHIVE_FORMATS: dict[str, tuple[str, str]] = {
    ".tar": ("tar", "w"),
    ".tar.gz": ("tar.gz", "w:gz"),
    ".tgz": ("tar.gz", "w:gz"),
    ".tar.bz2": ("tar.bz2", "w:bz2"),
    ".tbz2": ("tar.bz2", "w:bz2"),
    ".tar.xz": ("tar.xz", "w:xz"),
    ".txz": ("tar.xz", "w:xz"),
}
_SINGLE_COMPRESSED_FORMATS: dict[str, tuple[str, str]] = {
    ".gz": ("gzip", "gzip"),
    ".bz2": ("bzip2", "bzip2"),
    ".xz": ("xz", "xz"),
}


_CONCRETE_RESOURCE_ID_FORBIDDEN_CHARS = frozenset(' \t\r\n<>{}[]()"\'`?*,;')
_CONCRETE_RESOURCE_ID_SENTINELS = frozenset(
    {
        "-",
        "--",
        "id",
        "message_id",
        "placeholder",
        "example",
        "sample",
        "returned_id",
        "result_id",
        "none",
        "null",
    }
)


def _invalid_concrete_resource_id_reason(value: object) -> str:
    if not isinstance(value, str):
        return "missing"
    resource_id = value.strip()
    if not resource_id:
        return "missing"
    if resource_id.casefold() in _CONCRETE_RESOURCE_ID_SENTINELS:
        return "placeholder"
    if any(char in _CONCRETE_RESOURCE_ID_FORBIDDEN_CHARS for char in resource_id):
        return "placeholder"
    if len(resource_id) > 512:
        return "too_long"
    return ""


def _invalid_concrete_resource_name_reason(value: object) -> str:
    if not isinstance(value, str):
        return "missing"
    resource_name = value.strip()
    if not resource_name:
        return "missing"
    if resource_name.startswith("__") and resource_name.endswith("__"):
        return "placeholder"
    if len(resource_name) > 512:
        return "too_long"
    return ""


def _invalid_email_message_id_result(invocation: "ToolInvocation", raw_id: object) -> "ToolResult":
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={
            "reason": "invalid_message_id",
            "id": raw_id if isinstance(raw_id, str) else "",
            "required_source_tool": "email_search",
            "next_step": "Search email and pass one returned result id to email_read.",
        },
        error="Invalid email message id: use an id returned by email_search.",
    )


def _email_read_source_required_failure_result(
    invocation: "ToolInvocation",
    raw_id: object,
    exc: BaseException | None = None,
) -> "ToolResult":
    output = _account_tool_failure_output(
        invocation.principal_id,
        message_id=raw_id.strip() if isinstance(raw_id, str) else "",
    )
    output.update(
        {
            "reason": "invalid_message_id",
            "required_source_tool": "email_search",
            "next_step": "Search email and pass one returned result id to email_read.",
        }
    )
    if isinstance(exc, urllib.error.HTTPError):
        output["http_status"] = exc.code
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output=output,
        error=str(exc) if exc is not None else "Invalid email message id: use an id returned by email_search.",
    )


def _invalid_structured_identifier_result(
    invocation: "ToolInvocation",
    *,
    field_name: str,
    raw_value: object,
    reason: str,
    required_source_tool: str,
) -> "ToolResult":
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={
            "reason": reason,
            field_name: raw_value if isinstance(raw_value, str) else "",
            "required_source_tool": required_source_tool,
        },
        error=f"Invalid {field_name}: use an exact id returned by {required_source_tool}.",
    )
def _email_search_input_schema() -> dict[str, object]:
    return {
        "type": "object",
        "properties": {
            "query": {
                "type": "string",
                "minLength": 1,
                "description": "Provider-supported email search query.",
            },
            "limit": {
                "type": "integer",
                "minimum": 1,
                "maximum": 10,
                "description": "Maximum number of matching messages to return.",
            },
            "provider_id": {
                "type": "string",
                "description": "Optional configured connector/provider id when multiple email sources are active.",
            },
        },
        "required": ["query"],
        "additionalProperties": False,
    }


def _email_read_input_schema() -> dict[str, object]:
    return {
        "type": "object",
        "properties": {
            "id": {
                "type": "string",
                "minLength": 1,
                "description": "Exact id copied from a prior email_search results item.",
            },
            "provider_id": {
                "type": "string",
                "description": "Optional configured connector/provider id matching the source search result.",
            },
        },
        "required": ["id"],
        "additionalProperties": False,
    }


def _email_attachment_read_input_schema() -> dict[str, object]:
    return {
        "type": "object",
        "properties": {
            "message_id": {
                "type": "string",
                "minLength": 1,
                "description": "Exact email_read.message.id value for the message that contains the attachment.",
            },
            "attachment_id": {
                "type": "string",
                "minLength": 1,
                "description": "Exact email_read.message.attachments[].attachmentId value.",
            },
            "filename": {
                "type": "string",
                "description": "Optional filename copied from email_read.message.attachments[].filename.",
            },
            "mime_type": {
                "type": "string",
                "description": "Optional MIME type copied from email_read.message.attachments[].mimeType.",
            },
            "provider_id": {
                "type": "string",
                "description": "Optional configured connector/provider id matching the source message.",
            },
        },
        "required": ["message_id", "attachment_id"],
        "additionalProperties": False,
    }


_BROWSER_IMAGE_CONTENT_MIN_BYTES = 4_000
_BROWSER_IMAGE_CONTENT_MIN_WIDTH = 180
_BROWSER_IMAGE_CONTENT_MIN_HEIGHT = 120
_BROWSER_IMAGE_CONTENT_MIN_PIXELS = 40_000
_BROWSER_IMAGE_CONTENT_MAX_ASPECT_RATIO = 4.5
_BROWSER_IMAGE_CONTENT_MIN_LUMA_STDDEV = 3.0
_BROWSER_IMAGE_DIRECTIVE_MEDIA_TYPES = frozenset(
    {
        "image/png",
        "image/jpeg",
        "image/gif",
        "image/webp",
        "image/avif",
        "image/bmp",
    }
)
_BROWSER_IMAGE_FORMAT_SUFFIXES = {
    "PNG": ".png",
    "JPEG": ".jpg",
    "GIF": ".gif",
}
_LEGACY_GLOBAL_PERMISSION_PRINCIPALS = ("operator", "workspace:workspace_admin")


def _calendar_list_input_schema(*, include_provider_id: bool = True) -> dict[str, object]:
    properties: dict[str, object] = {
        "start": {"type": "string", "description": "Inclusive ISO-8601 start datetime."},
        "end": {"type": "string", "description": "Exclusive ISO-8601 end datetime."},
        "max": {"type": "integer", "minimum": 1, "description": "Maximum number of events to return."},
        "query": {
            "type": "string",
            "description": "Optional exact title, subject, or structured search text to mark matching returned events.",
        },
    }
    if include_provider_id:
        properties["provider_id"] = {
            "type": "string",
            "description": "Optional connector provider id. Defaults to an active calendar-capable connector.",
        }
    return {
        "type": "object",
        "properties": properties,
        "required": ["start", "end"],
        "additionalProperties": False,
    }


def _calendar_write_provider_property() -> dict[str, object]:
    return {
        "type": "string",
        "description": "Optional connector provider id. Defaults to an active write-capable calendar connector.",
    }


def _calendar_id_property() -> dict[str, object]:
    return {
        "type": "string",
        "description": "Optional Google Calendar id. Defaults to primary.",
    }


def _calendar_send_updates_property() -> dict[str, object]:
    return {
        "type": "string",
        "enum": ["all", "externalOnly", "none"],
        "description": "Who should receive Google Calendar update notifications. Defaults to all.",
    }


def _calendar_event_time_properties() -> dict[str, object]:
    return {
        "start": {
            "type": "string",
            "description": "Event start as ISO-8601 dateTime or YYYY-MM-DD all-day date.",
        },
        "end": {
            "type": "string",
            "description": "Event end as ISO-8601 dateTime or YYYY-MM-DD all-day date.",
        },
        "time_zone": {
            "type": "string",
            "description": "Optional IANA timezone for dateTime values, such as America/New_York.",
        },
    }


def _calendar_create_input_schema(*, include_provider_id: bool = True) -> dict[str, object]:
    properties: dict[str, object] = {
        "summary": {"type": "string", "description": "Event title/summary."},
        **_calendar_event_time_properties(),
        "location": {"type": "string", "description": "Optional event location."},
        "description": {"type": "string", "description": "Optional event description or notes."},
        "attendees": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Optional attendee email addresses.",
        },
        "calendar_id": _calendar_id_property(),
        "send_updates": _calendar_send_updates_property(),
    }
    if include_provider_id:
        properties["provider_id"] = _calendar_write_provider_property()
    return {
        "type": "object",
        "properties": properties,
        "required": ["summary", "start", "end"],
        "additionalProperties": False,
    }


def _calendar_update_input_schema(*, include_provider_id: bool = True) -> dict[str, object]:
    properties: dict[str, object] = {
        "event_id": {
            "type": "string",
            "description": "Exact event id returned by calendar_list.results[].id or a prior calendar write result.",
        },
        "summary": {"type": "string", "description": "Optional replacement event title/summary."},
        **_calendar_event_time_properties(),
        "location": {"type": "string", "description": "Optional replacement event location."},
        "description": {"type": "string", "description": "Optional replacement event description or notes."},
        "calendar_id": _calendar_id_property(),
        "send_updates": _calendar_send_updates_property(),
    }
    if include_provider_id:
        properties["provider_id"] = _calendar_write_provider_property()
    return {
        "type": "object",
        "properties": properties,
        "required": ["event_id"],
        "additionalProperties": False,
    }


def _calendar_respond_input_schema(*, include_provider_id: bool = True) -> dict[str, object]:
    properties: dict[str, object] = {
        "event_id": {
            "type": "string",
            "description": "Exact event id returned by calendar_list.results[].id or a prior calendar write result.",
        },
        "response_status": {
            "type": "string",
            "enum": ["accepted", "declined", "tentative", "needsAction"],
            "description": "Calendar RSVP response.",
        },
        "attendee_email": {
            "type": "string",
            "description": "Optional attendee email to update. If omitted, the connector attempts to update the current user's attendee response.",
        },
        "calendar_id": _calendar_id_property(),
        "send_updates": _calendar_send_updates_property(),
    }
    if include_provider_id:
        properties["provider_id"] = _calendar_write_provider_property()
    return {
        "type": "object",
        "properties": properties,
        "required": ["event_id", "response_status"],
        "additionalProperties": False,
    }


def _calendar_delete_input_schema(*, include_provider_id: bool = True) -> dict[str, object]:
    properties: dict[str, object] = {
        "event_id": {
            "type": "string",
            "description": "Exact event id returned by calendar_list.results[].id or a prior calendar write result.",
        },
        "calendar_id": _calendar_id_property(),
        "send_updates": _calendar_send_updates_property(),
    }
    if include_provider_id:
        properties["provider_id"] = _calendar_write_provider_property()
    return {
        "type": "object",
        "properties": properties,
        "required": ["event_id"],
        "additionalProperties": False,
    }


def _tool_grant_principal_candidates(principal_id: str | None) -> set[str]:
    principal_text = str(principal_id or "").strip()
    candidates = {candidate for candidate in {principal_text, permission_scope_principal(principal_text)} if candidate}
    candidates.update({
        GLOBAL_PERMISSION_PRINCIPAL,
        OPERATOR_PERMISSION_PRINCIPAL,
        *_LEGACY_GLOBAL_PERMISSION_PRINCIPALS,
    })
    return candidates


def _tool_blocked_for_principal(principal_id: str | None, tool_name: str | None) -> bool:
    try:
        from nullion.users import tool_blocked_for_principal

        return tool_blocked_for_principal(principal_id, tool_name)
    except Exception:
        logger.debug("Could not load user tool block policy", exc_info=True)
        return False


logger = logging.getLogger(__name__)

_SCHEDULER_CREATION_CAPSULE_LOCK = threading.Lock()
_SCHEDULER_CREATION_CAPSULE_TTL_SECONDS = 6 * 60 * 60
_SCHEDULER_CREATION_CAPSULE_LIMIT = 2048
_SCHEDULER_CREATION_CAPSULES: dict[str, float] = {}


class _SafeWebFetchRedirectHandler(urllib.request.HTTPRedirectHandler):
    """Follow redirects, but only to HTTP/HTTPS URLs (blocks data:, file:, etc).

    Also rejects redirects that resolve to private/loopback addresses to
    prevent SSRF via redirect chains (e.g. http://legit.com → http://169.254.x.x).
    """

    max_redirections = _WEB_FETCH_MAX_REDIRECTS

    def redirect_request(self, req, fp, code, msg, headers, newurl):
        parsed = urlparse(newurl)
        if parsed.scheme not in {"http", "https"}:
            # Refuse non-HTTP redirect targets (data:, file:, ftp:, etc.)
            return None
        # Refuse redirects to private/loopback hosts
        host = parsed.hostname or ""
        try:
            from ipaddress import ip_address as _ip
            addr = _ip(host)
            if not addr.is_global:
                return None
        except ValueError:
            # Hostname — allow it; DNS resolution happens later
            pass
        return super().redirect_request(req, fp, code, msg, headers, newurl)


_WEB_FETCH_RESOLUTION_LOCK = threading.RLock()


@dataclass(frozen=True)
class _WebFetchResolution:
    host: str
    address_infos: tuple[tuple[int, int, int, str, tuple], ...]


@contextmanager
def _pinned_web_fetch_resolution(resolution: _WebFetchResolution | None):
    if resolution is None:
        yield
        return

    with _WEB_FETCH_RESOLUTION_LOCK:
        original_getaddrinfo = socket.getaddrinfo

        def pinned_getaddrinfo(host, port, family=0, type=0, proto=0, flags=0):
            if isinstance(host, str) and host.strip().lower() == resolution.host:
                adjusted: list[tuple[int, int, int, str, tuple]] = []
                for info_family, info_type, info_proto, info_canonname, sockaddr in resolution.address_infos:
                    if family not in {0, info_family}:
                        continue
                    if type not in {0, info_type}:
                        continue
                    if proto not in {0, info_proto}:
                        continue
                    if not isinstance(sockaddr, tuple) or not sockaddr:
                        continue
                    if info_family == socket.AF_INET and len(sockaddr) >= 2:
                        adjusted_sockaddr = (sockaddr[0], port or sockaddr[1])
                    elif info_family == socket.AF_INET6 and len(sockaddr) >= 4:
                        adjusted_sockaddr = (sockaddr[0], port or sockaddr[1], sockaddr[2], sockaddr[3])
                    else:
                        adjusted_sockaddr = sockaddr
                    adjusted.append((info_family, info_type, info_proto, info_canonname, adjusted_sockaddr))
                if adjusted:
                    return adjusted
            return original_getaddrinfo(host, port, family, type, proto, flags)

        socket.getaddrinfo = pinned_getaddrinfo
        try:
            yield
        finally:
            socket.getaddrinfo = original_getaddrinfo


def _resolve_web_fetch_resolution(raw_url: str) -> _WebFetchResolution | None:
    parsed = urlparse(raw_url)
    if parsed.scheme not in {"http", "https"}:
        return None
    host = parsed.hostname
    if not isinstance(host, str) or not host:
        return None
    lowered = host.strip().lower()
    if lowered == "localhost":
        return None
    try:
        parsed_host = ip_address(lowered)
    except ValueError:
        try:
            address_infos = socket.getaddrinfo(lowered, parsed.port or None, proto=socket.IPPROTO_TCP)
        except OSError:
            return None
        filtered_infos: list[tuple[int, int, int, str, tuple]] = []
        for family, socktype, proto, canonname, sockaddr in address_infos:
            if family not in {socket.AF_INET, socket.AF_INET6}:
                continue
            if not isinstance(sockaddr, tuple) or not sockaddr:
                continue
            candidate = sockaddr[0]
            if not isinstance(candidate, str):
                continue
            if not ip_address(candidate).is_global:
                return None
            filtered_infos.append((family, socktype, proto, canonname, sockaddr))
        if not filtered_infos:
            return None
        return _WebFetchResolution(host=lowered, address_infos=tuple(filtered_infos))
    if not parsed_host.is_global:
        return None
    return None


def _principal_workspace_file_roots(principal_id: str | None) -> tuple[Path, ...]:
    principal = str(principal_id or "").strip()
    if not (
        principal in {"telegram_chat", "operator"}
        or principal.startswith(("user:", "workspace:", "web:", "telegram:", "slack:", "discord:"))
    ):
        return ()
    try:
        from nullion.workspace_storage import workspace_file_roots_for_principal

        return tuple(Path(root).resolve() for root in workspace_file_roots_for_principal(principal))
    except Exception:
        return ()


def _flow_context_is_scheduled_task_run(flow_context: object) -> bool:
    if not isinstance(flow_context, dict):
        return False
    if flow_context.get("scheduled_task_run"):
        return True
    trigger_flow = flow_context.get(FLOW_TRIGGER_CONTEXT_KEY)
    return isinstance(trigger_flow, dict) and bool(trigger_flow.get("scheduled_task_run"))


_FILESYSTEM_PATH_ARGUMENTS_BY_TOOL = {
    "archive_create": ("output_path", "source_dir"),
    "archive_extract": ("path", "output_dir", "manifest_output_path"),
    "audio_transcribe": ("path",),
    "file_download": ("output_path",),
    "file_read": ("path",),
    "file_write": ("path",),
    "file_patch": ("path",),
    "image_extract_text": ("path",),
    "image_generate": ("source_path", "output_path"),
    "document_create": ("output_path",),
    "presentation_create": ("output_path",),
    "spreadsheet_create": ("output_path",),
}


def _resolve_virtual_workspace_path(raw_path: str, *, principal_id: str | None) -> str:
    try:
        from nullion.workspace_storage import resolve_virtual_workspace_path_for_principal

        return str(resolve_virtual_workspace_path_for_principal(raw_path, principal_id))
    except Exception:
        return raw_path


def _with_resolved_virtual_workspace_paths(invocation: ToolInvocation) -> ToolInvocation:
    path_keys = _FILESYSTEM_PATH_ARGUMENTS_BY_TOOL.get(invocation.tool_name)
    if not path_keys:
        return invocation
    arguments = dict(invocation.arguments)
    changed = False
    for key in path_keys:
        value = arguments.get(key)
        if not isinstance(value, str) or not value:
            continue
        resolved_value = _resolve_virtual_workspace_path(value, principal_id=invocation.principal_id)
        if resolved_value == value:
            continue
        arguments[key] = resolved_value
        changed = True
    if not changed:
        return invocation
    return ToolInvocation(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        principal_id=invocation.principal_id,
        arguments=arguments,
        capsule_id=invocation.capsule_id,
        trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        flow_context=invocation.flow_context,
    )


def _effective_filesystem_roots(
    *,
    invocation: ToolInvocation,
    resolved_root: Path | None,
    resolved_allowed_roots: tuple[Path, ...] | None,
    include_principal_workspace: bool,
) -> tuple[Path, ...]:
    roots: list[Path] = []
    if resolved_root is not None:
        roots.append(resolved_root)
    if resolved_allowed_roots is not None:
        roots.extend(resolved_allowed_roots)
    if include_principal_workspace:
        roots.extend(_principal_workspace_file_roots(invocation.principal_id))
    return tuple(dict.fromkeys(root.resolve() for root in roots))


def _path_within_any_root(path: Path, roots: tuple[Path, ...]) -> bool:
    return any(_is_within_allowed_root(path, root) for root in roots)


def _safe_generated_artifact_filename(raw_name: object, *, suffix: str, fallback_stem: str) -> str:
    normalized_suffix = suffix if suffix.startswith(".") else f".{suffix}"
    name = Path(str(raw_name or "")).name.strip()
    if name and Path(name).suffix.lower() == normalized_suffix:
        stem = Path(name).stem
    elif name and not Path(name).suffix:
        stem = Path(name).name
    else:
        stem = str(fallback_stem or "nullion-artifact")
    stem = re.sub(r"[^A-Za-z0-9._-]+", "-", stem.strip()).strip("-._") or "nullion-artifact"
    return f"{stem[:80]}{normalized_suffix}"


def _workspace_generated_artifact_path(
    invocation: ToolInvocation,
    *,
    raw_path: object,
    suffix: str,
    stem: str,
) -> Path:
    normalized_suffix = suffix if suffix.startswith(".") else f".{suffix}"
    artifact_root = _workspace_artifact_root_for_principal(invocation.principal_id)
    if artifact_root is not None:
        try:
            raw_text = str(raw_path or "").strip()
            if raw_text:
                candidate = Path(raw_text).expanduser().resolve()
                if candidate.suffix.lower() == normalized_suffix and _path_within_any_root(candidate, (artifact_root,)):
                    return candidate
                filename = _safe_generated_artifact_filename(
                    candidate.name,
                    suffix=normalized_suffix,
                    fallback_stem=stem,
                )
                return (artifact_root / filename).resolve()
        except (OSError, RuntimeError, ValueError):
            pass
    from nullion.artifacts import artifact_path_for_generated_workspace_file

    return artifact_path_for_generated_workspace_file(
        principal_id=invocation.principal_id,
        suffix=normalized_suffix,
        stem=stem,
    ).resolve()


def _explicit_output_path_or_generated_artifact_path(
    invocation: ToolInvocation,
    *,
    raw_path: object,
    suffix: str,
    stem: str,
    roots: tuple[Path, ...],
) -> Path:
    normalized_suffix = suffix if suffix.startswith(".") else f".{suffix}"
    if isinstance(raw_path, str) and raw_path.strip():
        try:
            candidate = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser().resolve()
        except (OSError, RuntimeError, ValueError):
            candidate = Path(str(raw_path)).expanduser().resolve()
        artifact_root = _workspace_artifact_root_for_principal(invocation.principal_id)
        if (
            candidate.suffix.lower() == normalized_suffix
            and artifact_root is not None
            and not _path_within_any_root(candidate, (artifact_root,))
            and _generated_output_path_should_use_workspace_artifacts(candidate)
        ):
            return _workspace_generated_artifact_path(
                invocation,
                raw_path=raw_path,
                suffix=normalized_suffix,
                stem=stem,
            )
        if candidate.suffix.lower() == normalized_suffix and (
            _path_within_any_root(candidate, roots)
            or _is_approved_filesystem_path(candidate, invocation.trusted_filesystem_selectors)
        ):
            return candidate
    return _workspace_generated_artifact_path(
        invocation,
        raw_path=raw_path,
        suffix=normalized_suffix,
        stem=stem,
    )


def _generated_output_path_should_use_workspace_artifacts(path: Path) -> bool:
    try:
        from nullion.artifacts import nullion_data_home

        data_home = nullion_data_home()
    except Exception:
        return False
    return _path_within_any_root(path, (data_home,))


def _file_write_artifact_path_should_use_workspace_artifacts(
    path: Path,
    *,
    principal_id: str | None,
    effective_roots: tuple[Path, ...],
) -> bool:
    if path.suffix.lower() not in VALID_ATTACHMENT_EXTENSIONS:
        return False
    artifact_root = _workspace_artifact_root_for_principal(principal_id)
    if artifact_root is not None and _path_within_any_root(path, (artifact_root,)):
        return False
    try:
        from nullion.artifacts import nullion_data_home

        if _path_within_any_root(path, (nullion_data_home(),)):
            return True
    except Exception:
        pass
    return not _path_within_any_root(path, effective_roots)


def _resolve_local_workspace_file_input(
    raw_path: str,
    *,
    principal_id: str | None,
    effective_roots: tuple[Path, ...],
    trusted_filesystem_selectors: tuple[str, ...] = (),
) -> Path | None:
    path_text = str(raw_path or "").strip()
    if not path_text:
        return None
    resolved_text = _resolve_virtual_workspace_path(path_text, principal_id=principal_id)
    candidate = Path(resolved_text).expanduser()
    if candidate.is_absolute():
        resolved_candidate = candidate.resolve()
        if resolved_candidate.is_file() and (
            _path_within_any_root(resolved_candidate, effective_roots)
            or _is_approved_filesystem_path(resolved_candidate, trusted_filesystem_selectors)
        ):
            return resolved_candidate
        for root in effective_roots:
            if not _is_within_allowed_root(resolved_candidate, root):
                continue
            try:
                relative = resolved_candidate.relative_to(root)
            except ValueError:
                continue
            if relative.parts and relative.parts[0] in {"artifacts", "media", "files", "uploads", "scratch"}:
                relative = Path(*relative.parts[1:]) if len(relative.parts) > 1 else Path()
            if not relative.parts:
                continue
            for child_root_name in ("artifacts", "media", "files", "uploads", "scratch"):
                alternate = (root / child_root_name / relative).resolve()
                if alternate.is_file() and _path_within_any_root(alternate, effective_roots):
                    return alternate
        return None
    if any(part == ".." for part in candidate.parts):
        return None

    candidates: list[Path] = []
    for root in effective_roots:
        candidates.append(root / candidate)
        for child_root_name in ("artifacts", "media", "files", "uploads", "scratch"):
            candidates.append(root / child_root_name / candidate)

    seen: set[Path] = set()
    for possible_path in candidates:
        resolved_candidate = possible_path.resolve()
        if resolved_candidate in seen:
            continue
        seen.add(resolved_candidate)
        if resolved_candidate.is_file() and _path_within_any_root(resolved_candidate, effective_roots):
            return resolved_candidate
    return None


def _build_web_fetch_opener() -> urllib.request.OpenerDirector:
    https_handler = urllib.request.HTTPSHandler(context=_web_fetch_tls_context())
    return urllib.request.build_opener(
        urllib.request.ProxyHandler({}),
        https_handler,
        _SafeWebFetchRedirectHandler,
    )


def _web_fetch_tls_context() -> ssl.SSLContext:
    try:
        import certifi
    except ModuleNotFoundError:
        return ssl.create_default_context()
    try:
        return ssl.create_default_context(cafile=certifi.where())
    except Exception:
        logger.debug("Falling back to system TLS trust store for web fetches.", exc_info=True)
        return ssl.create_default_context()



def _is_global_literal_ip(host: str) -> bool:
    try:
        return ip_address(host.strip().lower()).is_global
    except ValueError:
        return False



_NETWORK_MODE_NONE = "none"
_NETWORK_MODE_LOCALHOST_ONLY = "localhost_only"
_NETWORK_MODE_APPROVED_ONLY = "approved_only"
_NETWORK_MODE_FULL = "full"
_VALID_NETWORK_MODES = {
    _NETWORK_MODE_NONE,
    _NETWORK_MODE_LOCALHOST_ONLY,
    _NETWORK_MODE_APPROVED_ONLY,
    _NETWORK_MODE_FULL,
}
_RESTRICTIVE_NETWORK_MODES = {
    _NETWORK_MODE_NONE,
    _NETWORK_MODE_LOCALHOST_ONLY,
    _NETWORK_MODE_APPROVED_ONLY,
}
_REQUIRED_RESTRICTIVE_TERMINAL_BACKEND_ATTESTED_CAPABILITIES = ("network_policy_enforced",)
_NETWORK_MODE_SCOPED_ATTESTED_CAPABILITIES = {
    _NETWORK_MODE_NONE: ("network_policy_enforced.none",),
    _NETWORK_MODE_LOCALHOST_ONLY: ("network_policy_enforced.localhost_only",),
    _NETWORK_MODE_APPROVED_ONLY: (
        "network_policy_enforced.approved_only",
        "approved_only_enforced_via_local_allowlist_proxy",
    ),
}
_EXPECTED_LAUNCHER_SUPPORTED_NETWORK_MODES = (
    _NETWORK_MODE_NONE,
    _NETWORK_MODE_LOCALHOST_ONLY,
    _NETWORK_MODE_APPROVED_ONLY,
)
_LAUNCHER_DESCRIPTOR_SCHEMA = "nullion.launcher.descriptor.v1"
_LAUNCHER_DESCRIPTOR_SOURCE_PROBE = "launcher_probe"
_LAUNCHER_DESCRIPTOR_SOURCE_LEGACY = "launcher_legacy"
_ALLOWED_LAUNCHER_DESCRIPTOR_SOURCES = {
    _LAUNCHER_DESCRIPTOR_SOURCE_PROBE,
    _LAUNCHER_DESCRIPTOR_SOURCE_LEGACY,
}
_EXPECTED_LAUNCHER_ENFORCEMENT_BY_MODE = {
    _NETWORK_MODE_NONE: "sandbox-exec deny network*",
    _NETWORK_MODE_LOCALHOST_ONLY: "sandbox-exec remote ip localhost:* only",
    _NETWORK_MODE_APPROVED_ONLY: "sandbox-exec localhost proxy only + in-process allowlist proxy",
}

_FILESYSTEM_BOUNDARY_DEFAULT = "default"
_FILESYSTEM_BOUNDARY_TRUSTED_ROOTS_ONLY = "trusted_roots_only"


def _required_attested_capabilities_for_network_mode(network_mode: str | None) -> tuple[str, ...]:
    if not _is_restrictive_network_mode(network_mode):
        return ()
    base = list(_REQUIRED_RESTRICTIVE_TERMINAL_BACKEND_ATTESTED_CAPABILITIES)
    for capability in _NETWORK_MODE_SCOPED_ATTESTED_CAPABILITIES.get(network_mode, ()):  # pragma: no branch
        if capability not in base:
            base.append(capability)
    return tuple(base)


def _serialize_boundary_fact(fact) -> dict[str, object]:
    return {
        "kind": fact.kind.value if hasattr(fact.kind, "value") else str(fact.kind),
        "target": fact.target,
        "host": fact.attributes.get("host"),
        "scheme": fact.attributes.get("scheme"),
        "address_class": fact.attributes.get("address_class"),
        "command_family": fact.attributes.get("command_family"),
    }


def _egress_attempts_for_invocation(invocation: ToolInvocation) -> list[dict[str, object]]:
    return [_serialize_boundary_fact(fact) for fact in extract_boundary_facts(invocation)]


def _selector_candidates_for_boundary_target(target: str) -> dict[str, str]:
    parsed = urlparse(target)
    if parsed.scheme and parsed.hostname:
        domain_selector = normalize_outbound_network_selector(target)
        return {
            "allow_once": domain_selector,
            "always_allow": domain_selector,
        }
    return {
        "allow_once": target,
        "always_allow": target,
    }


def _outbound_network_approval_context_from_result(result: ToolResult) -> dict[str, object] | None:
    if result.tool_name != "terminal_exec":
        return None
    output = result.output if isinstance(result.output, dict) else {}
    if output.get("reason") != "network_denied":
        return None
    if output.get("network_mode") not in {None, _NETWORK_MODE_NONE, _NETWORK_MODE_APPROVED_ONLY}:
        return None
    attempts = output.get("egress_attempts")
    if not isinstance(attempts, list):
        return None
    for attempt in attempts:
        if not isinstance(attempt, dict):
            continue
        if attempt.get("kind") != "outbound_network":
            continue
        target = attempt.get("target")
        if not isinstance(target, str) or not target:
            continue
        if attempt.get("address_class") != "public":
            continue
        return {
            "tool_name": "terminal_exec",
            "boundary_kind": "outbound_network",
            "target": target,
            "selector_candidates": _selector_candidates_for_boundary_target(target),
        }
    return None


def _boundary_approval_context_from_fact(fact) -> dict[str, object] | None:
    if fact.kind is BoundaryKind.OUTBOUND_NETWORK:
        if fact.attributes.get("address_class") != "public":
            return None
        return {
            "tool_name": fact.tool_name,
            "boundary_kind": fact.kind.value,
            "target": fact.target,
            "selector_candidates": _selector_candidates_for_boundary_target(fact.target),
        }
    if fact.kind is BoundaryKind.FILESYSTEM_ACCESS:
        return {
            "tool_name": fact.tool_name,
            "boundary_kind": fact.kind.value,
            "operation": fact.operation,
            "path": fact.attributes.get("path") or fact.target,
            "target": fact.target,
            "selector_candidates": {
                "allow_once": fact.target,
                "always_allow": fact.target,
            },
        }
    if fact.kind is BoundaryKind.ACCOUNT_ACCESS:
        account_type = fact.attributes.get("account_type", "unknown")
        return {
            "tool_name": fact.tool_name,
            "boundary_kind": fact.kind.value,
            "operation": fact.operation,
            "target": fact.target,
            "account_type": account_type,
            "selector_candidates": {
                "allow_once": fact.target,
                "always_allow": f"{account_type}:*",
            },
        }
    return None


def _selector_matches_target(*, selector: str, target: str) -> bool:
    parsed_target = urlparse(target)
    if parsed_target.scheme in {"http", "https"} and parsed_target.hostname:
        if normalize_outbound_network_selector(selector) == normalize_outbound_network_selector(target):
            return True
    if selector == target:
        return True
    if fnmatch(target, selector):
        return True
    if selector.endswith("/*") and target == selector[:-2]:
        return True
    if _selector_matches_www_family(selector=selector, target=target):
        return True
    return False


def _boundary_approval_match_key(*, context: dict[str, object], fallback_target: str) -> tuple[str, str] | None:
    boundary_kind = context.get("boundary_kind")
    if not isinstance(boundary_kind, str) or not boundary_kind:
        return None
    selectors = context.get("selector_candidates") if isinstance(context.get("selector_candidates"), dict) else {}
    candidates: list[object] = [
        selectors.get("always_allow") if isinstance(selectors, dict) else None,
        selectors.get("allow_once") if isinstance(selectors, dict) else None,
        context.get("target"),
        fallback_target,
    ]
    for candidate in candidates:
        if not isinstance(candidate, str) or not candidate:
            continue
        if boundary_kind == BoundaryKind.OUTBOUND_NETWORK.value:
            return (boundary_kind, normalize_outbound_network_selector(candidate))
        return (boundary_kind, candidate)
    return None


def _refresh_pending_approval_request(approval: ApprovalRequest, *, context: dict[str, object], resource: str) -> ApprovalRequest:
    merged_context = dict(approval.context or {})
    merged_context.update(context)
    return replace(
        approval,
        resource=resource,
        created_at=datetime.now(UTC),
        context=merged_context,
    )


def _plain_text_from_html_body(html_body: str) -> str:
    text = re.sub(r"(?is)<\s*br\s*/?\s*>", "\n", str(html_body or ""))
    text = re.sub(r"(?is)</\s*(?:p|div|section|article|tr|h[1-6])\s*>", "\n", text)
    text = re.sub(r"(?is)<\s*style\b[\s\S]*?</\s*style\s*>", " ", text)
    text = re.sub(r"(?is)<\s*script\b[\s\S]*?</\s*script\s*>", " ", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = unescape(text)
    lines = [" ".join(line.split()) for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def _email_html_body_from_invocation(invocation: ToolInvocation) -> str:
    inline_html = str(invocation.arguments.get("html_body") or "").strip()
    if inline_html:
        return inline_html
    html_path = str(invocation.arguments.get("html_path") or "").strip()
    if not html_path:
        return ""
    effective_roots = _principal_workspace_file_roots(invocation.principal_id)
    resolved = _resolve_local_workspace_file_input(
        html_path,
        principal_id=invocation.principal_id,
        effective_roots=effective_roots,
        trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
    )
    if resolved is None:
        resolved = Path(_resolve_virtual_workspace_path(html_path, principal_id=invocation.principal_id)).expanduser()
    if not resolved.exists() or not resolved.is_file():
        raise FileNotFoundError(f"HTML path does not exist or is not a file: {html_path}")
    return resolved.read_text(encoding="utf-8", errors="ignore").strip()


def _email_html_preview_path_for_invocation(invocation: ToolInvocation) -> str | None:
    try:
        html_path = str(invocation.arguments.get("html_path") or "").strip()
        if html_path:
            return html_path
        html_body = str(invocation.arguments.get("html_body") or "").strip()
        if not html_body:
            return None
        from nullion.artifacts import artifact_path_for_generated_workspace_file, normalize_html_document

        preview_path = artifact_path_for_generated_workspace_file(
            principal_id=invocation.principal_id,
            suffix=".html",
            stem="email-preview",
        )
        preview_path.parent.mkdir(parents=True, exist_ok=True)
        preview_path.write_text(normalize_html_document(html_body, title="Email Preview"), encoding="utf-8")
        return str(preview_path)
    except Exception:
        logger.debug("Could not create email HTML approval preview", exc_info=True)
        return None


def _selector_matches_www_family(*, selector: str, target: str) -> bool:
    selector_base_url = selector[:-2] if selector.endswith("/*") else selector
    parsed_selector = urlparse(selector_base_url)
    parsed_target = urlparse(target)
    if parsed_selector.scheme != parsed_target.scheme:
        return False
    selector_host = (parsed_selector.hostname or "").lower()
    target_host = (parsed_target.hostname or "").lower()
    if not selector_host or not target_host:
        return False
    selector_root = selector_host[4:] if selector_host.startswith("www.") else selector_host
    target_root = target_host[4:] if target_host.startswith("www.") else target_host
    if selector_root != target_root:
        return False
    if selector_host not in {selector_root, f"www.{selector_root}"}:
        return False
    if target_host not in {target_root, f"www.{target_root}"}:
        return False
    if selector.endswith("/*"):
        return True
    return parsed_selector.path.rstrip("/") == parsed_target.path.rstrip("/")


def _network_is_approved(*, target: str, approved_targets: object) -> bool:
    if not isinstance(approved_targets, (list, tuple, set, frozenset)):
        return False
    return any(isinstance(pattern, str) and _selector_matches_target(selector=pattern, target=target) for pattern in approved_targets)



def _normalize_network_mode(network_mode: object) -> str | None:
    if not isinstance(network_mode, str):
        return None
    if not network_mode:
        return None
    if network_mode in _VALID_NETWORK_MODES:
        return network_mode
    return None



def _network_attempt_allowed(*, attempt: dict[str, object], network_mode: str | None, approved_targets: object) -> bool:
    if network_mode in {None, _NETWORK_MODE_FULL}:
        return True
    if network_mode == _NETWORK_MODE_NONE:
        return False
    if network_mode == _NETWORK_MODE_LOCALHOST_ONLY:
        return attempt.get("address_class") in {"localhost", "loopback"}
    if network_mode == _NETWORK_MODE_APPROVED_ONLY:
        target = attempt.get("target")
        return isinstance(target, str) and _network_is_approved(target=target, approved_targets=approved_targets)
    return True



def _is_restrictive_network_mode(network_mode: str | None) -> bool:
    return network_mode in _RESTRICTIVE_NETWORK_MODES



def _missing_backend_attested_capabilities(
    descriptor: TerminalBackendDescriptor,
    *,
    required_capabilities: Iterable[str],
) -> tuple[str, ...]:
    return tuple(
        capability
        for capability in required_capabilities
        if capability not in descriptor.attested_capabilities
    )



def _requires_probe_derived_launcher_attestation(required_capabilities: Iterable[str]) -> bool:
    return any(isinstance(capability, str) and capability.startswith("network_policy_enforced") for capability in required_capabilities)



def _launcher_command_requires_probe_attestation(command: object) -> bool:
    if not isinstance(command, str):
        return False
    normalized_command = command.strip()
    if not normalized_command:
        return False
    return Path(normalized_command).parent != Path()



def _verify_launcher_v2_evidence(
    descriptor: TerminalBackendDescriptor,
    *,
    required_capabilities: Iterable[str],
) -> TerminalAttestationVerificationResult | None:
    evidence = descriptor.evidence
    if evidence is None or evidence.format != "nullion.launcher.v2":
        return None
    if descriptor.mode != "launcher":
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence requires launcher backend mode",
            metadata={"descriptor_mode": descriptor.mode},
        )
    payload = evidence.payload if isinstance(evidence.payload, dict) else {}
    descriptor_source = descriptor.metadata.get("descriptor_source")
    evidence_source = evidence.metadata.get("descriptor_source")
    if (
        descriptor_source not in _ALLOWED_LAUNCHER_DESCRIPTOR_SOURCES
        or evidence_source not in _ALLOWED_LAUNCHER_DESCRIPTOR_SOURCES
        or descriptor_source != evidence_source
    ):
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence provenance metadata mismatch",
            metadata={
                "descriptor_source": descriptor_source,
                "evidence_source": evidence_source,
            },
        )
    if (
        _requires_probe_derived_launcher_attestation(required_capabilities)
        and descriptor_source != _LAUNCHER_DESCRIPTOR_SOURCE_PROBE
        and _launcher_command_requires_probe_attestation(
            descriptor.metadata.get("launcher_command") or payload.get("launcher_command")
        )
    ):
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="restrictive launcher modes require probe-derived attestation",
            metadata={
                "descriptor_source": descriptor_source,
                "launcher_command": descriptor.metadata.get("launcher_command") or payload.get("launcher_command"),
            },
        )
    descriptor_schema = descriptor.metadata.get("descriptor_schema")
    evidence_schema = evidence.metadata.get("descriptor_schema")
    payload_schema = payload.get("schema")
    if (
        descriptor_schema != _LAUNCHER_DESCRIPTOR_SCHEMA
        or evidence_schema != _LAUNCHER_DESCRIPTOR_SCHEMA
        or payload_schema != _LAUNCHER_DESCRIPTOR_SCHEMA
    ):
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence schema metadata mismatch",
            metadata={
                "descriptor_schema": descriptor_schema,
                "evidence_schema": evidence_schema,
                "payload_schema": payload_schema,
            },
        )
    payload_mode = payload.get("mode")
    if payload_mode != "launcher":
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence payload mode mismatch",
            metadata={"payload_mode": payload_mode},
        )
    supported_network_modes = payload.get("supported_network_modes")
    if tuple(supported_network_modes) != _EXPECTED_LAUNCHER_SUPPORTED_NETWORK_MODES:
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence supported_network_modes mismatch",
            metadata={"supported_network_modes": supported_network_modes},
        )
    if payload.get("requires_sandbox_exec_for_restrictive_modes") is not True:
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence must require sandbox-exec for restrictive modes",
            metadata={"requires_sandbox_exec_for_restrictive_modes": payload.get("requires_sandbox_exec_for_restrictive_modes")},
        )
    if payload.get("proxy_bind_scope") != "127.0.0.1":
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence must bind approved_only proxy to 127.0.0.1",
            metadata={"proxy_bind_scope": payload.get("proxy_bind_scope")},
        )
    enforcement_by_mode = payload.get("enforcement_by_mode")
    if not isinstance(enforcement_by_mode, dict):
        return TerminalAttestationVerificationResult(
            is_valid=False,
            failure_reason="launcher evidence enforcement_by_mode missing",
            metadata={"enforcement_by_mode": enforcement_by_mode},
        )
    for mode, expected in _EXPECTED_LAUNCHER_ENFORCEMENT_BY_MODE.items():
        if enforcement_by_mode.get(mode) != expected:
            failure_reason = (
                "approved_only evidence must declare proxy-backed allowlist enforcement"
                if mode == _NETWORK_MODE_APPROVED_ONLY
                else f"launcher evidence enforcement mismatch for {mode}"
            )
            return TerminalAttestationVerificationResult(
                is_valid=False,
                failure_reason=failure_reason,
                metadata={"mode": mode, "enforcement": enforcement_by_mode.get(mode)},
            )
    return TerminalAttestationVerificationResult(is_valid=True)


class _DefaultTerminalAttestationVerifier:
    def verify(
        self,
        descriptor: TerminalBackendDescriptor,
        *,
        required_capabilities: Iterable[str],
    ) -> TerminalAttestationVerificationResult:
        missing_capabilities = _missing_backend_attested_capabilities(
            descriptor,
            required_capabilities=required_capabilities,
        )
        if missing_capabilities:
            return TerminalAttestationVerificationResult(
                is_valid=False,
                failure_reason="missing required attested capabilities: " + ", ".join(missing_capabilities),
                metadata={"missing_capabilities": list(missing_capabilities)},
            )
        launcher_v2_verification = _verify_launcher_v2_evidence(
            descriptor,
            required_capabilities=required_capabilities,
        )
        if launcher_v2_verification is not None:
            return launcher_v2_verification
        return TerminalAttestationVerificationResult(is_valid=True)



def verify_terminal_backend_attestation(
    descriptor: TerminalBackendDescriptor,
    *,
    required_capabilities: Iterable[str],
    verifier: TerminalAttestationVerifier | None = None,
) -> TerminalAttestationVerificationResult:
    active_verifier = verifier or _DefaultTerminalAttestationVerifier()
    return active_verifier.verify(descriptor, required_capabilities=required_capabilities)



_CANONICAL_TOOL_STATUSES = {
    "completed",
    "failed",
    "denied",
    "nonterminal",
    "unknown",
}
_COMPLETED_TOOL_STATUS_ALIASES = {"completed", "complete", "success", "succeeded", "ok", "done"}
_FAILED_TOOL_STATUS_ALIASES = {"failed", "failure", "error", "errored", "timeout", "timed_out", "partial"}
_DENIED_TOOL_STATUS_ALIASES = {"denied", "blocked", "approval_required", "capability_denied", "capability_not_granted"}
_NONTERMINAL_TOOL_STATUS_ALIASES = {"running", "pending", "started", "in_progress", "queued"}

class ToolRiskLevel(str, Enum):
    LOW = "low"
    MEDIUM = "medium"
    HIGH = "high"


_TOOL_RISK_SCORES = {
    ToolRiskLevel.LOW: 2,
    ToolRiskLevel.MEDIUM: 5,
    ToolRiskLevel.HIGH: 8,
}


class ToolSideEffectClass(str, Enum):
    READ = "read"
    WRITE = "write"
    EXTERNAL_WRITE = "external_write"
    ACCOUNT_WRITE = "account_write"
    DANGEROUS_EXEC = "dangerous_exec"


@dataclass(slots=True)
class ToolSpec:
    name: str
    description: str
    risk_level: ToolRiskLevel
    side_effect_class: ToolSideEffectClass
    requires_approval: bool
    timeout_seconds: int
    filesystem_boundary_policy: str = _FILESYSTEM_BOUNDARY_DEFAULT
    permission_scope: str = "global"
    input_schema: dict[str, object] | None = None
    capability_tags: tuple[str, ...] = ()
    continuation_tools: tuple[str, ...] = ()


@dataclass(slots=True)
class ToolInvocation:
    invocation_id: str
    tool_name: str
    principal_id: str
    arguments: dict[str, object]
    capsule_id: str | None = None
    trusted_filesystem_selectors: tuple[str, ...] = ()
    flow_context: dict[str, object] | None = None


@dataclass(slots=True)
class ToolResult:
    invocation_id: str
    tool_name: str
    status: str
    output: dict[str, object]
    error: str | None = None


ToolHandler = Callable[[ToolInvocation], ToolResult]
ToolCleanupHook = Callable[[str | None], None]


def _tool_handler_timeout_seconds(spec: ToolSpec) -> float | None:
    try:
        timeout = float(spec.timeout_seconds)
    except (TypeError, ValueError):
        return None
    if timeout <= 0:
        return None
    return timeout


def _tool_handler_should_use_timeout_worker(spec: ToolSpec) -> bool:
    return spec.side_effect_class == ToolSideEffectClass.READ


def _tool_handler_timeout_result(invocation: ToolInvocation, *, timeout_seconds: float) -> ToolResult:
    timeout_display = f"{timeout_seconds:g}"
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={
            "reason": "handler_timeout",
            "timeout_seconds": timeout_seconds,
            "message": f"Tool timed out after {timeout_display}s.",
        },
        error=f"Tool timed out after {timeout_display}s.",
    )


def _invoke_tool_handler_with_timeout(
    spec: ToolSpec,
    handler: ToolHandler,
    invocation: ToolInvocation,
) -> ToolResult:
    timeout_seconds = _tool_handler_timeout_seconds(spec)
    if timeout_seconds is None or not _tool_handler_should_use_timeout_worker(spec):
        return handler(invocation)

    complete = threading.Event()
    result_box: dict[str, object] = {}

    def _run_handler() -> None:
        try:
            result_box["result"] = handler(invocation)
        except BaseException as exc:  # pragma: no cover - re-raised on caller thread
            result_box["exception"] = exc
        finally:
            complete.set()

    thread = threading.Thread(
        target=_run_handler,
        name=f"nullion-tool-{spec.name}-timeout",
        daemon=True,
    )
    thread.start()
    if not complete.wait(timeout_seconds):
        return _tool_handler_timeout_result(invocation, timeout_seconds=timeout_seconds)
    if "exception" in result_box:
        raise result_box["exception"]  # type: ignore[misc]
    result = result_box.get("result")
    if isinstance(result, ToolResult):
        return result
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={"reason": "invalid_handler_result"},
        error=f"Tool handler returned invalid result: {type(result).__name__}",
    )


def _clear_deep_agent_profile_cache() -> None:
    try:
        from nullion.deep_agent_profiles import clear_deep_agent_profile_caches

        clear_deep_agent_profile_caches()
        from nullion.cron_delivery import clear_cron_execution_metadata_caches

        clear_cron_execution_metadata_caches()
    except Exception:
        logger.debug("Could not clear Deep Agents profile cache", exc_info=True)


_TEXT_FILE_WRITE_BLOCKED_EXTENSIONS = frozenset({
    ".doc",
    ".docx",
    ".pdf",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
})
_TERMINAL_DELIVERABLE_ARTIFACT_EXTENSIONS = frozenset(VALID_ATTACHMENT_EXTENSIONS) | _TEXT_FILE_WRITE_BLOCKED_EXTENSIONS
_MAX_TERMINAL_DISCOVERED_ARTIFACTS = 25
_MAX_TERMINAL_ARTIFACT_SCAN_ENTRIES = 10000

_CRON_TOOL_PROPERTIES: dict[str, dict[str, str]] = {
    "id": {"type": "string", "description": "Cron job id. Required for update operations."},
    "name": {"type": "string", "description": "Human-readable cron name."},
    "schedule": {"type": "string", "description": "Cron schedule expression."},
    "task": {"type": "string", "description": "Task instructions to run when the cron fires."},
    "enabled": {"type": "boolean", "description": "Whether the cron is active."},
    "workspace_id": {"type": "string", "description": "Workspace id that owns the cron."},
    "delivery_channel": {"type": "string", "description": "Delivery channel such as web, telegram, slack, or discord."},
    "delivery_target": {"type": "string", "description": "Channel-specific delivery target."},
    "html_image_delivery_mode": {
        "type": "string",
        "description": "HTML image delivery mode: linked, auto, or self_contained.",
    },
}


def _cron_tool_properties() -> dict[str, object]:
    return {name: dict(schema) for name, schema in _CRON_TOOL_PROPERTIES.items()}


def _default_input_schema_for_tool(tool_name: str) -> dict[str, object]:
    def cron_tool_properties() -> dict[str, object]:
        return {name: dict(schema) for name, schema in _CRON_TOOL_PROPERTIES.items()}

    schemas: dict[str, dict[str, object]] = {
        "create_cron": {
            "type": "object",
            "properties": cron_tool_properties(),
            "required": ["name", "schedule", "task"],
            "additionalProperties": False,
        },
        "list_crons": {
            "type": "object",
            "properties": {
                "workspace_id": {"type": "string", "description": "Workspace id to list. Defaults to the current workspace."},
                "include_all_workspaces": {"type": "boolean", "description": "Whether to include crons from every workspace."},
            },
            "additionalProperties": False,
        },
        "update_cron": {
            "type": "object",
            "properties": cron_tool_properties(),
            "required": ["id"],
            "additionalProperties": False,
        },
        "delete_cron": {
            "type": "object",
            "properties": {"id": {"type": "string", "description": "Cron job id to delete."}},
            "required": ["id"],
            "additionalProperties": False,
        },
        "delete_reminder": {
            "type": "object",
            "properties": {"task_id": {"type": "string", "description": "Reminder task id to delete."}},
            "required": ["task_id"],
            "additionalProperties": False,
        },
        "update_reminder": {
            "type": "object",
            "properties": {
                "task_id": {"type": "string", "description": "Reminder task id to update."},
                "text": {"type": "string", "description": "New reminder message. Omit to keep the current text."},
                "due_at": {
                    "type": "string",
                    "description": "New absolute ISO 8601 due time. Include timezone offset when known.",
                },
                "due_in_seconds": {
                    "type": "number",
                    "description": "New relative delay from the current moment, in seconds.",
                },
            },
            "required": ["task_id"],
            "additionalProperties": False,
        },
        "toggle_cron": {
            "type": "object",
            "properties": {
                "id": {"type": "string", "description": "Single cron job id to enable or disable."},
                "ids": {
                    "type": "array",
                    "description": "Exact cron ids from list_crons to enable or disable as one bulk operation.",
                    "items": {"type": "string"},
                },
                "all_current_workspace": {
                    "type": "boolean",
                    "description": (
                        "Enable or disable every cron in the current workspace. Use only when the structured "
                        "request applies to the whole current workspace."
                    ),
                },
                "enabled": {"type": "boolean", "description": "True to enable the cron, false to disable it."},
            },
            "required": ["enabled"],
            "additionalProperties": False,
        },
        "run_cron": {
            "type": "object",
            "properties": {
                "id": {"type": "string", "description": "Cron job id to run immediately."},
                "ids": {
                    "type": "array",
                    "description": "Exact cron ids to run immediately as one bulk operation after list_crons.",
                    "items": {"type": "string"},
                },
                "all_enabled": {
                    "type": "boolean",
                    "description": "Run every enabled cron in the current workspace as one bulk operation.",
                },
                "name": {"type": "string", "description": "Cron job name to run immediately when id is unavailable."},
            },
            "additionalProperties": False,
        },
        "file_read": {
            "type": "object",
            "properties": {"path": {"type": "string", "description": "Absolute or workspace-relative file path to read."}},
            "required": ["path"],
            "additionalProperties": False,
        },
        "archive_extract": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": (
                        "Absolute or workspace-relative archive/compressed file path to extract. The completed "
                        "result includes full CSV, JSON, and XLSX manifest paths for the extracted file list."
                    ),
                },
                "output_dir": {
                    "type": "string",
                    "description": "Optional absolute or workspace-relative output directory. Defaults to a sibling extraction folder.",
                },
                "manifest_output_path": {
                    "type": "string",
                    "description": (
                        "Optional absolute or workspace-relative .xlsx path for the generated extracted-file manifest. "
                        "Use this when the user requested a specific final workbook filename."
                    ),
                },
                "include_metadata": {
                    "type": "boolean",
                    "description": "Set true only when macOS/archive metadata sidecar entries are explicitly needed.",
                },
            },
            "required": ["path"],
            "additionalProperties": False,
        },
        "archive_create": {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": "Absolute or workspace-relative archive/compressed file path to create.",
                },
                "source_paths": {
                    "type": "array",
                    "description": "Files or directories to include in the archive.",
                    "items": {"type": "string"},
                },
                "source_dir": {
                    "type": "string",
                    "description": "Optional directory whose safe, non-metadata contents should be archived.",
                },
                "include_metadata": {
                    "type": "boolean",
                    "description": "Set true only when macOS/archive metadata sidecar entries are explicitly needed.",
                },
            },
            "required": ["output_path"],
            "additionalProperties": False,
        },
        "file_write": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute or workspace-relative file path to write."},
                "content": {
                    "type": "string",
                    "description": (
                        "Text content to write. Do not fabricate placeholder/example source URLs in delivered "
                        "artifacts; if current public source data is needed, use web/browser tools first. "
                        "For self-contained HTML, set img src values to local image artifact paths; file_write "
                        "will inline those local images as data URIs unless inline_local_html_images is false."
                    ),
                },
                "inline_local_html_images": {
                    "type": "boolean",
                    "description": (
                        "HTML only. Defaults to true for self-contained HTML. Set false for linked/non-self-contained "
                        "HTML so local image paths remain relative/sibling references instead of being inlined."
                    ),
                },
                "disallow_html_data_images": {
                    "type": "boolean",
                    "description": (
                        "HTML only. Set true for linked/non-self-contained HTML to reject data:image sources."
                    ),
                },
            },
            "required": ["path", "content"],
            "additionalProperties": False,
        },
        "document_create": {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": (
                        "Optional destination .docx path. If .html/.htm is provided, Nullion writes an HTML "
                        "document from the same structured content. If .pdf is provided, Nullion writes a real "
                        "PDF from the same structured content. If omitted, Nullion creates a .docx file in the "
                        "artifact directory."
                    ),
                },
                "title": {
                    "type": "string",
                    "description": (
                        "Optional document title used for the heading and default filename. "
                        "The generated document uses a report-quality layout profile with styled headings, readable spacing, and verified media embeds."
                    ),
                },
                "paragraphs": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional body paragraphs to include in order. Plain http(s) URLs are converted into clickable links.",
                },
                "sections": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string", "description": "Section heading."},
                            "body": {"type": "string", "description": "Section body text. Plain http(s) URLs are converted into clickable links."},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Optional bullet text for this section. Plain http(s) URLs are converted into clickable links.",
                            },
                            "image_paths": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": (
                                    "Optional existing content image artifact paths for this section. "
                                    "Browser screenshot artifacts must be supplied through screenshot_paths."
                                ),
                            },
                            "screenshot_paths": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": (
                                    "Optional browser/page screenshot artifact paths for this section."
                                ),
                            },
                        },
                        "additionalProperties": False,
                    },
                    "description": "Structured document sections.",
                },
                "tables": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "headers": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Column headers for a real Word table.",
                            },
                            "rows": {
                                "type": "array",
                                "items": {
                                    "type": "array",
                                    "items": {
                                        "anyOf": [
                                            {"type": "string"},
                                            {"type": "number"},
                                            {"type": "integer"},
                                            {"type": "boolean"},
                                            {"type": "null"},
                                        ],
                                    },
                                },
                                "description": "Table rows aligned to headers.",
                            },
                        },
                        "required": ["headers", "rows"],
                        "additionalProperties": False,
                    },
                    "description": "Optional real Word tables to include after introductory paragraphs.",
                },
                "image_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional existing content image artifact paths to embed in the document. "
                        "Browser screenshot artifacts must be supplied through screenshot_paths."
                    ),
                },
                "screenshot_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional browser/page screenshot artifact paths to embed in the document.",
                },
                "inline_local_html_images": {
                    "type": "boolean",
                    "description": "HTML output only. Defaults to true; set false for linked/non-self-contained HTML.",
                },
                "disallow_html_data_images": {
                    "type": "boolean",
                    "description": "HTML output only. Set true for linked/non-self-contained HTML to reject data:image sources.",
                },
            },
            "additionalProperties": False,
        },
        "spreadsheet_create": {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": (
                        "Optional destination path. .xlsx creates an Excel workbook with embedded media/charts; "
                        ".csv and .tsv create delimited text from the same rows/columns. If omitted, Nullion "
                        "creates an .xlsx file in the artifact directory."
                    ),
                },
                "title": {"type": "string", "description": "Optional workbook title used for the default filename."},
                "sheet_name": {"type": "string", "description": "Optional worksheet name."},
                "columns": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional ordered column names. If omitted, object row keys are used.",
                },
                "rows": {
                    "type": "array",
                    "items": {
                        "anyOf": [
                            {"type": "object"},
                            {
                                "type": "array",
                                "items": {
                                    "anyOf": [
                                        {"type": "string"},
                                        {"type": "number"},
                                        {"type": "integer"},
                                        {"type": "boolean"},
                                        {"type": "null"},
                                        {"type": "object"},
                                        {"type": "array", "items": {}},
                                    ],
                                },
                            },
                        ],
                    },
                    "description": (
                        "Rows as objects keyed by column name, or arrays matching the columns. "
                        "Use strings beginning with '=' for real Excel formula cells when formulas are requested; "
                        "static calculated numbers or prose about formula assumptions do not satisfy a formula request. "
                        "When the request or structured plan specifies an item/row count, make rows match that "
                        "total count; do not multiply a total count across sources/categories unless the request "
                        "explicitly asks for a count per source/category."
                    ),
                },
                "sheets": {
                    "type": "array",
                    "minItems": 1,
                    "items": {
                        "type": "object",
                        "properties": {
                            "sheet_name": {
                                "type": "string",
                                "description": "Worksheet name.",
                            },
                            "columns": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Ordered column names for this worksheet.",
                            },
                            "rows": {
                                "type": "array",
                                "items": {
                                    "anyOf": [
                                        {"type": "object"},
                                        {"type": "array", "items": {}},
                                    ]
                                },
                                "description": "Structured data rows for this worksheet.",
                            },
                        },
                        "required": ["sheet_name", "rows"],
                        "additionalProperties": False,
                    },
                    "description": (
                        "Optional multi-sheet workbook contents. Use this instead of top-level sheet_name, columns, "
                        "and rows whenever the requested workbook needs more than one worksheet, including a "
                        "separate Summary sheet. The first entry is the primary data sheet."
                    ),
                },
                "expected_rows": {
                    "type": "integer",
                    "minimum": 0,
                    "description": (
                        "Optional total data-row count required by the current request or structured plan. "
                        "Set this when the artifact has an explicit item/row count so the tool can reject a "
                        "wrong-sized table before delivery."
                    ),
                },
                "formulas_required": {
                    "type": "boolean",
                    "description": (
                        "Set true when the user requested formulas. The tool fails unless at least one row value "
                        "is a real Excel formula string beginning with '='."
                    ),
                },
                "image_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional existing content image artifact paths to embed, aligned to rows when possible. "
                        "Browser screenshot artifacts must be supplied through screenshot_paths."
                    ),
                },
                "screenshot_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional browser/page screenshot artifact paths to embed in a separate Screenshot column, "
                        "aligned to rows when possible."
                    ),
                },
                "charts": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Chart title."},
                            "sheet_name": {
                                "type": "string",
                                "description": (
                                    "Worksheet containing the category and value columns. Defaults to the first "
                                    "sheet; set this for charts based on another sheet in a multi-sheet workbook."
                                ),
                            },
                            "type": {
                                "type": "string",
                                "enum": ["bar", "line", "pie"],
                                "description": "Chart type.",
                            },
                            "categories_column": {
                                "type": "string",
                                "description": "Column name to use for category labels.",
                            },
                            "values_column": {
                                "type": "string",
                                "description": "Column name containing numeric chart values.",
                            },
                            "anchor": {
                                "type": "string",
                                "description": "Optional Excel anchor cell, for example H2.",
                            },
                        },
                        "required": ["type", "categories_column", "values_column"],
                        "additionalProperties": False,
                    },
                    "description": "Optional real Excel charts generated from columns on the selected worksheet.",
                },
                "conditional_formats": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "type": {
                                "type": "string",
                                "enum": ["data_bar", "color_scale"],
                                "description": "Conditional-format rule type.",
                            },
                            "column": {
                                "type": "string",
                                "description": "Column name to format.",
                            },
                            "sheet_name": {
                                "type": "string",
                                "description": (
                                    "Worksheet containing the column. Defaults to the first sheet; set this for "
                                    "conditional formatting on another sheet in a multi-sheet workbook."
                                ),
                            },
                            "color": {
                                "type": "string",
                                "description": "Optional six-character hex color for data_bar rules.",
                            },
                        },
                        "required": ["type", "column"],
                        "additionalProperties": False,
                    },
                    "description": "Optional real Excel conditional formatting for numeric columns.",
                },
            },
            "additionalProperties": False,
        },
        "presentation_create": {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": "Optional destination .pptx path. If omitted, Nullion creates one in the artifact directory.",
                },
                "title": {
                    "type": "string",
                    "description": (
                        "Optional deck title used for the default filename. "
                        "The generated deck uses a report-quality slide profile with styled titles, readable text, and aspect-ratio-safe media placement."
                    ),
                },
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide title."},
                            "body": {"type": "string", "description": "Short body text for the slide."},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Optional bullet text for the slide.",
                            },
                            "image_paths": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": (
                                    "Optional existing content image artifact paths for this slide. "
                                    "Browser screenshot artifacts must be supplied through screenshot_paths."
                                ),
                            },
                            "screenshot_paths": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Optional browser/page screenshot artifact paths for this slide.",
                            },
                        },
                        "additionalProperties": False,
                    },
                    "description": "Structured slide contents. If omitted, image_paths are placed one per slide.",
                },
                "image_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional existing content image artifact paths to place into slides. "
                        "Browser screenshot artifacts must be supplied through screenshot_paths."
                    ),
                },
                "screenshot_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional browser/page screenshot artifact paths to place into slides.",
                },
            },
            "additionalProperties": False,
        },
        "pdf_create": {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": "Destination .pdf path. If omitted, Nullion creates one in the artifact directory.",
                },
                "image_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Existing content image artifact paths to place into the PDF, one image per page. "
                        "When paired with text_pages, provide text_pages in the same order so each page's text "
                        "matches the corresponding image. Browser screenshot artifacts must be supplied through screenshot_paths."
                    ),
                },
                "screenshot_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Browser/page screenshot artifact paths to place into the PDF, one screenshot per page."
                    ),
                },
                "text_pages": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional report text pages to render into the PDF with extractable text and clickable URL links. "
                        "For reports/tables/cards that include names, prices, citations, listing links, or other "
                        "readable content, put that content here; image_paths alone creates an image-only PDF. "
                        "For multi-item reports with images, prefer one text page per image in matching order. "
                        "The generated PDF uses a report-quality layout profile; do not use browser screenshots as a substitute for readable report content."
                    ),
                },
                "html_pages": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional static HTML page fragments for designed PDFs that need custom visual layout, real tables, "
                        "cards, columns, image grids, or stronger styling than the default report text layout. "
                        "Use semantic HTML/CSS with readable text; scripts are ignored. "
                        "Use text_pages for plain report PDFs and html_pages for intentionally designed PDFs."
                    ),
                },
                "media_alignment": {
                    "type": "string",
                    "enum": ["auto", "align_pages", "preserve_text_pages"],
                    "description": (
                        "Controls how text_pages and media are paired. Use preserve_text_pages by default. "
                        "Use align_pages when each image/screenshot should share a page with its matching text section. "
                        "Use preserve_text_pages when the requested layout intentionally keeps the supplied text pages unchanged, "
                        "including one-page reports that contain multiple images."
                    ),
                },
                "title": {
                    "type": "string",
                    "description": "Optional title used for metadata, visible report headers, and default filename.",
                },
                "page_size": {
                    "type": "string",
                    "enum": ["letter", "a4"],
                    "description": "Optional page size. Defaults to letter.",
                },
            },
            "additionalProperties": False,
        },
        "pdf_edit": {
            "type": "object",
            "properties": {
                "input_path": {"type": "string", "description": "Existing PDF path to edit."},
                "output_path": {
                    "type": "string",
                    "description": "Destination .pdf path. If omitted, Nullion creates one in the artifact directory.",
                },
                "page_numbers": {
                    "type": "array",
                    "items": {"type": "integer"},
                    "description": "Optional 1-based page numbers to keep or reorder.",
                },
                "rotate_degrees": {
                    "type": "integer",
                    "enum": [0, 90, 180, 270],
                    "description": "Optional clockwise rotation applied to kept pages.",
                },
                "append_pdf_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional existing PDFs to append after the kept input pages.",
                },
                "append_image_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional content image paths to append as new PDF pages. Browser screenshot artifacts "
                        "must be supplied through append_screenshot_paths."
                    ),
                },
                "append_screenshot_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional browser/page screenshot artifact paths to append as new PDF pages."
                    ),
                },
                "append_text_pages": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional plain text pages to append.",
                },
                "title": {"type": "string", "description": "Optional title used for metadata and default filename."},
                "page_size": {
                    "type": "string",
                    "enum": ["letter", "a4"],
                    "description": "Optional page size for appended image/text pages. Defaults to letter.",
                },
            },
            "required": ["input_path"],
            "additionalProperties": False,
        },
        "terminal_exec": {
            "type": "object",
            "properties": {
                "command": {"type": "string", "description": "Shell command to execute."},
                "shell": {
                    "type": "string",
                    "enum": ["auto", "sh", "bash", "zsh", "powershell", "pwsh", "cmd"],
                    "description": "Optional shell family for this command. Defaults to the platform shell.",
                },
                "network_mode": {
                    "type": "string",
                    "enum": sorted(_VALID_NETWORK_MODES),
                    "description": "Optional network policy for the command.",
                },
                "timeout_seconds": {
                    "type": "integer",
                    "minimum": 1,
                    "maximum": 300,
                    "description": "Optional execution timeout in seconds.",
                },
            },
            "required": ["command"],
            "additionalProperties": False,
        },
        "web_fetch": {
            "type": "object",
            "properties": {"url": {"type": "string", "description": "HTTP or HTTPS URL to fetch."}},
            "required": ["url"],
            "additionalProperties": False,
        },
        "file_download": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "HTTP or HTTPS URL to download as bytes."},
                "output_path": {
                    "type": "string",
                    "description": "Optional absolute or workspace-relative destination path. Defaults to the workspace artifact directory.",
                },
                "filename": {
                    "type": "string",
                    "description": "Optional safe destination filename when output_path is omitted.",
                },
                "max_bytes": {
                    "type": "integer",
                    "minimum": 1,
                    "description": "Optional maximum response bytes to download.",
                },
                "timeout_seconds": {
                    "type": "number",
                    "minimum": 2,
                    "maximum": 60,
                    "description": "Optional total download timeout in seconds.",
                },
            },
            "required": ["url"],
            "additionalProperties": False,
        },
        "browser_image_collect": {
            "type": "object",
            "properties": {
                "page_url": {
                    "type": "string",
                    "description": "Optional HTTP/HTTPS page URL whose HTML should be inspected for image assets.",
                },
                "html": {
                    "type": "string",
                    "description": "Optional page HTML already obtained from browser/page tooling.",
                },
                "image_urls": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional explicit image asset URLs to fetch and save as local artifact files.",
                },
                "max_images": {
                    "type": "integer",
                    "minimum": 1,
                    "maximum": _BROWSER_IMAGE_COLLECT_MAX_IMAGES,
                    "description": "Maximum number of image files to save. Defaults to 10.",
                },
                "output_stem": {
                    "type": "string",
                    "description": "Optional filename stem for generated local image artifacts.",
                },
                "timeout_seconds": {
                    "type": "number",
                    "minimum": 2,
                    "maximum": 30,
                    "description": "Optional total image collection time budget. Defaults to 12 seconds.",
                },
                "quality_profile": {
                    "type": "string",
                    "enum": ["content", "any"],
                    "description": (
                        "Image validation profile. Use content for document/report media and any only when "
                        "small icons, logos, or other low-detail assets are intentionally requested."
                    ),
                },
            },
            "additionalProperties": False,
        },
        "weather_forecast": {
            "type": "object",
            "properties": {
                "location_text": {
                    "type": "string",
                    "description": "Optional place name, address, city, or postal code to resolve before fetching forecast data.",
                },
                "latitude": {"type": "number", "description": "Optional decimal latitude."},
                "longitude": {"type": "number", "description": "Optional decimal longitude."},
                "forecast_days": {
                    "type": "integer",
                    "minimum": 1,
                    "maximum": 7,
                    "description": "Number of forecast days to return. Defaults to 3.",
                },
                "timezone": {
                    "type": "string",
                    "description": "Optional IANA timezone. Defaults to auto from Open-Meteo.",
                },
                "include_hourly": {
                    "type": "boolean",
                    "description": "Whether to include hour-by-hour forecast rows in addition to current and daily forecast data.",
                },
                "hourly_hours": {
                    "type": "integer",
                    "minimum": 1,
                    "maximum": 168,
                    "description": "Maximum number of hourly forecast rows to return when include_hourly is true.",
                },
            },
            "additionalProperties": False,
        },
        "market_quote": {
            "type": "object",
            "properties": {
                "symbols": {
                    "type": "array",
                    "items": {"type": "string"},
                    "minItems": 1,
                    "maxItems": 16,
                    "description": "Ticker symbols to look up from public market quote data.",
                }
            },
            "required": ["symbols"],
            "additionalProperties": False,
        },
        "connector_request": {
            "type": "object",
            "properties": {
                "provider_id": {
                    "type": "string",
                    "description": "Workspace connection provider id.",
                },
                "url": {
                    "type": "string",
                    "description": (
                        "Full HTTP(S) gateway URL from the installed connector skill, under that "
                        "provider's configured base URL. Do not use generic public web URLs here."
                    ),
                },
                "params": {
                    "type": "object",
                    "description": "Optional query parameters for the connector request.",
                },
                "method": {
                    "type": "string",
                    "enum": ["GET", "HEAD", "POST", "PUT", "PATCH", "DELETE"],
                    "description": "HTTP method. Write methods require the connection permission mode to be read_write.",
                },
                "json": {
                    "type": "object",
                    "description": "Optional JSON object body for POST, PUT, PATCH, or DELETE requests.",
                },
                "body": {
                    "type": "string",
                    "description": "Optional raw request body for POST, PUT, PATCH, or DELETE requests.",
                },
            },
            "required": ["provider_id", "url"],
            "additionalProperties": False,
        },
        "email_send": {
            "type": "object",
            "properties": {
                "to": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Recipient email address(es).",
                },
                "subject": {"type": "string", "description": "Email subject."},
                "body": {
                    "type": "string",
                    "description": "Plain text email body. Provide this as the fallback text when sending HTML.",
                },
                "html_body": {
                    "type": "string",
                    "description": (
                        "Optional HTML email body. Use this when the reviewed draft is HTML or styled email content; "
                        "the send tool will deliver it as a text/html alternative."
                    ),
                },
                "html_path": {
                    "type": "string",
                    "description": "Optional local HTML artifact path to send as the email HTML body.",
                },
                "cc": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional CC email address(es).",
                },
                "bcc": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional BCC email address(es).",
                },
                "attachment_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional local artifact/media paths to attach.",
                },
            },
            "required": ["to", "subject"],
            "additionalProperties": False,
        },
        "skill_pack_read": {
            "type": "object",
            "properties": {
                "pack_id": {"type": "string", "description": "Installed skill pack id, such as owner/pack."},
                "path": {
                    "type": "string",
                    "description": "Relative reference file path shown in the enabled skill pack prompt.",
                },
            },
            "required": ["pack_id", "path"],
            "additionalProperties": False,
        },
        "web_search": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query."},
                "limit": {"type": "integer", "minimum": 1, "description": "Maximum number of results to return."},
            },
            "required": ["query"],
            "additionalProperties": False,
        },
        "file_search": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "Case-insensitive text or filename fragment to search for."},
                "root": {
                    "type": "string",
                    "description": "Optional absolute or workspace-relative folder to search. It must be inside the configured workspace or allowed roots.",
                },
                "search_contents": {
                    "type": "boolean",
                    "description": "Also search readable text file contents. Filename search always runs.",
                },
                "limit": {"type": "integer", "minimum": 1, "description": "Maximum number of matches to return."},
            },
            "required": ["pattern"],
            "additionalProperties": False,
        },
        "file_patch": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute or workspace-relative file path to edit."},
                "old_string": {"type": "string", "description": "Exact text to replace."},
                "new_string": {"type": "string", "description": "Replacement text."},
            },
            "required": ["path", "old_string", "new_string"],
            "additionalProperties": False,
        },
        "workspace_summary": {
            "type": "object",
            "properties": {},
            "additionalProperties": False,
        },
        "email_search": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Email search query."},
                "limit": {"type": "integer", "minimum": 1, "description": "Maximum number of messages to return."},
            },
            "required": ["query"],
            "additionalProperties": False,
        },
        "email_read": {
            "type": "object",
            "properties": {
                "id": {
                    "type": "string",
                    "description": "Exact email message id returned in email_search.results[].id. Do not infer ids or pass placeholders.",
                }
            },
            "required": ["id"],
            "additionalProperties": False,
        },
        "email_attachment_read": _email_attachment_read_input_schema(),
        "calendar_list": _calendar_list_input_schema(),
        "calendar_create": _calendar_create_input_schema(),
        "calendar_update": _calendar_update_input_schema(),
        "calendar_respond": _calendar_respond_input_schema(),
        "calendar_delete": _calendar_delete_input_schema(),
        "browser_navigate": {
            "type": "object",
            "properties": {
                "url": {
                    "type": "string",
                    "description": "HTTP/HTTPS URL, or a local HTML file path/file URL inside this workspace, to open in the browser.",
                }
            },
            "required": ["url"],
            "additionalProperties": False,
        },
        "browser_extract_items": {
            "type": "object",
            "properties": {
                "max_items": {
                    "type": "integer",
                    "minimum": 1,
                    "maximum": 100,
                    "description": "Maximum compact item rows to return. Defaults to 30.",
                },
                "selector": {
                    "type": "string",
                    "description": "Optional CSS selector limiting extraction to a result/list/table region.",
                },
            },
            "additionalProperties": False,
        },
        "browser_extract_detail": {
            "type": "object",
            "properties": {
                "selector": {
                    "type": "string",
                    "description": "Optional CSS selector limiting extraction to the verified detail region.",
                },
            },
            "additionalProperties": False,
        },
        "audio_transcribe": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the audio file."},
                "language": {"type": "string", "description": "Optional language code."},
            },
            "required": ["path"],
            "additionalProperties": False,
        },
        "image_extract_text": {
            "type": "object",
            "properties": {"path": {"type": "string", "description": "Path to the image file."}},
            "required": ["path"],
            "additionalProperties": False,
        },
        "image_generate": {
            "type": "object",
            "properties": {
                "prompt": {"type": "string", "description": "Image generation or edit instruction."},
                "output_path": {"type": "string", "description": "Destination path for the generated image file."},
                "size": {"type": "string", "description": "Optional output size, such as 1024x1024."},
                "source_path": {"type": "string", "description": "Optional source image path for image edits."},
            },
            "required": ["prompt", "output_path"],
            "additionalProperties": False,
        },
    }
    return schemas.get(
        tool_name,
        {
            "type": "object",
            "properties": {},
            "additionalProperties": True,
        },
    )


_MARKET_QUOTE_SYMBOL_RE = re.compile(r"^[A-Z][A-Z0-9.\-]{0,9}$")


def _normalize_market_quote_symbols(value: object) -> tuple[str, ...]:
    if isinstance(value, str):
        raw_values: Iterable[object] = re.split(r"[,\s]+", value)
    elif isinstance(value, Iterable) and not isinstance(value, (bytes, bytearray, dict)):
        raw_values = value
    else:
        raw_values = ()
    symbols: list[str] = []
    for raw in raw_values:
        symbol = str(raw or "").strip().upper()
        if not symbol or not _MARKET_QUOTE_SYMBOL_RE.match(symbol):
            continue
        if symbol not in symbols:
            symbols.append(symbol)
        if len(symbols) >= 16:
            break
    return tuple(symbols)


def _coerce_market_quote_price(value: object) -> float | None:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    if not (0 < number < 1_000_000):
        return None
    return number


def _default_market_quote_fetcher(symbol: str, timeout_seconds: int) -> dict[str, object]:
    encoded_symbol = quote(symbol, safe="")
    url = f"https://query1.finance.yahoo.com/v8/finance/chart/{encoded_symbol}?range=1d&interval=1m"
    request = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 Nullion"})
    with urllib.request.urlopen(request, timeout=timeout_seconds) as response:
        payload = json.loads(response.read(_WEB_FETCH_MAX_BODY_BYTES).decode("utf-8", errors="replace"))
    chart = payload.get("chart") if isinstance(payload, dict) else None
    results = chart.get("result") if isinstance(chart, dict) else None
    first = results[0] if isinstance(results, list) and results and isinstance(results[0], dict) else {}
    meta = first.get("meta") if isinstance(first, dict) else None
    if not isinstance(meta, dict):
        raise ValueError("Market quote response did not include chart metadata.")
    price = _coerce_market_quote_price(meta.get("regularMarketPrice"))
    if price is None:
        raise ValueError("Market quote response did not include a usable regularMarketPrice.")
    regular_market_time = meta.get("regularMarketTime")
    timestamp: str | None = None
    if isinstance(regular_market_time, (int, float)):
        timestamp = datetime.fromtimestamp(float(regular_market_time), tz=UTC).isoformat(timespec="seconds")
    previous_close = _coerce_market_quote_price(
        meta.get("previousClose") or meta.get("chartPreviousClose")
    )
    change = round(price - previous_close, 4) if previous_close is not None else None
    change_percent = (
        round(((price - previous_close) / previous_close) * 100, 4)
        if previous_close is not None
        else None
    )
    indicators = first.get("indicators") if isinstance(first, dict) else None
    quote_rows = indicators.get("quote") if isinstance(indicators, dict) else None
    quote_row = quote_rows[0] if isinstance(quote_rows, list) and quote_rows and isinstance(quote_rows[0], dict) else {}

    def _finite_market_values(key: str) -> list[float]:
        values = quote_row.get(key)
        if not isinstance(values, list):
            return []
        normalized: list[float] = []
        for value in values:
            try:
                number = float(value)
            except (TypeError, ValueError):
                continue
            if number >= 0:
                normalized.append(number)
        return normalized

    highs = [value for value in _finite_market_values("high") if value > 0]
    lows = [value for value in _finite_market_values("low") if value > 0]
    volumes = _finite_market_values("volume")
    return {
        "symbol": str(meta.get("symbol") or symbol).upper(),
        "name": str(meta.get("longName") or meta.get("shortName") or symbol).strip(),
        "price": round(price, 4),
        "previous_close": round(previous_close, 4) if previous_close is not None else None,
        "change": change,
        "change_percent": change_percent,
        "day_high": round(max(highs), 4) if highs else None,
        "day_low": round(min(lows), 4) if lows else None,
        "volume": int(sum(volumes)) if volumes else None,
        "fifty_two_week_high": _coerce_market_quote_price(meta.get("fiftyTwoWeekHigh")),
        "fifty_two_week_low": _coerce_market_quote_price(meta.get("fiftyTwoWeekLow")),
        "currency": str(meta.get("currency") or "").strip() or "USD",
        "exchange": str(meta.get("fullExchangeName") or meta.get("exchangeName") or "").strip(),
        "regular_market_time": timestamp,
        "regular_market_date": timestamp[:10] if timestamp else None,
        "source_url": url,
        "provider": "yahoo-finance-chart",
    }


def _build_market_quote_handler(
    quote_fetcher: Callable[[str, int], dict[str, object]],
) -> ToolHandler:
    def _handler(invocation: ToolInvocation) -> ToolResult:
        symbols = _normalize_market_quote_symbols(invocation.arguments.get("symbols"))
        if not symbols:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {"reason": "missing_symbols"},
                "market_quote requires one or more valid ticker symbols.",
            )
        quotes: list[dict[str, object]] = []
        errors: list[dict[str, str]] = []
        timeout_seconds = min(12, max(3, int(invocation.arguments.get("timeout_seconds") or 8)))

        def _fetch(symbol: str) -> tuple[str, dict[str, object] | None, str | None]:
            try:
                return symbol, quote_fetcher(symbol, timeout_seconds), None
            except Exception as exc:
                return symbol, None, str(exc)

        # The tool contract accepts a bounded batch of explicit symbols. Fetch
        # that batch concurrently so a healthy 16-symbol request cannot exceed
        # the handler deadline merely because each independent HTTP request was
        # performed serially. ``executor.map`` preserves the requested order.
        with ThreadPoolExecutor(max_workers=len(symbols), thread_name_prefix="nullion-market-quote") as executor:
            outcomes = list(executor.map(_fetch, symbols))
        for symbol, quote_row, error in outcomes:
            if quote_row is not None:
                quotes.append(quote_row)
            else:
                errors.append({"symbol": symbol, "error": error or "quote unavailable"})
        return ToolResult(
            invocation.invocation_id,
            invocation.tool_name,
            "completed" if quotes else "failed",
            {
                "provider": "yahoo-finance-chart",
                "symbols": list(symbols),
                "quotes": quotes,
                "errors": errors,
            },
            None if quotes else "No requested market quotes could be fetched.",
        )

    return _handler


@dataclass(frozen=True, slots=True)
class TerminalExecutionPolicy:
    network_mode: str | None
    approved_targets: tuple[str, ...]


@dataclass(frozen=True, slots=True)
class TerminalExecutionResult:
    exit_code: int
    stdout: str
    stderr: str
    shell: str = "unknown"
    argv: tuple[str, ...] = ()


@dataclass(frozen=True, slots=True)
class TerminalShellInvocation:
    shell: str
    argv: tuple[str, ...]


_FIND_PATH_OPTION_ARITY = {
    "-H": 0,
    "-L": 0,
    "-P": 0,
    "-O0": 0,
    "-O1": 0,
    "-O2": 0,
    "-O3": 0,
    "-D": 1,
}
_FIND_EXPRESSION_STARTERS = {
    "(",
    "!",
    "-",
}


def _terminal_filesystem_safety_denial(
    command: str,
    *,
    allowed_roots: Iterable[Path] = (),
    principal_id: str | None = None,
) -> dict[str, object] | None:
    """Reject shell commands likely to trigger broad local data traversal."""
    try:
        lexer = shlex.shlex(command, posix=True, punctuation_chars=True)
        lexer.whitespace_split = True
        tokens = list(lexer)
    except ValueError:
        tokens = shlex.split(command, posix=True)
    except Exception:
        return None
    if not tokens:
        return None

    home = str(Path.home())
    resolved_allowed_roots = tuple(Path(root).expanduser().resolve() for root in allowed_roots)
    home_variants = {home, "$HOME", "${HOME}", "~"}
    broad_root_variants = {"/", "/Users"}
    protected_home_roots = {
        str(Path(home) / "Library"),
        str(Path(home) / "Library" / "Application Support"),
        str(Path(home) / "Library" / "Containers"),
        str(Path(home) / "Library" / "Group Containers"),
    }
    unknown_workspace_denial = _unknown_workspace_storage_denial(tokens, principal_id=principal_id)
    if unknown_workspace_denial is not None:
        return unknown_workspace_denial

    for index, token in enumerate(tokens):
        if token != "find":
            continue
        path_tokens = _find_path_tokens(tokens[index + 1 :])
        if not path_tokens:
            path_tokens = ["."]
        maxdepth = _find_maxdepth(tokens[index + 1 :])
        for raw_path in path_tokens:
            normalized_path = _normalize_find_root(raw_path, home=home)
            if normalized_path in broad_root_variants:
                return {
                    "reason": "filesystem_traversal_denied",
                    "command_family": "find",
                    "path": raw_path,
                    "message": "Refusing to run a broad filesystem search from the system or users root.",
                }
            if _path_is_allowed_by_config(normalized_path, resolved_allowed_roots):
                continue
            if normalized_path in home_variants and maxdepth != 1:
                return {
                    "reason": "filesystem_traversal_denied",
                    "command_family": "find",
                    "path": raw_path,
                    "message": (
                        "Refusing to recursively search the home folder because it can cross "
                        "macOS-protected app data. Search a specific subfolder instead."
                    ),
                }
            if any(_path_is_at_or_under(normalized_path, root) for root in protected_home_roots) and maxdepth != 1:
                return {
                    "reason": "filesystem_traversal_denied",
                    "command_family": "find",
                    "path": raw_path,
                    "message": (
                        "Refusing to recursively search protected home Library data. "
                        "Use a specific file path or a narrow app folder with explicit approval."
                    ),
                }
    workspace_denial = _terminal_unknown_workspace_denial(tokens, resolved_allowed_roots, principal_id=principal_id)
    if workspace_denial is not None:
        return workspace_denial
    return None


_TERMINAL_MUTATION_BOUNDARY_OPERATIONS = frozenset({"write", "delete", "metadata_write"})
_SCHEDULED_TASK_HELPER_SCRIPT_EXTENSIONS = frozenset(
    {
        ".bat",
        ".bash",
        ".cjs",
        ".cmd",
        ".go",
        ".java",
        ".jl",
        ".js",
        ".kt",
        ".kts",
        ".lua",
        ".mjs",
        ".php",
        ".pl",
        ".ps1",
        ".py",
        ".pyw",
        ".r",
        ".rb",
        ".rs",
        ".scala",
        ".sh",
        ".swift",
        ".ts",
        ".tsx",
        ".zsh",
    }
)


def _workspace_scratch_root_for_principal(principal_id: str | None) -> Path | None:
    try:
        from nullion.workspace_storage import workspace_storage_roots_for_principal

        return workspace_storage_roots_for_principal(principal_id).scratch.resolve()
    except Exception:
        return None


def _workspace_artifact_root_for_principal(principal_id: str | None) -> Path | None:
    try:
        from nullion.workspace_storage import workspace_storage_roots_for_principal

        return workspace_storage_roots_for_principal(principal_id).artifacts.resolve()
    except Exception:
        return None


def _runtime_artifact_root() -> Path | None:
    try:
        from nullion.artifacts import nullion_data_home

        return (nullion_data_home() / "artifacts").resolve()
    except Exception:
        return None


def _path_within_root_text(path_text: object, root: Path | None) -> bool:
    if root is None:
        return False
    try:
        path = Path(str(path_text or "")).expanduser().resolve()
        path.relative_to(root)
        return True
    except (OSError, ValueError):
        return False


def _scheduled_task_helper_script_path(path_text: object) -> bool:
    try:
        return Path(str(path_text or "").strip().strip("'\"")).suffix.lower() in _SCHEDULED_TASK_HELPER_SCRIPT_EXTENSIONS
    except Exception:
        return False


def _terminal_filesystem_mutation_approval_denial(
    invocation: ToolInvocation,
    *,
    execution_cwd: Path | None,
    allowed_roots: Iterable[Path] = (),
) -> dict[str, object] | None:
    """Require a filesystem boundary grant before local shell writes/deletes."""
    resolved_allowed_roots = tuple(Path(root).expanduser().resolve() for root in allowed_roots)
    for fact in extract_boundary_facts(invocation):
        if fact.kind is not BoundaryKind.FILESYSTEM_ACCESS:
            continue
        if fact.operation not in _TERMINAL_MUTATION_BOUNDARY_OPERATIONS:
            continue
        try:
            raw_target = Path(str(fact.target)).expanduser()
            target_path = raw_target if raw_target.is_absolute() else (execution_cwd or Path.cwd()) / raw_target
            resolved_target = target_path.resolve()
        except OSError:
            resolved_target = Path(str(fact.target)).expanduser()
        if _path_within_any_root(resolved_target, resolved_allowed_roots):
            continue
        if _is_approved_filesystem_path(resolved_target, invocation.trusted_filesystem_selectors):
            continue
        return {
            "reason": "filesystem_mutation_approval_required",
            "boundary_kind": BoundaryKind.FILESYSTEM_ACCESS.value,
            "operation": fact.operation,
            "target": str(resolved_target),
            "command_family": fact.attributes.get("command_family"),
            "requires_approval": True,
            "message": (
                "Refusing to run a local shell command that would modify files outside the trusted workspace "
                "without an explicit filesystem approval."
            ),
        }
    return None


def _terminal_resolved_filesystem_target(fact, *, execution_cwd: Path | None) -> str:
    try:
        raw_target = Path(str(fact.target)).expanduser()
        target_path = raw_target if raw_target.is_absolute() else (execution_cwd or Path.cwd()) / raw_target
        return str(target_path.resolve())
    except OSError:
        return str(Path(str(fact.target)).expanduser())


def _terminal_resolved_target_text(raw_target: object, *, execution_cwd: Path | None) -> str:
    try:
        raw_text = str(raw_target or "").strip().strip("'\"")
        raw_path = Path(raw_text).expanduser()
        target_path = raw_path if raw_path.is_absolute() else (execution_cwd or Path.cwd()) / raw_path
        return str(target_path.resolve())
    except OSError:
        return str(Path(str(raw_target or "")).expanduser())


def _terminal_shell_tokens(command: str) -> list[str]:
    try:
        lexer = shlex.shlex(command, posix=True, punctuation_chars=True)
        lexer.whitespace_split = True
        return list(lexer)
    except Exception:
        try:
            return shlex.split(command, posix=True)
        except Exception:
            return str(command or "").split()


def _terminal_find_delete_targets(command: str, *, execution_cwd: Path | None) -> tuple[dict[str, object], ...]:
    tokens = _terminal_shell_tokens(command)
    targets: list[dict[str, object]] = []
    control_tokens = {";", "&&", "||", "|"}
    expression_starters = {"(", "!", "-name", "-iname", "-path", "-ipath", "-type", "-mtime", "-mmin", "-maxdepth", "-mindepth"}
    for index, token in enumerate(tokens):
        if Path(str(token).strip("'\"")).name.lower() != "find":
            continue
        segment: list[str] = []
        for candidate in tokens[index + 1 :]:
            if candidate in control_tokens:
                break
            segment.append(candidate)
        if "-delete" not in segment and not any(
            item == "-exec"
            and any(Path(str(part).strip("'\"")).name.lower() in {"rm", "rm.exe", "rmdir", "rmdir.exe"} for part in segment[pos + 1 :])
            for pos, item in enumerate(segment)
        ):
            continue
        raw_roots: list[str] = []
        for item in segment:
            if item in expression_starters or item.startswith("-"):
                break
            raw_roots.append(item)
        if not raw_roots:
            raw_roots = ["."]
        for raw_root in raw_roots:
            targets.append({
                "target": _terminal_resolved_target_text(raw_root, execution_cwd=execution_cwd),
                "operation": "delete",
                "command_family": "find",
            })
    return tuple(targets)


def _terminal_script_delete_targets(command: str, *, execution_cwd: Path | None) -> tuple[dict[str, object], ...]:
    lowered = str(command or "").lower()
    if not any(marker in lowered for marker in _TERMINAL_DESTRUCTIVE_SCRIPT_MARKERS):
        return ()
    token_families = {
        Path(str(token).strip("'\"")).name.lower()
        for token in _terminal_shell_tokens(command)[:6]
        if str(token).strip()
    }
    if not token_families.intersection(_TERMINAL_DESTRUCTIVE_SCRIPT_COMMANDS):
        return ()
    raw_paths: list[str] = []
    for match in _TERMINAL_PATH_LITERAL_RE.finditer(str(command or "")):
        raw_path = match.group("path").strip()
        if raw_path and raw_path not in raw_paths:
            raw_paths.append(raw_path)
    if not raw_paths and execution_cwd is not None:
        raw_paths = [str(execution_cwd)]
    return tuple(
        {
            "target": _terminal_resolved_target_text(raw_path, execution_cwd=execution_cwd),
            "operation": "delete",
            "command_family": "script",
        }
        for raw_path in raw_paths
    )


def _terminal_supplemental_delete_targets(
    command: str,
    *,
    execution_cwd: Path | None,
) -> tuple[dict[str, object], ...]:
    return (
        *_terminal_find_delete_targets(command, execution_cwd=execution_cwd),
        *_terminal_script_delete_targets(command, execution_cwd=execution_cwd),
    )


def _terminal_delete_targets_for_invocation(
    invocation: ToolInvocation,
    *,
    execution_cwd: Path | None,
) -> tuple[dict[str, object], ...]:
    targets: list[dict[str, object]] = []
    seen: set[str] = set()
    for fact in extract_boundary_facts(invocation):
        if fact.kind is not BoundaryKind.FILESYSTEM_ACCESS:
            continue
        if fact.operation != "delete":
            continue
        target = _terminal_resolved_filesystem_target(fact, execution_cwd=execution_cwd)
        if target in seen:
            continue
        seen.add(target)
        targets.append({
            "target": target,
            "operation": fact.operation,
            "command_family": fact.attributes.get("command_family"),
        })
    raw_command = invocation.arguments.get("command")
    if isinstance(raw_command, str) and raw_command.strip():
        for target in _terminal_supplemental_delete_targets(raw_command, execution_cwd=execution_cwd):
            target_path = str(target.get("target") or "")
            if not target_path or target_path in seen:
                continue
            seen.add(target_path)
            targets.append(dict(target))
    return tuple(targets)


def _terminal_execution_cwd_from_registry(registry: object) -> Path | None:
    roots = ()
    try:
        roots = tuple(registry.filesystem_allowed_roots())
    except Exception:
        roots = ()
    for root in roots:
        try:
            resolved = Path(root).expanduser().resolve()
        except (OSError, RuntimeError, ValueError):
            continue
        if resolved.is_dir():
            return resolved
    return None


def _terminal_destructive_permission_token(command: str, targets: Iterable[dict[str, object]]) -> str:
    payload = {
        "command": str(command or ""),
        "targets": [
            {
                "target": str(target.get("target") or ""),
                "operation": str(target.get("operation") or ""),
                "command_family": str(target.get("command_family") or ""),
            }
            for target in targets
        ],
    }
    encoded = json.dumps(payload, sort_keys=True, separators=(",", ":")).encode("utf-8")
    return hashlib.sha256(encoded).hexdigest()[:24]


def _terminal_destructive_permission(command: str, targets: Iterable[dict[str, object]]) -> str:
    return f"{TERMINAL_DESTRUCTIVE_PERMISSION_PREFIX}{_terminal_destructive_permission_token(command, targets)}"


def _terminal_destructive_approval_context(
    invocation: ToolInvocation,
    *,
    command: str,
    targets: tuple[dict[str, object], ...],
) -> dict[str, object]:
    token = _terminal_destructive_permission_token(command, targets)
    preview_targets = [str(target.get("target") or "") for target in targets[:_TERMINAL_DESTRUCTIVE_PREVIEW_LIMIT]]
    return {
        "tool_name": invocation.tool_name,
        "tool_description": "Run a local terminal command that deletes files.",
        "tool_risk_level": "high",
        "tool_side_effect_class": "delete",
        "tool_permission_scope": "exact_invocation",
        "tool_arguments": redact_value(dict(invocation.arguments or {})),
        "command": command,
        "operation": "delete",
        "target": f"{len(targets)} path{'s' if len(targets) != 1 else ''}",
        "destructive_targets": [dict(target) for target in targets],
        "destructive_preview_targets": preview_targets,
        "destructive_target_count": len(targets),
        "destructive_permission_token": token,
    }


def _terminal_output_artifact_root(
    *,
    execution_cwd: Path | None,
    allowed_roots: Iterable[Path],
    principal_id: str | None,
) -> Path | None:
    artifact_root = _workspace_artifact_root_for_principal(principal_id)
    if artifact_root is not None:
        return artifact_root
    if execution_cwd is not None:
        return execution_cwd / ".nullion-artifacts"
    for root in allowed_roots:
        try:
            resolved = Path(root).expanduser().resolve()
        except (OSError, RuntimeError, ValueError):
            continue
        return resolved / ".nullion-artifacts"
    return None


def _terminal_output_text_for_attachment(output: dict[str, object]) -> str:
    stdout = str(output.get("stdout") or "")
    stderr = str(output.get("stderr") or "")
    exit_code = output.get("exit_code", output.get("returncode"))
    parts: list[str] = []
    if stdout:
        parts.append("STDOUT\n" + stdout.rstrip())
    if stderr:
        parts.append("STDERR\n" + stderr.rstrip())
    if exit_code is not None:
        parts.append(f"EXIT CODE\n{exit_code}")
    return "\n\n".join(parts).strip()


def _materialize_terminal_output_attachment(
    output: dict[str, object],
    *,
    invocation: ToolInvocation,
    execution_cwd: Path | None,
    allowed_roots: Iterable[Path],
) -> str | None:
    text = _terminal_output_text_for_attachment(output)
    if len(text) <= _TERMINAL_OUTPUT_ATTACHMENT_THRESHOLD_CHARS:
        return None
    root = _terminal_output_artifact_root(
        execution_cwd=execution_cwd,
        allowed_roots=allowed_roots,
        principal_id=invocation.principal_id,
    )
    if root is None:
        return None
    try:
        root.mkdir(parents=True, exist_ok=True)
        path = root / f"terminal-output-{invocation.invocation_id[:12]}.txt"
        path.write_text(text + "\n", encoding="utf-8")
        return str(path.resolve())
    except OSError:
        logger.debug("Could not materialize terminal output attachment", exc_info=True)
        return None


def _unknown_workspace_storage_denial(tokens: list[str], *, principal_id: str | None = None) -> dict[str, object] | None:
    try:
        from nullion.workspace_storage import sanitize_workspace_id, workspace_storage_base

        storage_base = workspace_storage_base()
    except Exception:
        return None
    allowed_workspaces = {
        sanitize_workspace_id(workspace_id)
        for workspace_id in _registered_or_existing_workspace_ids(storage_base)
    }
    for token in tokens:
        if not token or token.startswith("-") or token in {";", "&&", "||", "|"}:
            continue
        try:
            candidate = Path(os.path.expandvars(token)).expanduser().resolve(strict=False)
        except Exception:
            continue
        try:
            relative = candidate.relative_to(storage_base)
        except ValueError:
            continue
        if not relative.parts:
            continue
        workspace_id = sanitize_workspace_id(relative.parts[0])
        if not _workspace_access_allowed(workspace_id, principal_id=principal_id, allowed_workspaces=allowed_workspaces):
            return {
                "reason": "unknown_workspace_denied",
                "workspace_id": workspace_id,
                "path": str(candidate),
                "message": f"Refusing to create or modify unknown workspace storage: {workspace_id}.",
            }
    return None


def _find_path_tokens(tokens: list[str]) -> list[str]:
    paths: list[str] = []
    index = 0
    while index < len(tokens):
        token = tokens[index]
        if token in {";", "&&", "||", "|"}:
            break
        if token in _FIND_PATH_OPTION_ARITY:
            index += 1 + _FIND_PATH_OPTION_ARITY[token]
            continue
        if any(token.startswith(prefix) for prefix in _FIND_EXPRESSION_STARTERS):
            break
        paths.append(token)
        index += 1
    return paths


def _find_maxdepth(tokens: list[str]) -> int | None:
    for index, token in enumerate(tokens):
        if token == "-maxdepth" and index + 1 < len(tokens):
            try:
                return int(tokens[index + 1])
            except ValueError:
                return None
    return None


def _normalize_find_root(raw_path: str, *, home: str) -> str:
    if raw_path in {"$HOME", "${HOME}", "~"}:
        return raw_path
    try:
        expanded = Path(os.path.expandvars(raw_path)).expanduser()
        return str(expanded.resolve(strict=False))
    except Exception:
        return raw_path


def _terminal_path_candidates(token: str) -> tuple[Path, ...]:
    stripped = str(token or "").strip().strip("'\"`")
    if not stripped:
        return ()
    candidates: list[str] = [stripped]
    for separator in ("=", ":"):
        if separator in stripped:
            tail = stripped.split(separator, 1)[1].strip()
            if tail:
                candidates.append(tail)
    paths: list[Path] = []
    for candidate in candidates:
        cleaned = candidate.strip("'\"`")
        if not cleaned:
            continue
        if cleaned.startswith(("~", "/", "$HOME", "${HOME}")) or "/.nullion/workspaces/" in cleaned:
            try:
                paths.append(Path(os.path.expandvars(cleaned)).expanduser().resolve(strict=False))
            except Exception:
                continue
    return tuple(dict.fromkeys(paths))


def _registered_or_existing_workspace_ids(base: Path) -> set[str]:
    ids: set[str] = set()
    try:
        from nullion.users import registered_workspace_ids

        ids.update(registered_workspace_ids())
    except Exception:
        ids.add("workspace_admin")
    try:
        if base.is_dir():
            ids.update(path.name for path in base.iterdir() if path.is_dir())
    except OSError:
        pass
    ids.add("workspace_admin")
    return ids


def _workspace_access_allowed(
    workspace_id: str,
    *,
    principal_id: str | None,
    allowed_workspaces: set[str],
) -> bool:
    try:
        from nullion.connections import workspace_id_for_principal
        from nullion.workspace_storage import sanitize_workspace_id

        principal_workspace = sanitize_workspace_id(workspace_id_for_principal(principal_id))
        target_workspace = sanitize_workspace_id(workspace_id)
    except Exception:
        principal_workspace = "workspace_admin"
        target_workspace = str(workspace_id or "workspace_admin")
    if principal_workspace == target_workspace:
        return True
    if principal_workspace == "workspace_admin" and target_workspace in allowed_workspaces:
        return True
    return False


def _terminal_unknown_workspace_denial(
    tokens: list[str],
    allowed_roots: Iterable[Path],
    *,
    principal_id: str | None = None,
) -> dict[str, object] | None:
    _ = allowed_roots
    try:
        from nullion.workspace_storage import workspace_storage_base

        base = workspace_storage_base()
    except Exception:
        return None
    allowed = _registered_or_existing_workspace_ids(base)
    for token in tokens:
        for path in _terminal_path_candidates(token):
            try:
                relative = path.relative_to(base)
            except ValueError:
                continue
            if not relative.parts:
                continue
            workspace_id = relative.parts[0]
            if _workspace_access_allowed(workspace_id, principal_id=principal_id, allowed_workspaces=allowed):
                continue
            return {
                "reason": "unknown_workspace_denied",
                "workspace_id": workspace_id,
                "path": str(path),
                "message": (
                    "Refusing to use or create an unregistered workspace directory. "
                    "Verify the person/workspace exists in the Users settings first."
                ),
            }
    return None


def _path_is_at_or_under(path: str, root: str) -> bool:
    try:
        normalized_path = Path(path)
        normalized_root = Path(root)
        return normalized_path == normalized_root or normalized_root in normalized_path.parents
    except Exception:
        return False


def _path_is_allowed_by_config(path: str, allowed_roots: Iterable[Path]) -> bool:
    if path in {"$HOME", "${HOME}", "~"}:
        path = str(Path.home())
    try:
        resolved_path = Path(os.path.expandvars(path)).expanduser().resolve(strict=False)
    except Exception:
        return False
    return any(_is_within_allowed_root(resolved_path, root) for root in allowed_roots)


@dataclass(frozen=True, slots=True)
class TerminalAttestationEvidence:
    format: str
    payload: dict[str, object]
    metadata: dict[str, object] = field(default_factory=dict)


@dataclass(frozen=True, slots=True)
class TerminalAttestationVerificationResult:
    is_valid: bool
    failure_reason: str | None = None
    metadata: dict[str, object] = field(default_factory=dict)


@dataclass(frozen=True, slots=True)
class TerminalBackendDescriptor:
    mode: str
    attested_capabilities: tuple[str, ...]
    metadata: dict[str, object]
    evidence: TerminalAttestationEvidence | None = None


class TerminalAttestationVerifier(Protocol):
    def verify(
        self,
        descriptor: TerminalBackendDescriptor,
        *,
        required_capabilities: Iterable[str],
    ) -> TerminalAttestationVerificationResult: ...


class TerminalExecutorBackend(Protocol):
    def describe(self) -> TerminalBackendDescriptor: ...

    def run(
        self,
        command: str,
        *,
        cwd: str | None,
        timeout: int,
        policy: TerminalExecutionPolicy,
        shell: str | None = None,
    ) -> TerminalExecutionResult: ...


_POSIX_SHELL_CHOICES = {
    "sh": ("sh", "-lc"),
    "bash": ("bash", "-lc"),
    "zsh": ("zsh", "-lc"),
}
_WINDOWS_SHELL_CHOICES = {
    "pwsh": ("pwsh", "-NoLogo", "-NoProfile", "-NonInteractive", "-ExecutionPolicy", "Bypass", "-Command"),
    "powershell": ("powershell.exe", "-NoLogo", "-NoProfile", "-NonInteractive", "-ExecutionPolicy", "Bypass", "-Command"),
    "cmd": ("cmd.exe", "/d", "/s", "/c"),
}


def _normalize_terminal_shell(raw_shell: object) -> str | None:
    if not isinstance(raw_shell, str):
        return None
    shell = raw_shell.strip().lower()
    return shell or None


def _terminal_shell_from_env() -> str | None:
    raw_shell = os.environ.get("NULLION_TERMINAL_SHELL")
    if not raw_shell:
        return None
    shell = raw_shell.strip()
    if not shell:
        return None
    shell_name = Path(shell).name.lower()
    if shell_name in {"pwsh", "pwsh.exe"}:
        return "pwsh"
    if shell_name in {"powershell", "powershell.exe"}:
        return "powershell"
    if shell_name in {"cmd", "cmd.exe"}:
        return "cmd"
    if shell_name in _POSIX_SHELL_CHOICES:
        return shell_name
    return shell


def _which_shell(executable: str) -> str | None:
    if Path(executable).parent != Path():
        return executable if Path(executable).exists() else None
    return shutil.which(executable)


def _resolve_posix_terminal_shell(command: str, requested_shell: str | None) -> TerminalShellInvocation:
    shell = requested_shell if requested_shell in _POSIX_SHELL_CHOICES else None
    if shell is None and requested_shell in _WINDOWS_SHELL_CHOICES:
        shell = "sh"
    if shell is None:
        env_shell = _terminal_shell_from_env()
        if env_shell and env_shell not in _POSIX_SHELL_CHOICES:
            resolved_env_shell = _which_shell(env_shell)
            if resolved_env_shell is not None:
                return TerminalShellInvocation(
                    shell=Path(resolved_env_shell).name.lower(),
                    argv=(resolved_env_shell, "-lc", command),
                )
        shell = env_shell if env_shell in _POSIX_SHELL_CHOICES else "sh"
    executable, flag = _POSIX_SHELL_CHOICES.get(shell, ("sh", "-lc"))
    resolved = _which_shell(executable) or "/bin/sh"
    return TerminalShellInvocation(shell=shell, argv=(resolved, flag, command))


def _resolve_windows_terminal_shell(command: str, requested_shell: str | None) -> TerminalShellInvocation:
    candidates: list[str] = []
    if requested_shell and requested_shell != "auto":
        candidates.append(requested_shell)
    env_shell = _terminal_shell_from_env()
    if env_shell:
        candidates.append(env_shell)
    candidates.extend(("pwsh", "powershell", "cmd"))

    seen: set[str] = set()
    for shell in candidates:
        if shell in seen:
            continue
        seen.add(shell)
        template = _WINDOWS_SHELL_CHOICES.get(shell)
        if template is None:
            continue
        executable = _which_shell(template[0])
        if executable is None:
            continue
        return TerminalShellInvocation(shell=shell, argv=(executable, *template[1:], command))

    comspec = os.environ.get("COMSPEC")
    cmd = comspec if comspec else "cmd.exe"
    return TerminalShellInvocation(shell="cmd", argv=(cmd, "/d", "/s", "/c", command))


def _resolve_terminal_shell_invocation(command: str, *, shell: str | None = None) -> TerminalShellInvocation:
    requested_shell = _normalize_terminal_shell(shell) or "auto"
    if os.name == "nt":
        return _resolve_windows_terminal_shell(command, requested_shell)
    return _resolve_posix_terminal_shell(command, requested_shell)


def _terminal_timeout_seconds(raw_timeout: object, *, default: int = 20, maximum: int = 300) -> int:
    if isinstance(raw_timeout, bool):
        return default
    try:
        timeout = int(raw_timeout) if raw_timeout is not None else default
    except (TypeError, ValueError):
        return default
    return max(1, min(maximum, timeout))


class SubprocessTerminalExecutorBackend:
    _INTERRUPTED_EXIT_CODE = 130
    _INTERRUPTED_MESSAGE = "Interrupted by user."

    def describe(self) -> TerminalBackendDescriptor:
        return TerminalBackendDescriptor(
            mode="subprocess",
            attested_capabilities=(),
            metadata={},
        )

    def run(
        self,
        command: str,
        *,
        cwd: str | None,
        timeout: int,
        policy: TerminalExecutionPolicy,
        shell: str | None = None,
    ) -> TerminalExecutionResult:
        del policy
        invocation = _resolve_terminal_shell_invocation(command, shell=shell)
        try:
            completed = subprocess.run(
                list(invocation.argv),
                shell=False,
                capture_output=True,
                text=True,
                timeout=timeout,
                cwd=cwd,
                check=False,
            )
        except KeyboardInterrupt:
            return TerminalExecutionResult(
                exit_code=self._INTERRUPTED_EXIT_CODE,
                stdout="",
                stderr=self._INTERRUPTED_MESSAGE,
                shell=invocation.shell,
                argv=invocation.argv,
            )
        except OSError as exc:
            return TerminalExecutionResult(
                exit_code=127,
                stdout="",
                stderr=f"Shell startup failed: {exc}",
                shell=invocation.shell,
                argv=invocation.argv,
            )
        return TerminalExecutionResult(
            exit_code=completed.returncode,
            stdout=completed.stdout,
            stderr=completed.stderr,
            shell=invocation.shell,
            argv=invocation.argv,
        )


@dataclass(frozen=True, slots=True)
class SandboxLauncherInvocation:
    argv: tuple[str, ...]
    env: dict[str, str]


class SandboxLauncherTerminalExecutorBackend:
    _INTERRUPTED_EXIT_CODE = 130
    _INTERRUPTED_MESSAGE = "Interrupted by user."

    def __init__(
        self,
        *,
        launcher_command: str,
        launcher_args: Iterable[str] = (),
    ) -> None:
        self._launcher_command = launcher_command
        self._launcher_args = tuple(arg for arg in launcher_args if isinstance(arg, str) and arg)
        self._launcher_argv = (launcher_command, *self._launcher_args)
        self._descriptor_cache: TerminalBackendDescriptor | None = None

    def _legacy_descriptor(self) -> TerminalBackendDescriptor:
        return TerminalBackendDescriptor(
            mode="launcher",
            attested_capabilities=(
                "network_policy_enforced",
                "network_policy_enforced.none",
                "network_policy_enforced.localhost_only",
                "network_policy_enforced.approved_only",
                "approved_only_enforced_via_local_allowlist_proxy",
            ),
            metadata={
                "launcher_command": self._launcher_command,
                "launcher_args": self._launcher_args,
                "supported_network_modes": ("none", "localhost_only", "approved_only"),
                "descriptor_source": _LAUNCHER_DESCRIPTOR_SOURCE_LEGACY,
                "descriptor_schema": _LAUNCHER_DESCRIPTOR_SCHEMA,
            },
            evidence=TerminalAttestationEvidence(
                format="nullion.launcher.v2",
                payload={
                    "schema": _LAUNCHER_DESCRIPTOR_SCHEMA,
                    "mode": "launcher",
                    "launcher_command": self._launcher_command,
                    "launcher_args": list(self._launcher_args),
                    "supported_network_modes": ["none", "localhost_only", "approved_only"],
                    "enforcement_by_mode": {
                        "none": "sandbox-exec deny network*",
                        "localhost_only": "sandbox-exec remote ip localhost:* only",
                        "approved_only": "sandbox-exec localhost proxy only + in-process allowlist proxy",
                    },
                    "requires_sandbox_exec_for_restrictive_modes": True,
                    "proxy_bind_scope": "127.0.0.1",
                },
                metadata={
                    "descriptor_source": _LAUNCHER_DESCRIPTOR_SOURCE_LEGACY,
                    "descriptor_schema": _LAUNCHER_DESCRIPTOR_SCHEMA,
                },
            ),
        )

    def _probe_launcher_descriptor(self) -> TerminalBackendDescriptor | None:
        command_path = Path(self._launcher_command)
        if command_path.parent == Path() and shutil.which(self._launcher_command) is None:
            return None
        try:
            completed = subprocess.run(
                [*self._launcher_argv, "--describe"],
                shell=False,
                capture_output=True,
                text=True,
                timeout=5,
                check=False,
            )
        except OSError:
            return None
        if completed.returncode != 0:
            return None
        try:
            payload = json.loads(completed.stdout)
        except json.JSONDecodeError:
            return None
        if not isinstance(payload, dict):
            return None
        if payload.get("schema") != _LAUNCHER_DESCRIPTOR_SCHEMA:
            return None
        if payload.get("mode") != "launcher":
            return None
        supported_network_modes = tuple(
            mode for mode in payload.get("supported_network_modes", ()) if isinstance(mode, str) and mode
        )
        attested_capabilities = tuple(
            capability
            for capability in payload.get("attested_capabilities", ())
            if isinstance(capability, str) and capability
        )
        enforcement_by_mode_raw = payload.get("enforcement_by_mode")
        enforcement_by_mode = {
            mode: enforcement
            for mode, enforcement in enforcement_by_mode_raw.items()
            if isinstance(mode, str) and isinstance(enforcement, str)
        } if isinstance(enforcement_by_mode_raw, dict) else {}
        proxy_bind_scope = payload.get("proxy_bind_scope")
        return TerminalBackendDescriptor(
            mode="launcher",
            attested_capabilities=attested_capabilities,
            metadata={
                "launcher_command": self._launcher_command,
                "launcher_args": self._launcher_args,
                "supported_network_modes": supported_network_modes,
                "descriptor_source": _LAUNCHER_DESCRIPTOR_SOURCE_PROBE,
                "descriptor_schema": _LAUNCHER_DESCRIPTOR_SCHEMA,
            },
            evidence=TerminalAttestationEvidence(
                format="nullion.launcher.v2",
                payload={
                    "schema": _LAUNCHER_DESCRIPTOR_SCHEMA,
                    "mode": "launcher",
                    "launcher_command": self._launcher_command,
                    "launcher_args": list(self._launcher_args),
                    "supported_network_modes": list(supported_network_modes),
                    "enforcement_by_mode": enforcement_by_mode,
                    "requires_sandbox_exec_for_restrictive_modes": payload.get(
                        "requires_sandbox_exec_for_restrictive_modes"
                    ),
                    "proxy_bind_scope": proxy_bind_scope,
                },
                metadata={
                    "descriptor_source": _LAUNCHER_DESCRIPTOR_SOURCE_PROBE,
                    "descriptor_schema": _LAUNCHER_DESCRIPTOR_SCHEMA,
                },
            ),
        )

    def describe(self) -> TerminalBackendDescriptor:
        if self._descriptor_cache is None:
            self._descriptor_cache = self._probe_launcher_descriptor() or self._legacy_descriptor()
        return self._descriptor_cache

    def _build_invocation(self, command: str, *, policy: TerminalExecutionPolicy) -> SandboxLauncherInvocation:
        network_mode = policy.network_mode or _NETWORK_MODE_FULL
        argv = [*self._launcher_argv, "--network-mode", network_mode]
        for target in policy.approved_targets:
            argv.extend(("--approved-target", target))
        argv.extend(("--", command))
        env = {
            "NULLION_TERMINAL_NETWORK_MODE": network_mode,
            "NULLION_TERMINAL_APPROVED_TARGETS": ",".join(policy.approved_targets),
        }
        return SandboxLauncherInvocation(argv=tuple(argv), env=env)

    def run(
        self,
        command: str,
        *,
        cwd: str | None,
        timeout: int,
        policy: TerminalExecutionPolicy,
        shell: str | None = None,
    ) -> TerminalExecutionResult:
        del shell
        invocation = self._build_invocation(command, policy=policy)
        env = {**os.environ, **invocation.env}
        try:
            completed = subprocess.run(
                list(invocation.argv),
                shell=False,
                capture_output=True,
                text=True,
                timeout=timeout,
                cwd=cwd,
                check=False,
                env=env,
            )
        except KeyboardInterrupt:
            return TerminalExecutionResult(
                exit_code=self._INTERRUPTED_EXIT_CODE,
                stdout="",
                stderr=self._INTERRUPTED_MESSAGE,
                shell="launcher",
                argv=invocation.argv,
            )
        except OSError as exc:
            return TerminalExecutionResult(
                exit_code=127,
                stdout="",
                stderr=f"Terminal launcher startup failed: {exc}",
                shell="launcher",
                argv=invocation.argv,
            )
        return TerminalExecutionResult(
            exit_code=completed.returncode,
            stdout=completed.stdout,
            stderr=completed.stderr,
            shell="launcher",
            argv=invocation.argv,
        )



def normalize_tool_status(status: str | None) -> str:
    if not isinstance(status, str):
        return "unknown"
    normalized = status.strip().lower()
    if not normalized:
        return "unknown"
    if normalized in _CANONICAL_TOOL_STATUSES:
        return normalized
    if normalized in _COMPLETED_TOOL_STATUS_ALIASES:
        return "completed"
    if normalized in _FAILED_TOOL_STATUS_ALIASES:
        return "failed"
    if normalized in _DENIED_TOOL_STATUS_ALIASES:
        return "denied"
    if normalized in _NONTERMINAL_TOOL_STATUS_ALIASES:
        return "nonterminal"
    return "unknown"



def normalize_tool_result(result: ToolResult) -> ToolResult:
    return ToolResult(
        invocation_id=result.invocation_id,
        tool_name=result.tool_name,
        status=normalize_tool_status(result.status),
        output=dict(result.output),
        error=result.error,
    )


_ACTION_RECEIPT_DETAIL_KEYS = (
    "path",
    "artifact_path",
    "output_path",
    "url",
    "id",
    "name",
    "task_id",
    "status",
    "delivery_status",
    "cron_delivery_status",
    "started_count",
    "failed_count",
    "result_text",
)
_ACTION_RECEIPT_SKIP_DETAIL_KEYS = {
    "content",
    "body",
    "html",
    "text",
    "raw",
    "headers",
    "authorization",
    "token",
    "api_key",
    "password",
    "secret",
    "_body_bytes",
}


def _receipt_value_summary(value: object, *, max_length: int = 260) -> str:
    if isinstance(value, bool):
        return "true" if value else "false"
    text = str(value if value is not None else "").strip()
    if not text:
        return ""
    text = " ".join(text.split())
    if len(text) <= max_length:
        return text
    return text[: max_length - 1].rstrip() + "..."


def _action_receipt(
    *,
    action: str,
    object_type: str,
    object_id: str = "",
    object_name: str = "",
    summary: str,
    details: list[str] | None = None,
    changes: list[dict[str, object]] | None = None,
) -> dict[str, object]:
    return {
        "type": "action_receipt",
        "action": action,
        "object_type": object_type,
        "object_id": object_id,
        "object_name": object_name,
        "summary": summary,
        "details": [detail for detail in (details or []) if str(detail or "").strip()],
        "changes": changes or [],
    }


def _connector_request_is_read_only(invocation: ToolInvocation) -> bool:
    if invocation.tool_name != "connector_request":
        return False
    method = str((invocation.arguments or {}).get("method") or "GET").strip().upper()
    return method in {"GET", "HEAD", "OPTIONS"}


def _action_object_type(spec: ToolSpec) -> str:
    tags = tuple(str(tag) for tag in (getattr(spec, "capability_tags", ()) or ()))
    for candidate in ("cron", "reminder", "email", "calendar", "media", "connector"):
        if candidate in tags:
            return candidate
    name = str(spec.name or "").strip()
    if "_" in name:
        return name.rsplit("_", 1)[0].replace("_", " ")
    return name or "tool action"


def _generic_action_receipt_details(output: dict[str, object]) -> list[str]:
    details: list[str] = []
    for key in _ACTION_RECEIPT_DETAIL_KEYS:
        value = output.get(key)
        if value in (None, "", [], {}):
            continue
        if key in _ACTION_RECEIPT_SKIP_DETAIL_KEYS:
            continue
        label = key.replace("_", " ").title()
        details.append(f"{label}: {_receipt_value_summary(value)}.")
    return details


def _with_default_action_receipt(
    *,
    spec: ToolSpec,
    invocation: ToolInvocation,
    result: ToolResult,
) -> ToolResult:
    if result.status != "completed" or not isinstance(result.output, dict):
        return result
    if result.output.get("action_receipt"):
        return result
    if (
        spec.side_effect_class is ToolSideEffectClass.READ
        or spec.side_effect_class is ToolSideEffectClass.DANGEROUS_EXEC
        or _connector_request_is_read_only(invocation)
    ):
        return result

    output = dict(result.output)
    summary = str(output.get("message") or output.get("summary") or "").strip()
    if not summary:
        summary = f"Completed action: {spec.name}."
    object_id = _receipt_value_summary(output.get("id") or output.get("task_id") or output.get("path") or "")
    object_name = _receipt_value_summary(output.get("name") or output.get("path") or output.get("url") or "")
    output["action_receipt"] = _action_receipt(
        action=str(output.get("action") or spec.name),
        object_type=_action_object_type(spec),
        object_id=object_id,
        object_name=object_name,
        summary=summary,
        details=_generic_action_receipt_details(output),
    )
    return ToolResult(
        invocation_id=result.invocation_id,
        tool_name=result.tool_name,
        status=result.status,
        output=output,
        error=result.error,
    )



class ToolRegistry:
    def __init__(
        self,
        *,
        plugin_registration_allowed: bool = True,
        extension_registration_allowed: bool | None = None,
        filesystem_allowed_roots: Iterable[Path] | None = None,
    ) -> None:
        if extension_registration_allowed is not None:
            plugin_registration_allowed = extension_registration_allowed
        self._specs: dict[str, ToolSpec] = {}
        self._handlers: dict[str, ToolHandler] = {}
        self._cleanup_hooks: list[ToolCleanupHook] = []
        self._plugin_registration_allowed = plugin_registration_allowed
        self._installed_plugins: set[str] = set()
        self._filesystem_allowed_roots = tuple(Path(root).resolve() for root in (filesystem_allowed_roots or ()))
        self._registry_revision = 0

    def register(self, spec: ToolSpec, handler: ToolHandler) -> None:
        if spec.name in self._specs:
            raise ValueError(f"Tool already registered: {spec.name}")
        self._specs[spec.name] = spec
        self._handlers[spec.name] = handler
        self._registry_revision += 1
        _clear_deep_agent_profile_cache()

    def unregister(self, name: str) -> None:
        self._specs.pop(name, None)
        self._handlers.pop(name, None)
        self._registry_revision += 1
        _clear_deep_agent_profile_cache()

    @property
    def registry_revision(self) -> int:
        return self._registry_revision

    def get_spec(self, name: str) -> ToolSpec:
        return self._specs[name]

    def list_specs(self) -> list[ToolSpec]:
        return [self._specs[name] for name in sorted(self._specs)]

    def can_invoke_tool(self, name: str) -> bool:
        return str(name or "") in self._specs and str(name or "") in self._handlers

    def list_tool_definitions(self) -> list[dict[str, object]]:
        definitions: list[dict[str, object]] = []
        registered_tool_names = set(self._specs).intersection(self._handlers)
        for spec in self.list_specs():
            if spec.name not in registered_tool_names:
                continue
            definition = {
                "name": spec.name,
                "description": spec.description,
                "input_schema": getattr(spec, "input_schema", None)
                or _default_input_schema_for_tool(spec.name),
                "capability_tags": list(getattr(spec, "capability_tags", ()) or ()),
                "side_effect_class": str(
                    getattr(spec.side_effect_class, "value", spec.side_effect_class)
                ),
                "risk_level": str(getattr(spec.risk_level, "value", spec.risk_level)),
                "requires_approval": bool(spec.requires_approval),
            }
            continuation_tools = [
                str(name)
                for name in (getattr(spec, "continuation_tools", ()) or ())
                if str(name or "") in registered_tool_names and str(name or "") != spec.name
            ]
            if continuation_tools:
                definition["continuation_tools"] = continuation_tools
            definitions.append(definition)
        return definitions

    def filesystem_allowed_roots(self) -> tuple[Path, ...]:
        return self._filesystem_allowed_roots

    def set_filesystem_allowed_roots(self, roots: Iterable[Path]) -> None:
        self._filesystem_allowed_roots = tuple(Path(root).resolve() for root in roots)

    def mark_plugin_installed(self, plugin_name: str) -> None:
        normalized = plugin_name.strip()
        if normalized:
            self._installed_plugins.add(normalized)

    def unmark_plugin_installed(self, plugin_name: str) -> None:
        normalized = plugin_name.strip()
        if normalized:
            self._installed_plugins.discard(normalized)

    def list_installed_plugins(self) -> list[str]:
        return sorted(self._installed_plugins)

    def is_plugin_installed(self, plugin_name: str) -> bool:
        return plugin_name in self._installed_plugins

    @staticmethod
    def _missing_required_schema_arguments(schema: dict[str, object], arguments: dict[str, object]) -> list[str]:
        missing: list[str] = []
        for raw_name in schema.get("required", ()) or ():
            name = str(raw_name)
            value = arguments.get(name)
            if value is None or (isinstance(value, str) and not value.strip()):
                missing.append(name)
        return missing

    def _preflight_schema_result(self, spec: ToolSpec, invocation: ToolInvocation) -> ToolResult | None:
        schema = spec.input_schema or _default_input_schema_for_tool(spec.name)
        if not isinstance(schema, dict):
            return None
        missing = self._missing_required_schema_arguments(schema, invocation.arguments or {})
        if not missing:
            return None
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="failed",
            output={
                "reason": "invalid_tool_arguments",
                "missing_required": missing,
                "suppress_activity": True,
            },
            error=f"Missing required arguments: {', '.join(missing)}",
        )

    def invoke(self, invocation: ToolInvocation) -> ToolResult:
        invocation = _with_resolved_virtual_workspace_paths(invocation)
        handler = self._handlers.get(invocation.tool_name)
        if handler is None:
            raise KeyError(f"Unknown tool: {invocation.tool_name}")
        spec = self._specs[invocation.tool_name]
        preflight = self._preflight_schema_result(spec, invocation)
        if preflight is not None:
            return preflight
        return _invoke_tool_handler_with_timeout(spec, handler, invocation)

    def register_cleanup_hook(self, hook: ToolCleanupHook) -> None:
        self._cleanup_hooks.append(hook)

    def run_cleanup_hooks(self, *, scope_id: str | None = None) -> None:
        for hook in tuple(self._cleanup_hooks):
            try:
                hook(scope_id)
            except Exception:
                logger.debug("Tool cleanup hook failed", exc_info=True)

    def require_extension_registration_allowed(self) -> None:
        if not self._plugin_registration_allowed:
            raise ValueError("Tool registry only accepts core tools")

    def require_plugin_registration_allowed(self) -> None:
        self.require_extension_registration_allowed()


def _load_sentinel_policy() -> SentinelPolicy:
    try:
        from nullion.preferences import load_preferences

        prefs = load_preferences()
        risk_threshold = 3 if prefs.auto_mode is False and prefs.sentinel_mode == "risk_based" else prefs.sentinel_risk_level
        return SentinelPolicy.from_values(
            mode=prefs.sentinel_mode,
            risk_threshold=risk_threshold,
        )
    except Exception:
        return SentinelPolicy()


def _env_flag(name: str, *, default: bool = True) -> bool:
    raw = os.environ.get(name)
    if raw is None or raw.strip() == "":
        return default
    return raw.strip().lower() not in {"0", "false", "no", "off"}


def _env_int(name: str, default: int, *, minimum: int = 0) -> int:
    raw = os.environ.get(name)
    if raw is None or raw.strip() == "":
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return max(minimum, value)


def _connector_access_enabled() -> bool:
    def _has_connector_configuration() -> bool:
        if str(os.environ.get("NULLION_CONNECTOR_GATEWAY") or "").strip():
            return True
        try:
            from nullion.connections import load_connection_registry

            return any(
                _connector_provider_id_looks_external(getattr(connection, "provider_id", ""))
                for connection in load_connection_registry().connections
                if getattr(connection, "active", True)
            )
        except Exception:
            return False

    raw = os.environ.get("NULLION_CONNECTOR_ACCESS_ENABLED")
    if raw is not None and raw.strip():
        return _env_flag("NULLION_CONNECTOR_ACCESS_ENABLED") and _has_connector_configuration()
    return _has_connector_configuration()


def _connector_provider_id_looks_external(provider_id: object) -> bool:
    normalized = str(provider_id or "").strip().lower()
    return normalized.startswith("skill_pack_connector_") or normalized.endswith("_connector_provider")


def _connector_provider_id_for_pack_id(pack_id: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", str(pack_id or "").strip().lower()).strip("_") or "custom_skill"
    return f"skill_pack_connector_{slug}"


def _installed_connector_skill_pack_id_for_provider(provider_id: str) -> str | None:
    try:
        from nullion.skill_pack_installer import list_installed_skill_packs

        for pack in list_installed_skill_packs():
            pack_id = str(getattr(pack, "pack_id", "") or "").strip()
            if pack_id and _connector_provider_id_for_pack_id(pack_id) == provider_id:
                return pack_id
    except Exception:
        return None
    return None


def _available_connector_provider_summaries(principal_id: str | None) -> list[dict[str, object]]:
    try:
        from nullion.connections import connection_for_principal, load_connection_registry

        summaries: list[dict[str, object]] = []
        seen: set[str] = set()
        for connection in load_connection_registry().connections:
            provider_id = str(getattr(connection, "provider_id", "") or "").strip()
            if (
                not provider_id
                or provider_id in seen
                or not getattr(connection, "active", True)
                or not _connector_provider_id_looks_external(provider_id)
            ):
                continue
            if connection_for_principal(principal_id, provider_id) is None:
                continue
            seen.add(provider_id)
            summary = {
                "provider_id": provider_id,
                "display_name": str(getattr(connection, "display_name", "") or provider_id),
                "credential_scope": str(getattr(connection, "credential_scope", "") or "workspace"),
            }
            skill_pack_id = _installed_connector_skill_pack_id_for_provider(provider_id)
            if skill_pack_id:
                summary["skill_pack_id"] = skill_pack_id
            base_urls = _connector_allowed_base_urls(connection, provider_id)
            if base_urls:
                summary["base_urls"] = list(base_urls)
            summaries.append(summary)
        return summaries
    except Exception:
        return []


def _account_tool_failure_output(
    principal_id: str | None,
    *,
    query: str | None = None,
    message_id: str | None = None,
) -> dict[str, object]:
    output: dict[str, object] = {}
    if query is not None:
        output["query"] = query
    if message_id is not None:
        output["id"] = message_id
    connector_providers = _available_connector_provider_summaries(principal_id)
    if connector_providers:
        output["available_connector_providers"] = connector_providers
        output["next_step"] = (
            "A native provider failed. If an enabled connector skill covers this task, "
            "consult that skill's instructions and try connector_request with a listed provider_id "
            "before concluding account access is unavailable."
        )
    return output


def _connector_request_boundary_preapproved(invocation: "ToolInvocation", fact: BoundaryFact) -> bool:
    provider_id = str(invocation.arguments.get("provider_id") or "").strip()
    if not provider_id:
        if invocation.tool_name in {"email_search", "email_read", "email_attachment_read"}:
            try:
                from nullion.connections import infer_email_plugin_provider

                provider_id = infer_email_plugin_provider(principal_id=invocation.principal_id) or ""
            except Exception:
                provider_id = ""
            if not provider_id:
                try:
                    from nullion.connections import default_email_connector_provider_id

                    provider_id = default_email_connector_provider_id(invocation.principal_id) or ""
                except Exception:
                    provider_id = ""
        if not provider_id:
            return False
    try:
        from nullion.connections import connection_for_principal

        connection = connection_for_principal(invocation.principal_id, provider_id)
    except Exception:
        connection = None
    if connection is None or not getattr(connection, "active", True):
        return False
    if invocation.tool_name in {"email_search", "email_read", "email_attachment_read"}:
        return (
            fact.kind is BoundaryKind.ACCOUNT_ACCESS
            and fact.operation == "read"
            and fact.target == provider_id
        )
    if invocation.tool_name != "connector_request":
        return False
    if fact.kind is BoundaryKind.ACCOUNT_ACCESS:
        if fact.operation != "read":
            return False
        return fact.target == provider_id
    if fact.kind is BoundaryKind.OUTBOUND_NETWORK:
        raw_url = invocation.arguments.get("url")
        if not isinstance(raw_url, str) or not raw_url.strip():
            return False
        try:
            url = _connector_request_url(raw_url, invocation.arguments.get("params"), connection, provider_id)
        except Exception:
            return False
        allowed_bases = _connector_allowed_base_urls(connection, provider_id)
        return bool(allowed_bases) and any(_url_is_under_base(url, base_url) for base_url in allowed_bases)
    return False

def _boundary_risk_score(fact: BoundaryFact) -> int:
    if fact.kind is BoundaryKind.ACCOUNT_ACCESS:
        return 7
    if fact.kind is BoundaryKind.FILESYSTEM_ACCESS:
        return 6
    if fact.kind is BoundaryKind.OUTBOUND_NETWORK:
        address_class = fact.attributes.get("address_class")
        if address_class in {"private", "loopback", "localhost", "link_local", "reserved"}:
            return 9
        return 4
    return 5


def _boundary_policy_principal_for_fact(principal_id: str, fact: BoundaryFact) -> str:
    if fact.kind is BoundaryKind.FILESYSTEM_ACCESS:
        return permission_scope_principal(principal_id)
    return GLOBAL_PERMISSION_PRINCIPAL


def _matching_active_boundary_permits_for_invocation(store: RuntimeStore, invocation: "ToolInvocation"):
    principal_ids = {
        invocation.principal_id,
        permission_scope_principal(invocation.principal_id),
        GLOBAL_PERMISSION_PRINCIPAL,
        OPERATOR_PERMISSION_PRINCIPAL,
        *_LEGACY_GLOBAL_PERMISSION_PRINCIPALS,
    }
    now = datetime.now(UTC)
    facts = extract_boundary_facts(invocation)
    if not facts:
        return []
    matched = {}
    for permit in store.list_boundary_permits():
        if permit.principal_id not in principal_ids:
            continue
        if not is_boundary_permit_active(permit, now=now):
            continue
        for fact in facts:
            if fact.kind is not permit.boundary_kind:
                continue
            policy_principal = _boundary_policy_principal_for_fact(invocation.principal_id, fact)
            request = BoundaryPolicyRequest(principal_id=policy_principal, boundary=fact)
            permit_rule = BoundaryPolicyRule(
                rule_id=f"permit:{permit.permit_id}",
                principal_id=policy_principal,
                kind=permit.boundary_kind,
                mode=PolicyMode.ALLOW,
                selector=permit.selector,
                created_by=permit.granted_by,
                created_at=permit.granted_at,
                expires_at=permit.expires_at,
                revoked_at=permit.revoked_at,
                reason="boundary_permit",
            )
            if evaluate_boundary_request(request, rules=[permit_rule]) is PolicyDecision.ALLOW:
                key = (permit.principal_id, permit.boundary_kind, permit.selector)
                current = matched.get(key)
                if current is None or permit.granted_at >= current.granted_at:
                    matched[key] = permit
                break
    return list(matched.values())


def _record_wildcard_boundary_permit_accesses(
    store: RuntimeStore,
    invocation: "ToolInvocation",
    permits,
) -> None:
    facts = extract_boundary_facts(invocation)
    if not facts:
        return
    seen: set[tuple[str, str, str]] = set()
    for permit in permits:
        if permit.selector != "*":
            continue
        for fact in facts:
            if fact.kind is not permit.boundary_kind:
                continue
            domain = (
                normalize_outbound_network_selector(fact.target)
                if fact.kind is BoundaryKind.OUTBOUND_NETWORK
                else fact.target
            )
            key = (permit.permit_id, domain, fact.target)
            if key in seen:
                continue
            seen.add(key)
            payload = {
                "permit_id": permit.permit_id,
                "approval_id": permit.approval_id,
                "principal_id": invocation.principal_id,
                "permit_principal_id": permit.principal_id,
                "boundary_kind": permit.boundary_kind.value,
                "selector": permit.selector,
                "domain": domain,
                "target": fact.target,
                "operation": fact.operation,
                "tool_name": invocation.tool_name,
                "invocation_id": invocation.invocation_id,
                "capsule_id": invocation.capsule_id,
                "accessed_at": datetime.now(UTC).isoformat(),
            }
            store.add_event(make_event("boundary_permit.wildcard_access", invocation.principal_id, payload))
            store.add_audit_record(
                make_audit_record("boundary_permit.wildcard_access", invocation.principal_id, payload)
            )


def _archive_entry_is_metadata(name: str) -> bool:
    normalized = str(name or "").replace("\\", "/").lstrip("/")
    if not normalized:
        return True
    parts = [part for part in normalized.split("/") if part]
    if not parts:
        return True
    if parts[0] == "__MACOSX":
        return True
    base = parts[-1]
    return base == ".DS_Store" or base.startswith("._")


def _archive_suffix(path: Path) -> str | None:
    name = path.name.lower()
    suffixes = sorted(
        [".zip", *_TAR_ARCHIVE_FORMATS.keys(), *_SINGLE_COMPRESSED_FORMATS.keys()],
        key=len,
        reverse=True,
    )
    for suffix in suffixes:
        if name.endswith(suffix):
            return suffix
    return None


def _archive_format_for_path(path: Path) -> tuple[str, str, str | None] | None:
    suffix = _archive_suffix(path)
    if suffix is None:
        return None
    if suffix == ".zip":
        return ("zip", "archive", suffix)
    if suffix in _TAR_ARCHIVE_FORMATS:
        return (_TAR_ARCHIVE_FORMATS[suffix][0], "archive", suffix)
    if suffix in _SINGLE_COMPRESSED_FORMATS:
        return (_SINGLE_COMPRESSED_FORMATS[suffix][0], "compressed_file", suffix)
    return None


def _archive_base_name(path: Path) -> str:
    suffix = _archive_suffix(path)
    if suffix and path.name.lower().endswith(suffix):
        base = path.name[: -len(suffix)]
        return base or path.stem or "contents"
    return path.stem or "contents"


def _archive_default_output_dir(archive_path: Path) -> Path:
    return archive_path.with_name(f"{_archive_base_name(archive_path)}_extracted")


def _safe_archive_member_path(output_dir: Path, member_name: str) -> Path | None:
    normalized = str(member_name or "").replace("\\", "/").lstrip("/")
    if not normalized:
        return None
    destination = (output_dir / normalized).resolve()
    try:
        destination.relative_to(output_dir)
    except ValueError:
        return None
    return destination


def _archive_entry_payload(info: zipfile.ZipInfo) -> dict[str, object]:
    return {
        "name": info.filename,
        "type": "directory" if info.is_dir() else "file",
        "bytes": 0 if info.is_dir() else info.file_size,
        "compressed_bytes": info.compress_size,
        "media_type": "" if info.is_dir() else (mimetypes.guess_type(info.filename)[0] or ""),
        "modified": "%04d-%02d-%02dT%02d:%02d:%02d"
        % (
            info.date_time[0],
            info.date_time[1],
            info.date_time[2],
            info.date_time[3],
            info.date_time[4],
            info.date_time[5],
        ),
    }


def _tar_entry_payload(info: tarfile.TarInfo) -> dict[str, object]:
    modified = datetime.fromtimestamp(info.mtime, UTC).replace(microsecond=0).isoformat() if info.mtime else ""
    entry_type = "directory" if info.isdir() else "file" if info.isfile() else "other"
    return {
        "name": info.name,
        "type": entry_type,
        "bytes": 0 if info.isdir() else info.size,
        "compressed_bytes": None,
        "media_type": "" if info.isdir() else (mimetypes.guess_type(info.name)[0] or ""),
        "modified": modified,
    }


def _archive_extension_counts(names: Iterable[str]) -> dict[str, int]:
    counts = Counter((Path(name).suffix.lower().lstrip(".") or "<none>") for name in names)
    return dict(sorted(counts.items()))


def _inspect_zip_archive(path: Path, *, entry_limit: int = _ARCHIVE_READ_ENTRY_LIMIT) -> dict[str, object]:
    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
    real_infos = [info for info in infos if not _archive_entry_is_metadata(info.filename)]
    metadata_infos = [info for info in infos if _archive_entry_is_metadata(info.filename)]
    limited_infos = real_infos[:entry_limit]
    entry_lines = [
        f"{entry['name']} ({entry['type']}, {entry['bytes']} bytes)"
        for entry in (_archive_entry_payload(info) for info in limited_infos)
    ]
    return {
        "path": str(path),
        "type": "archive",
        "archive_format": "zip",
        "media_type": "application/zip",
        "archive": True,
        "entries": [_archive_entry_payload(info) for info in limited_infos],
        "entry_count": len(real_infos),
        "file_count": sum(1 for info in real_infos if not info.is_dir()),
        "directory_count": sum(1 for info in real_infos if info.is_dir()),
        "metadata_entry_count": len(metadata_infos),
        "metadata_entries_ignored": len(metadata_infos),
        "total_archive_entries": len(infos),
        "extension_counts": _archive_extension_counts(info.filename for info in real_infos if not info.is_dir()),
        "truncated": len(real_infos) > entry_limit,
        "content": "\n".join(entry_lines),
    }


def _inspect_tar_archive(path: Path, *, archive_format: str, entry_limit: int = _ARCHIVE_READ_ENTRY_LIMIT) -> dict[str, object]:
    with tarfile.open(path, "r:*") as archive:
        infos = archive.getmembers()
    real_infos = [info for info in infos if not _archive_entry_is_metadata(info.name)]
    metadata_infos = [info for info in infos if _archive_entry_is_metadata(info.name)]
    limited_infos = real_infos[:entry_limit]
    entry_lines = [
        f"{entry['name']} ({entry['type']}, {entry['bytes']} bytes)"
        for entry in (_tar_entry_payload(info) for info in limited_infos)
    ]
    return {
        "path": str(path),
        "type": "archive",
        "archive_format": archive_format,
        "media_type": mimetypes.guess_type(path.name)[0] or "application/x-tar",
        "archive": True,
        "entries": [_tar_entry_payload(info) for info in limited_infos],
        "entry_count": len(real_infos),
        "file_count": sum(1 for info in real_infos if info.isfile()),
        "directory_count": sum(1 for info in real_infos if info.isdir()),
        "metadata_entry_count": len(metadata_infos),
        "metadata_entries_ignored": len(metadata_infos),
        "total_archive_entries": len(infos),
        "extension_counts": _archive_extension_counts(info.name for info in real_infos if info.isfile()),
        "truncated": len(real_infos) > entry_limit,
        "content": "\n".join(entry_lines),
    }


def _single_compressed_output_name(path: Path) -> str:
    suffix = _archive_suffix(path) or path.suffix.lower()
    if suffix and path.name.lower().endswith(suffix):
        return path.name[: -len(suffix)] or path.stem or "decompressed"
    return path.stem or "decompressed"


def _single_compressed_open(path: Path, archive_format: str):
    if archive_format == "gzip":
        return gzip.open(path, "rb")
    if archive_format == "bzip2":
        return bz2.open(path, "rb")
    if archive_format == "xz":
        return lzma.open(path, "rb")
    raise ValueError(f"Unsupported compressed format: {archive_format}")


def _inspect_single_compressed_file(path: Path, *, archive_format: str) -> dict[str, object]:
    output_name = _single_compressed_output_name(path)
    entry = {
        "name": output_name,
        "type": "file",
        "bytes": None,
        "compressed_bytes": path.stat().st_size,
        "media_type": mimetypes.guess_type(output_name)[0] or "",
        "modified": datetime.fromtimestamp(path.stat().st_mtime, UTC).replace(microsecond=0).isoformat(),
    }
    return {
        "path": str(path),
        "type": "compressed_file",
        "archive_format": archive_format,
        "media_type": mimetypes.guess_type(path.name)[0] or "application/octet-stream",
        "archive": True,
        "entries": [entry],
        "entry_count": 1,
        "file_count": 1,
        "directory_count": 0,
        "metadata_entry_count": 0,
        "metadata_entries_ignored": 0,
        "total_archive_entries": 1,
        "extension_counts": _archive_extension_counts([output_name]),
        "truncated": False,
        "content": f"{output_name} (file, compressed {path.stat().st_size} bytes)",
    }


def _inspect_archive_file(path: Path, *, entry_limit: int = _ARCHIVE_READ_ENTRY_LIMIT) -> dict[str, object]:
    archive_details = _archive_format_for_path(path)
    if archive_details is None:
        raise ValueError("unsupported_archive_format")
    archive_format, archive_kind, _suffix = archive_details
    if archive_format == "zip":
        return _inspect_zip_archive(path, entry_limit=entry_limit)
    if archive_kind == "archive":
        return _inspect_tar_archive(path, archive_format=archive_format, entry_limit=entry_limit)
    return _inspect_single_compressed_file(path, archive_format=archive_format)


def _archive_extract_error_result(invocation: ToolInvocation, *, path: Path | None, error: str, reason: str) -> ToolResult:
    output: dict[str, object] = {"reason": reason}
    if path is not None:
        output["path"] = str(path)
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output=output,
        error=error,
    )


def _archive_default_manifest_dir(archive_path: Path, output_dir: Path, *, workspace_root: Path | None = None) -> Path:
    if workspace_root is not None:
        artifact_dir = (workspace_root / "artifacts").resolve()
        if artifact_dir.exists() or archive_path.parent.name == "files":
            return artifact_dir
    if archive_path.parent.name == "files":
        sibling_artifacts = archive_path.parent.parent / "artifacts"
        if sibling_artifacts.exists():
            return sibling_artifacts.resolve()
    return archive_path.parent


def _archive_manifest_paths(
    archive_path: Path,
    output_dir: Path,
    *,
    manifest_output_path: Path | None = None,
    workspace_root: Path | None = None,
) -> tuple[Path, Path, Path]:
    if manifest_output_path is not None:
        xlsx_path = manifest_output_path
        base = xlsx_path.with_suffix("")
    else:
        manifest_dir = _archive_default_manifest_dir(archive_path, output_dir, workspace_root=workspace_root)
        base = manifest_dir / f"{output_dir.name}_manifest"
        xlsx_path = base.with_suffix(".xlsx")
    return base.with_suffix(".csv"), base.with_suffix(".json"), xlsx_path


def _archive_entry_basename(name: object) -> str:
    text = str(name or "").strip()
    return text.replace("\\", "/").rsplit("/", 1)[-1]


def _archive_entry_image_path(entry: dict[str, object]) -> Path | None:
    raw_path = entry.get("path")
    if not isinstance(raw_path, str) or not raw_path.strip():
        return None
    image_path = Path(raw_path)
    media_type = str(entry.get("media_type") or "").strip().lower()
    if not media_type.startswith("image/") and image_path.suffix.lower() not in _ARCHIVE_MANIFEST_IMAGE_SUFFIXES:
        return None
    if not image_path.is_file():
        return None
    embeddable_path = _embeddable_image_path(image_path)
    if embeddable_path.suffix.lower() not in _ARCHIVE_MANIFEST_IMAGE_SUFFIXES:
        return None
    if _local_raster_image_too_small(embeddable_path):
        return None
    return embeddable_path


def _write_archive_entry_manifests(
    archive_path: Path,
    output_dir: Path,
    entries: list[dict[str, object]],
    *,
    manifest_output_path: Path | None = None,
    workspace_root: Path | None = None,
) -> dict[str, object]:
    csv_path, json_path, xlsx_path = _archive_manifest_paths(
        archive_path,
        output_dir,
        manifest_output_path=manifest_output_path,
        workspace_root=workspace_root,
    )
    row_payloads: list[dict[str, object]] = []
    has_images = False
    for index, entry in enumerate(entries, start=1):
        filename = _archive_entry_basename(entry.get("name"))
        image_path = _archive_entry_image_path(entry)
        has_images = has_images or image_path is not None
        row_payloads.append(
            {
                "index": index,
                "name": str(entry.get("name") or ""),
                "filename": filename,
                "extension": Path(filename).suffix.lower(),
                "media_type": str(entry.get("media_type") or ""),
                "bytes": int(entry.get("bytes") or 0),
                "path": str(entry.get("path") or ""),
                "image_path": image_path,
            }
        )
    columns = [
        "#",
        "Name",
        "Filename",
        "Extension",
        "Media Type",
        "Bytes",
        "Extracted Path",
        *(["Image"] if has_images else []),
    ]
    rows: list[dict[str, object]] = []
    for payload in row_payloads:
        row: dict[str, object] = {
            "#": payload["index"],
            "Name": payload["name"],
            "Filename": payload["filename"],
            "Extension": payload["extension"],
            "Media Type": payload["media_type"],
            "Bytes": payload["bytes"],
            "Extracted Path": payload["path"],
        }
        if has_images:
            image_path = payload.get("image_path")
            row["Image"] = str(image_path) if isinstance(image_path, Path) else ""
        rows.append(row)
    csv_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    xlsx_path.parent.mkdir(parents=True, exist_ok=True)
    with csv_path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=columns)
        writer.writeheader()
        writer.writerows(rows)
    with json_path.open("w", encoding="utf-8") as handle:
        json.dump(
            {
                "columns": columns,
                "rows": rows,
                "entry_count": len(rows),
            },
            handle,
            ensure_ascii=False,
            indent=2,
        )
    manifest: dict[str, object] = {
        "manifest_csv_path": str(csv_path),
        "manifest_json_path": str(json_path),
        "manifest_row_count": len(rows),
        "manifest_columns": columns,
        "embedded_images": [],
        "source_image_paths": [],
        "manifest_embedded_image_count": 0,
    }
    try:
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as WorksheetImage
        from openpyxl.styles import Font
    except ModuleNotFoundError:
        return manifest
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Archive Items"
    sheet.append(columns)
    for cell in sheet[1]:
        cell.font = Font(bold=True)
    for row in rows:
        sheet.append(["" if column == "Image" else row[column] for column in columns])
    for column_cells in sheet.columns:
        max_width = min(
            80,
            max(len(str(cell.value or "")) for cell in column_cells) + 2,
        )
        sheet.column_dimensions[column_cells[0].column_letter].width = max(10, max_width)
    if has_images:
        image_column_index = columns.index("Image") + 1
        sheet.column_dimensions[sheet.cell(row=1, column=image_column_index).column_letter].width = 22
        embedded_images: list[str] = []
        optimized_image_paths: list[str] = []
        skipped_images: list[str] = []
        for row_index, payload in enumerate(row_payloads, start=2):
            image_path = payload.get("image_path")
            if not isinstance(image_path, Path):
                continue
            try:
                embed_path = optimized_embedded_image_path(image_path)
                worksheet_image = WorksheetImage(str(embed_path))
                if worksheet_image.width > 140:
                    scale = 140 / float(worksheet_image.width)
                    worksheet_image.width = int(worksheet_image.width * scale)
                    worksheet_image.height = int(worksheet_image.height * scale)
                if worksheet_image.height > 120:
                    scale = 120 / float(worksheet_image.height)
                    worksheet_image.width = int(worksheet_image.width * scale)
                    worksheet_image.height = int(worksheet_image.height * scale)
                sheet.add_image(
                    worksheet_image,
                    f"{sheet.cell(row=row_index, column=image_column_index).coordinate}",
                )
                sheet.row_dimensions[row_index].height = max(sheet.row_dimensions[row_index].height or 15, 92)
                embedded_images.append(str(image_path))
                if embed_path != image_path:
                    optimized_image_paths.append(str(embed_path))
            except Exception:
                skipped_images.append(str(image_path))
        manifest["embedded_images"] = embedded_images
        manifest["source_image_paths"] = embedded_images
        manifest["manifest_embedded_image_count"] = len(embedded_images)
        if optimized_image_paths:
            manifest["optimized_image_paths"] = optimized_image_paths
        if skipped_images:
            manifest["skipped_images"] = skipped_images
    workbook.save(xlsx_path)
    manifest["manifest_xlsx_path"] = str(xlsx_path)
    return manifest


def _build_archive_extract_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_path = invocation.arguments.get("path")
        if not isinstance(raw_path, str) or not raw_path:
            return _archive_extract_error_result(
                invocation,
                path=None,
                error="Missing required argument: path",
                reason="missing_path",
            )
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return _archive_extract_error_result(
                invocation,
                path=None,
                error="Archive extraction requires workspace_root or allowed_roots",
                reason="missing_workspace_roots",
            )
        archive_path = Path(raw_path).expanduser().resolve()
        if not _path_within_any_root(archive_path, effective_roots) and not _is_approved_filesystem_path(
            archive_path, invocation.trusted_filesystem_selectors
        ):
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"Path is outside workspace root: {archive_path}",
                reason="path_outside_workspace",
            )
        archive_details = _archive_format_for_path(archive_path)
        if archive_details is None:
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"Unsupported archive format: {archive_path.suffix or archive_path.name}",
                reason="unsupported_archive_format",
            )
        if not archive_path.is_file():
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"File not found: {archive_path}",
                reason="file_not_found",
            )
        raw_output_dir = invocation.arguments.get("output_dir")
        if isinstance(raw_output_dir, str) and raw_output_dir.strip():
            requested_output = Path(raw_output_dir).expanduser()
            output_dir = requested_output.resolve() if requested_output.is_absolute() else (archive_path.parent / requested_output).resolve()
        else:
            output_dir = _archive_default_output_dir(archive_path).resolve()
        if not _path_within_any_root(output_dir, effective_roots) and not _is_approved_filesystem_path(
            output_dir, invocation.trusted_filesystem_selectors
        ):
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"Output directory is outside workspace root: {output_dir}",
                reason="output_dir_outside_workspace",
            )
        manifest_output_path: Path | None = None
        raw_manifest_output_path = invocation.arguments.get("manifest_output_path")
        if isinstance(raw_manifest_output_path, str) and raw_manifest_output_path.strip():
            requested_manifest = Path(raw_manifest_output_path).expanduser()
            manifest_output_path = (
                requested_manifest.resolve()
                if requested_manifest.is_absolute()
                else (
                    _archive_default_manifest_dir(archive_path, output_dir, workspace_root=resolved_root)
                    / requested_manifest
                ).resolve()
            )
            if manifest_output_path.suffix.lower() != ".xlsx":
                return _archive_extract_error_result(
                    invocation,
                    path=archive_path,
                    error="manifest_output_path must end in .xlsx",
                    reason="invalid_manifest_output_path",
                )
            if not _path_within_any_root(manifest_output_path, effective_roots) and not _is_approved_filesystem_path(
                manifest_output_path, invocation.trusted_filesystem_selectors
            ):
                return _archive_extract_error_result(
                    invocation,
                    path=archive_path,
                    error=f"Manifest output path is outside workspace root: {manifest_output_path}",
                    reason="manifest_output_path_outside_workspace",
                )
        include_metadata = invocation.arguments.get("include_metadata") is True
        extracted_entries: list[dict[str, object]] = []
        skipped_metadata = 0
        skipped_unsafe = 0
        archive_format, archive_kind, _archive_suffix_text = archive_details
        try:
            output_dir.mkdir(parents=True, exist_ok=True)
            if archive_format == "zip":
                with zipfile.ZipFile(archive_path) as archive:
                    infos = archive.infolist()
                    if len(infos) > _ARCHIVE_EXTRACT_ENTRY_LIMIT:
                        return _archive_extract_error_result(
                            invocation,
                            path=archive_path,
                            error=f"Archive has too many entries to extract safely: {len(infos)}",
                            reason="archive_entry_limit_exceeded",
                        )
                    for info in infos:
                        if _archive_entry_is_metadata(info.filename) and not include_metadata:
                            skipped_metadata += 1
                            continue
                        destination = _safe_archive_member_path(output_dir, info.filename)
                        if destination is None:
                            skipped_unsafe += 1
                            continue
                        if info.is_dir():
                            destination.mkdir(parents=True, exist_ok=True)
                            continue
                        destination.parent.mkdir(parents=True, exist_ok=True)
                        with archive.open(info) as source, destination.open("wb") as target:
                            shutil.copyfileobj(source, target)
                        extracted_entries.append({
                            "name": info.filename,
                            "path": str(destination),
                            "bytes": destination.stat().st_size,
                            "media_type": mimetypes.guess_type(info.filename)[0] or "",
                        })
            elif archive_kind == "archive":
                with tarfile.open(archive_path, "r:*") as archive:
                    infos = archive.getmembers()
                    if len(infos) > _ARCHIVE_EXTRACT_ENTRY_LIMIT:
                        return _archive_extract_error_result(
                            invocation,
                            path=archive_path,
                            error=f"Archive has too many entries to extract safely: {len(infos)}",
                            reason="archive_entry_limit_exceeded",
                        )
                    for info in infos:
                        if _archive_entry_is_metadata(info.name) and not include_metadata:
                            skipped_metadata += 1
                            continue
                        destination = _safe_archive_member_path(output_dir, info.name)
                        if destination is None:
                            skipped_unsafe += 1
                            continue
                        if info.isdir():
                            destination.mkdir(parents=True, exist_ok=True)
                            continue
                        if not info.isfile():
                            skipped_unsafe += 1
                            continue
                        source = archive.extractfile(info)
                        if source is None:
                            skipped_unsafe += 1
                            continue
                        destination.parent.mkdir(parents=True, exist_ok=True)
                        with source, destination.open("wb") as target:
                            shutil.copyfileobj(source, target)
                        extracted_entries.append({
                            "name": info.name,
                            "path": str(destination),
                            "bytes": destination.stat().st_size,
                            "media_type": mimetypes.guess_type(info.name)[0] or "",
                        })
            else:
                output_name = _single_compressed_output_name(archive_path)
                destination = _safe_archive_member_path(output_dir, output_name)
                if destination is None:
                    skipped_unsafe += 1
                else:
                    destination.parent.mkdir(parents=True, exist_ok=True)
                    with _single_compressed_open(archive_path, archive_format) as source, destination.open("wb") as target:
                        shutil.copyfileobj(source, target)
                    extracted_entries.append({
                        "name": output_name,
                        "path": str(destination),
                        "bytes": destination.stat().st_size,
                        "media_type": mimetypes.guess_type(output_name)[0] or "",
                    })
        except (zipfile.BadZipFile, tarfile.TarError, EOFError, lzma.LZMAError):
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"File is not a valid supported archive: {archive_path}",
                reason="bad_archive_file",
            )
        except OSError as exc:
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"Could not extract archive: {exc}",
                reason="archive_extract_failed",
            )
        manifest_output: dict[str, object] = {}
        try:
            manifest_output = _write_archive_entry_manifests(
                archive_path,
                output_dir,
                extracted_entries,
                manifest_output_path=manifest_output_path,
                workspace_root=resolved_root,
            )
        except OSError as exc:
            return _archive_extract_error_result(
                invocation,
                path=archive_path,
                error=f"Could not write archive manifest: {exc}",
                reason="archive_manifest_write_failed",
            )
        manifest_artifact_paths = [
            str(path)
            for key in ("manifest_xlsx_path", "manifest_csv_path", "manifest_json_path")
            if isinstance(path := manifest_output.get(key), str) and path
        ]
        deliverable_manifest_paths = [
            str(path)
            for key in ("manifest_xlsx_path",)
            if manifest_output_path is not None and isinstance(path := manifest_output.get(key), str) and path
        ]
        artifact_descriptors = [
            artifact_output_descriptor(archive_path, role=ARTIFACT_ROLE_SOURCE, kind="archive"),
            *(
                artifact_output_descriptor(path, role=ARTIFACT_ROLE_DELIVERABLE, kind="archive_manifest")
                for path in deliverable_manifest_paths
            ),
            *(
                artifact_output_descriptor(path, role=ARTIFACT_ROLE_INTERMEDIATE, kind="archive_manifest_sidecar")
                for path in manifest_artifact_paths
                if path not in deliverable_manifest_paths
            ),
        ]
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "path": str(archive_path),
                "archive_path": str(archive_path),
                "archive_format": archive_format,
                "output_dir": str(output_dir),
                "entries": extracted_entries[:_ARCHIVE_READ_ENTRY_LIMIT],
                "entry_count": len(extracted_entries),
                "file_count": len(extracted_entries),
                "metadata_entries_ignored": skipped_metadata,
                "unsafe_entries_ignored": skipped_unsafe,
                "truncated": len(extracted_entries) > _ARCHIVE_READ_ENTRY_LIMIT,
                **manifest_output,
                "manifest_columns": manifest_output.get(
                    "manifest_columns",
                    ["#", "Name", "Media Type", "Bytes", "Extracted Path"],
                ),
                "manifest_artifact_paths": manifest_artifact_paths,
                "artifact_path": deliverable_manifest_paths[0] if deliverable_manifest_paths else "",
                "artifact_paths": deliverable_manifest_paths,
                "artifact_descriptors": artifact_descriptors,
            },
            error=None,
        )

    return handler


def _archive_create_error_result(invocation: ToolInvocation, *, path: Path | None, error: str, reason: str) -> ToolResult:
    output: dict[str, object] = {"reason": reason}
    if path is not None:
        output["output_path"] = str(path)
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output=output,
        error=error,
    )


def _build_archive_create_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def _resolve_input_path(raw_path: object, roots: tuple[Path, ...]) -> Path | None:
        if not isinstance(raw_path, str) or not raw_path.strip():
            return None
        candidate = Path(raw_path).expanduser()
        if candidate.is_absolute():
            candidates = (candidate,)
        else:
            candidates = tuple(root / candidate for root in roots)
        for possible in candidates:
            try:
                resolved = possible.resolve()
            except OSError:
                continue
            if _path_within_any_root(resolved, roots):
                return resolved
        return None

    def _iter_source_files(source: Path, *, include_metadata: bool) -> Iterable[tuple[Path, str]]:
        if source.is_file():
            if include_metadata or not _archive_entry_is_metadata(source.name):
                yield source, source.name
            return
        if not source.is_dir():
            return
        for current_dir, dirnames, filenames in os.walk(source, topdown=True, followlinks=False):
            current_path = Path(current_dir)
            dirnames[:] = [
                dirname
                for dirname in sorted(dirnames)
                if include_metadata or not _archive_entry_is_metadata(str((current_path / dirname).relative_to(source)))
            ]
            for filename in sorted(filenames):
                file_path = current_path / filename
                try:
                    resolved_file = file_path.resolve()
                except OSError:
                    continue
                if not resolved_file.is_file():
                    continue
                arcname = str(file_path.relative_to(source)).replace(os.sep, "/")
                if not include_metadata and _archive_entry_is_metadata(arcname):
                    continue
                yield resolved_file, arcname

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_output_path = invocation.arguments.get("output_path")
        if not isinstance(raw_output_path, str) or not raw_output_path.strip():
            return _archive_create_error_result(
                invocation,
                path=None,
                error="Missing required argument: output_path",
                reason="missing_output_path",
            )
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return _archive_create_error_result(
                invocation,
                path=None,
                error="Archive creation requires workspace_root or allowed_roots",
                reason="missing_workspace_roots",
            )
        output_path = Path(raw_output_path).expanduser()
        if not output_path.is_absolute():
            output_path = (effective_roots[0] / output_path)
        output_path = output_path.resolve()
        output_archive_details = _archive_format_for_path(output_path)
        if output_archive_details is None:
            return _archive_create_error_result(
                invocation,
                path=output_path,
                error="archive_create output_path must use a supported archive/compression extension.",
                reason="invalid_output_extension",
            )
        if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
            output_path, invocation.trusted_filesystem_selectors
        ):
            return _archive_create_error_result(
                invocation,
                path=output_path,
                error=f"Output path is outside workspace root: {output_path}",
                reason="output_path_outside_workspace",
            )
        include_metadata = invocation.arguments.get("include_metadata") is True
        raw_source_paths = invocation.arguments.get("source_paths")
        source_paths: list[Path] = []
        raw_source_dir = invocation.arguments.get("source_dir")
        source_dir = _resolve_input_path(raw_source_dir, effective_roots)
        if source_dir is not None:
            source_paths.append(source_dir)
        if isinstance(raw_source_paths, (list, tuple)):
            for raw_source_path in raw_source_paths:
                resolved_source = _resolve_input_path(raw_source_path, effective_roots)
                if resolved_source is not None:
                    source_paths.append(resolved_source)
        source_paths = list(dict.fromkeys(source_paths))
        if not source_paths:
            return _archive_create_error_result(
                invocation,
                path=output_path,
                error="archive_create requires source_dir or source_paths inside the workspace.",
                reason="missing_sources",
            )
        entries: list[tuple[Path, str]] = []
        output_identity = output_path.resolve()
        for source_path in source_paths:
            if not source_path.exists():
                continue
            for file_path, arcname in _iter_source_files(source_path, include_metadata=include_metadata):
                if file_path.resolve() == output_identity:
                    continue
                entries.append((file_path, arcname))
                if len(entries) > _ARCHIVE_EXTRACT_ENTRY_LIMIT:
                    return _archive_create_error_result(
                        invocation,
                        path=output_path,
                        error=f"Too many files to archive safely: {len(entries)}",
                        reason="archive_entry_limit_exceeded",
                    )
        deduped_entries: list[tuple[Path, str]] = []
        seen_names: set[str] = set()
        for file_path, arcname in entries:
            clean_name = arcname.replace("\\", "/").lstrip("/")
            if not clean_name or clean_name in seen_names:
                continue
            seen_names.add(clean_name)
            deduped_entries.append((file_path, clean_name))
        if not deduped_entries:
            return _archive_create_error_result(
                invocation,
                path=output_path,
                error="No archiveable files found in the requested sources.",
                reason="no_archiveable_files",
            )
        archive_format, archive_kind, suffix_text = output_archive_details
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            if archive_format == "zip":
                with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
                    for file_path, arcname in deduped_entries:
                        archive.write(file_path, arcname)
            elif archive_kind == "archive":
                mode = _TAR_ARCHIVE_FORMATS[suffix_text or ".tar"][1]
                with tarfile.open(output_path, mode) as archive:
                    for file_path, arcname in deduped_entries:
                        archive.add(file_path, arcname=arcname, recursive=False)
            else:
                if len(deduped_entries) != 1:
                    return _archive_create_error_result(
                        invocation,
                        path=output_path,
                        error="Single-file compression formats require exactly one source file.",
                        reason="single_file_compression_requires_one_source",
                    )
                source_file, _arcname = deduped_entries[0]
                if archive_format == "gzip":
                    opener = gzip.open
                elif archive_format == "bzip2":
                    opener = bz2.open
                elif archive_format == "xz":
                    opener = lzma.open
                else:
                    raise ValueError(f"Unsupported compressed format: {archive_format}")
                with source_file.open("rb") as source, opener(output_path, "wb") as target:
                    shutil.copyfileobj(source, target)
        except (OSError, tarfile.TarError, ValueError) as exc:
            return _archive_create_error_result(
                invocation,
                path=output_path,
                error=f"Could not create archive: {exc}",
                reason="archive_create_failed",
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "artifact_descriptors": [
                    artifact_output_descriptor(output_path, role=ARTIFACT_ROLE_DELIVERABLE, kind="archive")
                ],
                "archive_format": archive_format,
                "bytes_written": output_path.stat().st_size,
                "entry_count": len(deduped_entries),
                "entries": [
                    {
                        "name": arcname,
                        "path": str(file_path),
                        "bytes": file_path.stat().st_size,
                        "media_type": mimetypes.guess_type(arcname)[0] or "",
                    }
                    for file_path, arcname in deduped_entries[:_ARCHIVE_READ_ENTRY_LIMIT]
                ],
                "truncated": len(deduped_entries) > _ARCHIVE_READ_ENTRY_LIMIT,
            },
            error=None,
        )

    return handler


class ToolExecutor:
    def __init__(
        self,
        *,
        store: RuntimeStore,
        registry: ToolRegistry,
        allowed_tool_names: Iterable[str] | None = None,
        denied_tool_names: Iterable[str] | None = None,
        sentinel_policy: SentinelPolicy | None = None,
    ) -> None:
        self._store = store
        self._registry = registry
        self._allowed_tool_names = tuple(sorted(set(allowed_tool_names or ())))
        self._denied_tool_names = tuple(sorted(set(denied_tool_names or ())))
        self._sentinel_policy = sentinel_policy or _load_sentinel_policy()

    def _record_egress_event(self, event_type: str, invocation: ToolInvocation, payload: dict[str, object]) -> None:
        details = {
            "invocation_id": invocation.invocation_id,
            "tool_name": invocation.tool_name,
            "principal_id": invocation.principal_id,
            "capsule_id": invocation.capsule_id,
            **payload,
        }
        self._store.add_event(make_event(event_type, invocation.principal_id, details))
        self._store.add_audit_record(make_audit_record(event_type, invocation.principal_id, details))

    def _record_detected_egress_attempts(self, invocation: ToolInvocation, egress_attempts: list[dict[str, object]]) -> None:
        for attempt in egress_attempts:
            self._record_egress_event("tool.egress.detected", invocation, dict(attempt))

    def _record_egress_outcome(self, invocation: ToolInvocation, result: ToolResult, egress_attempts: list[dict[str, object]]) -> None:
        if not egress_attempts:
            return
        network_mode = result.output.get("network_mode") if isinstance(result.output, dict) else None
        outcome_event = None
        if result.output.get("reason") == "network_denied":
            outcome_event = "tool.egress.denied"
        elif result.status in {"completed", "failed", "nonterminal"}:
            outcome_event = "tool.egress.allowed"
        if outcome_event is None:
            return
        for attempt in egress_attempts:
            payload = dict(attempt)
            if network_mode is not None:
                payload["network_mode"] = network_mode
            payload["tool_status"] = result.status
            self._record_egress_event(outcome_event, invocation, payload)

    def _deny_invocation(
        self,
        invocation: ToolInvocation,
        *,
        reason: str,
        error: str,
        output: dict[str, object],
    ) -> ToolResult:
        denied_details = {
            "invocation_id": invocation.invocation_id,
            "tool_name": invocation.tool_name,
            "principal_id": invocation.principal_id,
            "capsule_id": invocation.capsule_id,
            "status": "denied",
            "reason": reason,
        }
        self._store.add_event(make_event("tool.invocation.denied", invocation.principal_id, denied_details))
        self._store.add_audit_record(make_audit_record("tool.invocation.denied", invocation.principal_id, denied_details))
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="denied",
            output=output,
            error=error,
        )

    def _has_active_grant(self, *, principal_id: str, permissions: Iterable[str]) -> bool:
        required = set(permissions)
        principal_ids = _tool_grant_principal_candidates(principal_id)
        for grant in self._store.list_permission_grants():
            if grant.principal_id not in principal_ids:
                continue
            if grant.permission not in required:
                continue
            if not is_permission_grant_active(grant):
                continue
            return True
        return False

    def _has_required_tool_grant(self, invocation: ToolInvocation) -> bool:
        if _tool_blocked_for_principal(invocation.principal_id, invocation.tool_name):
            return False
        if invocation.tool_name == "email_send":
            return bool(self._matching_email_send_review_grants(invocation))
        if self._matching_exact_tool_review_grants(invocation):
            return True
        if invocation.tool_name == "terminal_exec" and self._matching_terminal_destructive_grants(invocation):
            return True
        return self._has_active_grant(
            principal_id=invocation.principal_id,
            permissions=(
                f"tool:{invocation.tool_name}",
                f"tool.{invocation.tool_name}",
                invocation.tool_name,
            ),
        )

    def _matching_terminal_destructive_grants(self, invocation: ToolInvocation):
        if invocation.tool_name != "terminal_exec":
            return []
        raw_command = invocation.arguments.get("command")
        if not isinstance(raw_command, str) or not raw_command.strip():
            return []
        execution_cwd = _terminal_execution_cwd_from_registry(self._registry)
        targets = _terminal_delete_targets_for_invocation(invocation, execution_cwd=execution_cwd)
        if not targets:
            return []
        permission = _terminal_destructive_permission(raw_command, targets)
        principal_ids = _tool_grant_principal_candidates(invocation.principal_id)
        grants = []
        for grant in self._store.list_permission_grants():
            if grant.principal_id not in principal_ids:
                continue
            if grant.permission != permission:
                continue
            if is_permission_grant_active(grant):
                grants.append(grant)
        return grants

    def _revoke_terminal_destructive_grants(self, invocation: ToolInvocation) -> None:
        for grant in self._matching_terminal_destructive_grants(invocation):
            try:
                self._store.add_permission_grant(
                    revoke_permission_grant_record(
                        grant,
                        revoked_by="runtime",
                        revoked_at=datetime.now(UTC),
                        reason="Terminal destructive action approval consumed.",
                    )
                )
            except Exception:
                logger.debug("Could not revoke consumed terminal destructive approval grant", exc_info=True)

    def _matching_exact_tool_review_grants(self, invocation: ToolInvocation):
        current_arguments = redact_value(dict(invocation.arguments or {}))
        matching_approval_ids: set[str] = set()
        for approval in self._store.list_approval_requests():
            if getattr(getattr(approval, "status", None), "value", getattr(approval, "status", "")) != "approved":
                continue
            if approval.requested_by != invocation.principal_id:
                continue
            if approval.action != "use_tool":
                continue
            context = approval.context if isinstance(approval.context, dict) else {}
            approved_tool = str(context.get("tool_name") or approval.resource or "").strip()
            if approved_tool != invocation.tool_name:
                continue
            if context.get("tool_arguments") == current_arguments:
                matching_approval_ids.add(approval.approval_id)
        if not matching_approval_ids:
            return []
        grants = []
        principal_ids = _tool_grant_principal_candidates(invocation.principal_id)
        permission_names = {f"tool:{invocation.tool_name}", f"tool.{invocation.tool_name}", invocation.tool_name}
        for grant in self._store.list_permission_grants():
            if grant.approval_id not in matching_approval_ids:
                continue
            if grant.principal_id not in principal_ids:
                continue
            if grant.permission not in permission_names:
                continue
            if is_permission_grant_active(grant):
                grants.append(grant)
        return grants

    def _matching_email_send_review_grants(self, invocation: ToolInvocation):
        if invocation.tool_name != "email_send":
            return []
        return self._matching_exact_tool_review_grants(invocation)

    def _revoke_email_send_review_grants(self, invocation: ToolInvocation) -> None:
        for grant in self._matching_exact_tool_review_grants(invocation):
            try:
                self._store.add_permission_grant(
                    revoke_permission_grant_record(
                        grant,
                        revoked_by="runtime",
                        revoked_at=datetime.now(UTC),
                        reason="Email send review approval consumed.",
                    )
                )
            except Exception:
                logger.debug("Could not revoke consumed email_send approval grant", exc_info=True)

    def _find_pending_tool_approval(self, invocation: ToolInvocation):
        current_arguments = redact_value(dict(invocation.arguments or {})) if invocation.tool_name == "email_send" else None
        for approval in self._store.list_approval_requests():
            if approval.status.value != "pending":
                continue
            if approval.requested_by != invocation.principal_id:
                continue
            if approval.action != "use_tool":
                continue
            if approval.resource != invocation.tool_name:
                continue
            if current_arguments is not None:
                context = approval.context if isinstance(approval.context, dict) else {}
                if context.get("tool_arguments") != current_arguments:
                    continue
            return approval
        return None

    def _find_pending_boundary_policy_approval(self, invocation: ToolInvocation, *, context: dict[str, object]):
        target = context.get("target")
        if not isinstance(target, str) or not target:
            return None
        requested_workspace = permission_scope_principal(invocation.principal_id)
        requested_key = _boundary_approval_match_key(context=context, fallback_target=target)
        for approval in self._store.list_approval_requests():
            if approval.status.value != "pending":
                continue
            if permission_scope_principal(approval.requested_by) != requested_workspace:
                continue
            if approval.request_kind != "boundary_policy":
                continue
            if approval.action != "allow_boundary":
                continue
            approval_context = approval.context if isinstance(approval.context, dict) else {}
            approval_key = _boundary_approval_match_key(context=approval_context, fallback_target=approval.resource)
            if requested_key is None:
                if approval.resource != target:
                    continue
            elif approval_key != requested_key:
                continue
            return approval
        return None

    def _find_pending_terminal_destructive_approval(self, invocation: ToolInvocation, *, context: dict[str, object]):
        token = context.get("destructive_permission_token")
        if not isinstance(token, str) or not token:
            return None
        for approval in self._store.list_approval_requests():
            if approval.status.value != "pending":
                continue
            if approval.requested_by != invocation.principal_id:
                continue
            if approval.request_kind != TERMINAL_DESTRUCTIVE_ACTION_REQUEST_KIND:
                continue
            if approval.action != "use_tool" or approval.resource != invocation.tool_name:
                continue
            approval_context = approval.context if isinstance(approval.context, dict) else {}
            if approval_context.get("destructive_permission_token") == token:
                return approval
        return None

    def _ensure_terminal_destructive_approval_request(
        self,
        invocation: ToolInvocation,
        *,
        command: str,
        targets: tuple[dict[str, object], ...],
    ):
        workspace_id = "workspace_admin"
        try:
            from nullion.connections import workspace_id_for_principal

            workspace_id = workspace_id_for_principal(invocation.principal_id)
        except Exception:
            pass
        approval_context = {
            **_terminal_destructive_approval_context(invocation, command=command, targets=targets),
            "workspace_id": workspace_id,
            FLOW_TRIGGER_CONTEXT_KEY: dict(invocation.flow_context)
            if isinstance(invocation.flow_context, dict)
            else build_trigger_flow_context(
                principal_id=invocation.principal_id,
                invocation_id=invocation.invocation_id,
                capsule_id=invocation.capsule_id,
                flow_kind="terminal_destructive_action",
            ),
        }
        existing = self._find_pending_terminal_destructive_approval(invocation, context=approval_context)
        if existing is not None:
            refreshed = _refresh_pending_approval_request(
                existing,
                context=approval_context,
                resource=invocation.tool_name,
            )
            self._store.add_approval_request(refreshed)
            return refreshed
        approval = create_approval_request(
            requested_by=invocation.principal_id,
            action="use_tool",
            resource=invocation.tool_name,
            request_kind=TERMINAL_DESTRUCTIVE_ACTION_REQUEST_KIND,
            context=approval_context,
        )
        self._store.add_approval_request(approval)
        return approval

    def _terminal_destructive_preapproval_result(self, invocation: ToolInvocation) -> ToolResult | None:
        if invocation.tool_name != "terminal_exec":
            return None
        raw_command = invocation.arguments.get("command")
        if not isinstance(raw_command, str) or not raw_command.strip():
            return None
        execution_cwd = _terminal_execution_cwd_from_registry(self._registry)
        targets = _terminal_delete_targets_for_invocation(invocation, execution_cwd=execution_cwd)
        if not targets:
            return None
        if self._matching_terminal_destructive_grants(invocation):
            return None
        approval = self._ensure_terminal_destructive_approval_request(
            invocation,
            command=raw_command,
            targets=targets,
        )
        context = approval.context if isinstance(approval.context, dict) else {}
        return self._deny_invocation(
            invocation,
            reason="terminal_destructive_action_confirmation_required",
            error="Approval required before deleting files with terminal_exec",
            output={
                **context,
                "reason": "terminal_destructive_action_confirmation_required",
                "requires_approval": True,
                "approval_id": approval.approval_id,
            },
        )

    def _ensure_tool_approval_request(self, invocation: ToolInvocation):
        existing = self._find_pending_tool_approval(invocation)
        if existing is not None:
            return existing
        spec = self._registry.get_spec(invocation.tool_name)
        workspace_id = "workspace_admin"
        try:
            from nullion.connections import workspace_id_for_principal

            workspace_id = workspace_id_for_principal(invocation.principal_id)
        except Exception:
            pass
        tool_arguments = redact_value(dict(invocation.arguments or {}))
        approval_context = {
            "workspace_id": workspace_id,
            FLOW_TRIGGER_CONTEXT_KEY: dict(invocation.flow_context)
            if isinstance(invocation.flow_context, dict)
            else build_trigger_flow_context(
                principal_id=invocation.principal_id,
                invocation_id=invocation.invocation_id,
                capsule_id=invocation.capsule_id,
                flow_kind="tool_invocation",
            ),
            "tool_name": invocation.tool_name,
            "tool_description": spec.description,
            "tool_risk_level": str(getattr(spec.risk_level, "value", spec.risk_level)),
            "tool_side_effect_class": str(getattr(spec.side_effect_class, "value", spec.side_effect_class)),
            "requires_approval": spec.requires_approval,
            "tool_permission_scope": spec.permission_scope,
            "tool_arguments": tool_arguments,
        }
        if invocation.tool_name == "email_send":
            preview_path = _email_html_preview_path_for_invocation(invocation)
            if preview_path:
                approval_context["html_preview_path"] = preview_path
        approval = create_approval_request(
            requested_by=invocation.principal_id,
            action="use_tool",
            resource=invocation.tool_name,
            context=approval_context,
        )
        self._store.add_approval_request(approval)
        return approval

    def _ensure_boundary_policy_approval_request(self, invocation: ToolInvocation, *, context: dict[str, object]):
        target = context.get("target")
        if not isinstance(target, str) or not target:
            return None
        workspace_id = "workspace_admin"
        try:
            from nullion.connections import workspace_id_for_principal

            workspace_id = workspace_id_for_principal(invocation.principal_id)
        except Exception:
            pass
        approval_context = {
            **context,
            "workspace_id": workspace_id,
            FLOW_TRIGGER_CONTEXT_KEY: dict(invocation.flow_context)
            if isinstance(invocation.flow_context, dict)
            else build_trigger_flow_context(
                principal_id=invocation.principal_id,
                invocation_id=invocation.invocation_id,
                capsule_id=invocation.capsule_id,
                flow_kind="boundary_policy",
            ),
        }
        if invocation.tool_name == "email_send":
            # Boundary approvals must carry the reviewed email draft so a later
            # account pause cannot render an empty send-review card.
            approval_context.setdefault("tool_arguments", redact_value(dict(invocation.arguments or {})))
            preview_path = _email_html_preview_path_for_invocation(invocation)
            if preview_path:
                approval_context["html_preview_path"] = preview_path
        existing = self._find_pending_boundary_policy_approval(invocation, context=approval_context)
        if existing is not None:
            refreshed = _refresh_pending_approval_request(existing, context=approval_context, resource=target)
            self._store.add_approval_request(refreshed)
            return refreshed
        approval = create_approval_request(
            requested_by=invocation.principal_id,
            action="allow_boundary",
            resource=target,
            request_kind="boundary_policy",
            context=approval_context,
        )
        self._store.add_approval_request(approval)
        return approval

    def _boundary_rules_for_fact(self, invocation: ToolInvocation, *, fact) -> list[BoundaryPolicyRule]:
        policy_principal = _boundary_policy_principal_for_fact(invocation.principal_id, fact)
        principal_ids = {
            invocation.principal_id,
            policy_principal,
            permission_scope_principal(invocation.principal_id),
            GLOBAL_PERMISSION_PRINCIPAL,
            OPERATOR_PERMISSION_PRINCIPAL,
            *_LEGACY_GLOBAL_PERMISSION_PRINCIPALS,
        }
        rules = [
            replace(rule, principal_id=policy_principal)
            if rule.principal_id in {OPERATOR_PERMISSION_PRINCIPAL, *_LEGACY_GLOBAL_PERMISSION_PRINCIPALS}
            else rule
            for rule in self._store.list_boundary_policy_rules()
            if rule.principal_id in principal_ids and rule.kind is fact.kind
        ]
        now = datetime.now(UTC)
        for permit in self._store.list_boundary_permits():
            if permit.principal_id not in principal_ids:
                continue
            if permit.boundary_kind is not fact.kind:
                continue
            if not is_boundary_permit_active(permit, now=now):
                continue
            rules.append(
                BoundaryPolicyRule(
                    rule_id=f"permit:{permit.permit_id}",
                    principal_id=policy_principal,
                    kind=permit.boundary_kind,
                    mode=PolicyMode.ALLOW,
                    selector=permit.selector,
                    created_by=permit.granted_by,
                    created_at=permit.granted_at,
                    expires_at=permit.expires_at,
                    revoked_at=permit.revoked_at,
                    reason="boundary_permit",
                )
            )
        return rules

    def _filesystem_boundary_roots(self, principal_id: str | None = None) -> tuple[Path, ...]:
        return tuple(dict.fromkeys((
            *self._registry.filesystem_allowed_roots(),
            *_principal_workspace_file_roots(principal_id),
        )))

    def _filesystem_boundary_allowed(self, target: str, *, principal_id: str | None = None) -> bool:
        allowed_roots = self._filesystem_boundary_roots(principal_id)
        if not allowed_roots:
            return True
        try:
            path = Path(target).expanduser().resolve()
        except OSError:
            return False
        for root in allowed_roots:
            try:
                path.relative_to(root)
                return True
            except ValueError:
                continue
        return False

    def _scheduled_task_terminal_filesystem_boundary_denial(self, invocation: ToolInvocation, fact) -> ToolResult | None:
        if invocation.tool_name != "terminal_exec":
            return None
        if not _flow_context_is_scheduled_task_run(invocation.flow_context):
            return None
        if fact.kind is not BoundaryKind.FILESYSTEM_ACCESS:
            return None
        if fact.operation not in _TERMINAL_MUTATION_BOUNDARY_OPERATIONS:
            return None
        allowed_roots = self._filesystem_boundary_roots(invocation.principal_id)
        scratch_root = _workspace_scratch_root_for_principal(invocation.principal_id)
        output: dict[str, object] = {
            "reason": "scheduled_task_workspace_path_required",
            "boundary_kind": fact.kind.value,
            "operation": fact.operation,
            "target": fact.target,
            "command_family": fact.attributes.get("command_family"),
            "requires_approval": False,
            "retryable": True,
            "allowed_roots": [str(root) for root in allowed_roots],
            "message": (
                "Scheduled task terminal writes must stay inside trusted workspace storage. "
                "Retry helper scripts under suggested_scratch_root or another allowed workspace root."
            ),
        }
        if scratch_root is not None:
            output["suggested_scratch_root"] = str(scratch_root)
        return self._deny_invocation(
            invocation,
            reason="scheduled_task_workspace_path_required",
            error="Scheduled task terminal filesystem mutation outside trusted workspace",
            output=output,
        )

    def _scheduled_task_terminal_helper_artifact_denial(self, invocation: ToolInvocation, fact) -> ToolResult | None:
        if invocation.tool_name != "terminal_exec":
            return None
        if not _flow_context_is_scheduled_task_run(invocation.flow_context):
            return None
        if fact.kind is not BoundaryKind.FILESYSTEM_ACCESS:
            return None
        if fact.operation not in _TERMINAL_MUTATION_BOUNDARY_OPERATIONS:
            return None
        if not _scheduled_task_helper_script_path(fact.target):
            return None
        artifact_root = _workspace_artifact_root_for_principal(invocation.principal_id)
        if not _path_within_root_text(fact.target, artifact_root):
            return None
        scratch_root = _workspace_scratch_root_for_principal(invocation.principal_id)
        output: dict[str, object] = {
            "reason": "scheduled_task_helper_script_scratch_path_required",
            "boundary_kind": fact.kind.value,
            "operation": fact.operation,
            "target": fact.target,
            "command_family": fact.attributes.get("command_family"),
            "requires_approval": False,
            "retryable": True,
            "message": (
                "Scheduled task helper scripts are internal scratch files. "
                "Retry the script write under suggested_scratch_root instead of the delivery directory."
            ),
        }
        if scratch_root is not None:
            output["suggested_scratch_root"] = str(scratch_root)
        if artifact_root is not None:
            output["blocked_delivery_root"] = str(artifact_root)
        return self._deny_invocation(
            invocation,
            reason="scheduled_task_helper_script_scratch_path_required",
            error="Scheduled task helper script was written to the delivery directory",
            output=output,
        )

    def _scheduled_task_terminal_filesystem_preapproval_result(self, invocation: ToolInvocation) -> ToolResult | None:
        if invocation.tool_name != "terminal_exec":
            return None
        if not _flow_context_is_scheduled_task_run(invocation.flow_context):
            return None
        for fact in extract_boundary_facts(invocation):
            if fact.kind is not BoundaryKind.FILESYSTEM_ACCESS:
                continue
            if fact.operation not in _TERMINAL_MUTATION_BOUNDARY_OPERATIONS:
                continue
            helper_artifact_denial = self._scheduled_task_terminal_helper_artifact_denial(invocation, fact)
            if helper_artifact_denial is not None:
                return helper_artifact_denial
            if self._filesystem_boundary_allowed(fact.target, principal_id=invocation.principal_id):
                continue
            denial = self._scheduled_task_terminal_filesystem_boundary_denial(invocation, fact)
            if denial is not None:
                return denial
        return None

    def _media_filesystem_boundary_decision(self, fact) -> PolicyDecision | None:
        try:
            spec = self._registry.get_spec(fact.tool_name)
        except KeyError:
            return None
        if spec.filesystem_boundary_policy != _FILESYSTEM_BOUNDARY_TRUSTED_ROOTS_ONLY:
            return None
        return PolicyDecision.ALLOW if self._filesystem_boundary_allowed(fact.target) else PolicyDecision.DENY

    def _preflight_boundary_policy_result(self, invocation: ToolInvocation):
        facts = extract_boundary_facts(invocation)
        if invocation.tool_name == "terminal_exec":
            facts = [fact for fact in facts if fact.kind is BoundaryKind.FILESYSTEM_ACCESS]
            if not facts:
                return None
        for fact in facts:
            if _connector_request_boundary_preapproved(invocation, fact):
                continue
            if (
                invocation.tool_name != "terminal_exec"
                and fact.kind in {BoundaryKind.ACCOUNT_ACCESS, BoundaryKind.OUTBOUND_NETWORK}
                and self._matching_exact_tool_review_grants(invocation)
            ):
                # The exact reviewed tool approval already represents
                # consent for this invocation's own account/network boundary.
                continue
            if (
                invocation.tool_name != "terminal_exec"
                and fact.kind is BoundaryKind.FILESYSTEM_ACCESS
                and fact.operation == "read"
                and self._matching_exact_tool_review_grants(invocation)
            ):
                # Reviewed tool approvals may read the exact local files listed
                # in the approved invocation, such as email attachments.
                continue
            if (
                invocation.tool_name in _CALENDAR_ACCOUNT_WRITE_TOOLS
                and fact.kind is BoundaryKind.ACCOUNT_ACCESS
                and self._has_required_tool_grant(invocation)
            ):
                # The exact reviewed calendar mutation approval already
                # represents the user-visible consent for this account write.
                continue
            policy_principal = _boundary_policy_principal_for_fact(invocation.principal_id, fact)
            if fact.kind is BoundaryKind.OUTBOUND_NETWORK:
                request = BoundaryPolicyRequest(principal_id=policy_principal, boundary=fact)
                decision = evaluate_boundary_request(request, rules=self._boundary_rules_for_fact(invocation, fact=fact))
            elif fact.kind is BoundaryKind.FILESYSTEM_ACCESS:
                helper_artifact_denial = self._scheduled_task_terminal_helper_artifact_denial(invocation, fact)
                if helper_artifact_denial is not None:
                    return helper_artifact_denial
                if self._filesystem_boundary_allowed(fact.target, principal_id=invocation.principal_id):
                    continue
                scheduled_terminal_denial = self._scheduled_task_terminal_filesystem_boundary_denial(invocation, fact)
                if scheduled_terminal_denial is not None:
                    return scheduled_terminal_denial
                rules = self._boundary_rules_for_fact(invocation, fact=fact)
                media_decision = self._media_filesystem_boundary_decision(fact)
                if media_decision == PolicyDecision.ALLOW:
                    # File is within trusted roots — fast-path allow without rule lookup.
                    decision = media_decision
                else:
                    matching_rules = [
                        rule
                        for rule in rules
                        if rule.selector == fact.target
                        or fnmatch(fact.target, rule.selector)
                        or (rule.selector.endswith("/*") and fact.target == rule.selector[:-2])
                    ]
                    if matching_rules:
                        # Explicit user-granted rules take precedence even for media tools.
                        request = BoundaryPolicyRequest(principal_id=policy_principal, boundary=fact)
                        decision = evaluate_boundary_request(request, rules=matching_rules)
                    elif media_decision == PolicyDecision.DENY:
                        # Media tool with trusted-roots-only policy, file outside trusted roots,
                        # and no explicit grant — hard deny (no approval flow).
                        decision = PolicyDecision.DENY
                    else:
                        decision = PolicyDecision.ALLOW if self._filesystem_boundary_allowed(fact.target) else PolicyDecision.REQUIRE_APPROVAL
            elif fact.kind is BoundaryKind.ACCOUNT_ACCESS:
                rules = self._boundary_rules_for_fact(invocation, fact=fact)
                if rules:
                    request = BoundaryPolicyRequest(principal_id=policy_principal, boundary=fact)
                    decision = evaluate_boundary_request(request, rules=rules)
                else:
                    # Default: require approval before any account-scoped operation
                    decision = PolicyDecision.REQUIRE_APPROVAL
            else:
                continue
            decision = self._sentinel_policy.decision_for_risk(_boundary_risk_score(fact), baseline=decision)
            if decision is PolicyDecision.ALLOW:
                continue
            if decision is PolicyDecision.REQUIRE_APPROVAL:
                context = _boundary_approval_context_from_fact(fact)
                if context is None:
                    continue
                approval = self._ensure_boundary_policy_approval_request(invocation, context=context)
                if approval is None:
                    continue
                approval_context = approval.context if isinstance(approval.context, dict) else context
                return self._deny_invocation(
                    invocation,
                    reason="approval_required",
                    error="Approval required for outbound network boundary policy",
                    output={
                        **approval_context,
                        "reason": "approval_required",
                        "requires_approval": True,
                        "approval_id": approval.approval_id,
                    },
                )
            return self._deny_invocation(
                invocation,
                reason="boundary_denied",
                error=f"Boundary denied for tool: {invocation.tool_name}",
                output={"reason": "boundary_denied", "target": fact.target, "boundary_kind": fact.kind.value},
            )
        return None

    def _with_trusted_filesystem_selectors(self, invocation: ToolInvocation) -> ToolInvocation:
        selectors = list(invocation.trusted_filesystem_selectors)
        policy_principal = permission_scope_principal(invocation.principal_id)
        for fact in extract_boundary_facts(invocation):
            if fact.kind is not BoundaryKind.FILESYSTEM_ACCESS:
                continue
            if self._filesystem_boundary_allowed(fact.target, principal_id=invocation.principal_id):
                continue
            rules = self._boundary_rules_for_fact(invocation, fact=fact)
            matching_rules = [
                rule
                for rule in rules
                if rule.selector == fact.target
                or fnmatch(fact.target, rule.selector)
                or (rule.selector.endswith("/*") and fact.target == rule.selector[:-2])
            ]
            if not matching_rules:
                continue
            request = BoundaryPolicyRequest(principal_id=policy_principal, boundary=fact)
            if evaluate_boundary_request(request, rules=matching_rules) is PolicyDecision.ALLOW:
                selectors.extend(rule.selector for rule in matching_rules if rule.mode is PolicyMode.ALLOW)
        trusted = tuple(dict.fromkeys(selectors))
        if trusted == invocation.trusted_filesystem_selectors:
            return invocation
        return ToolInvocation(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            principal_id=invocation.principal_id,
            arguments=dict(invocation.arguments),
            capsule_id=invocation.capsule_id,
            trusted_filesystem_selectors=trusted,
            flow_context=invocation.flow_context,
        )

    def _rewrite_network_denied_result_as_boundary_approval_required(
        self,
        invocation: ToolInvocation,
        result: ToolResult,
    ) -> ToolResult:
        context = _outbound_network_approval_context_from_result(result)
        if context is None:
            return result
        approval = self._ensure_boundary_policy_approval_request(invocation, context=context)
        if approval is None:
            return result
        approval_context = approval.context if isinstance(approval.context, dict) else context
        return ToolResult(
            invocation_id=result.invocation_id,
            tool_name=result.tool_name,
            status="denied",
            output={
                **approval_context,
                "reason": "approval_required",
                "requires_approval": True,
                "approval_id": approval.approval_id,
            },
            error="Approval required for outbound network boundary policy",
        )

    def _record_prompt_injection_scan(self, invocation: ToolInvocation, result: ToolResult) -> None:
        if result.status != "completed":
            return
        scan = scan_tool_output(result.tool_name, result.output)
        if not scan.detected:
            return
        payload = {
            "invocation_id": result.invocation_id,
            "tool_name": result.tool_name,
            "principal_id": invocation.principal_id,
            "capsule_id": invocation.capsule_id,
            "severity": scan.severity,
            "findings": [finding.to_dict() for finding in scan.findings],
        }
        self._store.add_event(make_event("security.prompt_injection.detected", invocation.principal_id, payload))
        self._store.add_audit_record(
            make_audit_record("security.prompt_injection.detected", invocation.principal_id, payload)
        )

    def invoke(self, invocation: ToolInvocation) -> ToolResult:
        invocation = _with_resolved_virtual_workspace_paths(invocation)
        if _tool_blocked_for_principal(invocation.principal_id, invocation.tool_name):
            return self._deny_invocation(
                invocation,
                reason="user_tool_blocked",
                error=f"Tool is blocked for this user: {invocation.tool_name}",
                output={"reason": "user_tool_blocked", "tool_name": invocation.tool_name},
            )
        if invocation.tool_name in self._denied_tool_names:
            return self._deny_invocation(
                invocation,
                reason="capability_denied",
                error=f"Capability denied for tool: {invocation.tool_name}",
                output={"reason": "capability_denied", "denied_tools": list(self._denied_tool_names)},
            )

        if self._allowed_tool_names and invocation.tool_name not in self._allowed_tool_names:
            return self._deny_invocation(
                invocation,
                reason="capability_not_granted",
                error=f"Capability not granted for tool: {invocation.tool_name}",
                output={"reason": "capability_not_granted", "allowed_tools": list(self._allowed_tool_names)},
            )

        scheduled_terminal_denial = self._scheduled_task_terminal_filesystem_preapproval_result(invocation)
        if scheduled_terminal_denial is not None:
            return scheduled_terminal_denial

        destructive_terminal_approval = self._terminal_destructive_preapproval_result(invocation)
        if destructive_terminal_approval is not None:
            return destructive_terminal_approval

        spec = self._registry.get_spec(invocation.tool_name)

        tool_decision = (
            self._sentinel_policy.decision_for_risk(
                _TOOL_RISK_SCORES.get(spec.risk_level, 5),
                baseline=PolicyDecision.REQUIRE_APPROVAL,
            )
            if (spec.requires_approval or self._sentinel_policy.mode.value == "ask_all")
            else PolicyDecision.ALLOW
        )
        if tool_decision is PolicyDecision.REQUIRE_APPROVAL and not self._has_required_tool_grant(invocation):
            approval = self._ensure_tool_approval_request(invocation)
            approval_context = approval.context if isinstance(approval.context, dict) else {}
            return self._deny_invocation(
                invocation,
                reason="approval_required",
                error=f"Approval required for tool: {invocation.tool_name}",
                output={
                    **approval_context,
                    "reason": "approval_required",
                    "requires_approval": True,
                    "approval_id": approval.approval_id,
                },
            )

        preflight_result = self._preflight_boundary_policy_result(invocation)
        if preflight_result is not None:
            return preflight_result
        invocation = self._with_trusted_filesystem_selectors(invocation)

        egress_attempts = _egress_attempts_for_invocation(invocation)
        self._record_detected_egress_attempts(invocation, egress_attempts)

        details = {
            "invocation_id": invocation.invocation_id,
            "tool_name": invocation.tool_name,
            "principal_id": invocation.principal_id,
            "capsule_id": invocation.capsule_id,
            "status": "started",
        }
        self._store.add_event(make_event("tool.invocation.started", invocation.principal_id, details))
        self._store.add_audit_record(make_audit_record("tool.invocation.started", invocation.principal_id, details))

        started_at = perf_counter()
        try:
            result = self._registry.invoke(invocation)
        except Exception as exc:
            result = ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"reason": "handler_exception"},
                error=str(exc),
            )
        duration_ms = (perf_counter() - started_at) * 1000
        result = normalize_tool_result(result)
        result = _with_default_action_receipt(spec=spec, invocation=invocation, result=result)
        if invocation.tool_name == "email_send" and result.status != "denied":
            self._revoke_email_send_review_grants(invocation)
        if invocation.tool_name == "terminal_exec" and result.status != "denied":
            self._revoke_terminal_destructive_grants(invocation)
        self._record_egress_outcome(invocation, result, egress_attempts)
        result = self._rewrite_network_denied_result_as_boundary_approval_required(invocation, result)
        if invocation.tool_name != "terminal_exec" and result.status != "denied":
            active_permits = _matching_active_boundary_permits_for_invocation(self._store, invocation)
            _record_wildcard_boundary_permit_accesses(self._store, invocation, active_permits)
            for permit in active_permits:
                if permit.uses_remaining <= 0:
                    continue
                self._store.add_boundary_permit(consume_boundary_permit_record(permit))
        self._record_prompt_injection_scan(invocation, result)
        completed_details = {
            "invocation_id": result.invocation_id,
            "tool_name": result.tool_name,
            "principal_id": invocation.principal_id,
            "capsule_id": invocation.capsule_id,
            "status": result.status,
        }
        if result.error:
            completed_details["error"] = str(result.error).strip().replace("\n", " ")[:240]
        output = result.output if isinstance(result.output, dict) else {}
        for output_key, detail_key in (
            ("reason", "reason"),
            ("summary", "summary"),
            ("message", "message"),
            ("url", "url"),
            ("path", "path"),
        ):
            value = output.get(output_key)
            if value:
                completed_details[detail_key] = str(value).strip().replace("\n", " ")[:240]
        self._store.add_event(make_event("tool.invocation.completed", invocation.principal_id, completed_details))
        self._store.add_audit_record(make_audit_record("tool.invocation.completed", invocation.principal_id, completed_details))
        try:
            from nullion.builder_routes import build_route_observation

            observation = build_route_observation(
                invocation=invocation,
                result=result,
                duration_ms=duration_ms,
                capability_tags=getattr(spec, "capability_tags", ()) or (),
            )
            if observation is not None:
                add_observation = getattr(self._store, "add_builder_route_observation", None)
                if callable(add_observation):
                    add_observation(observation)
                logger.info(
                    "tool route timing tool=%s source=%s status=%s duration_ms=%.1f reason=%s capsule_id=%s principal_id=%s",
                    observation.get("tool_name"),
                    observation.get("source_domain") or "unknown",
                    observation.get("status"),
                    float(observation.get("duration_ms") or 0.0),
                    observation.get("reason") or "",
                    observation.get("capsule_id") or "",
                    observation.get("principal_id") or "",
                )
        except Exception:
            logger.debug("Builder route observation failed", exc_info=True)
        return result



def _is_within_allowed_root(path: Path, allowed_root: Path) -> bool:
    try:
        path.relative_to(allowed_root)
    except ValueError:
        return False
    return True



def _is_approved_filesystem_path(path: Path, selectors: Iterable[str]) -> bool:
    target = str(path)
    for selector in selectors:
        if not isinstance(selector, str) or not selector:
            continue
        try:
            resolved_selector = str(Path(selector).expanduser().resolve())
        except OSError:
            resolved_selector = str(Path(selector).expanduser())
        if (
            selector == target
            or resolved_selector == target
            or fnmatch(target, selector)
            or fnmatch(target, resolved_selector)
            or (selector.endswith("/*") and target == selector[:-2])
            or (resolved_selector.endswith("/*") and target == resolved_selector[:-2])
        ):
            return True
    return False


_FILE_WALK_PRUNED_DIR_NAMES = frozenset(
    {
        ".git",
        ".hg",
        ".mypy_cache",
        ".pytest_cache",
        ".ruff_cache",
        ".svn",
        ".venv",
        "__pycache__",
        "brave-debug",
        "chrome-debug",
        "chromium-debug",
        "node_modules",
        "venv",
    }
)

_MACOS_PROTECTED_APP_DATA_SUFFIXES = (
    ("library", "application support", "addressbook"),
    ("library", "calendars"),
    ("library", "containers"),
    ("library", "group containers"),
    ("library", "mail"),
    ("library", "messages"),
    ("library", "safari"),
)


def _has_path_suffix(path: Path, suffix: tuple[str, ...]) -> bool:
    parts = tuple(part.lower() for part in path.parts)
    if len(parts) < len(suffix):
        return False
    return any(parts[index : index + len(suffix)] == suffix for index in range(len(parts) - len(suffix) + 1))


def _should_prune_filesystem_walk_dir(path: Path) -> bool:
    name = path.name.lower()
    if name in _FILE_WALK_PRUNED_DIR_NAMES:
        return True
    return any(_has_path_suffix(path, suffix) for suffix in _MACOS_PROTECTED_APP_DATA_SUFFIXES)


def _build_file_read_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_path = invocation.arguments.get("path")
        if not isinstance(raw_path, str) or not raw_path:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: path",
            )
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="File access requires workspace_root or allowed_roots",
            )

        resolved_file = _resolve_local_workspace_file_input(
            raw_path,
            principal_id=invocation.principal_id,
            effective_roots=effective_roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
        path = resolved_file or Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser().resolve()
        if not _path_within_any_root(path, effective_roots) and not _is_approved_filesystem_path(
            path, invocation.trusted_filesystem_selectors
        ):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Path is outside workspace root: {path}",
            )

        if path.is_dir():
            try:
                children = sorted(path.iterdir(), key=lambda child: child.name.lower())
            except OSError as exc:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"path": str(path), "type": "directory"},
                    error=f"Could not list directory: {exc}",
                )
            entry_limit = 100
            entries = []
            for child in children[:entry_limit]:
                entry_type = "directory" if child.is_dir() else "file"
                try:
                    size = None if child.is_dir() else child.stat().st_size
                except OSError:
                    size = None
                entries.append(
                    {
                        "name": child.name,
                        "path": str(child),
                        "type": entry_type,
                        "bytes": size,
                        "media_type": "" if child.is_dir() else (mimetypes.guess_type(child.name)[0] or ""),
                    }
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "path": str(path),
                    "type": "directory",
                    "is_directory": True,
                    "entries": entries,
                    "entry_count": len(children),
                    "file_count": sum(1 for child in children if child.is_file()),
                    "directory_count": sum(1 for child in children if child.is_dir()),
                    "truncated": len(children) > entry_limit,
                },
                error=None,
            )

        archive_details = _archive_format_for_path(path)
        if archive_details is not None:
            try:
                archive_output = _inspect_archive_file(path)
            except FileNotFoundError:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={},
                    error=f"File not found: {path}",
                )
            except (zipfile.BadZipFile, tarfile.TarError, EOFError, lzma.LZMAError, OSError) as exc:
                archive_format = archive_details[0]
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(path),
                        "type": "archive",
                        "archive_format": archive_format,
                    },
                    error=f"Could not inspect archive: {exc}",
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output=archive_output,
                error=None,
            )

        media_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
        if path.suffix.lower() in {".xlsx", ".xlsm"}:
            try:
                from openpyxl import load_workbook

                workbook = load_workbook(path, read_only=True, data_only=False)
                worksheets: list[dict[str, object]] = []
                content_sections: list[str] = []
                workbook_truncated = False
                max_rows = 200
                max_columns = 50
                max_content_chars = 24_000
                for worksheet in workbook.worksheets:
                    rows: list[list[object]] = []
                    worksheet_truncated = bool(
                        int(worksheet.max_row or 0) > max_rows
                        or int(worksheet.max_column or 0) > max_columns
                    )
                    for row in worksheet.iter_rows(
                        min_row=1,
                        max_row=min(int(worksheet.max_row or 0), max_rows),
                        max_col=min(int(worksheet.max_column or 0), max_columns),
                        values_only=True,
                    ):
                        normalized_row: list[object] = []
                        for value in row:
                            if value is None or isinstance(value, (str, int, float, bool)):
                                normalized_row.append(value)
                            else:
                                normalized_row.append(str(value))
                        while normalized_row and normalized_row[-1] is None:
                            normalized_row.pop()
                        rows.append(normalized_row)
                    while rows and not rows[-1]:
                        rows.pop()
                    worksheets.append(
                        {
                            "name": worksheet.title,
                            "rows": rows,
                            "row_count": int(worksheet.max_row or 0),
                            "column_count": int(worksheet.max_column or 0),
                            "truncated": worksheet_truncated,
                        }
                    )
                    lines = [f"Sheet: {worksheet.title}"]
                    lines.extend("\t".join("" if value is None else str(value) for value in row) for row in rows)
                    content_sections.append("\n".join(lines))
                    workbook_truncated = workbook_truncated or worksheet_truncated
                workbook.close()
                content = "\n\n".join(content_sections)
                if len(content) > max_content_chars:
                    content = content[:max_content_chars].rstrip() + "\n[truncated]"
                    workbook_truncated = True
            except FileNotFoundError:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={},
                    error=f"File not found: {path}",
                )
            except ModuleNotFoundError as exc:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"path": str(path), "type": "spreadsheet", "media_type": media_type},
                    error=f"Spreadsheet reading requires the shipped openpyxl runtime dependency: {exc}",
                )
            except Exception as exc:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"path": str(path), "type": "spreadsheet", "media_type": media_type},
                    error=f"Could not extract spreadsheet contents: {exc}",
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "path": str(path),
                    "type": "spreadsheet",
                    "media_type": media_type,
                    "content": content,
                    "sheet_count": len(worksheets),
                    "worksheets": worksheets,
                    "truncated": workbook_truncated,
                    "extraction_method": "openpyxl",
                },
                error=None,
            )
        if path.suffix.lower() == ".pdf" or media_type == "application/pdf":
            try:
                data = path.read_bytes()
                content, page_count, truncated = _extract_pdf_text(data)
            except FileNotFoundError:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={},
                    error=f"File not found: {path}",
                )
            except Exception as exc:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"path": str(path), "type": "pdf", "media_type": media_type},
                    error=f"Could not extract PDF text: {exc}",
                )
            if not content:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(path),
                        "type": "pdf",
                        "media_type": media_type,
                        "page_count": page_count,
                    },
                    error=f"PDF contains no extractable text: {path}",
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "path": str(path),
                    "type": "pdf",
                    "media_type": media_type,
                    "content": content,
                    "page_count": page_count,
                    "truncated": truncated,
                    "extraction_method": "pypdf",
                },
                error=None,
            )

        try:
            content = path.read_text(encoding="utf-8")
        except FileNotFoundError:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"File not found: {path}",
            )
        except UnicodeDecodeError:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "path": str(path),
                    "type": "file",
                    "media_type": media_type,
                },
                error=f"File is not readable as UTF-8 text: {path}",
            )

        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"path": str(path), "content": content},
            error=None,
        )

    return handler



def _build_file_write_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_path = invocation.arguments.get("path")
        raw_content = invocation.arguments.get("content")
        if not isinstance(raw_path, str) or not raw_path:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: path",
            )
        if not isinstance(raw_content, str):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: content",
            )
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="File access requires workspace_root or allowed_roots",
            )

        path = Path(raw_path).expanduser().resolve()
        if _file_write_artifact_path_should_use_workspace_artifacts(
            path,
            principal_id=invocation.principal_id,
            effective_roots=effective_roots,
        ):
            path = _workspace_generated_artifact_path(
                invocation,
                raw_path=raw_path,
                suffix=path.suffix.lower(),
                stem=path.stem or "file",
            )
        if not _path_within_any_root(path, effective_roots) and not _is_approved_filesystem_path(
            path, invocation.trusted_filesystem_selectors
        ):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Path is outside workspace root: {path}",
            )
        if path.suffix.lower() in _TEXT_FILE_WRITE_BLOCKED_EXTENSIONS:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"path": str(path)},
                error=(
                    f"file_write is text-only and cannot create {path.suffix.lower()} artifacts. "
                    "Use a dedicated artifact tool or verified generator for that file type."
                ),
            )
        placeholder_errors = _text_artifact_placeholder_url_errors(raw_content)
        if placeholder_errors:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "path": str(path),
                    "reason": "placeholder_source_urls",
                    "placeholder_urls": placeholder_errors[:10],
                },
                error=(
                    "file_write content contains placeholder/example source URLs. Use verified source URLs "
                    "from tool results or request web/browser scope before writing the artifact."
                ),
            )
        content = raw_content
        embedded_html_images: list[dict[str, object]] = []
        inline_local_html_images = invocation.arguments.get("inline_local_html_images")
        if not isinstance(inline_local_html_images, bool):
            inline_local_html_images = True
        disallow_html_data_images = invocation.arguments.get("disallow_html_data_images") is True
        if path.suffix.lower() in {".html", ".htm"} and inline_local_html_images:
            content, embedded_html_images = _inline_html_local_images(
                raw_content,
                principal_id=invocation.principal_id,
                output_path=path,
                effective_roots=effective_roots,
                trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
            )
        if path.suffix.lower() in {".html", ".htm"} and disallow_html_data_images:
            parser = _HtmlImageSrcParser()
            try:
                parser.feed(content)
                parser.close()
            except Exception:
                parser.sources = []
            data_image_sources = [
                source for source in parser.sources if source.strip().lower().startswith("data:image/")
            ]
            if data_image_sources:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(path),
                        "reason": "html_data_images_disallowed",
                        "data_image_count": len(data_image_sources),
                    },
                    error=(
                        "file_write linked HTML cannot contain data:image sources when "
                        "disallow_html_data_images is true. Reference local sibling image files instead."
                    ),
                )
        html_image_errors = _html_embedded_image_errors(content) if path.suffix.lower() in {".html", ".htm"} else []
        if html_image_errors:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "path": str(path),
                    "reason": "invalid_embedded_html_images",
                    "invalid_images": html_image_errors[:10],
                },
                error=(
                    "file_write HTML contains malformed or truncated embedded image data URIs. "
                    "Use complete local image bytes or browser-collected image artifacts before writing the HTML."
                ),
            )
        if path.suffix.lower() in {".html", ".htm"}:
            normalized_html = content.strip().casefold()
            missing_document_parts = [
                token
                for token in ("<html", "<body", "</body>", "</html>")
                if token not in normalized_html
            ]
            if missing_document_parts:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(path),
                        "reason": "incomplete_html_document",
                        "missing_document_parts": missing_document_parts,
                    },
                    error=(
                        "file_write HTML is incomplete or truncated. Write a complete HTML document with html/body "
                        "opening and closing tags before treating it as a deliverable artifact."
                    ),
                )
        required_content_tokens = _flow_context_required_artifact_content_tokens(
            invocation.flow_context,
            path.suffix.lower(),
        )
        if required_content_tokens:
            from nullion.artifact_validation import missing_required_artifact_content_tokens

            missing_content_tokens = missing_required_artifact_content_tokens(
                content,
                suffix=path.suffix.lower(),
                required_tokens=required_content_tokens,
            )
            if missing_content_tokens:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(path),
                        "reason": "artifact_required_content_missing",
                        "missing_required_content_tokens": list(missing_content_tokens),
                        "required_content_token_count": len(required_content_tokens),
                    },
                    error=(
                        "The artifact is incomplete because required structured content keys are missing. "
                        "Include every missing key in the visible artifact content before delivery."
                    ),
                )
        if (
            path.suffix.lower() in {".html", ".htm"}
            and inline_local_html_images
            and _flow_context_requires_embedded_media_for_extension(invocation.flow_context, path.suffix.lower())
            and _html_embedded_raster_data_image_count(content) <= 0
        ):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "path": str(path),
                    "reason": "html_embedded_raster_image_required",
                },
                error=(
                    "file_write HTML has an embedded-media delivery contract but contains no embedded raster "
                    "image data. Use browser-collected or downloaded PNG/JPG/WebP/GIF assets and reference "
                    "their local paths before writing the self-contained HTML."
                ),
            )

        path.parent.mkdir(parents=True, exist_ok=True)
        binary_image_decoded = False
        image_suffix = path.suffix.lower()
        if image_suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
            candidate = content.strip()
            data_uri_match = re.fullmatch(r"data:image/[^;,]+;base64,(?P<data>[A-Za-z0-9+/=\s]+)", candidate)
            if data_uri_match:
                candidate = data_uri_match.group("data")
            if re.fullmatch(r"[A-Za-z0-9+/=\s]+", candidate or ""):
                try:
                    image_bytes = base64.b64decode(re.sub(r"\s+", "", candidate), validate=True)
                except Exception:
                    image_bytes = b""
                if image_bytes.startswith((b"\x89PNG\r\n\x1a\n", b"\xff\xd8\xff", b"GIF87a", b"GIF89a", b"RIFF")):
                    path.write_bytes(image_bytes)
                    binary_image_decoded = True
        if not binary_image_decoded:
            path.write_text(content, encoding="utf-8")
        bytes_written = path.stat().st_size if path.exists() else len(content.encode("utf-8"))
        output: dict[str, object] = {"path": str(path), "bytes_written": bytes_written}
        if path.suffix.lower() in VALID_ATTACHMENT_EXTENSIONS:
            output["artifact_path"] = str(path)
            output["artifact_paths"] = [str(path)]
        if binary_image_decoded:
            output["binary_image_decoded"] = True
        if embedded_html_images:
            output["embedded_html_images"] = embedded_html_images
        if path.suffix.lower() in {".html", ".htm"}:
            output["html_image_mode"] = "self_contained" if inline_local_html_images else "linked"
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output=output,
            error=None,
        )

    return handler


_TEXT_ARTIFACT_URL_RE = re.compile(r"https?://[^\s\"'<>),;]+", flags=re.IGNORECASE)


def _text_artifact_placeholder_url_errors(content: str) -> list[str]:
    errors: list[str] = []
    for raw_url in _TEXT_ARTIFACT_URL_RE.findall(content or ""):
        url = raw_url.rstrip(".,)")
        parsed = urlparse(url)
        host = (parsed.netloc or "").lower()
        path = unquote(parsed.path or "").lower()
        if host in {"example.com", "example.org", "example.net"} or host.endswith(".example.com"):
            errors.append(url)
            continue
        if re.search(r"(?:^|[-_/])example(?:[-_/]|$)", path):
            errors.append(url)
    return list(dict.fromkeys(errors))


class _HtmlImageSrcParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.sources: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() != "img":
            return
        attrs_dict = {str(key or "").lower(): value for key, value in attrs}
        source = attrs_dict.get("src")
        if isinstance(source, str) and source.strip():
            self.sources.append(source.strip())


_HTML_IMG_SRC_ATTR_RE = re.compile(
    r"(<img\b[^>]*?\bsrc\s*=\s*)([\"'])(.*?)(\2)",
    flags=re.IGNORECASE | re.DOTALL,
)


def _html_local_image_source_path(source: str) -> str | None:
    source_text = str(source or "").strip()
    if not source_text or source_text.startswith("#"):
        return None
    parsed = urlparse(source_text)
    if parsed.scheme.lower() in {"http", "https", "data", "javascript", "mailto", "cid", "blob"}:
        return None
    if parsed.scheme and parsed.scheme.lower() != "file":
        return None
    if parsed.scheme.lower() == "file":
        return unquote(parsed.path or "")
    without_fragment = source_text.split("#", 1)[0].split("?", 1)[0]
    return unquote(without_fragment)


def _image_data_uri_for_file(path: Path) -> str | None:
    path = optimized_embedded_image_path(path)
    mime = mimetypes.guess_type(str(path))[0]
    if mime is None:
        suffix = path.suffix.lower()
        if suffix in {".jpg", ".jpeg"}:
            mime = "image/jpeg"
        elif suffix == ".png":
            mime = "image/png"
        elif suffix == ".gif":
            mime = "image/gif"
        elif suffix == ".webp":
            mime = "image/webp"
        elif suffix == ".svg":
            mime = "image/svg+xml"
    if not mime or not mime.lower().startswith("image/"):
        return None
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{data}"


def _inline_html_local_images(
    content: str,
    *,
    principal_id: str | None,
    output_path: Path,
    effective_roots: tuple[Path, ...],
    trusted_filesystem_selectors: tuple[str, ...],
) -> tuple[str, list[dict[str, object]]]:
    roots = tuple(dict.fromkeys((*effective_roots, output_path.parent.resolve())))
    embedded: list[dict[str, object]] = []

    def replace(match: re.Match[str]) -> str:
        raw_source = unescape(match.group(3).strip())
        source_path = _html_local_image_source_path(raw_source)
        if source_path is None:
            return match.group(0)
        image_path = _resolve_local_workspace_file_input(
            source_path,
            principal_id=principal_id,
            effective_roots=roots,
            trusted_filesystem_selectors=trusted_filesystem_selectors,
        )
        if image_path is None:
            return match.group(0)
        data_uri = _image_data_uri_for_file(image_path)
        if data_uri is None:
            return match.group(0)
        embedded.append({"source": raw_source, "path": str(image_path), "bytes": image_path.stat().st_size})
        return f"{match.group(1)}{match.group(2)}{data_uri}{match.group(4)}"

    return _HTML_IMG_SRC_ATTR_RE.sub(replace, content or ""), embedded


def _html_embedded_image_errors(content: str) -> list[str]:
    parser = _HtmlImageSrcParser()
    try:
        parser.feed(content or "")
    except Exception:
        return []
    errors: list[str] = []
    for index, source in enumerate(parser.sources, start=1):
        if not source.lower().startswith("data:image/"):
            continue
        if "," not in source:
            errors.append(f"image {index}: malformed data URI")
            continue
        header, payload = source.split(",", 1)
        media_type = header[5:].split(";", 1)[0].lower()
        if ";base64" not in header.lower():
            decoded = unquote_to_bytes(payload)
        else:
            compact_payload = re.sub(r"\s+", "", payload)
            try:
                decoded = base64.b64decode(compact_payload, validate=True)
            except Exception:
                errors.append(f"image {index}: invalid base64 data URI")
                continue
        if media_type != "image/svg+xml" and len(decoded) < 1024:
            errors.append(f"image {index}: embedded raster image is too small ({len(decoded)} bytes)")
    return errors


def _html_embedded_raster_data_image_count(content: str) -> int:
    parser = _HtmlImageSrcParser()
    try:
        parser.feed(content or "")
        parser.close()
    except Exception:
        return 0
    count = 0
    for source in parser.sources:
        if not source.lower().startswith("data:image/") or "," not in source:
            continue
        header = source.split(",", 1)[0]
        media_type = header[5:].split(";", 1)[0].lower()
        if media_type in {"image/png", "image/jpeg", "image/jpg", "image/webp", "image/gif"}:
            count += 1
    return count


def _normalized_extension_values(value: object) -> set[str]:
    raw_values = value if isinstance(value, (list, tuple, set)) else (value,)
    normalized: set[str] = set()
    for raw in raw_values:
        text = str(raw or "").strip().lower()
        if not text:
            continue
        if not text.startswith("."):
            text = f".{text}"
        normalized.add(".html" if text == ".htm" else text)
    return normalized


def _flow_context_requires_embedded_media_for_extension(flow_context: object, extension: str) -> bool:
    if not isinstance(flow_context, dict):
        return False
    normalized_extension = ".html" if extension == ".htm" else extension
    required_extensions: set[str] = set()
    for key in ("required_embedded_media_extensions", "embedded_media_artifact_extensions"):
        required_extensions.update(_normalized_extension_values(flow_context.get(key)))
    return normalized_extension in required_extensions


def _flow_context_requires_artifact_delivery_for_extension(flow_context: object, extension: str = "") -> bool:
    if not isinstance(flow_context, dict):
        return False
    if flow_context.get("requires_artifact_delivery"):
        return True
    required_extensions: set[str] = set()
    for key in ("artifact_extensions", "required_artifact_extensions", "requested_artifact_extensions"):
        required_extensions.update(_normalized_extension_values(flow_context.get(key)))
    if not required_extensions:
        return False
    normalized_extension = ".html" if extension == ".htm" else str(extension or "").strip().lower()
    return not normalized_extension or normalized_extension in required_extensions


def _flow_context_required_artifact_content_tokens(flow_context: object, extension: str) -> tuple[str, ...]:
    if not isinstance(flow_context, dict):
        return ()
    required_extensions: set[str] = set()
    for key in ("artifact_extensions", "required_artifact_extensions", "requested_artifact_extensions"):
        required_extensions.update(_normalized_extension_values(flow_context.get(key)))
    normalized_extension = ".html" if extension == ".htm" else str(extension or "").strip().lower()
    if required_extensions and normalized_extension not in required_extensions:
        return ()
    from nullion.artifact_validation import normalize_required_artifact_content_tokens

    return normalize_required_artifact_content_tokens(
        flow_context.get("required_artifact_content_tokens")
    )


def _existing_xlsx_embedded_media_count(path: Path) -> int:
    if path.suffix.lower() != ".xlsx":
        return 0
    try:
        if not path.is_file() or path.stat().st_size <= 0:
            return 0
    except OSError:
        return 0
    try:
        with zipfile.ZipFile(path) as package:
            return sum(1 for member in package.namelist() if member.startswith("xl/media/"))
    except Exception:
        return 0


def _document_output_path(
    invocation: ToolInvocation,
    *,
    raw_path: object,
    title: object,
    roots: tuple[Path, ...],
) -> Path:
    if isinstance(raw_path, str) and raw_path.strip():
        raw_suffix = Path(str(raw_path)).suffix.lower()
        suffix = raw_suffix if raw_suffix in {".docx", ".html", ".htm", ".pdf"} else ".docx"
        return _explicit_output_path_or_generated_artifact_path(
            invocation,
            raw_path=raw_path,
            suffix=suffix,
            stem=_safe_pdf_stem(str(title or "document")),
            roots=roots,
        )
    try:
        from nullion.artifacts import artifact_path_for_generated_workspace_file

        return artifact_path_for_generated_workspace_file(
            principal_id=invocation.principal_id,
            suffix=".docx",
            stem=_safe_pdf_stem(str(title or "document")),
        ).resolve()
    except Exception:
        return (roots[0] / f"{_safe_pdf_stem(str(title or 'document'))}.docx").resolve()


def _document_sections(raw_sections: object) -> tuple[list[dict[str, object]], str | None]:
    if raw_sections is None:
        return [], None
    if not isinstance(raw_sections, list):
        return [], "sections must be a list"
    sections: list[dict[str, object]] = []
    for index, raw_section in enumerate(raw_sections, start=1):
        if not isinstance(raw_section, dict):
            return [], "sections entries must be objects"
        bullets, bullet_error = _coerce_string_list(raw_section.get("bullets"), field=f"sections[{index}].bullets")
        if bullet_error is not None:
            return [], bullet_error
        image_paths, image_error = _coerce_string_list(raw_section.get("image_paths"), field=f"sections[{index}].image_paths")
        if image_error is not None:
            return [], image_error
        screenshot_paths, screenshot_error = _coerce_string_list(
            raw_section.get("screenshot_paths"),
            field=f"sections[{index}].screenshot_paths",
        )
        if screenshot_error is not None:
            return [], screenshot_error
        sections.append(
            {
                "heading": str(raw_section.get("heading") or "").strip(),
                "body": str(raw_section.get("body") or "").strip(),
                "bullets": bullets,
                "image_paths": image_paths,
                "screenshot_paths": screenshot_paths,
            }
        )
    return sections, None


def _document_tables(raw_tables: object) -> tuple[list[dict[str, object]], str | None]:
    if raw_tables is None:
        return [], None
    if not isinstance(raw_tables, list):
        return [], "tables must be a list"
    tables: list[dict[str, object]] = []
    for index, raw_table in enumerate(raw_tables, start=1):
        if not isinstance(raw_table, dict):
            return [], "tables entries must be objects"
        headers, header_error = _coerce_string_list(raw_table.get("headers"), field=f"tables[{index}].headers")
        if header_error is not None:
            return [], header_error
        if not headers:
            return [], f"tables[{index}].headers must include at least one column"
        raw_rows = raw_table.get("rows") or []
        if not isinstance(raw_rows, list):
            return [], f"tables[{index}].rows must be a list"
        rows: list[list[str]] = []
        for row_index, raw_row in enumerate(raw_rows, start=1):
            if not isinstance(raw_row, (list, tuple)):
                return [], f"tables[{index}].rows[{row_index}] must be a list"
            row_values = [str(value if value is not None else "") for value in raw_row[: len(headers)]]
            if len(row_values) < len(headers):
                row_values.extend([""] * (len(headers) - len(row_values)))
            rows.append(row_values)
        tables.append({"headers": headers, "rows": rows})
    return tables, None


_DOCUMENT_URL_RE = re.compile(r"https?://[^\s\"'<>]+", flags=re.IGNORECASE)


def _document_add_hyperlink(paragraph: object, url: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    part = paragraph.part
    relationship_id = part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)

    run = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    properties.append(color)
    properties.append(underline)
    run.append(properties)

    text = OxmlElement("w:t")
    text.text = url
    run.append(text)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _document_add_paragraph_with_links(document: object, text: str, *, style: str | None = None) -> object:
    paragraph = document.add_paragraph(style=style)
    cursor = 0
    for match in _DOCUMENT_URL_RE.finditer(text or ""):
        if match.start() > cursor:
            paragraph.add_run(text[cursor : match.start()])
        raw_url = match.group(0)
        trailing = ""
        while raw_url and raw_url[-1] in ".,);]":
            trailing = raw_url[-1] + trailing
            raw_url = raw_url[:-1]
        if raw_url:
            _document_add_hyperlink(paragraph, raw_url)
        if trailing:
            paragraph.add_run(trailing)
        cursor = match.end()
    if cursor < len(text or ""):
        paragraph.add_run((text or "")[cursor:])
    return paragraph


def _resolve_document_image_paths(
    raw_paths: list[str],
    *,
    roots: tuple[Path, ...],
    invocation: ToolInvocation,
    allow_browser_screenshots: bool = False,
) -> tuple[list[Path], list[str], str | None]:
    resolved: list[Path] = []
    skipped: list[str] = []
    for raw_path in raw_paths:
        if _spreadsheet_http_url(raw_path):
            skipped.append(raw_path)
            return resolved, skipped, (
                "document_create image_paths and screenshot_paths must be local artifact file paths. "
                "Fetch remote image URLs first, then pass the saved local paths."
            )
        image_path = _resolve_local_workspace_file_input(
            raw_path,
            principal_id=invocation.principal_id,
            effective_roots=roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
        if image_path is None:
            candidate = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser()
            if candidate.is_absolute():
                image_path = candidate.resolve()
                if not _path_within_any_root(image_path, roots) and not _is_approved_filesystem_path(
                    image_path,
                    invocation.trusted_filesystem_selectors,
                ):
                    return resolved, skipped, f"Image path is outside workspace root: {image_path}"
            else:
                skipped.append(raw_path)
                continue
        if not image_path.is_file():
            skipped.append(raw_path)
            continue
        if _screenshot_image_rejected(image_path, allow_browser_screenshots=allow_browser_screenshots):
            skipped.append(raw_path)
            continue
        embeddable_path = _embeddable_image_path(image_path)
        if not allow_browser_screenshots and _local_raster_image_too_small(embeddable_path):
            skipped.append(raw_path)
            continue
        resolved.append(embeddable_path)
    return resolved, skipped, None


def _document_html_image_src(image_path: Path, *, output_path: Path, inline_local_html_images: bool) -> str:
    if inline_local_html_images:
        return str(image_path)
    try:
        return os.path.relpath(image_path, output_path.parent)
    except ValueError:
        return str(image_path)


def _write_document_html_artifact(
    *,
    invocation: ToolInvocation,
    output_path: Path,
    title: str,
    paragraphs: list[str],
    sections: list[dict[str, object]],
    tables: list[dict[str, object]],
    image_paths: list[str],
    screenshot_paths: list[str],
    effective_roots: tuple[Path, ...],
) -> ToolResult:
    inline_local_html_images = invocation.arguments.get("inline_local_html_images")
    if not isinstance(inline_local_html_images, bool):
        inline_local_html_images = True
    disallow_html_data_images = invocation.arguments.get("disallow_html_data_images") is True
    embedded_images: list[str] = []
    embedded_screenshots: list[str] = []
    skipped_images: list[str] = []

    def resolve_sources(raw_paths: list[str], *, allow_browser_screenshots: bool, embedded_target: list[str]) -> list[str]:
        resolved, skipped, error = _resolve_document_image_paths(
            raw_paths,
            roots=effective_roots,
            invocation=invocation,
            allow_browser_screenshots=allow_browser_screenshots,
        )
        skipped_images.extend(skipped)
        if error is not None:
            raise ValueError(error)
        sources: list[str] = []
        for image_path in resolved:
            embedded_target.append(str(image_path))
            embed_path = optimized_embedded_image_path(image_path)
            sources.append(_document_html_image_src(embed_path, output_path=output_path, inline_local_html_images=inline_local_html_images))
        return sources

    try:
        top_image_sources = resolve_sources(image_paths, allow_browser_screenshots=False, embedded_target=embedded_images)
        top_screenshot_sources = resolve_sources(
            screenshot_paths,
            allow_browser_screenshots=True,
            embedded_target=embedded_screenshots,
        )
        section_sources: list[tuple[list[str], list[str]]] = []
        for section in sections:
            section_sources.append(
                (
                    resolve_sources(
                        [str(path) for path in section.get("image_paths") or ()],
                        allow_browser_screenshots=False,
                        embedded_target=embedded_images,
                    ),
                    resolve_sources(
                        [str(path) for path in section.get("screenshot_paths") or ()],
                        allow_browser_screenshots=True,
                        embedded_target=embedded_screenshots,
                    ),
                )
            )
    except ValueError as exc:
        return ToolResult(
            invocation.invocation_id,
            invocation.tool_name,
            "failed",
            {
                "path": str(output_path),
                "reason": "artifact_media_inputs_failed",
                "paragraphs": len(paragraphs),
                "sections": len(sections),
                "tables": len(tables),
                "embedded_images": embedded_images,
                "embedded_screenshots": embedded_screenshots,
                "skipped_images": skipped_images,
            },
            str(exc),
        )
    if skipped_images:
        return ToolResult(
            invocation.invocation_id,
            invocation.tool_name,
            "failed",
            {
                "path": str(output_path),
                "reason": "artifact_media_embed_failed",
                "paragraphs": len(paragraphs),
                "sections": len(sections),
                "tables": len(tables),
                "embedded_images": embedded_images,
                "embedded_screenshots": embedded_screenshots,
                "skipped_images": skipped_images,
            },
            (
                "document_create could not embed all requested image/screenshot paths; "
                f"skipped {len(skipped_images)} path(s). Use existing raster image files such as .png, .jpg, or .jpeg."
            ),
        )

    def paragraph_html(text: str) -> str:
        return f"<p>{html.escape(text)}</p>" if text else ""

    chunks = [
        "<!doctype html>",
        "<html>",
        "<head>",
        '<meta charset="utf-8">',
        f"<title>{html.escape(title)}</title>",
        "<style>body{font-family:Arial,sans-serif;margin:2rem;line-height:1.45;color:#111827}img{max-width:100%;height:auto;border-radius:6px;margin:0.5rem 0 1rem}table{border-collapse:collapse;width:100%;margin:1rem 0}td,th{border:1px solid #d1d5db;padding:0.45rem;text-align:left}th{background:#f3f4f6}.section{margin-top:1.25rem}</style>",
        "</head>",
        "<body>",
        f"<h1>{html.escape(title)}</h1>",
    ]
    chunks.extend(paragraph_html(paragraph) for paragraph in paragraphs)
    for source in [*top_image_sources, *top_screenshot_sources]:
        chunks.append(f'<img src="{html.escape(source, quote=True)}" alt="">')
    for table_spec in tables:
        headers = [str(header) for header in table_spec["headers"]]
        chunks.append("<table><thead><tr>")
        chunks.extend(f"<th>{html.escape(header)}</th>" for header in headers)
        chunks.append("</tr></thead><tbody>")
        for row_values in table_spec["rows"]:
            chunks.append("<tr>")
            chunks.extend(f"<td>{html.escape(str(value))}</td>" for value in row_values)
            chunks.append("</tr>")
        chunks.append("</tbody></table>")
    for section, (image_sources, screenshot_sources) in zip(sections, section_sources, strict=False):
        heading = str(section.get("heading") or "").strip()
        body = str(section.get("body") or "").strip()
        chunks.append('<section class="section">')
        if heading:
            chunks.append(f"<h2>{html.escape(heading)}</h2>")
        if body:
            chunks.append(paragraph_html(body))
        bullets = [str(bullet) for bullet in section.get("bullets") or ()]
        if bullets:
            chunks.append("<ul>")
            chunks.extend(f"<li>{html.escape(bullet)}</li>" for bullet in bullets)
            chunks.append("</ul>")
        for source in [*image_sources, *screenshot_sources]:
            chunks.append(f'<img src="{html.escape(source, quote=True)}" alt="">')
        chunks.append("</section>")
    chunks.extend(["</body>", "</html>"])
    content = "\n".join(chunks)
    if inline_local_html_images:
        content, html_embeds = _inline_html_local_images(
            content,
            principal_id=invocation.principal_id,
            output_path=output_path,
            effective_roots=effective_roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
    else:
        html_embeds = []
    if disallow_html_data_images:
        parser = _HtmlImageSrcParser()
        try:
            parser.feed(content)
            parser.close()
        except Exception:
            parser.sources = []
        if any(source.strip().lower().startswith("data:image/") for source in parser.sources):
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {"path": str(output_path), "reason": "html_data_images_disallowed"},
                "document_create linked HTML cannot contain data:image sources when disallow_html_data_images is true.",
            )
    required_content_tokens = _flow_context_required_artifact_content_tokens(
        invocation.flow_context,
        output_path.suffix.lower(),
    )
    if required_content_tokens:
        from nullion.artifact_validation import missing_required_artifact_content_tokens

        missing_content_tokens = missing_required_artifact_content_tokens(
            content,
            suffix=output_path.suffix.lower(),
            required_tokens=required_content_tokens,
        )
        if missing_content_tokens:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "reason": "artifact_required_content_missing",
                    "missing_required_content_tokens": list(missing_content_tokens),
                    "required_content_token_count": len(required_content_tokens),
                },
                (
                    "The artifact is incomplete because required structured content keys are missing. "
                    "Include every missing key in the visible artifact content before delivery."
                ),
            )
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(content, encoding="utf-8")
    output = {
        "path": str(output_path),
        "artifact_path": str(output_path),
        "artifact_paths": [str(output_path)],
        "paragraphs": len(paragraphs),
        "sections": len(sections),
        "tables": len(tables),
        "embedded_images": embedded_images,
        "embedded_screenshots": embedded_screenshots,
        "skipped_images": skipped_images,
        "quality_profile": "html_document_v1",
        "html_image_mode": "self_contained" if inline_local_html_images else "linked",
    }
    if html_embeds:
        output["embedded_html_images"] = html_embeds
    return ToolResult(invocation.invocation_id, invocation.tool_name, "completed", output, None)


def _document_pdf_html_page(
    *,
    title: str,
    paragraphs: list[str],
    sections: list[dict[str, object]],
    tables: list[dict[str, object]],
) -> str:
    chunks = [
        f"<h1>{html.escape(title)}</h1>",
        *[f"<p>{html.escape(paragraph)}</p>" for paragraph in paragraphs if paragraph],
    ]
    for table_spec in tables:
        headers = [str(header) for header in table_spec["headers"]]
        chunks.append("<table><thead><tr>")
        chunks.extend(f"<th>{html.escape(header)}</th>" for header in headers)
        chunks.append("</tr></thead><tbody>")
        for row_values in table_spec["rows"]:
            chunks.append("<tr>")
            chunks.extend(f"<td>{html.escape(str(value))}</td>" for value in row_values)
            chunks.append("</tr>")
        chunks.append("</tbody></table>")
    for section in sections:
        heading = str(section.get("heading") or "").strip()
        body = str(section.get("body") or "").strip()
        bullets = [str(bullet) for bullet in section.get("bullets") or ()]
        chunks.append("<section>")
        if heading:
            chunks.append(f"<h2>{html.escape(heading)}</h2>")
        if body:
            chunks.append(f"<p>{html.escape(body)}</p>")
        if bullets:
            chunks.append("<ul>")
            chunks.extend(f"<li>{html.escape(bullet)}</li>" for bullet in bullets)
            chunks.append("</ul>")
        chunks.append("</section>")
    return "\n".join(chunks)


def _document_create_pdf_arguments(
    *,
    invocation: ToolInvocation,
    output_path: Path,
    title: str,
    paragraphs: list[str],
    sections: list[dict[str, object]],
    tables: list[dict[str, object]],
    image_paths: list[str],
    screenshot_paths: list[str],
) -> dict[str, object]:
    section_image_paths = [
        str(path)
        for section in sections
        for path in (section.get("image_paths") or ())
        if str(path or "").strip()
    ]
    section_screenshot_paths = [
        str(path)
        for section in sections
        for path in (section.get("screenshot_paths") or ())
        if str(path or "").strip()
    ]
    arguments: dict[str, object] = {
        "output_path": str(output_path),
        "title": title,
        "html_pages": [
            _document_pdf_html_page(
                title=title,
                paragraphs=paragraphs,
                sections=sections,
                tables=tables,
            )
        ],
    }
    if image_paths or section_image_paths:
        arguments["image_paths"] = [*image_paths, *section_image_paths]
    if screenshot_paths or section_screenshot_paths:
        arguments["screenshot_paths"] = [*screenshot_paths, *section_screenshot_paths]
    for key in ("page_size",):
        value = invocation.arguments.get(key)
        if value is not None:
            arguments[key] = value
    return arguments


def _build_document_create_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "document_create requires workspace_root or allowed_roots")

        title = str(invocation.arguments.get("title") or "Document").strip() or "Document"
        paragraphs, paragraph_error = _coerce_string_list(invocation.arguments.get("paragraphs"), field="paragraphs")
        if paragraph_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, paragraph_error)
        sections, section_error = _document_sections(invocation.arguments.get("sections"))
        if section_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, section_error)
        tables, table_error = _document_tables(invocation.arguments.get("tables"))
        if table_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, table_error)
        image_paths, image_error = _coerce_string_list(invocation.arguments.get("image_paths"), field="image_paths")
        if image_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, image_error)
        screenshot_paths, screenshot_error = _coerce_string_list(
            invocation.arguments.get("screenshot_paths"),
            field="screenshot_paths",
        )
        if screenshot_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, screenshot_error)
        output_path = _document_output_path(
            invocation,
            raw_path=invocation.arguments.get("output_path"),
            title=title,
            roots=effective_roots,
        )
        if output_path.suffix.lower() == ".pdf":
            pdf_handler = _build_pdf_create_handler(
                workspace_root=resolved_root,
                allowed_roots=resolved_allowed_roots,
                include_principal_workspace=include_principal_workspace,
            )
            pdf_result = pdf_handler(
                ToolInvocation(
                    invocation_id=invocation.invocation_id,
                    tool_name="pdf_create",
                    principal_id=invocation.principal_id,
                    arguments=_document_create_pdf_arguments(
                        invocation=invocation,
                        output_path=output_path,
                        title=title,
                        paragraphs=paragraphs,
                        sections=sections,
                        tables=tables,
                        image_paths=image_paths,
                        screenshot_paths=screenshot_paths,
                    ),
                    capsule_id=invocation.capsule_id,
                    trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
                    flow_context=invocation.flow_context,
                )
            )
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                pdf_result.status,
                {
                    **(pdf_result.output if isinstance(pdf_result.output, dict) else {}),
                    "delegated_tool_name": "pdf_create",
                },
                pdf_result.error,
            )
        if output_path.suffix.lower() not in {".docx", ".html", ".htm"}:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {"path": str(output_path)}, "document_create output_path must end in .docx, .html, .htm, or .pdf")
        if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
            output_path,
            invocation.trusted_filesystem_selectors,
        ):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Path is outside workspace root: {output_path}")
        if output_path.suffix.lower() in {".html", ".htm"}:
            return _write_document_html_artifact(
                invocation=invocation,
                output_path=output_path,
                title=title,
                paragraphs=paragraphs,
                sections=sections,
                tables=tables,
                image_paths=image_paths,
                screenshot_paths=screenshot_paths,
                effective_roots=effective_roots,
            )

        try:
            from docx import Document
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
            from docx.shared import Inches, Pt, RGBColor
        except ModuleNotFoundError as exc:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "reason": "missing_dependency",
                    "dependency_id": "python-docx",
                    "dependency": "python-docx",
                    "package": "python-docx",
                    "requirement": "python-docx>=1.1,<2",
                    "license": "MIT",
                    "install_command": "python -m pip install 'python-docx>=1.1,<2'",
                },
                f"document_create requires python-docx: {exc}",
            )

        document = Document()
        section = document.sections[0]
        section.top_margin = Inches(0.65)
        section.bottom_margin = Inches(0.65)
        section.left_margin = Inches(0.72)
        section.right_margin = Inches(0.72)
        styles = document.styles
        normal_font = styles["Normal"].font
        normal_font.name = "Arial"
        normal_font.size = Pt(10.5)
        heading_1 = styles["Heading 1"].font
        heading_1.name = "Arial"
        heading_1.size = Pt(22)
        heading_1.bold = True
        heading_1.color.rgb = RGBColor(17, 24, 39)
        heading_2 = styles["Heading 2"].font
        heading_2.name = "Arial"
        heading_2.size = Pt(15)
        heading_2.bold = True
        heading_2.color.rgb = RGBColor(31, 41, 55)
        try:
            document.core_properties.title = title
        except Exception:
            pass
        title_paragraph = document.add_heading(title, level=1)
        title_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        for paragraph in paragraphs:
            paragraph_obj = _document_add_paragraph_with_links(document, paragraph)
            paragraph_obj.paragraph_format.space_after = Pt(7)
        for table_spec in tables:
            headers = list(table_spec["headers"])
            table_rows = list(table_spec["rows"])
            table = document.add_table(rows=1, cols=len(headers))
            table.style = "Table Grid"
            for cell, header in zip(table.rows[0].cells, headers):
                cell.text = str(header)
                for run in cell.paragraphs[0].runs:
                    run.bold = True
            for row_values in table_rows:
                cells = table.add_row().cells
                for cell, value in zip(cells, row_values):
                    cell.text = str(value)
            document.add_paragraph()

        embedded_images: list[str] = []
        embedded_screenshots: list[str] = []
        optimized_image_paths: list[str] = []
        skipped_images: list[str] = []

        def add_images(
            raw_paths: list[str],
            *,
            allow_browser_screenshots: bool = False,
            embedded_target: list[str] = embedded_images,
        ) -> None:
            resolved, skipped, error = _resolve_document_image_paths(
                raw_paths,
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=allow_browser_screenshots,
            )
            skipped_images.extend(skipped)
            if error is not None:
                raise ValueError(error)
            for image_path in resolved:
                try:
                    embed_path = optimized_embedded_image_path(image_path)
                    document.add_picture(str(embed_path), width=Inches(5.8))
                    embedded_target.append(str(image_path))
                    if embed_path != image_path:
                        optimized_image_paths.append(str(embed_path))
                except Exception:
                    skipped_images.append(str(image_path))

        try:
            add_images(image_paths)
            add_images(
                screenshot_paths,
                allow_browser_screenshots=True,
                embedded_target=embedded_screenshots,
            )
            for section in sections:
                heading = str(section.get("heading") or "").strip()
                body = str(section.get("body") or "").strip()
                if heading:
                    heading_paragraph = document.add_heading(heading, level=2)
                    heading_paragraph.paragraph_format.space_before = Pt(10)
                    heading_paragraph.paragraph_format.space_after = Pt(4)
                if body:
                    body_paragraph = _document_add_paragraph_with_links(document, body)
                    body_paragraph.paragraph_format.space_after = Pt(7)
                for bullet in section.get("bullets") or ():
                    bullet_paragraph = _document_add_paragraph_with_links(document, str(bullet), style="List Bullet")
                    bullet_paragraph.paragraph_format.space_after = Pt(3)
                add_images([str(path) for path in section.get("image_paths") or ()])
                add_images(
                    [str(path) for path in section.get("screenshot_paths") or ()],
                    allow_browser_screenshots=True,
                    embedded_target=embedded_screenshots,
                )
        except ValueError as exc:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "reason": "artifact_media_inputs_failed",
                    "paragraphs": len(paragraphs),
                    "sections": len(sections),
                    "tables": len(tables),
                    "embedded_images": embedded_images,
                    "embedded_screenshots": embedded_screenshots,
                    "skipped_images": skipped_images,
                },
                str(exc),
            )
        if skipped_images:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "reason": "artifact_media_embed_failed",
                    "paragraphs": len(paragraphs),
                    "sections": len(sections),
                    "tables": len(tables),
                    "embedded_images": embedded_images,
                    "embedded_screenshots": embedded_screenshots,
                    "skipped_images": skipped_images,
                },
                (
                    "document_create could not embed all requested image/screenshot paths; "
                    f"skipped {len(skipped_images)} path(s). Use existing raster image files such as .png, .jpg, or .jpeg."
                ),
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return ToolResult(
            invocation.invocation_id,
            invocation.tool_name,
            "completed",
            {
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "paragraphs": len(paragraphs),
                "sections": len(sections),
                "tables": len(tables),
                "embedded_images": embedded_images,
                "embedded_screenshots": embedded_screenshots,
                "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                "skipped_images": skipped_images,
                "quality_profile": "report_quality_v1",
                "layout_features": [
                    "styled_headings",
                    "readable_margins",
                    "hyperlinked_urls",
                    *(
                        ["verified_media_embeds"]
                        if embedded_images or embedded_screenshots
                        else []
                    ),
                ],
            },
            None,
        )

    return handler


_SPREADSHEET_IMAGE_KEYS = ("image_path", "image_paths", "image")
_SPREADSHEET_SCREENSHOT_KEYS = ("screenshot_path", "screenshot_paths", "screenshot")
_SPREADSHEET_IMAGE_KEY_SET = frozenset(_SPREADSHEET_IMAGE_KEYS)
_SPREADSHEET_SCREENSHOT_KEY_SET = frozenset(_SPREADSHEET_SCREENSHOT_KEYS)
_SPREADSHEET_LINK_KEYS = ("url", "link", "href", "uri", "source_url", "page_url")
_SPREADSHEET_LINK_LABEL_KEYS = ("title", "name", "label", "text", "caption")
_SPREADSHEET_MARKDOWN_LINK_RE = re.compile(r"^\s*\[([^\]\n]{1,200})\]\((https?://[^)\s]+)\)\s*$", flags=re.IGNORECASE)
_SPREADSHEET_AUTOLINK_RE = re.compile(r"^\s*<\s*(https?://[^>\s]+)\s*>\s*$", flags=re.IGNORECASE)
_SPREADSHEET_LOCAL_IMAGE_SUFFIXES = frozenset({".bmp", ".gif", ".jpeg", ".jpg", ".png"})
_SPREADSHEET_EMPTY_REQUIRED_TEXT = frozenset({"", "-", "\u2014"})
_BROWSER_SCREENSHOT_ARTIFACT_NAME_RE = re.compile(r"^screenshot-[0-9a-f]{12}\.png$", re.IGNORECASE)
_SCREENSHOT_ARTIFACT_PARENT_NAMES = frozenset({"artifacts", ".nullion-artifacts"})


def _json_scalar_for_spreadsheet(value: object) -> str | int | float | bool | None:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return json.dumps(value, ensure_ascii=False, sort_keys=True)


def _spreadsheet_http_url(value: object) -> str | None:
    if not isinstance(value, str):
        return None
    candidate = value.strip().strip("`\"'")
    if not candidate:
        return None
    parsed = urlparse(candidate)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        return None
    return candidate


def _spreadsheet_http_image_url(value: object) -> str | None:
    url = _spreadsheet_http_url(value)
    if not url:
        return None
    parsed = urlparse(url)
    if Path(parsed.path).suffix.lower() not in _SPREADSHEET_LOCAL_IMAGE_SUFFIXES:
        return None
    return url


def _spreadsheet_nested_row_url(value: str) -> str | None:
    outer = _spreadsheet_http_url(value)
    if not outer:
        return None
    parsed = urlparse(outer)
    base_url = f"{parsed.scheme}://{parsed.netloc}"
    try:
        decoded_outer = unquote(outer)
    except Exception:
        decoded_outer = outer
    for match in re.finditer(r"https?://", decoded_outer, flags=re.IGNORECASE):
        start = match.start()
        if start <= 0:
            continue
        nested = _spreadsheet_http_url(decoded_outer[start:])
        if nested and nested != outer:
            return nested
    for _key, param_value in parse_qsl(parsed.query, keep_blank_values=True):
        values = [param_value]
        try:
            decoded = unquote(param_value)
        except Exception:
            decoded = param_value
        if decoded not in values:
            values.append(decoded)
        for candidate in values:
            text = str(candidate or "").strip()
            if not text.startswith(("http://", "https://", "/")):
                continue
            nested = _spreadsheet_http_url(urljoin(base_url, text))
            if nested and nested != outer:
                return nested
    return None


def _spreadsheet_canonical_row_hyperlink(value: str) -> str:
    nested = _spreadsheet_nested_row_url(value)
    candidate = nested or value
    parsed = urlparse(candidate)
    if nested:
        parsed = parsed._replace(query="")
    return parsed._replace(fragment="").geturl()


def _spreadsheet_is_aggregate_row_url(url: str) -> bool:
    parsed = urlparse(url)
    path = (parsed.path or "").strip("/")
    path_parts = {part.lower() for part in path.split("/") if part}
    if path_parts.intersection({"search", "s"}):
        return True
    query_keys = {key.lower() for key in dict(parse_qsl(parsed.query, keep_blank_values=True))}
    return bool(query_keys.intersection({"k", "q", "query", "search"}))


def _spreadsheet_cell_value_and_hyperlink(value: object) -> tuple[object, str | None]:
    if isinstance(value, dict):
        link_candidate: str | None = None
        for key in _SPREADSHEET_LINK_KEYS:
            if key in value:
                link_candidate = _spreadsheet_http_url(value.get(key))
                if link_candidate:
                    break
        if link_candidate:
            hyperlink = _spreadsheet_canonical_row_hyperlink(link_candidate)
            display_value: object = link_candidate
            for key in _SPREADSHEET_LINK_LABEL_KEYS:
                label = value.get(key)
                if isinstance(label, str) and label.strip():
                    display_value = label.strip()
                    break
            if display_value == link_candidate:
                display_value = hyperlink
            return _json_scalar_for_spreadsheet(display_value), hyperlink
        return _json_scalar_for_spreadsheet(value), None

    if isinstance(value, str):
        markdown = _SPREADSHEET_MARKDOWN_LINK_RE.fullmatch(value)
        if markdown:
            label = markdown.group(1).strip()
            link = _spreadsheet_http_url(markdown.group(2))
            if link:
                hyperlink = _spreadsheet_canonical_row_hyperlink(link)
                return (label or hyperlink), hyperlink
        autolink = _SPREADSHEET_AUTOLINK_RE.fullmatch(value)
        if autolink:
            link = _spreadsheet_http_url(autolink.group(1))
            if link:
                hyperlink = _spreadsheet_canonical_row_hyperlink(link)
                return hyperlink, hyperlink
        link = _spreadsheet_http_url(value)
        if link:
            hyperlink = _spreadsheet_canonical_row_hyperlink(link)
            return hyperlink, hyperlink

    return _json_scalar_for_spreadsheet(value), None


def _spreadsheet_is_image_key(key: object) -> bool:
    return str(key or "").strip().lower() in _SPREADSHEET_IMAGE_KEY_SET


def _spreadsheet_is_screenshot_key(key: object) -> bool:
    return str(key or "").strip().lower() in _SPREADSHEET_SCREENSHOT_KEY_SET


def _spreadsheet_is_reserved_media_key(key: object) -> bool:
    return _spreadsheet_is_image_key(key) or _spreadsheet_is_screenshot_key(key)


def _spreadsheet_required_text_missing(value: object) -> bool:
    text = str(value if value is not None else "").strip()
    return text.lower() in _SPREADSHEET_EMPTY_REQUIRED_TEXT


def _spreadsheet_columns(raw_columns: object, rows: list[object]) -> list[str]:
    columns = [
        str(column).strip()
        for column in raw_columns or ()
        if str(column or "").strip() and not _spreadsheet_is_reserved_media_key(column)
    ] if isinstance(raw_columns, (list, tuple)) else []
    if columns:
        return list(dict.fromkeys(columns))
    discovered: list[str] = []
    for row in rows:
        if isinstance(row, dict):
            for key in row:
                if _spreadsheet_is_reserved_media_key(key):
                    continue
                column = str(key).strip()
                if column and column not in discovered:
                    discovered.append(column)
        elif isinstance(row, (list, tuple)):
            for index in range(len(row)):
                column = f"Column {index + 1}"
                if column not in discovered:
                    discovered.append(column)
    return discovered or ["Value"]


def _spreadsheet_row_values(row: object, columns: list[str]) -> list[tuple[object, str | None]]:
    if isinstance(row, dict):
        return [_spreadsheet_cell_value_and_hyperlink(row.get(column)) for column in columns]
    if isinstance(row, (list, tuple)):
        values = [_spreadsheet_cell_value_and_hyperlink(value) for value in row]
        if len(values) < len(columns):
            values.extend([(None, None)] * (len(columns) - len(values)))
        return values[: len(columns)]
    return [_spreadsheet_cell_value_and_hyperlink(row), *((None, None),) * max(0, len(columns) - 1)]


def _write_delimited_spreadsheet_rows(
    *,
    output_path: Path,
    columns: list[str],
    rows: list[object],
    delimiter: str,
) -> dict[str, object]:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    row_values_by_row = [_spreadsheet_row_values(row, columns) for row in rows]
    with output_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle, delimiter=delimiter)
        writer.writerow(columns)
        for row_values in row_values_by_row:
            writer.writerow([hyperlink or (value if value is not None else "") for value, hyperlink in row_values])
    return {
        "path": str(output_path),
        "artifact_path": str(output_path),
        "artifact_paths": [str(output_path)],
        "artifact_descriptors": [
            artifact_output_descriptor(output_path, role=ARTIFACT_ROLE_DELIVERABLE, kind="spreadsheet")
        ],
        "rows": len(rows),
        "columns": len(columns),
        "bytes_written": output_path.stat().st_size,
        "format": "tsv" if delimiter == "\t" else "csv",
    }


def _spreadsheet_image_candidates_from_value(value: object) -> list[object]:
    if isinstance(value, (list, tuple)):
        return list(value)
    return [value]


def _spreadsheet_invalid_image_value(value: object) -> str:
    text = _json_scalar_for_spreadsheet(value)
    return str(text if text is not None else value)[:300]


def _spreadsheet_row_media_path(
    row: object,
    fallback: object,
    *,
    is_media_key: Callable[[object], bool],
) -> tuple[str | None, str | None]:
    candidates: list[object] = []
    if isinstance(row, dict):
        for key, value in row.items():
            if is_media_key(key):
                candidates.extend(_spreadsheet_image_candidates_from_value(value))
    candidates.extend(_spreadsheet_image_candidates_from_value(fallback))
    for candidate in candidates:
        if candidate is None:
            continue
        if isinstance(candidate, str):
            if candidate.strip():
                return candidate.strip(), None
            continue
        return None, _spreadsheet_invalid_image_value(candidate)
    return None, None


def _spreadsheet_row_image_path(row: object, fallback: object) -> tuple[str | None, str | None]:
    path, invalid = _spreadsheet_row_media_path(row, fallback, is_media_key=_spreadsheet_is_image_key)
    if path or invalid:
        return path, invalid
    candidates: list[object] = []
    if isinstance(row, dict):
        candidates.extend(row.values())
    elif isinstance(row, (list, tuple)):
        candidates.extend(row)
    for candidate in candidates:
        if not isinstance(candidate, str):
            continue
        text = candidate.strip()
        if not text:
            continue
        remote_image_url = _spreadsheet_http_image_url(text)
        if remote_image_url:
            return remote_image_url, None
        if _spreadsheet_http_url(text):
            continue
        if Path(text).suffix.lower() in _SPREADSHEET_LOCAL_IMAGE_SUFFIXES:
            return text, None
    return None, None


def _spreadsheet_row_screenshot_path(row: object, fallback: object) -> tuple[str | None, str | None]:
    return _spreadsheet_row_media_path(row, fallback, is_media_key=_spreadsheet_is_screenshot_key)


def _spreadsheet_rows_have_formula(rows: list[object], columns: list[str] | None = None) -> bool:
    def value_has_formula(value: object) -> bool:
        if isinstance(value, str):
            return value.strip().startswith("=")
        if isinstance(value, dict):
            return any(value_has_formula(item) for item in value.values())
        if isinstance(value, (list, tuple)):
            return any(value_has_formula(item) for item in value)
        return False

    if columns is not None:
        return any(value_has_formula(value) for row in rows for value, _hyperlink in _spreadsheet_row_values(row, columns))
    return any(value_has_formula(row) for row in rows)


def _spreadsheet_chart_specs(
    raw_charts: object,
    columns: list[str],
    *,
    sheet_columns: Mapping[str, list[str]] | None = None,
    default_sheet_name: str = "Sheet1",
) -> tuple[list[dict[str, str]], str | None]:
    if raw_charts is None:
        return [], None
    if not isinstance(raw_charts, list):
        return [], "charts must be a list"
    chart_specs: list[dict[str, str]] = []
    for index, raw_chart in enumerate(raw_charts, start=1):
        if not isinstance(raw_chart, dict):
            return [], "charts entries must be objects"
        chart_type = str(raw_chart.get("type") or "").strip().lower()
        if chart_type not in {"bar", "line", "pie"}:
            return [], f"charts[{index}].type must be one of bar, line, or pie"
        sheet_name = str(raw_chart.get("sheet_name") or default_sheet_name).strip() or default_sheet_name
        target_columns = columns
        if sheet_columns is not None:
            if sheet_name not in sheet_columns:
                return [], f"charts[{index}].sheet_name must match a spreadsheet worksheet"
            target_columns = sheet_columns[sheet_name]
        column_set = set(target_columns)
        categories_column = str(raw_chart.get("categories_column") or "").strip()
        values_column = str(raw_chart.get("values_column") or "").strip()
        if categories_column not in column_set:
            return [], f"charts[{index}].categories_column must match a spreadsheet column"
        if values_column not in column_set:
            return [], f"charts[{index}].values_column must match a spreadsheet column"
        anchor = str(raw_chart.get("anchor") or "").strip().upper() or f"H{2 + ((index - 1) * 16)}"
        if not re.fullmatch(r"[A-Z]{1,3}[1-9][0-9]{0,6}", anchor):
            return [], f"charts[{index}].anchor must be an Excel cell reference"
        chart_specs.append(
            {
                "type": chart_type,
                "title": str(raw_chart.get("title") or "").strip(),
                "sheet_name": sheet_name,
                "categories_column": categories_column,
                "values_column": values_column,
                "anchor": anchor,
            }
        )
    return chart_specs, None


def _spreadsheet_conditional_format_specs(
    raw_formats: object,
    columns: list[str],
    *,
    sheet_columns: Mapping[str, list[str]] | None = None,
    default_sheet_name: str = "Sheet1",
) -> tuple[list[dict[str, str]], str | None]:
    if raw_formats is None:
        return [], None
    if not isinstance(raw_formats, list):
        return [], "conditional_formats must be a list"
    format_specs: list[dict[str, str]] = []
    for index, raw_format in enumerate(raw_formats, start=1):
        if not isinstance(raw_format, dict):
            return [], "conditional_formats entries must be objects"
        format_type = str(raw_format.get("type") or "").strip().lower()
        if format_type not in {"data_bar", "color_scale"}:
            return [], f"conditional_formats[{index}].type must be one of data_bar or color_scale"
        sheet_name = str(raw_format.get("sheet_name") or default_sheet_name).strip() or default_sheet_name
        target_columns = columns
        if sheet_columns is not None:
            if sheet_name not in sheet_columns:
                return [], f"conditional_formats[{index}].sheet_name must match a spreadsheet worksheet"
            target_columns = sheet_columns[sheet_name]
        column_set = set(target_columns)
        column = str(raw_format.get("column") or "").strip()
        if column not in column_set:
            return [], f"conditional_formats[{index}].column must match a spreadsheet column"
        color = re.sub(r"[^0-9A-Fa-f]", "", str(raw_format.get("color") or "").strip())[:6]
        format_specs.append(
            {"type": format_type, "column": column, "color": color, "sheet_name": sheet_name}
        )
    return format_specs, None


def _spreadsheet_row_has_media_contract(
    row: object,
    fallback: object,
    *,
    is_media_key: Callable[[object], bool],
) -> bool:
    if fallback is not None:
        return True
    if not isinstance(row, dict):
        return False
    return any(is_media_key(key) for key in row)


def _svg_text_preview(svg_path: Path) -> str:
    try:
        root = ElementTree.fromstring(svg_path.read_text(encoding="utf-8", errors="ignore"))
    except Exception:
        return svg_path.stem.replace("_", " ").replace("-", " ").strip()
    texts: list[str] = []
    for element in root.iter():
        if not str(element.tag).lower().endswith("text"):
            continue
        value = "".join(element.itertext()).strip()
        if value:
            texts.append(re.sub(r"\s+", " ", value))
        if len(texts) >= 3:
            break
    return " / ".join(texts).strip() or svg_path.stem.replace("_", " ").replace("-", " ").strip()


def _write_svg_fallback_raster(svg_path: Path) -> Path | None:
    raster_path = svg_path.with_suffix(".png")
    if raster_path.is_file():
        return raster_path
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception:
        return None
    try:
        preview = _svg_text_preview(svg_path)
        image = Image.new("RGB", (640, 360), "#f8fafc")
        draw = ImageDraw.Draw(image)
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Bold.ttf", 30)
            body_font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", 18)
        except Exception:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
        draw.rounded_rectangle((24, 24, 616, 336), radius=28, fill="#ffffff", outline="#cbd5e1", width=3)
        draw.rounded_rectangle((48, 52, 184, 188), radius=24, fill="#dbeafe", outline="#60a5fa", width=3)
        draw.ellipse((84, 86, 148, 150), fill="#2563eb")
        wrapped = textwrap.wrap(preview or "Generated visual", width=34)[:4]
        draw.text((220, 72), "Generated visual", fill="#0f172a", font=title_font)
        y = 126
        for line in wrapped:
            draw.text((220, y), line, fill="#334155", font=body_font)
            y += 28
        draw.line((220, 250, 560, 250), fill="#a78bfa", width=8)
        draw.line((220, 280, 500, 280), fill="#34d399", width=8)
        image.save(raster_path, format="PNG")
        return raster_path
    except Exception:
        return None


def _embeddable_image_path(image_path: Path) -> Path:
    if image_path.suffix.lower() != ".svg":
        return image_path
    for suffix in (".png", ".jpg", ".jpeg"):
        companion = image_path.with_suffix(suffix)
        if companion.is_file():
            return companion
    fallback = _write_svg_fallback_raster(image_path)
    if fallback is not None:
        return fallback
    return image_path


def optimized_embedded_image_path(
    image_path: Path,
    *,
    max_dimension: int = 2048,
    jpeg_quality: int = 95,
) -> Path:
    """Return a visually faithful derivative for embedding into generated artifacts."""
    source_path = _embeddable_image_path(image_path)
    if source_path.suffix.lower() not in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}:
        return source_path
    try:
        source_size = source_path.stat().st_size
    except OSError:
        return source_path
    try:
        from PIL import Image, ImageOps
    except Exception:
        return source_path
    try:
        with Image.open(source_path) as opened:
            image = ImageOps.exif_transpose(opened)
            width, height = image.size
            has_alpha = image.mode in {"RGBA", "LA"} or (
                image.mode == "P" and "transparency" in getattr(image, "info", {})
            )
            if width <= 0 or height <= 0:
                return source_path
            target_width, target_height = width, height
            largest = max(width, height)
            if largest > max_dimension:
                scale = max_dimension / float(largest)
                target_width = max(1, int(width * scale))
                target_height = max(1, int(height * scale))
                image = image.resize((target_width, target_height), Image.Resampling.LANCZOS)
            digest = hashlib.sha256(
                f"{source_path.resolve()}:{source_size}:{source_path.stat().st_mtime_ns}:{max_dimension}:{jpeg_quality}:baseline".encode(
                    "utf-8"
                )
            ).hexdigest()[:16]
            cache_dir = source_path.parent / ".nullion-optimized-media"
            cache_dir.mkdir(parents=True, exist_ok=True)
            output_suffix = ".png" if has_alpha else ".jpg"
            output_path = cache_dir / f"{source_path.stem}-{digest}-embed{output_suffix}"
            if output_path.is_file():
                return output_path
            if has_alpha:
                image.save(output_path, format="PNG", optimize=True)
            else:
                if image.mode != "RGB":
                    image = image.convert("RGB")
                image.save(
                    output_path,
                    format="JPEG",
                    quality=jpeg_quality,
                    subsampling=0,
                    optimize=True,
                    progressive=False,
                )
            if output_path.is_file() and (
                output_path.stat().st_size < source_size
                or (source_size <= 750_000 and source_path.suffix.lower() in {".jpg", ".jpeg"})
            ):
                return output_path
            try:
                output_path.unlink(missing_ok=True)
            except OSError:
                pass
            return source_path
    except Exception:
        return source_path


def _local_raster_image_too_small(image_path: Path, *, min_bytes: int = 128) -> bool:
    if image_path.suffix.lower() not in {".png", ".jpg", ".jpeg"}:
        return False
    try:
        return image_path.stat().st_size < min_bytes
    except OSError:
        return True


def _spreadsheet_numeric_cell_value(value: object) -> float | None:
    if isinstance(value, bool) or value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip()
    if not text:
        return None
    percent = text.endswith("%")
    text = text.strip("%").replace(",", "")
    text = re.sub(r"^[^\d.+-]+", "", text)
    try:
        number = float(text)
    except ValueError:
        return None
    return number / 100.0 if percent else number


def _spreadsheet_formula_cached_values(ws) -> dict[str, float]:
    try:
        from openpyxl.utils.cell import coordinate_to_tuple, range_boundaries
    except ModuleNotFoundError:
        return {}

    def cell_number(reference: str) -> float | None:
        try:
            row_index, column_index = coordinate_to_tuple(reference.replace("$", ""))
            return _spreadsheet_numeric_cell_value(ws.cell(row=row_index, column=column_index).value)
        except Exception:
            return None

    def range_numbers(reference: str) -> list[float]:
        try:
            min_col, min_row, max_col, max_row = range_boundaries(reference.replace("$", ""))
        except Exception:
            number = cell_number(reference)
            return [] if number is None else [number]
        values: list[float] = []
        for row_index in range(min_row, max_row + 1):
            for column_index in range(min_col, max_col + 1):
                number = _spreadsheet_numeric_cell_value(ws.cell(row=row_index, column=column_index).value)
                if number is not None:
                    values.append(number)
        return values

    def evaluate(formula: str) -> float | None:
        expression = formula.strip()
        if expression.startswith("="):
            expression = expression[1:].strip()

        def replace_function(match: re.Match[str]) -> str:
            function_name = match.group(1).upper()
            args = [
                number
                for raw_arg in match.group(2).split(",")
                for number in range_numbers(raw_arg.strip())
            ]
            if not args:
                return "0"
            if function_name == "SUM":
                return str(sum(args))
            if function_name == "AVERAGE":
                return str(sum(args) / len(args))
            if function_name == "MIN":
                return str(min(args))
            if function_name == "MAX":
                return str(max(args))
            if function_name == "COUNT":
                return str(len(args))
            return "0"

        expression = re.sub(
            r"\b(SUM|AVERAGE|MIN|MAX|COUNT)\s*\(([^()]*)\)",
            replace_function,
            expression,
            flags=re.IGNORECASE,
        )
        if ":" in expression:
            return None

        def replace_cell(match: re.Match[str]) -> str:
            number = cell_number(match.group(0))
            return "0" if number is None else str(number)

        expression = re.sub(r"\$?[A-Z]{1,3}\$?[1-9][0-9]{0,6}", replace_cell, expression)
        if not re.fullmatch(r"[0-9eE+\-*/().\s]+", expression):
            return None
        try:
            value = eval(expression, {"__builtins__": {}}, {})
        except Exception:
            return None
        return float(value) if isinstance(value, (int, float)) else None

    cached_values: dict[str, float] = {}
    for row in ws.iter_rows():
        for cell in row:
            value = cell.value
            if isinstance(value, str) and value.startswith("="):
                cached_value = evaluate(value)
                if cached_value is not None:
                    cached_values[cell.coordinate] = cached_value
    return cached_values


def _inject_spreadsheet_cached_formula_values(output_path: Path, cached_values: dict[str, float]) -> None:
    if not cached_values:
        return
    try:
        from openpyxl import load_workbook
    except ModuleNotFoundError:
        return
    try:
        workbook = load_workbook(output_path, read_only=True, data_only=False)
        sheet_names = list(workbook.sheetnames)
        workbook.close()
    except Exception:
        return
    if not sheet_names:
        return
    sheet_index = 1
    worksheet_member = f"xl/worksheets/sheet{sheet_index}.xml"
    namespace = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
    tmp_path = output_path.with_suffix(output_path.suffix + ".tmp")
    try:
        with zipfile.ZipFile(output_path, "r") as source, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as target:
            for member in source.infolist():
                data = source.read(member.filename)
                if member.filename == worksheet_member:
                    root = ElementTree.fromstring(data)
                    for cell_element in root.iter(f"{namespace}c"):
                        reference = cell_element.attrib.get("r")
                        if reference not in cached_values:
                            continue
                        value_element = cell_element.find(f"{namespace}v")
                        if value_element is None:
                            value_element = ElementTree.SubElement(cell_element, f"{namespace}v")
                        cached_value = cached_values[reference]
                        value_element.text = str(int(cached_value)) if cached_value.is_integer() else f"{cached_value:.12g}"
                    data = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
                target.writestr(member, data)
        tmp_path.replace(output_path)
    except Exception:
        try:
            tmp_path.unlink(missing_ok=True)
        except OSError:
            pass


def _is_browser_screenshot_image_path(image_path: Path) -> bool:
    candidate = image_path.expanduser().resolve(strict=False)
    if not _BROWSER_SCREENSHOT_ARTIFACT_NAME_RE.match(candidate.name):
        return False
    return any(part in _SCREENSHOT_ARTIFACT_PARENT_NAMES for part in candidate.parts)


def _screenshot_image_rejected(image_path: Path, *, allow_browser_screenshots: bool) -> bool:
    return not allow_browser_screenshots and _is_browser_screenshot_image_path(image_path)


def _spreadsheet_output_path(
    invocation: ToolInvocation,
    *,
    raw_path: object,
    title: object,
    roots: tuple[Path, ...],
) -> Path:
    if isinstance(raw_path, str) and raw_path.strip():
        stem = re.sub(r"[^A-Za-z0-9._-]+", "-", str(title or "spreadsheet").strip()).strip("-._")
        raw_suffix = Path(str(raw_path)).suffix.lower()
        suffix = raw_suffix if raw_suffix in {".xlsx", ".csv", ".tsv"} else ".xlsx"
        return _explicit_output_path_or_generated_artifact_path(
            invocation,
            raw_path=raw_path,
            suffix=suffix,
            stem=stem or "spreadsheet",
            roots=roots,
        )
    from nullion.artifacts import artifact_path_for_generated_workspace_file

    stem = re.sub(r"[^A-Za-z0-9._-]+", "-", str(title or "spreadsheet").strip()).strip("-._")
    return artifact_path_for_generated_workspace_file(
        principal_id=invocation.principal_id,
        suffix=".xlsx",
        stem=stem or "spreadsheet",
    )


def _spreadsheet_archive_manifest_row_count(output_path: Path) -> int | None:
    if output_path.suffix.lower() not in {".xlsx", ".csv", ".tsv"}:
        return None
    manifest_path = output_path.with_suffix(".json")
    if not manifest_path.is_file():
        return None
    try:
        payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    if not isinstance(payload, dict):
        return None
    raw_count = payload.get("entry_count")
    if isinstance(raw_count, bool) or not isinstance(raw_count, int) or raw_count < 0:
        return None
    return raw_count


def _build_spreadsheet_create_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "File access requires workspace_root or allowed_roots")

        Workbook = None
        WorksheetImage = None
        Font = None
        openpyxl_import_error: ModuleNotFoundError | None = None
        try:
            from openpyxl import Workbook
            from openpyxl.drawing.image import Image as WorksheetImage
            from openpyxl.styles import Font
        except ModuleNotFoundError as exc:
            openpyxl_import_error = exc

        raw_sheets = invocation.arguments.get("sheets")
        sheet_specs: list[dict[str, object]] = []
        if raw_sheets is not None:
            if not isinstance(raw_sheets, list) or not raw_sheets:
                return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "sheets must be a non-empty list")
            for index, raw_sheet in enumerate(raw_sheets, start=1):
                if not isinstance(raw_sheet, dict):
                    return ToolResult(
                        invocation.invocation_id,
                        invocation.tool_name,
                        "failed",
                        {},
                        f"sheets[{index}] must be an object",
                    )
                raw_sheet_rows = raw_sheet.get("rows")
                if not isinstance(raw_sheet_rows, list):
                    return ToolResult(
                        invocation.invocation_id,
                        invocation.tool_name,
                        "failed",
                        {},
                        f"sheets[{index}].rows must be a list",
                    )
                raw_sheet_name = str(raw_sheet.get("sheet_name") or "").strip()
                if not raw_sheet_name:
                    return ToolResult(
                        invocation.invocation_id,
                        invocation.tool_name,
                        "failed",
                        {},
                        f"sheets[{index}].sheet_name must be non-empty",
                    )
                sheet_specs.append(
                    {
                        "sheet_name": raw_sheet_name,
                        "columns": raw_sheet.get("columns"),
                        "rows": list(raw_sheet_rows),
                    }
                )
            rows = list(sheet_specs[0]["rows"])
            primary_columns = sheet_specs[0].get("columns")
            primary_sheet_name = str(sheet_specs[0]["sheet_name"])
        else:
            raw_rows = invocation.arguments.get("rows") or []
            if not isinstance(raw_rows, list):
                return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "rows must be a list")
            rows = list(raw_rows)
            primary_columns = invocation.arguments.get("columns")
            primary_sheet_name = str(invocation.arguments.get("sheet_name") or "Sheet1")
        if not isinstance(rows, list):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "rows must be a list")
        expected_rows_raw = invocation.arguments.get("expected_rows")
        expected_rows: int | None = None
        if expected_rows_raw is not None:
            if isinstance(expected_rows_raw, bool) or not isinstance(expected_rows_raw, int) or expected_rows_raw < 0:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {"reason": "invalid_expected_rows", "expected_rows": expected_rows_raw},
                    "expected_rows must be a non-negative integer when provided",
                )
            expected_rows = expected_rows_raw
            actual_row_count = (
                sum(len(list(sheet_spec.get("rows") or [])) for sheet_spec in sheet_specs)
                if sheet_specs
                else len(rows)
            )
            if actual_row_count != expected_rows:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {
                        "rows": actual_row_count,
                        "expected_rows": expected_rows,
                        "reason": "row_count_mismatch",
                    },
                    "spreadsheet_create rows must match expected_rows when a structured row-count contract is provided.",
                )
        columns = _spreadsheet_columns(primary_columns, rows)
        workbook_sheet_specs = (
            sheet_specs
            if sheet_specs
            else [
                {
                    "sheet_name": primary_sheet_name,
                    "columns": primary_columns,
                    "rows": rows,
                }
            ]
        )
        sheet_columns_by_name = {
            str(sheet_spec.get("sheet_name") or "").strip(): _spreadsheet_columns(
                sheet_spec.get("columns"),
                list(sheet_spec.get("rows") or []),
            )
            for sheet_spec in workbook_sheet_specs
        }
        formulas_required = bool(invocation.arguments.get("formulas_required"))
        workbook_has_formula = any(
            _spreadsheet_rows_have_formula(
                list(sheet_spec.get("rows") or []),
                sheet_columns_by_name.get(str(sheet_spec.get("sheet_name") or "").strip()),
            )
            for sheet_spec in workbook_sheet_specs
        )
        if formulas_required and not workbook_has_formula:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "rows": len(rows),
                    "columns": columns,
                    "reason": "missing_required_formulas",
                },
                "spreadsheet_create requires at least one row value in the selected columns to begin with '=' when formulas_required is true.",
            )
        image_paths = list(invocation.arguments.get("image_paths") or [])
        if not isinstance(image_paths, list):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "image_paths must be a list")
        screenshot_paths = list(invocation.arguments.get("screenshot_paths") or [])
        if not isinstance(screenshot_paths, list):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "screenshot_paths must be a list")
        row_image_specs = [
            _spreadsheet_row_image_path(row, image_paths[index] if index < len(image_paths) else None)
            for index, row in enumerate(rows)
        ]
        row_screenshot_specs = [
            _spreadsheet_row_screenshot_path(row, screenshot_paths[index] if index < len(screenshot_paths) else None)
            for index, row in enumerate(rows)
        ]
        row_image_paths = [path for path, _invalid in row_image_specs]
        row_image_invalid_values = [invalid for _path, invalid in row_image_specs]
        row_screenshot_paths = [path for path, _invalid in row_screenshot_specs]
        row_screenshot_invalid_values = [invalid for _path, invalid in row_screenshot_specs]
        row_has_image_contract = [
            _spreadsheet_row_has_media_contract(
                row,
                image_paths[index] if index < len(image_paths) else None,
                is_media_key=_spreadsheet_is_image_key,
            )
            for index, row in enumerate(rows)
        ]
        row_has_screenshot_contract = [
            _spreadsheet_row_has_media_contract(
                row,
                screenshot_paths[index] if index < len(screenshot_paths) else None,
                is_media_key=_spreadsheet_is_screenshot_key,
            )
            for index, row in enumerate(rows)
        ]
        media_contract_row_indices = {
            index
            for index, has_contract in enumerate(row_has_image_contract)
            if has_contract or row_image_paths[index] or row_image_invalid_values[index]
        } | {
            index
            for index, has_contract in enumerate(row_has_screenshot_contract)
            if has_contract or row_screenshot_paths[index] or row_screenshot_invalid_values[index]
        }
        include_image_column = any(row_image_paths) or any(row_image_invalid_values)
        include_screenshot_column = any(row_screenshot_paths) or any(row_screenshot_invalid_values)
        workbook_columns = [
            *columns,
            *(["Image"] if include_image_column else []),
            *(["Screenshot"] if include_screenshot_column else []),
        ]
        output_path = _spreadsheet_output_path(
            invocation,
            raw_path=invocation.arguments.get("output_path"),
            title=invocation.arguments.get("title"),
            roots=effective_roots,
        )
        if output_path.suffix.lower() not in {".xlsx", ".csv", ".tsv"}:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {"path": str(output_path)},
                "spreadsheet_create output_path must end in .xlsx, .csv, or .tsv",
            )
        if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
            output_path,
            invocation.trusted_filesystem_selectors,
        ):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Path is outside workspace root: {output_path}")
        archive_manifest_rows = _spreadsheet_archive_manifest_row_count(output_path)
        if archive_manifest_rows is not None and len(rows) != archive_manifest_rows:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "rows": len(rows),
                    "expected_rows": archive_manifest_rows,
                    "reason": "archive_manifest_row_count_mismatch",
                    "existing_manifest_path": str(output_path),
                },
                "spreadsheet_create cannot overwrite an archive manifest with a partial row set.",
            )
        chart_specs, chart_error = _spreadsheet_chart_specs(
            invocation.arguments.get("charts"),
            columns,
            sheet_columns=sheet_columns_by_name,
            default_sheet_name=primary_sheet_name,
        )
        if chart_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, chart_error)
        conditional_format_specs, conditional_format_error = _spreadsheet_conditional_format_specs(
            invocation.arguments.get("conditional_formats"),
            columns,
            sheet_columns=sheet_columns_by_name,
            default_sheet_name=primary_sheet_name,
        )
        if conditional_format_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, conditional_format_error)
        if output_path.suffix.lower() in {".csv", ".tsv"}:
            output = _write_delimited_spreadsheet_rows(
                output_path=output_path,
                columns=columns,
                rows=rows,
                delimiter="\t" if output_path.suffix.lower() == ".tsv" else ",",
            )
            return ToolResult(invocation.invocation_id, invocation.tool_name, "completed", output, None)

        row_values_by_row = [_spreadsheet_row_values(row, columns) for row in rows]
        if Workbook is None or WorksheetImage is None or Font is None:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "reason": "missing_dependency",
                    "dependency_id": "openpyxl",
                    "dependency": "openpyxl",
                    "package": "openpyxl",
                    "requirement": "openpyxl>=3.1,<4",
                    "license": "MIT",
                    "install_command": "python -m pip install 'openpyxl>=3.1,<4'",
                    "rich_features_requested": bool(
                        image_paths
                        or screenshot_paths
                        or any(row_image_paths)
                        or any(row_screenshot_paths)
                        or chart_specs
                        or conditional_format_specs
                    ),
                },
                f"spreadsheet_create requires the shipped openpyxl runtime dependency to create .xlsx artifacts: {openpyxl_import_error}",
            )
        missing_image_rows = [
            index + 2
            for index, image_path_text in enumerate(row_image_paths)
            if include_image_column
            and row_has_image_contract[index]
            and not image_path_text
            and not row_image_invalid_values[index]
        ]
        link_column_indices = [
            index
            for index in range(len(columns))
            if any(index < len(row_values) and row_values[index][1] for row_values in row_values_by_row)
        ]
        missing_link_cells = [
            {"row": row_index + 2, "column": columns[column_index]}
            for row_index, row_values in enumerate(row_values_by_row)
            for column_index in link_column_indices
            if row_index in media_contract_row_indices
            and column_index < len(row_values)
            and not row_values[column_index][1]
            and _spreadsheet_required_text_missing(row_values[column_index][0])
        ]
        if missing_image_rows or missing_link_cells:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "rows": len(rows),
                    "columns": len(workbook_columns),
                    "embedded_images": [],
                    "embedded_screenshots": [],
                    "reason": "incomplete_required_row_values",
                    "missing_image_rows": missing_image_rows,
                    "missing_link_cells": missing_link_cells,
                },
                "spreadsheet_create cannot attach a mixed/incomplete table; rows with requested image or hyperlink cells must provide row-specific values.",
            )

        wb = Workbook()
        ws = wb.active
        sheet_name = primary_sheet_name.strip() or "Sheet1"
        ws.title = re.sub(r"[\[\]:*?/\\]", " ", sheet_name)[:31] or "Sheet1"
        ws.append(workbook_columns)
        for cell in ws[1]:
            cell.font = Font(bold=True)

        embedded_images: list[str] = []
        embedded_screenshots: list[str] = []
        optimized_image_paths: list[str] = []
        skipped_images: list[str] = []
        remote_image_urls: list[str] = []
        aggregate_row_links: list[str] = []
        image_column_index = len(columns) + 1 if include_image_column else None
        screenshot_column_index = len(columns) + (1 if include_image_column else 0) + 1 if include_screenshot_column else None

        def add_row_image(
            *,
            raw_path: str,
            row_index: int,
            column_index: int,
            allow_browser_screenshots: bool,
            embedded_target: list[str],
        ) -> None:
            if _spreadsheet_http_url(raw_path):
                remote_image_urls.append(raw_path)
                skipped_images.append(raw_path)
                return
            image_path = _resolve_local_workspace_file_input(
                raw_path,
                principal_id=invocation.principal_id,
                effective_roots=effective_roots,
                trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
            )
            if image_path is None:
                skipped_images.append(raw_path)
                return
            if _screenshot_image_rejected(image_path, allow_browser_screenshots=allow_browser_screenshots):
                skipped_images.append(raw_path)
                return
            image_path = _embeddable_image_path(image_path)
            if not allow_browser_screenshots and _local_raster_image_too_small(image_path):
                skipped_images.append(raw_path)
                return
            try:
                embed_path = optimized_embedded_image_path(image_path)
                image = WorksheetImage(str(embed_path))
                if image.width > 140:
                    scale = 140 / float(image.width)
                    image.width = int(image.width * scale)
                    image.height = int(image.height * scale)
                if image.height > 120:
                    scale = 120 / float(image.height)
                    image.width = int(image.width * scale)
                    image.height = int(image.height * scale)
                ws.add_image(image, f"{ws.cell(row=row_index, column=column_index).coordinate}")
                ws.row_dimensions[row_index].height = max(ws.row_dimensions[row_index].height or 15, 92)
                embedded_target.append(str(image_path))
                if embed_path != image_path:
                    optimized_image_paths.append(str(embed_path))
            except Exception:
                skipped_images.append(raw_path)

        for row_index, row in enumerate(rows, start=2):
            values = row_values_by_row[row_index - 2]
            ws.append(
                [
                    *(value for value, _hyperlink in values),
                    *([""] if include_image_column else []),
                    *([""] if include_screenshot_column else []),
                ]
            )
            for col_index, (_value, hyperlink) in enumerate(values, start=1):
                if hyperlink:
                    if _spreadsheet_is_aggregate_row_url(hyperlink):
                        aggregate_row_links.append(hyperlink)
                    cell = ws.cell(row=row_index, column=col_index)
                    cell.hyperlink = hyperlink
                    cell.font = Font(color="0563C1", underline="single")
            image_path_text = row_image_paths[row_index - 2]
            if image_column_index is None or not image_path_text:
                invalid_image_value = row_image_invalid_values[row_index - 2]
                if invalid_image_value:
                    skipped_images.append(invalid_image_value)
            else:
                add_row_image(
                    raw_path=image_path_text,
                    row_index=row_index,
                    column_index=image_column_index,
                    allow_browser_screenshots=False,
                    embedded_target=embedded_images,
                )
            screenshot_path_text = row_screenshot_paths[row_index - 2]
            if screenshot_column_index is None or not screenshot_path_text:
                invalid_screenshot_value = row_screenshot_invalid_values[row_index - 2]
                if invalid_screenshot_value:
                    skipped_images.append(invalid_screenshot_value)
            else:
                add_row_image(
                    raw_path=screenshot_path_text,
                    row_index=row_index,
                    column_index=screenshot_column_index,
                    allow_browser_screenshots=True,
                    embedded_target=embedded_screenshots,
                )

        created_sheet_names = [ws.title]
        for sheet_index, sheet_spec in enumerate(sheet_specs[1:], start=2):
            extra_rows = list(sheet_spec.get("rows") or [])
            extra_columns = _spreadsheet_columns(sheet_spec.get("columns"), extra_rows)
            extra_sheet_name = re.sub(r"[\[\]:*?/\\]", " ", str(sheet_spec.get("sheet_name") or "").strip())[:31]
            if not extra_sheet_name:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {"reason": "invalid_sheet_name", "sheet_index": sheet_index},
                    f"sheets[{sheet_index}].sheet_name is invalid",
                )
            if extra_sheet_name in created_sheet_names:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {"reason": "duplicate_sheet_name", "sheet_name": extra_sheet_name},
                    "spreadsheet_create requires unique worksheet names",
                )
            extra_ws = wb.create_sheet(extra_sheet_name)
            created_sheet_names.append(extra_ws.title)
            extra_ws.append(extra_columns)
            for cell in extra_ws[1]:
                cell.font = Font(bold=True)
            for extra_row in extra_rows:
                extra_values = _spreadsheet_row_values(extra_row, extra_columns)
                extra_ws.append([value for value, _hyperlink in extra_values])
                current_row = extra_ws.max_row
                for column_index, (_value, hyperlink) in enumerate(extra_values, start=1):
                    if hyperlink:
                        cell = extra_ws.cell(row=current_row, column=column_index)
                        cell.hyperlink = hyperlink
                        cell.font = Font(color="0563C1", underline="single")
            for column_cells in extra_ws.columns:
                header = str(column_cells[0].value or "")
                max_length = max(len(str(cell.value or "")) for cell in column_cells[:100])
                extra_ws.column_dimensions[column_cells[0].column_letter].width = min(
                    max(max_length + 2, len(header) + 2, 12),
                    48,
                )

        existing_embedded_media_count = _existing_xlsx_embedded_media_count(output_path)
        embedded_media_required = _flow_context_requires_embedded_media_for_extension(
            invocation.flow_context,
            output_path.suffix.lower(),
        )
        if (
            embedded_media_required
            and existing_embedded_media_count > 0
            and not embedded_images
            and not embedded_screenshots
            and not skipped_images
        ):
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "reason": "stale_embedded_media_cannot_satisfy_regeneration",
                    "existing_embedded_media_count": existing_embedded_media_count,
                    "required_arguments": ["image_paths", "screenshot_paths"],
                    "message": (
                        "The existing workbook's embedded media is stale for this regeneration. "
                        "Use a current-turn image source tool, then retry spreadsheet_create with the resulting "
                        "local raster path in image_paths or screenshot_paths."
                    ),
                },
                "spreadsheet_create will not reuse stale embedded media for a regenerated visual workbook.",
            )

        if skipped_images:
            reason = "spreadsheet_embed_paths_failed"
            error_message = (
                "spreadsheet_create could not embed all requested image/screenshot paths; "
                f"skipped {len(skipped_images)} path(s). Use existing raster image files such as .png, .jpg, or .jpeg."
            )
            if remote_image_urls:
                reason = "remote_image_paths_not_supported"
                error_message = (
                    "spreadsheet_create image_paths and screenshot_paths must be local artifact file paths. "
                    "Fetch remote image URLs first, then pass the saved local paths."
                )
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "rows": len(rows),
                    "columns": len(workbook_columns),
                    "sheet_names": created_sheet_names,
                    "embedded_images": embedded_images,
                    "embedded_screenshots": embedded_screenshots,
                    "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                    "skipped_images": skipped_images,
                    "reason": reason,
                    "remote_image_urls": remote_image_urls,
                },
                error_message,
            )

        if (
            embedded_media_required
            and not embedded_images
            and not embedded_screenshots
        ):
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "rows": len(rows),
                    "columns": len(workbook_columns),
                    "embedded_images": embedded_images,
                    "embedded_screenshots": embedded_screenshots,
                    "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                    "skipped_images": skipped_images,
                    "artifact_extensions": [output_path.suffix.lower()],
                    "embedded_media_artifact_extensions": [output_path.suffix.lower()],
                    "required_arguments": ["image_paths", "screenshot_paths"],
                    "reason": "artifact_media_required_by_turn_contract",
                },
                (
                    "spreadsheet_create has an embedded-media delivery contract but no local image or "
                    "screenshot paths were embedded. Fetch, generate, or collect local raster image files "
                    "first, then retry with image_paths or screenshot_paths."
                ),
            )

        if aggregate_row_links and len(rows) > 1:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "rows": len(rows),
                    "columns": len(workbook_columns),
                    "embedded_images": embedded_images,
                    "embedded_screenshots": embedded_screenshots,
                    "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                    "reason": "non_row_specific_links",
                    "aggregate_row_links": list(dict.fromkeys(aggregate_row_links)),
                },
                "spreadsheet_create row links must be direct row-specific source/item URLs; aggregate homepage/search/result URLs cannot be used as per-row links.",
            )

        applied_conditional_formats: list[dict[str, str]] = []
        if conditional_format_specs:
            try:
                from openpyxl.formatting.rule import ColorScaleRule, DataBarRule
            except ModuleNotFoundError as exc:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {"reason": "missing_dependency", "dependency": "openpyxl.formatting.rule"},
                    f"spreadsheet_create requires openpyxl conditional formatting support: {exc}",
                )
            for format_spec in conditional_format_specs:
                target_sheet_name = re.sub(r"[\[\]:*?/\\]", " ", format_spec["sheet_name"])[:31] or "Sheet1"
                target_ws = wb[target_sheet_name]
                target_columns = sheet_columns_by_name[format_spec["sheet_name"]]
                target_rows = next(
                    list(sheet_spec.get("rows") or [])
                    for sheet_spec in workbook_sheet_specs
                    if str(sheet_spec.get("sheet_name") or "").strip() == format_spec["sheet_name"]
                )
                if not target_rows:
                    continue
                column_index = target_columns.index(format_spec["column"]) + 1
                column_letter = target_ws.cell(row=1, column=column_index).column_letter
                cell_range = f"{column_letter}2:{column_letter}{len(target_rows) + 1}"
                if format_spec["type"] == "data_bar":
                    rule = DataBarRule(
                        start_type="min",
                        end_type="max",
                        color=format_spec["color"] or "638EC6",
                        showValue=True,
                    )
                else:
                    rule = ColorScaleRule(
                        start_type="min",
                        start_color="F8696B",
                        mid_type="percentile",
                        mid_value=50,
                        mid_color="FFEB84",
                        end_type="max",
                        end_color="63BE7B",
                    )
                target_ws.conditional_formatting.add(cell_range, rule)
                applied_spec = {**format_spec, "range": cell_range}
                if format_spec["sheet_name"] == primary_sheet_name:
                    applied_spec.pop("sheet_name", None)
                applied_conditional_formats.append(applied_spec)

        embedded_charts: list[dict[str, str]] = []
        if chart_specs:
            try:
                from openpyxl.chart import BarChart, LineChart, PieChart, Reference
            except ModuleNotFoundError as exc:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {"reason": "missing_dependency", "dependency": "openpyxl.chart"},
                    f"spreadsheet_create requires openpyxl chart support: {exc}",
                )
            for chart_spec in chart_specs:
                target_sheet_name = re.sub(r"[\[\]:*?/\\]", " ", chart_spec["sheet_name"])[:31] or "Sheet1"
                target_ws = wb[target_sheet_name]
                target_columns = sheet_columns_by_name[chart_spec["sheet_name"]]
                target_rows = next(
                    list(sheet_spec.get("rows") or [])
                    for sheet_spec in workbook_sheet_specs
                    if str(sheet_spec.get("sheet_name") or "").strip() == chart_spec["sheet_name"]
                )
                if not target_rows:
                    continue
                category_column_index = target_columns.index(chart_spec["categories_column"]) + 1
                value_column_index = target_columns.index(chart_spec["values_column"]) + 1
                chart_type = chart_spec["type"]
                if chart_type == "line":
                    chart = LineChart()
                elif chart_type == "pie":
                    chart = PieChart()
                else:
                    chart = BarChart()
                chart.title = chart_spec["title"] or None
                data = Reference(
                    target_ws,
                    min_col=value_column_index,
                    min_row=1,
                    max_row=len(target_rows) + 1,
                )
                categories = Reference(
                    target_ws,
                    min_col=category_column_index,
                    min_row=2,
                    max_row=len(target_rows) + 1,
                )
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(categories)
                if chart_type != "pie":
                    chart.y_axis.title = chart_spec["values_column"]
                    chart.x_axis.title = chart_spec["categories_column"]
                target_ws.add_chart(chart, chart_spec["anchor"])
                applied_chart_spec = dict(chart_spec)
                if chart_spec["sheet_name"] == primary_sheet_name:
                    applied_chart_spec.pop("sheet_name", None)
                embedded_charts.append(applied_chart_spec)

        for column_cells in ws.columns:
            header = str(column_cells[0].value or "")
            max_length = max(len(str(cell.value or "")) for cell in column_cells[:100])
            ws.column_dimensions[column_cells[0].column_letter].width = min(max(max_length + 2, len(header) + 2, 12), 48)
        if image_column_index is not None:
            ws.column_dimensions[ws.cell(row=1, column=image_column_index).column_letter].width = 22
        if screenshot_column_index is not None:
            ws.column_dimensions[ws.cell(row=1, column=screenshot_column_index).column_letter].width = 22

        cached_formula_values = _spreadsheet_formula_cached_values(ws)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(output_path)
        _inject_spreadsheet_cached_formula_values(output_path, cached_formula_values)
        return ToolResult(
            invocation.invocation_id,
            invocation.tool_name,
            "completed",
            {
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "artifact_descriptors": [
                    artifact_output_descriptor(output_path, role=ARTIFACT_ROLE_DELIVERABLE, kind="spreadsheet")
                ],
                "rows": len(rows),
                "columns": len(workbook_columns),
                "sheet_names": created_sheet_names,
                "embedded_images": embedded_images,
                "embedded_screenshots": embedded_screenshots,
                "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                "embedded_charts": embedded_charts,
                "conditional_formats": applied_conditional_formats,
                "formulas_required": formulas_required,
                "skipped_images": skipped_images,
                "remote_image_urls": remote_image_urls,
            },
            None,
        )

    return handler


def _presentation_output_path(
    invocation: ToolInvocation,
    *,
    raw_path: object,
    title: object,
    roots: tuple[Path, ...],
) -> Path:
    if isinstance(raw_path, str) and raw_path.strip():
        return _explicit_output_path_or_generated_artifact_path(
            invocation,
            raw_path=raw_path,
            suffix=".pptx",
            stem=_safe_pdf_stem(str(title or "presentation")),
            roots=roots,
        )
    try:
        from nullion.artifacts import artifact_path_for_generated_workspace_file

        return artifact_path_for_generated_workspace_file(
            principal_id=invocation.principal_id,
            suffix=".pptx",
            stem=_safe_pdf_stem(str(title or "presentation")),
        ).resolve()
    except Exception:
        return (roots[0] / f"{_safe_pdf_stem(str(title or 'presentation'))}.pptx").resolve()


def _presentation_slide_specs(
    raw_slides: object,
    image_paths: list[str],
    screenshot_paths: list[str],
    *,
    title: str,
) -> tuple[list[dict[str, object]], str | None]:
    if raw_slides is not None and not isinstance(raw_slides, list):
        return [], "slides must be a list"
    slides: list[dict[str, object]] = []
    for index, raw_slide in enumerate(raw_slides or [], start=1):
        if not isinstance(raw_slide, dict):
            return [], "slides entries must be objects"
        bullets, bullet_error = _coerce_string_list(raw_slide.get("bullets"), field=f"slides[{index}].bullets")
        if bullet_error is not None:
            return [], bullet_error
        slide_images, image_error = _coerce_string_list(raw_slide.get("image_paths"), field=f"slides[{index}].image_paths")
        if image_error is not None:
            return [], image_error
        slide_screenshots, screenshot_error = _coerce_string_list(
            raw_slide.get("screenshot_paths"),
            field=f"slides[{index}].screenshot_paths",
        )
        if screenshot_error is not None:
            return [], screenshot_error
        slides.append(
            {
                "title": str(raw_slide.get("title") or f"Slide {index}").strip() or f"Slide {index}",
                "body": str(raw_slide.get("body") or "").strip(),
                "bullets": bullets,
                "image_paths": slide_images,
                "screenshot_paths": slide_screenshots,
            }
        )
    had_explicit_slides = bool(slides)
    if not slides:
        media_paths = [*image_paths, *screenshot_paths]
        if media_paths:
            slides = [
                {
                    "title": title or f"Image {index}",
                    "body": "",
                    "bullets": [],
                    "image_paths": [image_path] if index <= len(image_paths) else [],
                    "screenshot_paths": [image_path] if index > len(image_paths) else [],
                }
                for index, image_path in enumerate(media_paths, start=1)
            ]
        else:
            slides = [{"title": title or "Presentation", "body": "", "bullets": [], "image_paths": [], "screenshot_paths": []}]
    if not had_explicit_slides:
        return slides, None
    for target_key, media_paths in (("image_paths", image_paths), ("screenshot_paths", screenshot_paths)):
        if not media_paths or not slides:
            continue
        for index, image_path in enumerate(media_paths):
            target = slides[index] if index < len(slides) else slides[-1]
            target.setdefault(target_key, [])
            cast_images = target[target_key]
            if isinstance(cast_images, list):
                cast_images.append(image_path)
    return slides, None


def _resolve_presentation_image_paths(
    raw_paths: list[str],
    *,
    roots: tuple[Path, ...],
    invocation: ToolInvocation,
    allow_browser_screenshots: bool = False,
) -> tuple[list[Path], list[str], str | None]:
    resolved: list[Path] = []
    skipped: list[str] = []
    for raw_path in raw_paths:
        if _spreadsheet_http_url(raw_path):
            skipped.append(raw_path)
            return resolved, skipped, (
                "presentation_create image_paths and screenshot_paths must be local artifact file paths. "
                "Fetch remote image URLs first, then pass the saved local paths."
            )
        image_path = _resolve_local_workspace_file_input(
            raw_path,
            principal_id=invocation.principal_id,
            effective_roots=roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
        if image_path is None:
            candidate = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser()
            if candidate.is_absolute():
                image_path = candidate.resolve()
                if not _path_within_any_root(image_path, roots) and not _is_approved_filesystem_path(
                    image_path,
                    invocation.trusted_filesystem_selectors,
                ):
                    return resolved, skipped, f"Image path is outside workspace root: {image_path}"
            else:
                skipped.append(raw_path)
                continue
        if not image_path.is_file():
            skipped.append(raw_path)
            continue
        if _screenshot_image_rejected(image_path, allow_browser_screenshots=allow_browser_screenshots):
            skipped.append(raw_path)
            continue
        embeddable_path = _embeddable_image_path(image_path)
        if not allow_browser_screenshots and _local_raster_image_too_small(embeddable_path):
            skipped.append(raw_path)
            continue
        resolved.append(embeddable_path)
    return resolved, skipped, None


def _add_presentation_text(slide, *, title: str, body: str, bullets: list[str], has_images: bool) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(9.1), Inches(0.55))
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    paragraph.text = title[:120]
    paragraph.font.bold = True
    paragraph.font.size = Pt(28)
    paragraph.font.color.rgb = RGBColor(17, 24, 39)

    text_width = Inches(4.25 if has_images else 9.1)
    body_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), text_width, Inches(5.6))
    text_frame = body_box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    first = True
    for line in [body, *bullets]:
        text = str(line or "").strip()
        if not text:
            continue
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        paragraph.text = text[:700]
        paragraph.font.size = Pt(17 if first and body else 15)
        paragraph.font.color.rgb = RGBColor(31, 41, 55)
        paragraph.space_after = Pt(7)
        if not first or text in bullets:
            paragraph.level = 0
        first = False


def _add_presentation_images(slide, image_paths: list[Path]) -> tuple[list[str], list[str], list[str]]:
    from pptx.util import Inches
    from PIL import Image

    embedded: list[str] = []
    optimized: list[str] = []
    failed: list[str] = []
    if not image_paths:
        return embedded, optimized, failed
    max_width = Inches(4.25)
    max_height = Inches(4.85 if len(image_paths) == 1 else 2.25)
    left = Inches(5.25)
    top = Inches(1.18)
    for index, image_path in enumerate(image_paths[:2]):
        image_top = top + Inches(2.45 * index)
        try:
            embed_path = optimized_embedded_image_path(image_path)
            with Image.open(embed_path) as image:
                width_px, height_px = image.size
            if width_px <= 0 or height_px <= 0:
                raise ValueError("image has invalid dimensions")
            scale = min(float(max_width) / float(width_px), float(max_height) / float(height_px))
            picture_width = int(width_px * scale)
            picture_height = int(height_px * scale)
            picture_left = int(left + (max_width - picture_width) / 2)
            picture_top = int(image_top + (max_height - picture_height) / 2)
            slide.shapes.add_picture(str(embed_path), picture_left, picture_top, width=picture_width, height=picture_height)
            embedded.append(str(image_path))
            if embed_path != image_path:
                optimized.append(str(embed_path))
        except Exception:
            failed.append(str(image_path))
    return embedded, optimized, failed


def _build_presentation_create_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "presentation_create requires workspace_root or allowed_roots")
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "reason": "missing_dependency",
                    "dependency_id": "python-pptx",
                    "dependency": "python-pptx",
                    "package": "python-pptx",
                    "requirement": "python-pptx>=1.0,<2",
                    "license": "MIT",
                    "install_command": "python -m pip install 'python-pptx>=1.0,<2'",
                },
                f"presentation_create requires python-pptx: {exc}",
            )

        title = str(invocation.arguments.get("title") or "Nullion presentation").strip() or "Nullion presentation"
        image_paths, image_error = _coerce_string_list(invocation.arguments.get("image_paths"), field="image_paths")
        if image_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, image_error)
        screenshot_paths, screenshot_error = _coerce_string_list(
            invocation.arguments.get("screenshot_paths"),
            field="screenshot_paths",
        )
        if screenshot_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, screenshot_error)
        slides, slide_error = _presentation_slide_specs(
            invocation.arguments.get("slides"),
            image_paths,
            screenshot_paths,
            title=title,
        )
        if slide_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, slide_error)

        output_path = _presentation_output_path(
            invocation,
            raw_path=invocation.arguments.get("output_path"),
            title=title,
            roots=effective_roots,
        )
        if output_path.suffix.lower() != ".pptx":
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {"path": str(output_path)}, "presentation_create output_path must end in .pptx")
        if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
            output_path,
            invocation.trusted_filesystem_selectors,
        ):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Path is outside workspace root: {output_path}")

        deck = Presentation()
        blank_layout = deck.slide_layouts[6]
        embedded_images: list[str] = []
        embedded_screenshots: list[str] = []
        optimized_image_paths: list[str] = []
        skipped_images: list[str] = []
        for slide_spec in slides:
            slide = deck.slides.add_slide(blank_layout)
            slide_images, slide_skipped, image_error = _resolve_presentation_image_paths(
                list(slide_spec.get("image_paths") or []),
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=False,
            )
            skipped_images.extend(slide_skipped)
            if image_error is not None:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {
                        "path": str(output_path),
                        "reason": "artifact_media_inputs_failed",
                        "slide_count": len(slides),
                        "embedded_images": embedded_images,
                        "embedded_screenshots": embedded_screenshots,
                        "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                        "skipped_images": skipped_images,
                    },
                    image_error,
                )
            slide_screenshots, screenshot_skipped, image_error = _resolve_presentation_image_paths(
                list(slide_spec.get("screenshot_paths") or []),
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=True,
            )
            skipped_images.extend(screenshot_skipped)
            if image_error is not None:
                return ToolResult(
                    invocation.invocation_id,
                    invocation.tool_name,
                    "failed",
                    {
                        "path": str(output_path),
                        "reason": "artifact_media_inputs_failed",
                        "slide_count": len(slides),
                        "embedded_images": embedded_images,
                        "embedded_screenshots": embedded_screenshots,
                        "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                        "skipped_images": skipped_images,
                    },
                    image_error,
                )
            slide_media_paths = [*slide_images, *slide_screenshots]
            screenshot_identity = {str(path) for path in slide_screenshots}
            _add_presentation_text(
                slide,
                title=str(slide_spec.get("title") or title),
                body=str(slide_spec.get("body") or ""),
                bullets=[str(item) for item in slide_spec.get("bullets") or []],
                has_images=bool(slide_media_paths),
            )
            slide_embedded, slide_optimized, slide_failed = _add_presentation_images(slide, slide_media_paths)
            optimized_image_paths.extend(slide_optimized)
            for embedded_path in slide_embedded:
                if embedded_path in screenshot_identity:
                    embedded_screenshots.append(embedded_path)
                else:
                    embedded_images.append(embedded_path)
            skipped_images.extend(slide_failed)

        if skipped_images:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {
                    "path": str(output_path),
                    "reason": "artifact_media_embed_failed",
                    "slide_count": len(slides),
                    "embedded_images": embedded_images,
                    "embedded_screenshots": embedded_screenshots,
                    "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                    "skipped_images": skipped_images,
                },
                (
                    "presentation_create could not embed all requested image/screenshot paths; "
                    f"skipped {len(skipped_images)} path(s). Use existing raster image files such as .png, .jpg, or .jpeg."
                ),
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        deck.save(output_path)
        return ToolResult(
            invocation.invocation_id,
            invocation.tool_name,
            "completed",
            {
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "slide_count": len(slides),
                "embedded_images": embedded_images,
                "embedded_screenshots": embedded_screenshots,
                "optimized_image_paths": list(dict.fromkeys(optimized_image_paths)),
                "skipped_images": skipped_images,
                "bytes_written": output_path.stat().st_size,
                "quality_profile": "report_quality_v1",
                "layout_features": [
                    "styled_slide_titles",
                    "readable_text_layout",
                    "aspect_ratio_safe_media",
                    *(
                        ["verified_media_embeds"]
                        if embedded_images or embedded_screenshots
                        else []
                    ),
                ],
            },
            None,
        )

    return handler


_PDF_RENDER_DPI = 300.0
_PDF_JPEG_QUALITY = 95
_PDF_PAGE_SIZES = {
    "letter": (2550, 3300),
    "a4": (2480, 3508),
}


def _pdf_points_to_px(points: float) -> int:
    return max(1, int(round((float(points) * _PDF_RENDER_DPI) / 72.0)))


def _load_pdf_font(*, size_px: int, bold: bool = False):
    from PIL import ImageFont

    candidates = ("DejaVuSans-Bold.ttf", "Arial Bold.ttf", "Arial.ttf") if bold else ("DejaVuSans.ttf", "Arial.ttf")
    for name in candidates:
        try:
            return ImageFont.truetype(name, size=size_px)
        except Exception:
            continue
    return ImageFont.load_default()


def _pdf_default_output_path(invocation: ToolInvocation, *, title: str, roots: tuple[Path, ...]) -> Path:
    try:
        from nullion.artifacts import artifact_path_for_generated_workspace_file

        return artifact_path_for_generated_workspace_file(
            principal_id=invocation.principal_id,
            suffix=".pdf",
            stem=_safe_pdf_stem(title),
        ).resolve()
    except Exception:
        root = roots[0]
        return (root / f"{_safe_pdf_stem(title)}.pdf").resolve()


def _safe_pdf_stem(title: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9._-]+", "-", str(title or "nullion-artifact").strip().lower())
    cleaned = cleaned.strip(".-")
    return cleaned[:48] or "nullion-artifact"


def _coerce_string_list(value: object, *, field: str) -> tuple[list[str], str | None]:
    if value is None:
        return [], None
    if isinstance(value, str):
        text = value.strip()
        return ([text] if text else []), None
    if not isinstance(value, list):
        return [], f"{field} must be a list of strings"
    items = [str(item).strip() for item in value if isinstance(item, str) and str(item).strip()]
    return items, None


def _image_to_pdf_page(path: Path, *, page_size: tuple[int, int]):
    from PIL import Image, ImageOps

    with Image.open(path) as image:
        image = ImageOps.exif_transpose(image)
        if image.mode not in {"RGB", "L"}:
            background = Image.new("RGB", image.size, "white")
            if image.mode in {"RGBA", "LA"}:
                background.paste(image, mask=image.getchannel("A"))
                image = background
            else:
                image = image.convert("RGB")
        else:
            image = image.convert("RGB")
        fitted = ImageOps.contain(image, page_size)
        page = Image.new("RGB", page_size, "white")
        offset = ((page_size[0] - fitted.width) // 2, (page_size[1] - fitted.height) // 2)
        page.paste(fitted, offset)
        return page


def _text_to_pdf_page(text: str, *, title: str, page_size: tuple[int, int]):
    from PIL import Image, ImageDraw

    page = Image.new("RGB", page_size, "#f8fafc")
    draw = ImageDraw.Draw(page)
    body_font = _load_pdf_font(size_px=_pdf_points_to_px(12))
    title_font = _load_pdf_font(size_px=_pdf_points_to_px(18), bold=True)
    meta_font = _load_pdf_font(size_px=_pdf_points_to_px(9))
    margin = _pdf_points_to_px(36)
    header_height = _pdf_points_to_px(46)
    max_width = max(120, page_size[0] - margin * 2)
    body_line_height = max(22, int((draw.textbbox((0, 0), "Ag", font=body_font)[3]) * 1.35))
    title_line_height = max(30, int((draw.textbbox((0, 0), "Ag", font=title_font)[3]) * 1.25))

    def wrap_for_width(paragraph: str) -> list[str]:
        words = [word for word in paragraph.split() if word]
        if not words:
            return [""]
        lines: list[str] = []
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            if draw.textlength(candidate, font=body_font) <= max_width:
                current = candidate
            else:
                lines.append(current)
                if draw.textlength(word, font=body_font) <= max_width:
                    current = word
                else:
                    # Emergency fallback for oversized unbroken tokens.
                    hard_lines = textwrap.wrap(word, width=40) or [word]
                    lines.extend(hard_lines[:-1])
                    current = hard_lines[-1]
        lines.append(current)
        return lines

    draw.rectangle((0, 0, page_size[0], header_height), fill="#111827")
    y = margin // 2
    if title.strip():
        draw.text((margin, y), title.strip()[:140], fill="#ffffff", font=title_font)
    y = header_height + margin
    draw.rounded_rectangle(
        (margin // 2, header_height + margin // 2, page_size[0] - margin // 2, page_size[1] - margin // 2),
        radius=_pdf_points_to_px(8),
        fill="#ffffff",
        outline="#e5e7eb",
        width=max(1, _pdf_points_to_px(1)),
    )
    for paragraph in str(text or "").splitlines() or [""]:
        lines = wrap_for_width(paragraph)
        for line in lines:
            if y > page_size[1] - margin * 2:
                return page
            draw.text((margin, y), line, fill="#111827", font=body_font)
            y += body_line_height
        y += _pdf_points_to_px(6)
    draw.text((margin, page_size[1] - margin), "Generated report", fill="#6b7280", font=meta_font)
    return page


def _pdf_chromium_executable() -> str | None:
    configured = os.environ.get("NULLION_CHROMIUM_EXECUTABLE")
    candidates = [
        configured,
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Brave Browser.app/Contents/MacOS/Brave Browser",
        shutil.which("google-chrome"),
        shutil.which("chromium"),
        shutil.which("chromium-browser"),
        shutil.which("chrome"),
    ]
    for candidate in candidates:
        if candidate and Path(candidate).is_file():
            return str(candidate)
    return None


def _normalize_pdf_report_text(text: str) -> str:
    normalized_lines: list[str] = []
    for raw_line in _strip_pdf_text_emoji(str(text or "")).splitlines():
        line = raw_line.replace("Â·", "•").replace("\u00c2\u00b7", "•")
        line = re.sub(r"^(\s*)(?:[-*•]\s+)?(?:b7|B7)(?=\s+)", r"\1- ", line)
        line = re.sub(r"(?<=\S)\s+(?:b7|B7)\s+(?=\S)", " · ", line)
        normalized_lines.append(line)
    return "\n".join(normalized_lines)


def _pdf_text_html(text: str) -> str:
    normalized = _normalize_pdf_report_text(text)

    def replace_url(match: re.Match[str]) -> str:
        url = match.group(0)
        trailing = ""
        while url and url[-1] in ".,);]":
            trailing = url[-1] + trailing
            url = url[:-1]
        safe_url = html.escape(url, quote=True)
        return f'<a href="{safe_url}">{safe_url}</a>{html.escape(trailing)}'

    def line_html(line: str) -> str:
        escaped = html.escape(line)
        return _TEXT_ARTIFACT_URL_RE.sub(replace_url, escaped)

    lines = [line.strip() for line in normalized.splitlines()]
    blocks: list[str] = []
    index = 0
    while index < len(lines or [""]):
        line = (lines or [""])[index]
        table_html, next_index = _pdf_markdown_table_html(lines, index, line_html=line_html)
        if table_html:
            blocks.append(table_html)
            index = next_index
            continue
        if line.startswith(("- ", "* ")):
            blocks.append(f"<p class=\"bullet\">{line_html(line[2:].strip()) or '&nbsp;'}</p>")
        elif re.match(r"^\d+[\).]\s+", line):
            blocks.append(f"<p class=\"bullet\">{line_html(line)}</p>")
        elif line.endswith(":") and len(line) <= 90:
            blocks.append(f"<h2>{line_html(line[:-1])}</h2>")
        else:
            blocks.append(f"<p>{line_html(line) or '&nbsp;'}</p>")
        index += 1
    return "".join(blocks)


def _pdf_markdown_table_html(
    lines: list[str],
    start_index: int,
    *,
    line_html: Callable[[str], str],
) -> tuple[str | None, int]:
    if start_index + 1 >= len(lines):
        return None, start_index
    header = lines[start_index]
    separator = lines[start_index + 1]
    if "|" not in header or "|" not in separator:
        return None, start_index
    separator_cells = [cell.strip() for cell in separator.strip().strip("|").split("|")]
    if not separator_cells or not all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in separator_cells):
        return None, start_index
    headers = [cell.strip() for cell in header.strip().strip("|").split("|")]
    if len(headers) != len(separator_cells):
        return None, start_index
    rows: list[list[str]] = []
    index = start_index + 2
    while index < len(lines):
        row = lines[index]
        if "|" not in row or not row.strip():
            break
        cells = [cell.strip() for cell in row.strip().strip("|").split("|")]
        if len(cells) != len(headers):
            break
        rows.append(cells)
        index += 1
    if not rows:
        return None, start_index
    thead = "".join(f"<th>{line_html(cell)}</th>" for cell in headers)
    tbody = "".join(
        "<tr>" + "".join(f"<td>{line_html(cell)}</td>" for cell in row) + "</tr>"
        for row in rows
    )
    return f'<table class="report-table"><thead><tr>{thead}</tr></thead><tbody>{tbody}</tbody></table>', index


_PDF_SECTION_MARKER_RE = re.compile(r"^\s*(?:[-*•]\s*)?\d+[\).]\s+\S+")


def _split_single_pdf_report_page_for_media(text_pages: list[str], *, media_count: int) -> tuple[list[str], bool]:
    if media_count <= 1 or len(text_pages) != 1:
        return text_pages, False

    text = str(text_pages[0] or "")
    lines = text.splitlines()
    section_starts = [index for index, line in enumerate(lines) if _PDF_SECTION_MARKER_RE.match(line.strip())]
    if len(section_starts) != media_count:
        return text_pages, False

    preamble = "\n".join(lines[: section_starts[0]]).strip()
    aligned_pages: list[str] = []
    for section_index, start in enumerate(section_starts):
        end = section_starts[section_index + 1] if section_index + 1 < len(section_starts) else len(lines)
        section = "\n".join(lines[start:end]).strip()
        if section_index == 0 and preamble:
            section = f"{preamble}\n\n{section}".strip()
        aligned_pages.append(section)
    if len(aligned_pages) != media_count or any(not page.strip() for page in aligned_pages):
        return text_pages, False
    return aligned_pages, True


def _pdf_image_data_uri(path: Path) -> str:
    path = optimized_embedded_image_path(path)
    mime = mimetypes.guess_type(str(path))[0] or "image/png"
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{data}"


def _optimized_image_paths_for_output(paths: Iterable[Path]) -> list[str]:
    optimized: list[str] = []
    for path in paths:
        embed_path = optimized_embedded_image_path(path)
        if embed_path != path:
            optimized.append(str(embed_path))
    return list(dict.fromkeys(optimized))


def _pdf_page_count(path: Path, *, fallback: int) -> int:
    try:
        from pypdf import PdfReader

        return len(PdfReader(str(path)).pages)
    except Exception:
        return fallback


def _resolve_pdf_image_sources(
    raw_paths: list[str],
    *,
    roots: tuple[Path, ...],
    invocation: ToolInvocation,
    allow_browser_screenshots: bool,
    screenshot_field: str,
) -> tuple[list[Path], str | None]:
    resolved: list[Path] = []
    for raw_path in raw_paths:
        if _spreadsheet_http_url(raw_path):
            return resolved, (
                "pdf image path inputs must be local artifact file paths. "
                "Fetch remote image URLs first, then pass the saved local paths."
            )
        image_path = _resolve_local_workspace_file_input(
            raw_path,
            principal_id=invocation.principal_id,
            effective_roots=roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
        if image_path is None:
            candidate = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser()
            if candidate.is_absolute():
                image_path = candidate.resolve()
                if not _path_within_any_root(image_path, roots) and not _is_approved_filesystem_path(
                    image_path, invocation.trusted_filesystem_selectors
                ):
                    return resolved, f"Image path is outside workspace root: {image_path}"
            else:
                return resolved, f"Image file not found: {raw_path}"
        if not image_path.is_file():
            return resolved, f"Image file not found: {image_path}"
        if _screenshot_image_rejected(image_path, allow_browser_screenshots=allow_browser_screenshots):
            return resolved, f"Browser screenshot artifacts must be supplied through {screenshot_field}."
        resolved.append(image_path)
    return resolved, None


def _save_text_pdf_with_chromium(
    path: Path,
    *,
    title: str,
    text_pages: list[str],
    image_paths: list[Path],
    page_size_name: str,
) -> bool:
    executable = _pdf_chromium_executable()
    if executable is None:
        return False
    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return False

    preserve_text_pages = bool(text_pages) and len(text_pages) < len(image_paths)
    page_count = max(len(text_pages), 1) if preserve_text_pages else max(len(text_pages), len(image_paths), 1)
    page_blocks: list[str] = []
    for index in range(page_count):
        text_html = _pdf_text_html(text_pages[index] if index < len(text_pages) else "")
        page_image_paths: list[Path] = []
        if preserve_text_pages:
            start = index
            end = len(image_paths) if index == page_count - 1 else index + 1
            page_image_paths = image_paths[start:end]
        elif index < len(image_paths):
            page_image_paths = [image_paths[index]]
        image_html = "".join(
            f'<figure><img class="report-image" src="{_pdf_image_data_uri(image_path)}" alt=""></figure>'
            for image_path in page_image_paths
        )
        if len(page_image_paths) > 1:
            image_html = f'<div class="report-image-grid">{image_html}</div>'
        page_blocks.append(
            "<section class=\"page\">"
            "<header>"
            f"<h1>{html.escape(title or 'Report')}</h1>"
            "</header>"
            "<main>"
            f"{image_html}"
            f"<div class=\"text\">{text_html}</div>"
            "</main>"
            "</section>"
        )
    html_doc = (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        f"<title>{html.escape(title or 'PDF')}</title>"
        "<style>"
        "@page{margin:0.45in;} body{font-family:Arial,Helvetica,sans-serif;color:#111827;margin:0;background:#f8fafc;}"
        ".page{break-after:page;page-break-after:always;background:white;border:1px solid #e5e7eb;border-radius:8px;overflow:hidden;min-height:9.9in;}"
        ".page:last-child{break-after:auto;page-break-after:auto;}"
        "header{background:#111827;color:white;padding:18px 22px;} h1{font-size:20px;line-height:1.15;margin:0;}"
        "main{padding:22px;} figure{float:right;margin:0 0 16px 22px;padding:8px;border:1px solid #e5e7eb;border-radius:6px;background:#fff;}"
        ".report-image-grid{float:right;width:2.55in;margin:0 0 16px 22px;display:grid;grid-template-columns:1fr;gap:7px;}"
        ".report-image-grid figure{float:none;margin:0;}"
        ".report-image-grid .report-image{max-height:1.38in;}"
        ".report-image{max-width:2.35in;max-height:2.2in;object-fit:contain;display:block;}"
        "h2{font-size:14px;margin:13px 0 7px;color:#1f2937;} p{font-size:11.5px;line-height:1.45;margin:0 0 8px;}"
        ".bullet{padding-left:14px;text-indent:-10px;} .bullet:before{content:'• ';color:#2563eb;font-weight:bold;}"
        ".report-table{width:100%;border-collapse:collapse;margin:10px 0 14px;font-size:10.5px;line-height:1.3;}"
        ".report-table th{background:#f1f5f9;color:#111827;text-align:left;font-weight:700;}"
        ".report-table th,.report-table td{border:1px solid #d1d5db;padding:5px 6px;vertical-align:top;}"
        "a{color:#0645ad;text-decoration:underline;word-break:break-word;}"
        "</style></head><body>"
        + "".join(page_blocks)
        + "</body></html>"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(headless=True, executable_path=executable)
        try:
            page = browser.new_page()
            page.set_content(html_doc, wait_until="load")
            page.pdf(
                path=str(path),
                format="A4" if page_size_name == "a4" else "Letter",
                print_background=True,
                prefer_css_page_size=False,
            )
        finally:
            browser.close()
    return True


_PDF_HTML_SCRIPT_RE = re.compile(r"(?is)<\s*script\b[^>]*>.*?<\s*/\s*script\s*>")
_PDF_HTML_BODY_RE = re.compile(r"(?is)<\s*body\b[^>]*>(.*?)<\s*/\s*body\s*>")
_PDF_HTML_STYLE_RE = re.compile(r"(?is)<\s*style\b[^>]*>(.*?)<\s*/\s*style\s*>")
_PDF_HTML_IMG_RE = re.compile(r"(?is)<\s*img\b")
_EMOJI_PRESENTATION_RE = re.compile(
    "["
    "\U0001F1E6-\U0001F1FF"
    "\U0001F300-\U0001FAFF"
    "\u2600-\u27BF"
    "\ufe0f"
    "\u200d"
    "]"
)


def _strip_pdf_text_emoji(text: str) -> str:
    return _EMOJI_PRESENTATION_RE.sub("", str(text or ""))


_PDF_HTML_PRINT_SAFETY_SCRIPT = """
() => {
  const risky = [];
  let pageRisk = false;
  for (const el of document.querySelectorAll('.designed-page *')) {
    const style = getComputedStyle(el);
    const page = el.closest('.designed-page');
    const rect = el.getBoundingClientRect();
    const pageRect = page ? page.getBoundingClientRect() : null;
    const bottomAnchored = style.bottom !== 'auto' || style.insetBlockEnd !== 'auto';
    const outOfPage = pageRect && (rect.bottom > pageRect.bottom - 2 || rect.top < pageRect.top - 2);
    if (page && page.scrollHeight > page.clientHeight + 4) {
      pageRisk = true;
    }
    if (
      style.position === 'fixed' ||
      style.position === 'sticky' ||
      (style.position === 'absolute' && (bottomAnchored || outOfPage))
    ) {
      el.dataset.nullionFlowSafe = '1';
      risky.push(el);
      pageRisk = true;
    }
  }
  if (pageRisk) {
    document.body.classList.add('nullion-print-flow-safe');
  }
  return risky.length;
}
"""


_PDF_HTML_MEDIA_REPAIR_SCRIPT = """
(sources) => {
  if (!Array.isArray(sources) || !sources.length) {
    return 0;
  }
  let repaired = 0;
  let sourceIndex = 0;
  for (const img of Array.from(document.images || [])) {
    const broken = !img.getAttribute('src') || img.naturalWidth === 0 || img.naturalHeight === 0;
    if (!broken) {
      continue;
    }
    img.src = sources[sourceIndex % sources.length];
    img.dataset.nullionRepairedMedia = '1';
    sourceIndex += 1;
    repaired += 1;
  }
  return repaired;
}
"""


_PDF_HTML_PRINT_SCALE_SCRIPT = """
(limitPx) => {
  const limit = Number(limitPx) || 0;
  if (!limit) {
    return 0.92;
  }
  let maxHeight = 0;
  for (const page of Array.from(document.querySelectorAll('.designed-page'))) {
    const inner = page.querySelector(':scope > .designed-page-inner') || page;
    const rect = inner.getBoundingClientRect();
    maxHeight = Math.max(maxHeight, inner.scrollHeight || 0, rect.height || 0);
  }
  if (!maxHeight) {
    return 0.92;
  }
  if (maxHeight > limit * 0.98) {
    return Math.max(0.55, Math.min(0.88, (limit / maxHeight) * 0.96));
  }
  return Math.max(0.55, Math.min(0.92, (limit / maxHeight) * 0.985));
}
"""


_PDF_HTML_APPLY_PRINT_SCALE_SCRIPT = """
(scale) => {
  const ratio = Number(scale) || 1;
  if (ratio >= 0.995) {
    return 0;
  }
  let scaled = 0;
  for (const page of Array.from(document.querySelectorAll('.designed-page'))) {
    const inner = page.querySelector(':scope > .designed-page-inner') || page;
    const originalHeight = Math.max(inner.scrollHeight || 0, inner.getBoundingClientRect().height || 0);
    inner.style.zoom = String(ratio);
    inner.dataset.nullionPrintScale = String(Math.round(ratio * 1000) / 1000);
    if (originalHeight) {
      page.style.minHeight = `${Math.ceil(originalHeight * ratio)}px`;
    }
    page.style.overflow = 'visible';
    scaled += 1;
  }
  return scaled;
}
"""


def _pdf_html_page_contract_dimensions(page_size_name: str) -> tuple[float, float, float]:
    page_width = 8.27 if page_size_name == "a4" else 8.5
    page_height = 11.69 if page_size_name == "a4" else 11.0
    margin = 0.42
    return max(1.0, page_width - (margin * 2)), max(1.0, page_height - (margin * 2)), margin


def _repair_pdf_html_media_sources(page: object, media_sources: list[str]) -> int:
    if not media_sources:
        return 0
    evaluate = getattr(page, "evaluate", None)
    if not callable(evaluate):
        return 0
    try:
        result = evaluate(_PDF_HTML_MEDIA_REPAIR_SCRIPT, media_sources)
    except Exception:
        return 0
    return result if isinstance(result, int) else 0


def _apply_pdf_html_print_safety(page: object) -> int:
    evaluate = getattr(page, "evaluate", None)
    if not callable(evaluate):
        return 0
    try:
        result = evaluate(_PDF_HTML_PRINT_SAFETY_SCRIPT)
    except Exception:
        return 0
    return result if isinstance(result, int) else 0


def _pdf_html_print_scale(page: object, *, page_size_name: str) -> float:
    evaluate = getattr(page, "evaluate", None)
    if not callable(evaluate):
        return 0.92
    _, content_height_inches, _ = _pdf_html_page_contract_dimensions(page_size_name)
    limit_px = max(1, int(round(content_height_inches * 96)))
    try:
        result = evaluate(_PDF_HTML_PRINT_SCALE_SCRIPT, limit_px)
    except Exception:
        return 0.92
    if isinstance(result, (int, float)):
        return max(0.55, min(0.92, float(result)))
    return 0.92


def _apply_pdf_html_print_scale(page: object, scale: float) -> int:
    if scale >= 0.995:
        return 0
    evaluate = getattr(page, "evaluate", None)
    if not callable(evaluate):
        return 0
    try:
        result = evaluate(_PDF_HTML_APPLY_PRINT_SCALE_SCRIPT, scale)
    except Exception:
        return 0
    return result if isinstance(result, int) else 0


def _pdf_html_fragment(raw_html: str) -> tuple[str, list[str]]:
    cleaned = _PDF_HTML_SCRIPT_RE.sub("", _strip_pdf_text_emoji(str(raw_html or "")))
    styles = [match.group(1) for match in _PDF_HTML_STYLE_RE.finditer(cleaned)]
    body_match = _PDF_HTML_BODY_RE.search(cleaned)
    if body_match:
        fragment = body_match.group(1)
    else:
        fragment = re.sub(r"(?is)<!doctype[^>]*>", "", cleaned)
        fragment = re.sub(r"(?is)</?\s*(?:html|head|body)\b[^>]*>", "", fragment)
        fragment = _PDF_HTML_STYLE_RE.sub("", fragment)
    return fragment.strip() or "<p></p>", styles


def _save_html_pdf_with_chromium(
    path: Path,
    *,
    title: str,
    html_pages: list[str],
    image_paths: list[Path] | tuple[Path, ...] = (),
    screenshot_paths: list[Path] | tuple[Path, ...] = (),
    page_size_name: str,
) -> bool:
    executable = _pdf_chromium_executable()
    if executable is None:
        return False
    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return False

    page_fragments: list[str] = []
    page_blocks: list[str] = []
    custom_styles: list[str] = []
    for raw_page in html_pages:
        fragment, styles = _pdf_html_fragment(raw_page)
        custom_styles.extend(styles)
        page_fragments.append(fragment)
    media_blocks: list[str] = []
    media_paths = [*image_paths, *screenshot_paths]
    content_width_inches, content_height_inches, page_margin_inches = _pdf_html_page_contract_dimensions(page_size_name)
    media_data: list[tuple[int, Path, str, str]] = []
    for index, media_path in enumerate(media_paths, start=1):
        media_src = _pdf_image_data_uri(media_path)
        media_label = "Screenshot" if media_path in screenshot_paths else "Image"
        media_data.append((index, media_path, media_src, media_label))
        media_blocks.append(
            "<figure>"
            f'<img src="{html.escape(media_src, quote=True)}" alt="{html.escape(media_label)} {index}">'
            f"<figcaption>{html.escape(media_label)} {index}: {html.escape(media_path.name)}</figcaption>"
            "</figure>"
        )
    if media_data and page_fragments and not any(_PDF_HTML_IMG_RE.search(fragment) for fragment in page_fragments):
        index, media_path, media_src, media_label = media_data[0]
        primary_media = (
            '<figure class="primary-media">'
            f'<img src="{html.escape(media_src, quote=True)}" alt="{html.escape(media_label)} {index}">'
            f"<figcaption>{html.escape(media_label)} {index}: {html.escape(media_path.name)}</figcaption>"
            "</figure>"
        )
        page_fragments[0] = f"{primary_media}{page_fragments[0]}"
    page_blocks = [f'<section class="designed-page"><div class="designed-page-inner">{fragment}</div></section>' for fragment in page_fragments]
    if media_blocks and not page_fragments:
        page_blocks.append(
            '<section class="designed-page media-appendix"><div class="designed-page-inner">'
            "<h2>Embedded Visual Assets</h2>"
            '<div class="media-grid">'
            + "".join(media_blocks)
            + "</div></div></section>"
        )
    html_doc = (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        f"<title>{html.escape(title or 'PDF')}</title>"
        "<style>"
        f"@page{{margin:{page_margin_inches:.2f}in;}}"
        f":root{{--nullion-pdf-content-width:{content_width_inches:.2f}in;--nullion-pdf-content-height:{content_height_inches:.2f}in;}}"
        "body{font-family:Arial,Helvetica,sans-serif;color:#111827;margin:0;background:#f8fafc;}"
        ".designed-page{break-after:page;page-break-after:always;min-height:9.95in;background:white;}"
        ".designed-page:last-child{break-after:auto;page-break-after:auto;}"
        "h1,h2,h3,p,figure,table,.designed-page-inner>*{break-inside:avoid;page-break-inside:avoid;}"
        "img{max-width:100%;height:auto;} table{width:100%;border-collapse:collapse;}"
        "th,td{border:1px solid #d1d5db;padding:6px 8px;vertical-align:top;} th{background:#f1f5f9;text-align:left;}"
        ".media-appendix h2{margin:0 0 18px;font-size:24px;}.media-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px;}"
        ".media-grid figure{margin:0;border:1px solid #e5e7eb;padding:10px;background:#f9fafb;break-inside:avoid;}"
        ".media-grid img{width:100%;max-height:3.4in;object-fit:contain;background:white;}.media-grid figcaption{font-size:10px;color:#4b5563;margin-top:6px;}"
        ".primary-media{float:right;width:38%;margin:0 0 14px 20px;border:1px solid #e5e7eb;padding:8px;background:#f9fafb;break-inside:avoid;}"
        ".primary-media img{width:100%;max-height:2.25in;object-fit:cover;background:white;}.primary-media figcaption{font-size:9px;color:#4b5563;margin-top:5px;}"
        ".designed-page-inner{transform-origin:top left;}"
        "a{color:#0645ad;text-decoration:underline;word-break:break-word;}"
        "</style>"
        + "".join(f"<style>{style}</style>" for style in custom_styles)
        + "<style>"
        ".designed-page{box-sizing:border-box!important;width:var(--nullion-pdf-content-width)!important;"
        "max-width:var(--nullion-pdf-content-width)!important;min-height:var(--nullion-pdf-content-height)!important;"
        "margin:0 auto!important;overflow:visible!important;}"
        ".designed-page-inner{box-sizing:border-box!important;width:100%!important;max-width:100%!important;"
        "min-height:auto!important;overflow:visible!important;}"
        ".designed-page-inner>:first-child{box-sizing:border-box!important;width:100%!important;max-width:100%!important;"
        "height:auto!important;min-height:auto!important;overflow:visible!important;}"
        ".designed-page-inner img,.designed-page-inner svg,.designed-page-inner canvas{max-width:100%!important;}"
        "body.nullion-print-flow-safe .designed-page{height:auto!important;min-height:auto!important;overflow:visible!important;}"
        "body.nullion-print-flow-safe .designed-page [data-nullion-flow-safe='1']{position:static!important;inset:auto!important;"
        "top:auto!important;right:auto!important;bottom:auto!important;left:auto!important;transform:none!important;width:auto!important;"
        "max-width:100%!important;margin:14px 0!important;}"
        "</style>"
        + "</head><body>"
        + "".join(page_blocks)
        + "</body></html>"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(headless=True, executable_path=executable)
        try:
            page = browser.new_page()
            page.set_content(html_doc, wait_until="load")
            _repair_pdf_html_media_sources(page, [media_src for _, _, media_src, _ in media_data])
            _apply_pdf_html_print_safety(page)
            print_scale = _pdf_html_print_scale(page, page_size_name=page_size_name)
            applied_print_scale = _apply_pdf_html_print_scale(page, print_scale)
            page.pdf(
                path=str(path),
                format="A4" if page_size_name == "a4" else "Letter",
                print_background=True,
                prefer_css_page_size=False,
                scale=0.98 if applied_print_scale else print_scale,
            )
        finally:
            browser.close()
    return True


def _build_pdf_pages(
    *,
    image_paths: list[str],
    screenshot_paths: list[str],
    text_pages: list[str],
    page_size: tuple[int, int],
    roots: tuple[Path, ...],
    invocation: ToolInvocation,
    title: str,
    screenshot_field: str = "screenshot_paths",
) -> tuple[list[object], list[str], list[str], str | None]:
    pages = []
    source_images: list[str] = []
    source_screenshots: list[str] = []

    def add_image_pages(raw_paths: list[str], *, allow_browser_screenshots: bool, target: list[str]) -> str | None:
        for raw_path in raw_paths:
            if _spreadsheet_http_url(raw_path):
                return (
                    "pdf image path inputs must be local artifact file paths. "
                    "Fetch remote image URLs first, then pass the saved local paths."
                )
            image_path = _resolve_local_workspace_file_input(
                raw_path,
                principal_id=invocation.principal_id,
                effective_roots=roots,
                trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
            )
            if image_path is None:
                candidate = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser()
                if candidate.is_absolute():
                    image_path = candidate.resolve()
                    if not _path_within_any_root(image_path, roots) and not _is_approved_filesystem_path(
                        image_path, invocation.trusted_filesystem_selectors
                    ):
                        return f"Image path is outside workspace root: {image_path}"
                else:
                    return f"Image file not found: {raw_path}"
            if not image_path.is_file():
                return f"Image file not found: {image_path}"
            if _screenshot_image_rejected(image_path, allow_browser_screenshots=allow_browser_screenshots):
                return f"Browser screenshot artifacts must be supplied through {screenshot_field}."
            try:
                pages.append(_image_to_pdf_page(optimized_embedded_image_path(image_path), page_size=page_size))
            except Exception as exc:
                return f"Could not load image file {image_path}: {exc}"
            target.append(str(image_path))
        return None

    page_error = add_image_pages(
        image_paths,
        allow_browser_screenshots=False,
        target=source_images,
    )
    if page_error is not None:
        return pages, source_images, source_screenshots, page_error
    page_error = add_image_pages(
        screenshot_paths,
        allow_browser_screenshots=True,
        target=source_screenshots,
    )
    if page_error is not None:
        return pages, source_images, source_screenshots, page_error
    for text in text_pages:
        pages.append(_text_to_pdf_page(text, title=title, page_size=page_size))
    return pages, source_images, source_screenshots, None


def _save_pdf_pages(path: Path, pages: list[object], *, title: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    pages[0].save(
        path,
        "PDF",
        save_all=True,
        append_images=pages[1:],
        resolution=_PDF_RENDER_DPI,
        quality=_PDF_JPEG_QUALITY,
        subsampling=0,
        optimize=True,
        title=title or None,
    )


def _build_pdf_create_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="pdf_create requires workspace_root or allowed_roots",
            )

        image_paths, image_error = _coerce_string_list(invocation.arguments.get("image_paths"), field="image_paths")
        if image_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, image_error)
        screenshot_paths, screenshot_error = _coerce_string_list(
            invocation.arguments.get("screenshot_paths"),
            field="screenshot_paths",
        )
        if screenshot_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, screenshot_error)
        text_pages, text_error = _coerce_string_list(invocation.arguments.get("text_pages"), field="text_pages")
        if text_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, text_error)
        html_pages, html_error = _coerce_string_list(invocation.arguments.get("html_pages"), field="html_pages")
        if html_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, html_error)
        if not image_paths and not screenshot_paths and not text_pages and not html_pages:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="pdf_create requires at least one image_paths, screenshot_paths, text_pages, or html_pages entry",
            )

        title = str(invocation.arguments.get("title") or "Nullion PDF").strip()
        raw_page_size = str(invocation.arguments.get("page_size") or "letter").strip().lower()
        page_size = _PDF_PAGE_SIZES.get(raw_page_size)
        if page_size is None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Unsupported page_size: {raw_page_size}",
            )

        raw_output_path = invocation.arguments.get("output_path")
        if isinstance(raw_output_path, str) and raw_output_path.strip():
            output_path = _explicit_output_path_or_generated_artifact_path(
                invocation,
                raw_path=raw_output_path,
                suffix=".pdf",
                stem=_safe_pdf_stem(title),
                roots=effective_roots,
            )
        else:
            output_path = _pdf_default_output_path(invocation, title=title, roots=effective_roots)

        if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
            output_path, invocation.trusted_filesystem_selectors
        ):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Output path is outside workspace root: {output_path}",
            )

        if html_pages:
            inlined_html_pages: list[str] = []
            embedded_html_images: list[dict[str, object]] = []
            unresolved_local_html_images: list[str] = []
            for html_page in html_pages:
                inlined_page, embedded = _inline_html_local_images(
                    html_page,
                    principal_id=invocation.principal_id,
                    output_path=output_path,
                    effective_roots=effective_roots,
                    trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
                )
                embedded_html_images.extend(embedded)
                parser = _HtmlImageSrcParser()
                try:
                    parser.feed(inlined_page)
                    parser.close()
                except Exception:
                    parser.sources = []
                for source in parser.sources:
                    if _html_local_image_source_path(source) is not None:
                        unresolved_local_html_images.append(source)
                inlined_html_pages.append(inlined_page)
            if unresolved_local_html_images:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(output_path),
                        "reason": "html_pdf_unresolved_local_images",
                        "unresolved_image_sources": list(dict.fromkeys(unresolved_local_html_images))[:10],
                    },
                    error=(
                        "pdf_create HTML pages contain local image references that could not be resolved. "
                        "Use local artifact image paths or pass image_paths separately."
                    ),
                )
            source_images, image_source_error = _resolve_pdf_image_sources(
                image_paths,
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=False,
                screenshot_field="screenshot_paths",
            )
            if image_source_error is not None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(output_path),
                        "reason": "artifact_media_inputs_failed",
                        "image_paths": image_paths,
                        "screenshot_paths": screenshot_paths,
                        "source_image_paths": [str(path) for path in source_images],
                        "source_screenshot_paths": [],
                    },
                    error=image_source_error,
                )
            source_screenshots, screenshot_source_error = _resolve_pdf_image_sources(
                screenshot_paths,
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=True,
                screenshot_field="screenshot_paths",
            )
            if screenshot_source_error is not None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(output_path),
                        "reason": "artifact_media_inputs_failed",
                        "image_paths": image_paths,
                        "screenshot_paths": screenshot_paths,
                        "source_image_paths": [str(path) for path in source_images],
                        "source_screenshot_paths": [str(path) for path in source_screenshots],
                    },
                    error=screenshot_source_error,
                )
            try:
                if _save_html_pdf_with_chromium(
                    output_path,
                    title=title,
                    html_pages=inlined_html_pages,
                    image_paths=source_images,
                    screenshot_paths=source_screenshots,
                    page_size_name=raw_page_size,
                ):
                    size_bytes = output_path.stat().st_size if output_path.exists() else 0
                    html_embedded_image_paths = [
                        str(item["path"])
                        for item in embedded_html_images
                        if isinstance(item.get("path"), str)
                    ]
                    has_media = bool(source_images or source_screenshots or html_embedded_image_paths)
                    actual_page_count = _pdf_page_count(output_path, fallback=max(len(html_pages), 1) + (1 if has_media else 0))
                    layout_features = [
                        "custom_html_pages",
                        "styled_tables",
                        "clickable_links",
                        "readable_text_layout",
                    ]
                    if has_media:
                        layout_features.append("verified_media_embeds")
                    return ToolResult(
                        invocation_id=invocation.invocation_id,
                        tool_name=invocation.tool_name,
                        status="completed",
                        output={
                            "path": str(output_path),
                            "artifact_path": str(output_path),
                            "artifact_paths": [str(output_path)],
                            "bytes_written": size_bytes,
                            "page_count": actual_page_count,
                            "html_pages": len(html_pages),
                            "source_image_paths": [str(path) for path in source_images],
                            "source_screenshot_paths": [str(path) for path in source_screenshots],
                            "optimized_image_paths": _optimized_image_paths_for_output([*source_images, *source_screenshots]),
                            "embedded_html_images": embedded_html_images,
                            "text_layer": True,
                            "quality_profile": "designed_pdf_v1",
                            "layout_features": layout_features,
                        },
                        error=None,
                    )
            except Exception as exc:
                logger.info("html_pdf_chromium_failed: %s", exc, exc_info=True)
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"path": str(output_path), "reason": "html_pdf_render_failed"},
                error="PDF creation failed while rendering html_pages.",
            )

        if text_pages:
            raw_media_alignment = str(invocation.arguments.get("media_alignment") or "preserve_text_pages").strip().lower()
            if raw_media_alignment not in {"auto", "align_pages", "preserve_text_pages"}:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={},
                    error=f"Unsupported media_alignment: {raw_media_alignment}",
                )
            source_images, image_source_error = _resolve_pdf_image_sources(
                image_paths,
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=False,
                screenshot_field="screenshot_paths",
            )
            if image_source_error is not None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(output_path),
                        "reason": "artifact_media_inputs_failed",
                        "image_paths": image_paths,
                        "screenshot_paths": screenshot_paths,
                        "source_image_paths": [str(path) for path in source_images],
                        "source_screenshot_paths": [],
                    },
                    error=image_source_error,
                )
            source_screenshots, screenshot_source_error = _resolve_pdf_image_sources(
                screenshot_paths,
                roots=effective_roots,
                invocation=invocation,
                allow_browser_screenshots=True,
                screenshot_field="screenshot_paths",
            )
            if screenshot_source_error is not None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "path": str(output_path),
                        "reason": "artifact_media_inputs_failed",
                        "image_paths": image_paths,
                        "screenshot_paths": screenshot_paths,
                        "source_image_paths": [str(path) for path in source_images],
                        "source_screenshot_paths": [str(path) for path in source_screenshots],
                    },
                    error=screenshot_source_error,
                )
            media_count = len(source_images) + len(source_screenshots)
            text_pages_for_render, media_alignment_applied = _split_single_pdf_report_page_for_media(
                text_pages,
                media_count=media_count,
            ) if raw_media_alignment == "align_pages" else (text_pages, False)
            try:
                if _save_text_pdf_with_chromium(
                    output_path,
                    title=title,
                    text_pages=text_pages_for_render,
                    image_paths=[*source_images, *source_screenshots],
                    page_size_name=raw_page_size,
                ):
                    size_bytes = output_path.stat().st_size if output_path.exists() else 0
                    expected_page_count = (
                        max(len(text_pages_for_render), media_count, 1)
                        if media_alignment_applied
                        else max(len(text_pages_for_render), 1)
                    )
                    actual_page_count = _pdf_page_count(output_path, fallback=expected_page_count)
                    layout_features = [
                        "styled_report_pages",
                        "clickable_links",
                        "readable_text_layout",
                    ]
                    if source_images or source_screenshots:
                        layout_features.append("verified_media_embeds")
                    if media_alignment_applied:
                        layout_features.append("media_aligned_report_pages")
                    return ToolResult(
                        invocation_id=invocation.invocation_id,
                        tool_name=invocation.tool_name,
                        status="completed",
                        output={
                            "path": str(output_path),
                            "artifact_path": str(output_path),
                            "artifact_paths": [str(output_path)],
                            "bytes_written": size_bytes,
                            "page_count": actual_page_count,
                            "text_pages": len(text_pages_for_render),
                            "source_image_paths": [str(path) for path in source_images],
                            "source_screenshot_paths": [str(path) for path in source_screenshots],
                            "optimized_image_paths": _optimized_image_paths_for_output([*source_images, *source_screenshots]),
                            "text_layer": True,
                            "quality_profile": "report_quality_v1",
                            "layout_features": layout_features,
                        },
                        error=None,
                    )
            except Exception as exc:
                logger.info("text_pdf_chromium_failed: %s", exc, exc_info=True)
            text_pages = text_pages_for_render

        pages, source_images, source_screenshots, page_error = _build_pdf_pages(
            image_paths=image_paths,
            screenshot_paths=screenshot_paths,
            text_pages=text_pages,
            page_size=page_size,
            roots=effective_roots,
            invocation=invocation,
            title=title,
            screenshot_field="screenshot_paths",
        )
        if page_error is not None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "path": str(output_path),
                    "reason": "artifact_media_inputs_failed",
                    "image_paths": image_paths,
                    "screenshot_paths": screenshot_paths,
                    "source_image_paths": source_images,
                    "source_screenshot_paths": source_screenshots,
                },
                error=page_error,
            )
        try:
            _save_pdf_pages(output_path, pages, title=title)
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"PDF creation failed: {exc}",
            )
        finally:
            for page in pages:
                try:
                    page.close()
                except Exception:
                    pass

        size_bytes = output_path.stat().st_size if output_path.exists() else 0
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "bytes_written": size_bytes,
                "page_count": len(source_images) + len(source_screenshots) + len(text_pages),
                "source_image_paths": source_images,
                "source_screenshot_paths": source_screenshots,
                "optimized_image_paths": _optimized_image_paths_for_output(
                    [*(Path(path) for path in source_images), *(Path(path) for path in source_screenshots)]
                ),
                "quality_profile": "report_quality_v1",
                "layout_features": [
                    "styled_report_pages",
                    "readable_text_layout",
                    *(
                        ["verified_media_embeds"]
                        if source_images or source_screenshots
                        else []
                    ),
                ],
            },
            error=None,
        )

    return handler


def _build_pdf_edit_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        try:
            from pypdf import PdfReader, PdfWriter
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"pdf_edit requires the pypdf package: {exc}",
            )

        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "pdf_edit requires workspace_root or allowed_roots")

        raw_input_path = invocation.arguments.get("input_path")
        if not isinstance(raw_input_path, str) or not raw_input_path.strip():
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "Missing required argument: input_path")
        input_path = _resolve_local_workspace_file_input(
            raw_input_path,
            principal_id=invocation.principal_id,
            effective_roots=effective_roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
        if input_path is None:
            input_path = Path(_resolve_virtual_workspace_path(raw_input_path, principal_id=invocation.principal_id)).expanduser().resolve()
            if not _path_within_any_root(input_path, effective_roots) and not _is_approved_filesystem_path(
                input_path, invocation.trusted_filesystem_selectors
            ):
                return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Input path is outside workspace root: {input_path}")
        if not input_path.is_file():
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"PDF file not found: {input_path}")

        append_pdf_paths, append_pdf_error = _coerce_string_list(invocation.arguments.get("append_pdf_paths"), field="append_pdf_paths")
        if append_pdf_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, append_pdf_error)
        append_image_paths, append_image_error = _coerce_string_list(invocation.arguments.get("append_image_paths"), field="append_image_paths")
        if append_image_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, append_image_error)
        append_screenshot_paths, append_screenshot_error = _coerce_string_list(
            invocation.arguments.get("append_screenshot_paths"),
            field="append_screenshot_paths",
        )
        if append_screenshot_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, append_screenshot_error)
        append_text_pages, append_text_error = _coerce_string_list(invocation.arguments.get("append_text_pages"), field="append_text_pages")
        if append_text_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, append_text_error)

        title = str(invocation.arguments.get("title") or input_path.stem or "Nullion PDF").strip()
        raw_page_size = str(invocation.arguments.get("page_size") or "letter").strip().lower()
        page_size = _PDF_PAGE_SIZES.get(raw_page_size)
        if page_size is None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Unsupported page_size: {raw_page_size}")

        raw_output_path = invocation.arguments.get("output_path")
        if isinstance(raw_output_path, str) and raw_output_path.strip():
            output_path = _explicit_output_path_or_generated_artifact_path(
                invocation,
                raw_path=raw_output_path,
                suffix=".pdf",
                stem=_safe_pdf_stem(f"{title}-edited"),
                roots=effective_roots,
            )
        else:
            output_path = _pdf_default_output_path(invocation, title=f"{title}-edited", roots=effective_roots)
        if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
            output_path, invocation.trusted_filesystem_selectors
        ):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Output path is outside workspace root: {output_path}")

        raw_page_numbers = invocation.arguments.get("page_numbers")
        if raw_page_numbers is not None and not isinstance(raw_page_numbers, list):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "page_numbers must be a list of 1-based integers")
        rotate_degrees = invocation.arguments.get("rotate_degrees")
        if rotate_degrees is None:
            rotate_degrees = 0
        if not isinstance(rotate_degrees, int) or rotate_degrees not in {0, 90, 180, 270}:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "rotate_degrees must be one of 0, 90, 180, 270")

        temp_paths: list[Path] = []
        try:
            reader = PdfReader(str(input_path))
            writer = PdfWriter()
            selected_pages = (
                [int(page_number) - 1 for page_number in raw_page_numbers if isinstance(page_number, int)]
                if isinstance(raw_page_numbers, list)
                else list(range(len(reader.pages)))
            )
            if isinstance(raw_page_numbers, list) and len(selected_pages) != len(raw_page_numbers):
                return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "page_numbers must contain only integers")
            for page_index in selected_pages:
                if page_index < 0 or page_index >= len(reader.pages):
                    return ToolResult(
                        invocation.invocation_id,
                        invocation.tool_name,
                        "failed",
                        {},
                        f"page_numbers contains out-of-range page: {page_index + 1}",
                    )
                page = reader.pages[page_index]
                if rotate_degrees:
                    page = page.rotate(rotate_degrees)
                writer.add_page(page)

            for raw_path in append_pdf_paths:
                append_path = _resolve_local_workspace_file_input(
                    raw_path,
                    principal_id=invocation.principal_id,
                    effective_roots=effective_roots,
                    trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
                )
                if append_path is None:
                    append_path = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser().resolve()
                    if not _path_within_any_root(append_path, effective_roots) and not _is_approved_filesystem_path(
                        append_path, invocation.trusted_filesystem_selectors
                    ):
                        return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Append PDF path is outside workspace root: {append_path}")
                if not append_path.is_file():
                    return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"Append PDF file not found: {append_path}")
                append_reader = PdfReader(str(append_path))
                for page in append_reader.pages:
                    writer.add_page(page)

            if append_image_paths or append_screenshot_paths or append_text_pages:
                pages, _source_images, _source_screenshots, page_error = _build_pdf_pages(
                    image_paths=append_image_paths,
                    screenshot_paths=append_screenshot_paths,
                    text_pages=append_text_pages,
                    page_size=page_size,
                    roots=effective_roots,
                    invocation=invocation,
                    title=title,
                    screenshot_field="append_screenshot_paths",
                )
                if page_error is not None:
                    return ToolResult(
                        invocation.invocation_id,
                        invocation.tool_name,
                        "failed",
                        {
                            "path": str(output_path),
                            "reason": "artifact_media_inputs_failed",
                            "input_path": str(input_path),
                            "image_paths": append_image_paths,
                            "screenshot_paths": append_screenshot_paths,
                            "source_image_paths": _source_images,
                            "source_screenshot_paths": _source_screenshots,
                        },
                        page_error,
                    )
                output_path.parent.mkdir(parents=True, exist_ok=True)
                temp_path = output_path.parent / f".{output_path.stem}-append-{uuid4().hex[:12]}.pdf"
                temp_paths.append(temp_path)
                try:
                    _save_pdf_pages(temp_path, pages, title=title)
                finally:
                    for page in pages:
                        try:
                            page.close()
                        except Exception:
                            pass
                append_reader = PdfReader(str(temp_path))
                for page in append_reader.pages:
                    writer.add_page(page)

            if len(writer.pages) == 0:
                return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "pdf_edit produced no pages")
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with output_path.open("wb") as handle:
                writer.write(handle)
            page_count = len(PdfReader(str(output_path)).pages)
        except Exception as exc:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, f"PDF edit failed: {exc}")
        finally:
            for temp_path in temp_paths:
                try:
                    temp_path.unlink(missing_ok=True)
                except Exception:
                    pass

        size_bytes = output_path.stat().st_size if output_path.exists() else 0
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "bytes_written": size_bytes,
                "page_count": page_count,
            },
            error=None,
        )

    return handler



def _build_file_patch_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_path = invocation.arguments.get("path")
        raw_old_string = invocation.arguments.get("old_string")
        raw_new_string = invocation.arguments.get("new_string")
        raw_replace_all = invocation.arguments.get("replace_all")
        if not isinstance(raw_path, str) or not raw_path:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: path",
            )
        if not isinstance(raw_old_string, str):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: old_string",
            )
        if raw_old_string == "":
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="old_string must be non-empty",
            )
        if not isinstance(raw_new_string, str):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: new_string",
            )
        if raw_replace_all is not None and not isinstance(raw_replace_all, bool):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="replace_all must be a boolean when provided",
            )

        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="file_patch requires workspace_root or allowed_roots",
            )

        path = Path(raw_path).expanduser().resolve()
        if not _path_within_any_root(path, effective_roots):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Path is outside workspace root: {path}",
            )

        try:
            content = path.read_text(encoding="utf-8")
        except FileNotFoundError:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"File not found: {path}",
            )

        replace_all = bool(raw_replace_all)
        replacement_count = content.count(raw_old_string)
        if replacement_count == 0:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Old string not found in file: {path}",
            )
        if not replace_all and replacement_count != 1:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Old string must match exactly once in file: {path}",
            )

        new_content = content.replace(raw_old_string, raw_new_string) if replace_all else content.replace(raw_old_string, raw_new_string, 1)
        path.write_text(new_content, encoding="utf-8")
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"path": str(path), "replacements": replacement_count if replace_all else 1},
            error=None,
        )

    return handler



def _build_workspace_summary_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def _resolve_candidate_within_scope(candidate: Path, root: Path) -> Path | None:
        try:
            resolved_candidate = candidate.resolve()
        except OSError:
            return None
        if not _is_within_allowed_root(resolved_candidate, root):
            return None
        return resolved_candidate

    def _display_path(*, root: Path, path: Path) -> str:
        relative_path = path.relative_to(root).as_posix()
        return relative_path

    def handler(invocation: ToolInvocation) -> ToolResult:
        _ = invocation.arguments  # arguments unused; cannot del a slotted dataclass attribute
        roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        max_entries = _env_int("NULLION_WORKSPACE_SUMMARY_MAX_ENTRIES", 20_000, minimum=1)
        max_sample_files = _env_int("NULLION_WORKSPACE_SUMMARY_SAMPLE_FILES", 40, minimum=0)
        max_extensions = _env_int("NULLION_WORKSPACE_SUMMARY_EXTENSIONS", 40, minimum=1)
        seen_directories: set[Path] = set()
        seen_files: set[Path] = set()
        extensions: Counter[str] = Counter()
        sample_files: list[str] = []
        scanned_entries = 0
        truncated = False

        for root in roots:
            if scanned_entries >= max_entries:
                truncated = True
                break
            for current_dir, dirnames, filenames in os.walk(root, topdown=True, followlinks=False):
                if scanned_entries >= max_entries:
                    truncated = True
                    dirnames[:] = []
                    break
                current_path = Path(current_dir)
                scoped_dirnames: list[str] = []
                for dirname in sorted(dirnames):
                    scanned_entries += 1
                    if scanned_entries > max_entries:
                        truncated = True
                        break
                    candidate_dir = current_path / dirname
                    if candidate_dir.is_symlink() or _should_prune_filesystem_walk_dir(candidate_dir):
                        continue
                    resolved_dir = _resolve_candidate_within_scope(candidate_dir, root)
                    if resolved_dir is None:
                        continue
                    scoped_dirnames.append(dirname)
                    if resolved_dir not in seen_directories:
                        seen_directories.add(resolved_dir)
                dirnames[:] = scoped_dirnames
                for filename in sorted(filenames):
                    scanned_entries += 1
                    if scanned_entries > max_entries:
                        truncated = True
                        break
                    candidate_file = current_path / filename
                    if candidate_file.is_symlink():
                        continue
                    resolved_file = _resolve_candidate_within_scope(candidate_file, root)
                    if resolved_file is None or resolved_file in seen_files:
                        continue
                    seen_files.add(resolved_file)
                    extensions[candidate_file.suffix.lower()] += 1
                    display_path = _display_path(root=root, path=candidate_file)
                    if len(roots) > 1:
                        display_path = f"{root}::{display_path}"
                    if len(sample_files) < max_sample_files:
                        sample_files.append(display_path)

        extension_rows = [
            {"extension": extension, "count": count}
            for extension, count in sorted(extensions.items(), key=lambda item: (-item[1], item[0]))
        ]

        output = {
            "roots": [str(root) for root in roots],
            "file_count": len(seen_files),
            "directory_count": len(seen_directories),
            "scanned_entries": scanned_entries,
            "truncated": truncated,
            "extensions": extension_rows[:max_extensions],
            "sample_files": sorted(sample_files),
        }
        if len(extension_rows) > max_extensions:
            output["extensions_truncated"] = {"shown": max_extensions, "total": len(extension_rows)}
        if len(seen_files) > max_sample_files:
            output["sample_files_truncated"] = {"shown": max_sample_files, "total": len(seen_files)}

        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output=output,
            error=None,
        )

    return handler



def _terminal_deliverable_artifact_paths_since(
    roots: Iterable[Path],
    *,
    since_timestamp: float,
) -> list[str]:
    candidates: list[tuple[float, str]] = []
    seen: set[str] = set()
    scanned_entries = 0
    for raw_root in roots:
        try:
            root = Path(raw_root).expanduser().resolve()
        except (OSError, RuntimeError, ValueError):
            continue
        if not root.is_dir():
            continue
        stack = [root]
        while stack and scanned_entries < _MAX_TERMINAL_ARTIFACT_SCAN_ENTRIES:
            current = stack.pop()
            try:
                with os.scandir(current) as entries:
                    for entry in entries:
                        scanned_entries += 1
                        if scanned_entries > _MAX_TERMINAL_ARTIFACT_SCAN_ENTRIES:
                            break
                        try:
                            if entry.is_dir(follow_symlinks=False):
                                stack.append(Path(entry.path))
                                continue
                            if not entry.is_file(follow_symlinks=False):
                                continue
                            path = Path(entry.path).resolve()
                            if path.suffix.lower() not in _TERMINAL_DELIVERABLE_ARTIFACT_EXTENSIONS:
                                continue
                            stat = path.stat()
                        except (OSError, RuntimeError, ValueError):
                            continue
                        if stat.st_size <= 0 or stat.st_mtime < (since_timestamp - 1.0):
                            continue
                        resolved = str(path)
                        if resolved in seen:
                            continue
                        seen.add(resolved)
                        candidates.append((stat.st_mtime, resolved))
            except OSError:
                continue
    candidates.sort(key=lambda item: (item[0], item[1]))
    return promote_supporting_asset_artifact_paths(
        [path for _mtime, path in candidates[-_MAX_TERMINAL_DISCOVERED_ARTIFACTS:]],
        artifact_roots=tuple(Path(root).expanduser().resolve() for root in roots),
    )


def _build_terminal_exec_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
) -> ToolHandler:
    cwd = Path(workspace_root).resolve() if workspace_root is not None else None
    execution_cwd = cwd if cwd is not None and cwd.is_dir() else None
    resolved_allowed_roots = (
        tuple(Path(root).expanduser().resolve() for root in allowed_roots)
        if allowed_roots is not None
        else (() if cwd is None else (cwd,))
    )
    backend = terminal_executor_backend or SubprocessTerminalExecutorBackend()

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_command = invocation.arguments.get("command")
        if not isinstance(raw_command, str) or not raw_command:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: command",
            )

        filesystem_denial = _terminal_filesystem_safety_denial(
            raw_command,
            allowed_roots=resolved_allowed_roots,
            principal_id=invocation.principal_id,
        )
        if filesystem_denial is not None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=filesystem_denial,
                error=str(filesystem_denial.get("message") or "Filesystem traversal denied"),
            )

        mutation_denial = _terminal_filesystem_mutation_approval_denial(
            invocation,
            execution_cwd=execution_cwd,
            allowed_roots=resolved_allowed_roots,
        )
        if mutation_denial is not None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="denied",
                output=mutation_denial,
                error=str(mutation_denial.get("message") or "Filesystem mutation approval required"),
            )

        egress_attempts = _egress_attempts_for_invocation(invocation)
        raw_network_mode = invocation.arguments.get("network_mode")
        network_mode = _normalize_network_mode(raw_network_mode)
        if raw_network_mode is not None and network_mode not in _VALID_NETWORK_MODES:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "reason": "invalid_network_mode",
                    "network_mode": raw_network_mode,
                    "allowed_network_modes": sorted(_VALID_NETWORK_MODES),
                },
                error=f"Invalid terminal execution network mode: {raw_network_mode}",
            )
        approved_targets = invocation.arguments.get("approved_targets")
        policy = TerminalExecutionPolicy(
            network_mode=network_mode,
            approved_targets=tuple(target for target in approved_targets if isinstance(target, str))
            if isinstance(approved_targets, (list, tuple, set, frozenset))
            else (),
        )
        shell = _normalize_terminal_shell(invocation.arguments.get("shell"))
        timeout_seconds = _terminal_timeout_seconds(invocation.arguments.get("timeout_seconds"))
        has_network_egress = bool(egress_attempts)

        if _is_restrictive_network_mode(network_mode):
            backend_descriptor = backend.describe()
            required_attested_capabilities = _required_attested_capabilities_for_network_mode(network_mode)
            if backend_descriptor.mode != "subprocess":
                missing_capabilities = _missing_backend_attested_capabilities(
                    backend_descriptor,
                    required_capabilities=required_attested_capabilities,
                )
                verification = verify_terminal_backend_attestation(
                    backend_descriptor,
                    required_capabilities=required_attested_capabilities,
                    verifier=terminal_attestation_verifier,
                )
                if missing_capabilities:
                    return ToolResult(
                        invocation_id=invocation.invocation_id,
                        tool_name=invocation.tool_name,
                        status="failed",
                        output={
                            "reason": "backend_attestation_missing",
                            "network_mode": network_mode,
                            "required_attested_capabilities": list(required_attested_capabilities),
                            "backend": {
                                "mode": backend_descriptor.mode,
                                "attested_capabilities": list(backend_descriptor.attested_capabilities),
                            },
                        },
                        error=(
                            "Restrictive terminal execution requires backend attestation: "
                            + ", ".join(missing_capabilities)
                        ),
                    )
                if not verification.is_valid:
                    failure_reason = verification.failure_reason or "verification failed"
                    backend_output: dict[str, object] = {
                        "mode": backend_descriptor.mode,
                        "attested_capabilities": list(backend_descriptor.attested_capabilities),
                    }
                    if backend_descriptor.evidence is not None:
                        backend_output["evidence_format"] = backend_descriptor.evidence.format
                    return ToolResult(
                        invocation_id=invocation.invocation_id,
                        tool_name=invocation.tool_name,
                        status="failed",
                        output={
                            "reason": "backend_attestation_invalid",
                            "network_mode": network_mode,
                            "required_attested_capabilities": list(required_attested_capabilities),
                            "backend": backend_output,
                            "attestation_verification": {
                                "is_valid": verification.is_valid,
                                "failure_reason": verification.failure_reason,
                                "metadata": dict(verification.metadata),
                            },
                        },
                        error=(
                            "Restrictive terminal execution requires verified backend attestation: "
                            + failure_reason
                        ),
                    )

        if has_network_egress and not all(
            _network_attempt_allowed(
                attempt=attempt,
                network_mode=network_mode,
                approved_targets=approved_targets,
            )
            for attempt in egress_attempts
        ):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"reason": "network_denied", "network_mode": network_mode, "egress_attempts": egress_attempts},
                error="Network egress denied by terminal execution policy",
            )

        try:
            started_at = datetime.now(UTC).timestamp()
            try:
                completed = backend.run(
                    raw_command,
                    cwd=str(execution_cwd) if execution_cwd is not None else None,
                    timeout=timeout_seconds,
                    policy=policy,
                    shell=shell,
                )
            except TypeError as exc:
                if "shell" not in str(exc):
                    raise
                completed = backend.run(
                    raw_command,
                    cwd=str(execution_cwd) if execution_cwd is not None else None,
                    timeout=timeout_seconds,
                    policy=policy,
                )
        except subprocess.TimeoutExpired:
            output = {"egress_attempts": egress_attempts} if has_network_egress else {}
            if network_mode is not None:
                output["network_mode"] = network_mode
            output["timeout_seconds"] = timeout_seconds
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=output,
                error=f"Command timed out after {timeout_seconds}s",
            )

        output = {
            "stdout": completed.stdout,
            "stderr": completed.stderr,
            "exit_code": completed.exit_code,
            "shell": completed.shell,
            "timeout_seconds": timeout_seconds,
        }
        if network_mode is not None:
            output["network_mode"] = network_mode
        if has_network_egress:
            output["egress_attempts"] = egress_attempts
        if completed.exit_code == 0:
            artifact_scan_roots = [
                *(Path(root).expanduser().resolve() for root in resolved_allowed_roots),
            ]
            runtime_artifact_root = _runtime_artifact_root()
            if runtime_artifact_root is not None:
                artifact_scan_roots.insert(0, runtime_artifact_root)
            principal_artifact_root = _workspace_artifact_root_for_principal(invocation.principal_id)
            if principal_artifact_root is not None:
                artifact_scan_roots.insert(0, principal_artifact_root)
            artifact_paths = _terminal_deliverable_artifact_paths_since(
                tuple(dict.fromkeys(artifact_scan_roots)),
                since_timestamp=started_at,
            )
            if artifact_paths:
                output["artifact_paths"] = artifact_paths
        terminal_output_path = _materialize_terminal_output_attachment(
            output,
            invocation=invocation,
            execution_cwd=execution_cwd,
            allowed_roots=resolved_allowed_roots,
        )
        if terminal_output_path:
            raw_artifact_paths = output.get("artifact_paths")
            artifact_paths = (
                list(raw_artifact_paths)
                if isinstance(raw_artifact_paths, (list, tuple, set, frozenset))
                else ([raw_artifact_paths] if raw_artifact_paths else [])
            )
            artifact_paths.append(terminal_output_path)
            output["artifact_paths"] = list(dict.fromkeys(str(path) for path in artifact_paths if str(path)))
            output["terminal_output_path"] = terminal_output_path
            descriptors = list(output.get("artifact_descriptors") or [])
            descriptors.append(
                artifact_output_descriptor(
                    terminal_output_path,
                    role=ARTIFACT_ROLE_DELIVERABLE,
                    kind="terminal_output",
                    label="Terminal output",
                )
            )
            output["artifact_descriptors"] = descriptors
        error = None
        if completed.exit_code != 0:
            if completed.exit_code == 127 and completed.stderr.startswith(
                ("Shell startup failed:", "Terminal launcher startup failed:")
            ):
                error = completed.stderr
            else:
                error = f"Command failed with exit code {completed.exit_code}"
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed" if completed.exit_code == 0 else "failed",
            output=output,
            error=error,
        )

    return handler



def _http_retry_url_for_https_transport_failure(url: str) -> str | None:
    parsed = urlparse(url)
    if parsed.scheme != "https" or parsed.username or parsed.password or parsed.port is not None:
        return None
    return parsed._replace(scheme="http").geturl()


def _is_https_transport_eof(exc: Exception) -> bool:
    reason = exc.reason if isinstance(exc, urllib.error.URLError) else exc
    if isinstance(reason, ssl.SSLEOFError):
        return True
    if isinstance(reason, ssl.SSLError):
        library, reason_text = reason.args[:2] if len(reason.args) >= 2 else ("", "")
        return "UNEXPECTED_EOF_WHILE_READING" in str(library) or "UNEXPECTED_EOF_WHILE_READING" in str(reason_text)
    return False


def _web_fetch_binary_suffix(url: str, content_type: str) -> str:
    parsed_suffix = Path(urlparse(url).path).suffix.lower()
    if parsed_suffix in VALID_ATTACHMENT_EXTENSIONS:
        return parsed_suffix
    guessed = mimetypes.guess_extension(content_type.split(";", 1)[0].strip())
    if guessed:
        return ".jpg" if guessed == ".jpe" else guessed
    return ".bin"


def _web_fetch_is_textual(content_type: str, data: bytes) -> bool:
    media_type = content_type.split(";", 1)[0].strip().lower()
    if media_type.startswith("text/") or media_type in {
        "application/json",
        "application/javascript",
        "application/xml",
        "application/xhtml+xml",
        "application/rss+xml",
        "application/atom+xml",
        "image/svg+xml",
    }:
        return True
    sample = data[:1024]
    if b"\x00" in sample:
        return False
    try:
        sample.decode("utf-8")
        return media_type in {"", "application/octet-stream"}
    except UnicodeDecodeError:
        return False


def _fetch_web_url_once(url: str, timeout_seconds: int) -> dict[str, object]:
    resolution = _resolve_web_fetch_resolution(url)
    parsed = urlparse(url)
    host = parsed.hostname
    if (
        parsed.scheme not in {"http", "https"}
        or not isinstance(host, str)
        or not host
        or (resolution is None and not _is_global_literal_ip(host))
    ):
        raise ValueError(f"Blocked URL for web_fetch: {url}")
    request = urllib.request.Request(
        url,
        headers={
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36"
            ),
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9",
        },
    )
    opener = _build_web_fetch_opener()
    with _pinned_web_fetch_resolution(resolution):
        with opener.open(request, timeout=timeout_seconds) as response:
            status_code = getattr(response, "status", 200)
            content_type = response.headers.get_content_type()
            data = response.read(_WEB_FETCH_MAX_BODY_BYTES + 1)
    truncated = len(data) > _WEB_FETCH_MAX_BODY_BYTES
    if truncated:
        data = data[:_WEB_FETCH_MAX_BODY_BYTES]
    if not _web_fetch_is_textual(content_type, data):
        return {
            "url": url,
            "status_code": status_code,
            "content_type": content_type,
            "content_kind": "binary",
            "body_size": len(data),
            "body_truncated": truncated,
            "suggested_extension": _web_fetch_binary_suffix(url, content_type),
            "_body_bytes": data,
        }
    body = data[:65536].decode("utf-8", "ignore")  # 64 KB of text
    title_match = re.search(r"<title[^>]*>(.*?)</title>", body, re.IGNORECASE | re.DOTALL)
    title = None if title_match is None else re.sub(r"\s+", " ", unescape(title_match.group(1))).strip() or None
    text = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", body)).strip()
    return {
        "url": url,
        "status_code": status_code,
        "content_type": content_type,
        "title": title,
        "text": text,
        "body": body,
    }


def _default_web_fetcher(url: str, timeout_seconds: int) -> dict[str, object]:
    try:
        return _fetch_web_url_once(url, timeout_seconds)
    except Exception as exc:
        retry_url = _http_retry_url_for_https_transport_failure(url)
        if retry_url is None or not _is_https_transport_eof(exc):
            raise
        response = _fetch_web_url_once(retry_url, timeout_seconds)
        response["requested_url"] = url
        response["transport_fallback"] = {
            "from_scheme": "https",
            "to_scheme": "http",
            "reason": "https_transport_eof",
        }
        return response


def _browser_image_url(value: object, *, base_url: str | None = None) -> str | None:
    text = str(value or "").strip()
    if not text:
        return None
    if text.startswith("data:image/"):
        return text
    resolved = urljoin(base_url, text) if base_url else text
    parsed = urlparse(resolved)
    if parsed.scheme in {"http", "https"} and parsed.netloc:
        return resolved
    return None


def _browser_image_srcset_urls(value: object, *, base_url: str | None = None) -> list[str]:
    urls: list[str] = []
    for candidate in str(value or "").split(","):
        raw_url = candidate.strip().split()[0] if candidate.strip() else ""
        image_url = _browser_image_url(raw_url, base_url=base_url)
        if image_url:
            urls.append(image_url)
    return urls


_PLACEHOLDER_IMAGE_HOSTS = frozenset(
    {
        "dummyimage.com",
        "httpbin.org",
        "placehold.co",
        "placeholder.com",
        "via.placeholder.com",
    }
)


def _is_placeholder_image_url(url: str) -> bool:
    parsed = urlparse(url)
    host = (parsed.hostname or "").strip().lower()
    if host in _PLACEHOLDER_IMAGE_HOSTS:
        return True
    return host.startswith("placeholder.") or host.endswith(".placeholder.com")


class _BrowserImageCandidateHTMLParser(HTMLParser):
    def __init__(self, *, base_url: str | None = None) -> None:
        super().__init__(convert_charrefs=True)
        self.base_url = base_url
        self.candidates: list[dict[str, object]] = []
        self._script_type: str | None = None
        self._script_chunks: list[str] = []

    def _append(self, raw_url: object, *, role: str, attributes: dict[str, str]) -> None:
        image_url = _browser_image_url(raw_url, base_url=self.base_url)
        if not image_url:
            return
        self.candidates.append(
            {
                "source_url": image_url,
                "page_url": self.base_url,
                "role": role,
                "alt": attributes.get("alt") or "",
                "title": attributes.get("title") or "",
                "width": attributes.get("width") or "",
                "height": attributes.get("height") or "",
            }
        )

    def _append_srcset(self, raw_srcset: object, *, role: str, attributes: dict[str, str]) -> None:
        for image_url in _browser_image_srcset_urls(raw_srcset, base_url=self.base_url):
            self._append(image_url, role=role, attributes=attributes)

    def _append_style_urls(self, raw_style: object, *, role: str, attributes: dict[str, str]) -> None:
        for match in re.finditer(r"url\(([^)]+)\)", str(raw_style or ""), flags=re.IGNORECASE):
            raw_url = match.group(1).strip().strip("\"'")
            self._append(raw_url, role=role, attributes=attributes)

    def _append_json_ld_images(self, value: object, *, attributes: dict[str, str]) -> None:
        if isinstance(value, str):
            self._append(value, role="json_ld_image", attributes=attributes)
            return
        if isinstance(value, list):
            for item in value:
                self._append_json_ld_images(item, attributes=attributes)
            return
        if not isinstance(value, dict):
            return
        for key, nested in value.items():
            normalized_key = str(key or "").strip().lower()
            if normalized_key in {"image", "thumbnail", "thumbnailurl", "contenturl", "url"}:
                self._append_json_ld_images(nested, attributes=attributes)
            elif isinstance(nested, (dict, list)):
                self._append_json_ld_images(nested, attributes=attributes)

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag_name = tag.lower()
        attributes = {str(name).lower(): str(value or "") for name, value in attrs}
        self._append_style_urls(attributes.get("style"), role="style_background", attributes=attributes)
        if tag_name == "img":
            for attr_name in (
                "src",
                "data-src",
                "data-original",
                "data-lazy-src",
                "data-image",
                "data-image-src",
                "data-full-src",
                "data-large-src",
                "data-original-src",
                "data-media-url",
                "data-thumb",
                "data-thumbnail",
            ):
                self._append(attributes.get(attr_name), role=f"img_{attr_name}", attributes=attributes)
            for attr_name in ("srcset", "data-srcset", "data-lazy-srcset"):
                self._append_srcset(attributes.get(attr_name), role=f"img_{attr_name}", attributes=attributes)
            return
        if tag_name == "source":
            for attr_name in ("srcset", "data-srcset", "data-lazy-srcset"):
                self._append_srcset(attributes.get(attr_name), role=f"source_{attr_name}", attributes=attributes)
            return
        if tag_name == "meta":
            key = (attributes.get("property") or attributes.get("name") or "").strip().lower()
            if key in {"og:image", "og:image:url", "twitter:image", "twitter:image:src"}:
                self._append(attributes.get("content"), role=key, attributes=attributes)
            return
        if tag_name == "link":
            rel_tokens = {
                token.strip().lower()
                for token in (attributes.get("rel") or "").replace(",", " ").split()
                if token.strip()
            }
            if "image_src" in rel_tokens:
                self._append(attributes.get("href"), role="link_image_src", attributes=attributes)
            return
        if tag_name == "script":
            script_type = (attributes.get("type") or "").strip().lower()
            if script_type in {"application/ld+json", "application/json+ld"}:
                self._script_type = script_type
                self._script_chunks = []

    def handle_data(self, data: str) -> None:
        if self._script_type:
            self._script_chunks.append(data)

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() != "script" or not self._script_type:
            return
        script_text = "".join(self._script_chunks).strip()
        self._script_type = None
        self._script_chunks = []
        if not script_text:
            return
        try:
            payload = json.loads(script_text)
        except json.JSONDecodeError:
            return
        self._append_json_ld_images(payload, attributes={})


def _browser_image_candidates_from_html(html: str, *, page_url: str | None = None) -> list[dict[str, object]]:
    parser = _BrowserImageCandidateHTMLParser(base_url=page_url)
    parser.feed(html)
    parser.close()
    return parser.candidates


def _dedupe_browser_image_candidates(candidates: Iterable[dict[str, object]]) -> list[dict[str, object]]:
    deduped: list[dict[str, object]] = []
    seen: set[str] = set()
    for candidate in candidates:
        source_url = candidate.get("source_url")
        if not isinstance(source_url, str) or not source_url.strip():
            continue
        key = source_url.strip()
        if key in seen:
            continue
        seen.add(key)
        deduped.append({**candidate, "source_url": key})
    return deduped


def _browser_image_candidate_declared_pixels(candidate: dict[str, object]) -> int:
    def dimension(value: object) -> int:
        match = re.search(r"\d+", str(value or ""))
        if not match:
            return 0
        try:
            return int(match.group(0))
        except ValueError:
            return 0

    return dimension(candidate.get("width")) * dimension(candidate.get("height"))


def _browser_image_candidate_sort_key(candidate: dict[str, object]) -> tuple[int, int, int]:
    role = str(candidate.get("role") or "")
    role_priority = 2 if role in {"explicit", "page_url_image"} else 1 if role.startswith(("og:", "twitter:")) else 0
    text_score = int(bool(str(candidate.get("alt") or "").strip())) + int(bool(str(candidate.get("title") or "").strip()))
    return (role_priority, _browser_image_candidate_declared_pixels(candidate), text_score)


def _browser_image_content_quality_error(
    *,
    data_bytes: int,
    width: int,
    height: int,
    luma_stddev: float,
) -> str | None:
    if data_bytes < _BROWSER_IMAGE_CONTENT_MIN_BYTES:
        return f"image is too small for document media ({data_bytes} bytes)"
    if width < _BROWSER_IMAGE_CONTENT_MIN_WIDTH or height < _BROWSER_IMAGE_CONTENT_MIN_HEIGHT:
        return f"image dimensions are too small for document media ({width}x{height})"
    if width * height < _BROWSER_IMAGE_CONTENT_MIN_PIXELS:
        return f"image area is too small for document media ({width * height} pixels)"
    shortest = max(1, min(width, height))
    aspect_ratio = max(width, height) / shortest
    if aspect_ratio > _BROWSER_IMAGE_CONTENT_MAX_ASPECT_RATIO:
        return f"image aspect ratio is too extreme for document media ({aspect_ratio:.2f}:1)"
    if luma_stddev < _BROWSER_IMAGE_CONTENT_MIN_LUMA_STDDEV:
        return f"image has too little visual detail for document media (luma stddev {luma_stddev:.2f})"
    return None


def _fetch_browser_image_binary(url: str, timeout_seconds: int) -> dict[str, object]:
    resolution = _resolve_web_fetch_resolution(url)
    parsed = urlparse(url)
    host = parsed.hostname
    if _is_placeholder_image_url(url):
        raise ValueError(f"Blocked placeholder image URL for browser_image_collect: {url}")
    if (
        parsed.scheme not in {"http", "https"}
        or not isinstance(host, str)
        or not host
        or (resolution is None and not _is_global_literal_ip(host))
    ):
        raise ValueError(f"Blocked URL for browser_image_collect: {url}")
    request = urllib.request.Request(
        url,
        headers={
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36"
            ),
            "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9",
        },
    )
    opener = _build_web_fetch_opener()
    with _pinned_web_fetch_resolution(resolution):
        with opener.open(request, timeout=timeout_seconds) as response:
            status_code = getattr(response, "status", 200)
            content_type = response.headers.get_content_type()
            data = response.read(_BROWSER_IMAGE_MAX_BYTES + 1)
            final_url = getattr(response, "url", None) or getattr(response, "geturl", lambda: url)()
    if len(data) > _BROWSER_IMAGE_MAX_BYTES:
        raise ValueError("image response exceeded maximum size")
    return {
        "source_url": url,
        "final_url": str(final_url or url),
        "status_code": status_code,
        "content_type": content_type,
        "data": data,
    }


def _browser_image_bytes_from_data_url(url: str) -> dict[str, object]:
    header, separator, payload = url.partition(",")
    if not separator or ";base64" not in header.lower():
        raise ValueError("Only base64 data image URLs are supported")
    media_type = header[5:].split(";", 1)[0].strip().lower()
    if media_type not in _BROWSER_IMAGE_DIRECTIVE_MEDIA_TYPES:
        raise ValueError(f"Unsupported data image type: {media_type or 'unknown'}")
    data = base64.b64decode(payload, validate=True)
    if len(data) > _BROWSER_IMAGE_MAX_BYTES:
        raise ValueError("image response exceeded maximum size")
    return {
        "source_url": url,
        "final_url": url,
        "status_code": 200,
        "content_type": media_type,
        "data": data,
    }


def _materialize_browser_image_artifact(
    *,
    invocation: ToolInvocation,
    image_payload: dict[str, object],
    output_stem: str,
    quality_profile: str = "content",
) -> dict[str, object]:
    data = image_payload.get("data")
    if not isinstance(data, bytes) or not data:
        raise ValueError("image response had no bytes")
    content_type = str(image_payload.get("content_type") or "").split(";", 1)[0].strip().lower()
    if content_type and content_type not in _BROWSER_IMAGE_DIRECTIVE_MEDIA_TYPES and not content_type.startswith("application/octet-stream"):
        raise ValueError(f"unsupported image content type: {content_type}")
    try:
        from PIL import Image
    except ModuleNotFoundError as exc:
        raise RuntimeError(f"browser_image_collect requires Pillow to validate image assets: {exc}") from exc

    with Image.open(BytesIO(data)) as image:
        width, height = image.size
        image_format = str(image.format or "").upper()
        if quality_profile == "content":
            from PIL import ImageStat

            preview = image.convert("L").resize((64, 64))
            luma_stddev = float(ImageStat.Stat(preview).stddev[0])
            quality_error = _browser_image_content_quality_error(
                data_bytes=len(data),
                width=width,
                height=height,
                luma_stddev=luma_stddev,
            )
            if quality_error is not None:
                raise ValueError(quality_error)
        if image_format in _BROWSER_IMAGE_FORMAT_SUFFIXES:
            suffix = _BROWSER_IMAGE_FORMAT_SUFFIXES[image_format]
            converted_bytes = data
        else:
            suffix = ".png"
            converted = image.convert("RGBA" if image.mode in {"RGBA", "LA", "P"} else "RGB")
            buffer = BytesIO()
            converted.save(buffer, format="PNG")
            converted_bytes = buffer.getvalue()

    from nullion.artifacts import artifact_path_for_generated_workspace_file

    path = artifact_path_for_generated_workspace_file(
        principal_id=invocation.principal_id,
        suffix=suffix,
        stem=_safe_pdf_stem(output_stem) or "browser-image",
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(converted_bytes)
    return {
        "local_path": str(path),
        "artifact_path": str(path),
        "source_url": image_payload.get("source_url"),
        "final_url": image_payload.get("final_url"),
        "content_type": content_type,
        "bytes": len(converted_bytes),
        "width": width,
        "height": height,
        "format": suffix.lstrip("."),
    }


def _build_browser_image_collect_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_max_images = invocation.arguments.get("max_images")
        try:
            max_images = int(raw_max_images) if raw_max_images is not None else 10
        except (TypeError, ValueError):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "max_images must be an integer")
        max_images = max(1, min(_BROWSER_IMAGE_COLLECT_MAX_IMAGES, max_images))
        output_stem = str(invocation.arguments.get("output_stem") or "browser-image").strip() or "browser-image"
        quality_profile = str(invocation.arguments.get("quality_profile") or "content").strip().lower()
        if quality_profile not in {"content", "any"}:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {},
                "quality_profile must be one of: content, any",
            )
        raw_timeout_seconds = invocation.arguments.get("timeout_seconds")
        try:
            total_timeout_seconds = float(raw_timeout_seconds) if raw_timeout_seconds is not None else 12.0
        except (TypeError, ValueError):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "timeout_seconds must be a number")
        total_timeout_seconds = max(2.0, min(30.0, total_timeout_seconds))
        deadline = perf_counter() + total_timeout_seconds
        page_url = _browser_image_url(invocation.arguments.get("page_url"))
        html = invocation.arguments.get("html")
        if html is not None and not isinstance(html, str):
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, "html must be a string")
        image_urls, image_url_error = _coerce_string_list(invocation.arguments.get("image_urls"), field="image_urls")
        if image_url_error is not None:
            return ToolResult(invocation.invocation_id, invocation.tool_name, "failed", {}, image_url_error)

        candidates: list[dict[str, object]] = []
        for image_url in image_urls:
            resolved_image_url = _browser_image_url(image_url, base_url=page_url)
            if resolved_image_url:
                candidates.append({"source_url": resolved_image_url, "page_url": page_url, "role": "explicit"})
        skipped: list[dict[str, object]] = []
        if page_url and not html:
            try:
                fetched_page = _default_web_fetcher(page_url, min(6, max(1, int(deadline - perf_counter()))))
            except Exception as exc:
                skipped.append({"source_url": page_url, "reason": f"page_fetch_failed: {exc}"})
            else:
                if fetched_page.get("content_kind") == "binary" and str(fetched_page.get("content_type") or "").startswith("image/"):
                    candidates.append({"source_url": page_url, "page_url": page_url, "role": "page_url_image"})
                else:
                    body = fetched_page.get("body")
                    if isinstance(body, str):
                        html = body
        if isinstance(html, str) and html.strip():
            candidates.extend(_browser_image_candidates_from_html(html, page_url=page_url))

        candidates = sorted(
            _dedupe_browser_image_candidates(candidates),
            key=_browser_image_candidate_sort_key,
            reverse=True,
        )[: max(max_images * 3, max_images)]
        if not candidates:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {"reason": "no_image_candidates", "image_paths": [], "images": [], "skipped_images": skipped},
                "browser_image_collect requires page_url, html, or image_urls that contain image asset URLs.",
            )

        images: list[dict[str, object]] = []
        for candidate in candidates:
            if len(images) >= max_images:
                break
            remaining_seconds = deadline - perf_counter()
            if remaining_seconds <= 0:
                skipped.append({"source_url": str(candidate.get("source_url") or ""), "reason": "image collection time budget exhausted"})
                break
            source_url = str(candidate.get("source_url") or "")
            try:
                image_payload = (
                    _browser_image_bytes_from_data_url(source_url)
                    if source_url.startswith("data:image/")
                    else _fetch_browser_image_binary(source_url, max(1, min(4, int(remaining_seconds))))
                )
                materialized = _materialize_browser_image_artifact(
                    invocation=invocation,
                    image_payload=image_payload,
                    output_stem=output_stem,
                    quality_profile=quality_profile,
                )
            except Exception as exc:
                skipped.append({"source_url": source_url, "reason": str(exc)})
                continue
            images.append(
                {
                    **{key: value for key, value in candidate.items() if key != "source_url"},
                    **materialized,
                }
            )

        image_paths = [str(image["local_path"]) for image in images if isinstance(image.get("local_path"), str)]
        output = {
            "image_paths": image_paths,
            "artifact_paths": image_paths,
            "artifact_descriptors": [
                artifact_output_descriptor(path, role=ARTIFACT_ROLE_SOURCE, kind="browser_image")
                for path in image_paths
            ],
            "images": images,
            "skipped_images": skipped,
            "candidate_count": len(candidates),
            "saved_count": len(image_paths),
            "timeout_seconds": total_timeout_seconds,
        }
        if not image_paths:
            output["reason"] = "image_collection_failed"
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                output,
                "browser_image_collect could not save any embeddable local image artifacts.",
            )
        return ToolResult(invocation.invocation_id, invocation.tool_name, "completed", output, None)

    return handler


def register_browser_image_collect_tool(registry: ToolRegistry) -> ToolRegistry:
    try:
        registry.get_spec("browser_image_collect")
        return registry
    except KeyError:
        pass
    registry.register(
        ToolSpec(
            name="browser_image_collect",
            description=(
                "Extract image asset URLs from page HTML or explicit browser/page image URLs, fetch the image bytes, "
                "and save embeddable local image artifact files. Use this before document_create, spreadsheet_create, "
                "presentation_create, or pdf_create when those artifacts need real content images instead of page screenshots."
            ),
            risk_level=ToolRiskLevel.MEDIUM,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=30,
            capability_tags=("public_web", "browser", "media", "image", "artifact"),
            continuation_tools=("document_create", "spreadsheet_create", "presentation_create", "pdf_create"),
        ),
        _build_browser_image_collect_handler(),
    )
    return registry


def register_web_fetch_tool(
    registry: ToolRegistry,
    web_fetcher: Callable[[str, int], dict[str, object]] | None = None,
) -> ToolRegistry:
    try:
        registry.get_spec("web_fetch")
        return registry
    except KeyError:
        pass
    registry.register(
        ToolSpec(
            name="web_fetch",
            description=(
                "Fetch a URL and return its content and response metadata. "
                "Use only when direct HTTP fetching is enabled by the runtime policy; otherwise use browser tools."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.READ,
            requires_approval=False,
            timeout_seconds=20,
        ),
        _build_web_fetch_handler(web_fetcher or _default_web_fetcher),
    )
    return registry


def _filename_from_content_disposition(value: str | None) -> str | None:
    if not isinstance(value, str) or not value.strip():
        return None
    for part in value.split(";"):
        key, sep, raw = part.strip().partition("=")
        if not sep:
            continue
        if key.strip().lower() in {"filename", "filename*"}:
            candidate = raw.strip().strip("\"'")
            if "''" in candidate:
                candidate = candidate.split("''", 1)[1]
            candidate = unquote(candidate).strip()
            if candidate:
                return candidate
    return None


def _safe_download_filename(value: object, *, url: str, content_type: str | None = None) -> str:
    raw_name = str(value or "").strip()
    if not raw_name:
        raw_name = Path(unquote(urlparse(url).path or "")).name
    if not raw_name:
        suffix = _web_fetch_binary_suffix(url, content_type or "application/octet-stream")
        raw_name = f"download{suffix}"
    name = Path(raw_name).name.strip().replace("\x00", "")
    name = re.sub(r"[^A-Za-z0-9._ -]+", "_", name).strip(" .")
    if not name:
        name = "download.bin"
    if "." not in name:
        suffix = _web_fetch_binary_suffix(url, content_type or "application/octet-stream")
        if suffix:
            name = f"{name}{suffix}"
    return name[:180]


def _file_download_error_result(
    invocation: ToolInvocation,
    *,
    error: str,
    reason: str,
    output: dict[str, object] | None = None,
) -> ToolResult:
    payload = dict(output or {})
    payload["reason"] = reason
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output=payload,
        error=error,
    )


def _file_download_timeout_seconds(raw_timeout: object) -> float:
    default_timeout = 45.0
    if raw_timeout is None:
        return default_timeout
    try:
        timeout = float(raw_timeout)
    except (TypeError, ValueError):
        return default_timeout
    return max(2.0, min(60.0, timeout))


def _try_set_response_socket_timeout(response: object, timeout_seconds: float) -> None:
    for attr_path in (
        ("fp", "raw", "_sock"),
        ("fp", "raw", "_fp", "fp", "raw", "_sock"),
    ):
        target = response
        try:
            for attr in attr_path:
                target = getattr(target, attr)
            settimeout = getattr(target, "settimeout", None)
            if callable(settimeout):
                settimeout(max(1.0, float(timeout_seconds)))
                return
        except Exception:
            continue


def _download_archive_integrity_error(path: Path, *, final_url: str, content_type: str) -> str | None:
    suffixes = {path.suffix.lower(), Path(unquote(urlparse(final_url).path or "")).suffix.lower()}
    media_type = str(content_type or "").split(";", 1)[0].strip().lower()
    looks_like_zip = ".zip" in suffixes or "zip" in media_type
    if looks_like_zip and not zipfile.is_zipfile(path):
        return f"Downloaded archive is not a valid ZIP file: {path}"
    return None


def _build_file_download_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_url = invocation.arguments.get("url")
        if not isinstance(raw_url, str) or not raw_url.strip():
            return _file_download_error_result(invocation, error="Missing required argument: url", reason="missing_url")
        url = raw_url.strip()
        resolution = _resolve_web_fetch_resolution(url)
        parsed = urlparse(url)
        host = parsed.hostname
        if (
            parsed.scheme not in {"http", "https"}
            or not isinstance(host, str)
            or not host
            or (resolution is None and not _is_global_literal_ip(host))
        ):
            return _file_download_error_result(
                invocation,
                error=f"Blocked URL for file_download: {url}",
                reason="blocked_url",
                output={"url": url},
            )
        effective_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not effective_roots:
            return _file_download_error_result(
                invocation,
                error="file_download requires workspace_root or allowed_roots",
                reason="missing_workspace_roots",
                output={"url": url},
            )
        raw_max_bytes = invocation.arguments.get("max_bytes")
        max_bytes = _FILE_DOWNLOAD_MAX_BYTES
        if isinstance(raw_max_bytes, int) and raw_max_bytes > 0:
            max_bytes = min(raw_max_bytes, _FILE_DOWNLOAD_MAX_BYTES)
        timeout_seconds = _file_download_timeout_seconds(invocation.arguments.get("timeout_seconds"))
        try:
            integrity_retry_attempt = int(invocation.arguments.get("_download_integrity_retry_attempt") or 0)
        except (TypeError, ValueError):
            integrity_retry_attempt = 0
        deadline = perf_counter() + timeout_seconds
        raw_output_path = invocation.arguments.get("output_path")
        output_path: Path | None = None
        if isinstance(raw_output_path, str) and raw_output_path.strip():
            candidate = Path(raw_output_path).expanduser()
            output_path = candidate.resolve() if candidate.is_absolute() else (effective_roots[0] / candidate).resolve()
            if not _path_within_any_root(output_path, effective_roots) and not _is_approved_filesystem_path(
                output_path, invocation.trusted_filesystem_selectors
            ):
                return _file_download_error_result(
                    invocation,
                    error=f"Output path is outside workspace root: {output_path}",
                    reason="output_path_outside_workspace",
                    output={"url": url, "output_path": str(output_path)},
                )
        request = urllib.request.Request(
            url,
            headers={
                "User-Agent": (
                    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/124.0.0.0 Safari/537.36"
                ),
                "Accept": "*/*",
                "Accept-Language": "en-US,en;q=0.9",
            },
        )
        try:
            opener = _build_web_fetch_opener()
            with _pinned_web_fetch_resolution(resolution):
                socket_timeout = min(10.0, max(1.0, timeout_seconds))
                with opener.open(request, timeout=socket_timeout) as response:
                    status_code = int(getattr(response, "status", 200) or 200)
                    final_url = response.geturl()
                    content_type = response.headers.get_content_type()
                    content_length = response.headers.get("Content-Length")
                    try:
                        expected_bytes = int(content_length) if content_length else None
                    except ValueError:
                        expected_bytes = None
                    if expected_bytes is not None and expected_bytes > max_bytes:
                        return _file_download_error_result(
                            invocation,
                            error=f"Download is too large: {expected_bytes} bytes exceeds {max_bytes} bytes",
                            reason="download_too_large",
                            output={
                                "url": url,
                                "final_url": final_url,
                                "content_type": content_type,
                                "content_length": expected_bytes,
                                "max_bytes": max_bytes,
                            },
                        )
                    if output_path is None:
                        header_name = _filename_from_content_disposition(response.headers.get("Content-Disposition"))
                        safe_name = _safe_download_filename(
                            invocation.arguments.get("filename") or header_name,
                            url=final_url or url,
                            content_type=content_type,
                        )
                        from nullion.artifacts import artifact_path_for_generated_workspace_file

                        suffix = Path(safe_name).suffix or _web_fetch_binary_suffix(final_url or url, content_type)
                        output_path = artifact_path_for_generated_workspace_file(
                            principal_id=invocation.principal_id,
                            suffix=suffix,
                            stem=Path(safe_name).stem or "download",
                        )
                    output_path.parent.mkdir(parents=True, exist_ok=True)
                    bytes_written = 0
                    with output_path.open("wb") as target:
                        while True:
                            remaining_seconds = deadline - perf_counter()
                            if remaining_seconds <= 0:
                                target.close()
                                try:
                                    output_path.unlink()
                                except OSError:
                                    pass
                                return _file_download_error_result(
                                    invocation,
                                    error=f"Download timed out after {timeout_seconds:g}s",
                                    reason="download_timeout",
                                    output={
                                        "url": url,
                                        "final_url": final_url,
                                        "content_type": content_type,
                                        "bytes_read": bytes_written,
                                        "timeout_seconds": timeout_seconds,
                                    },
                                )
                            _try_set_response_socket_timeout(response, min(socket_timeout, remaining_seconds))
                            chunk = response.read(1024 * 1024)
                            if not chunk:
                                break
                            bytes_written += len(chunk)
                            if perf_counter() > deadline:
                                target.close()
                                try:
                                    output_path.unlink()
                                except OSError:
                                    pass
                                return _file_download_error_result(
                                    invocation,
                                    error=f"Download timed out after {timeout_seconds:g}s",
                                    reason="download_timeout",
                                    output={
                                        "url": url,
                                        "final_url": final_url,
                                        "content_type": content_type,
                                        "bytes_read": bytes_written,
                                        "timeout_seconds": timeout_seconds,
                                    },
                                )
                            if bytes_written > max_bytes:
                                target.close()
                                try:
                                    output_path.unlink()
                                except OSError:
                                    pass
                                return _file_download_error_result(
                                    invocation,
                                    error=f"Download exceeded {max_bytes} bytes",
                                    reason="download_too_large",
                                    output={
                                        "url": url,
                                        "final_url": final_url,
                                        "content_type": content_type,
                                        "bytes_read": bytes_written,
                                        "max_bytes": max_bytes,
                                    },
                                )
                            target.write(chunk)
        except Exception as exc:
            return _file_download_error_result(
                invocation,
                error=f"Could not download file: {exc}",
                reason="download_failed",
                output={"url": url},
            )
        integrity_error = _download_archive_integrity_error(
            output_path,
            final_url=str(final_url or url),
            content_type=str(content_type or ""),
        )
        if integrity_error:
            bytes_written = output_path.stat().st_size if output_path.exists() else 0
            try:
                output_path.unlink()
            except OSError:
                pass
            if integrity_retry_attempt < 1:
                retry_args = dict(invocation.arguments)
                retry_args["_download_integrity_retry_attempt"] = integrity_retry_attempt + 1
                retry_args["timeout_seconds"] = max(timeout_seconds, 60.0)
                return handler(replace(invocation, arguments=retry_args))
            return _file_download_error_result(
                invocation,
                error=integrity_error,
                reason="download_integrity_failed",
                output={
                    "url": url,
                    "final_url": final_url,
                    "content_type": content_type,
                    "bytes_read": bytes_written,
                    "timeout_seconds": timeout_seconds,
                },
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "url": url,
                "final_url": final_url,
                "status_code": status_code,
                "content_type": content_type,
                "path": str(output_path),
                "artifact_path": str(output_path),
                "artifact_paths": [str(output_path)],
                "artifact_descriptors": [
                    artifact_output_descriptor(output_path, role=ARTIFACT_ROLE_SOURCE, kind="download")
                ],
                "bytes_written": output_path.stat().st_size,
                "max_bytes": max_bytes,
                "timeout_seconds": timeout_seconds,
            },
            error=None,
        )

    return handler


def _json_get(url: str, timeout_seconds: int, *, attempts: int = 3) -> object:
    request = urllib.request.Request(url, headers={"User-Agent": "Nullion/1.0"})
    last_error: Exception | None = None
    for _attempt in range(max(1, attempts)):
        try:
            opener = urllib.request.build_opener(
                urllib.request.ProxyHandler({}),
                urllib.request.HTTPSHandler(context=_web_fetch_tls_context()),
            )
            with opener.open(request, timeout=timeout_seconds) as response:
                payload = response.read(500_000)
            break
        except (OSError, TimeoutError, urllib.error.URLError) as exc:
            last_error = exc
    else:
        assert last_error is not None
        raise last_error
    return json.loads(payload.decode("utf-8"))


def _weather_json_get(url: str, timeout_seconds: int) -> object:
    return _json_get(url, timeout_seconds, attempts=1)


def _weather_http_timeout_seconds() -> int:
    raw = os.getenv("NULLION_WEATHER_HTTP_TIMEOUT_SECONDS", "").strip()
    if raw:
        try:
            return min(15, max(2, int(raw)))
        except ValueError:
            pass
    return 5


_WEATHER_LOCATION_CACHE_LOCK = threading.Lock()
_WEATHER_LOCATION_CACHE: dict[tuple[int, str], tuple[float, dict[str, object]]] = {}


def _weather_location_cache_ttl_seconds() -> int:
    raw = os.getenv("NULLION_WEATHER_LOCATION_CACHE_TTL_SECONDS", "").strip()
    if raw:
        try:
            return max(0, min(7 * 24 * 60 * 60, int(raw)))
        except ValueError:
            pass
    return 24 * 60 * 60


def _weather_location_cache_key(location_text: str, json_get: Callable[[str, int], object]) -> tuple[int, str]:
    normalized = " ".join(str(location_text or "").strip().lower().split())
    return (id(json_get), normalized)


def _cached_weather_location(
    location_text: str,
    *,
    json_get: Callable[[str, int], object],
) -> dict[str, object] | None:
    ttl_seconds = _weather_location_cache_ttl_seconds()
    if ttl_seconds <= 0:
        return None
    key = _weather_location_cache_key(location_text, json_get)
    now = perf_counter()
    with _WEATHER_LOCATION_CACHE_LOCK:
        cached = _WEATHER_LOCATION_CACHE.get(key)
        if not cached:
            return None
        cached_at, location = cached
        if now - cached_at > ttl_seconds:
            _WEATHER_LOCATION_CACHE.pop(key, None)
            return None
        return dict(location)


def _remember_weather_location(
    location_text: str,
    location: dict[str, object],
    *,
    json_get: Callable[[str, int], object],
) -> dict[str, object]:
    ttl_seconds = _weather_location_cache_ttl_seconds()
    if ttl_seconds <= 0:
        return location
    key = _weather_location_cache_key(location_text, json_get)
    with _WEATHER_LOCATION_CACHE_LOCK:
        _WEATHER_LOCATION_CACHE[key] = (perf_counter(), dict(location))
    return location


def _coerce_float_arg(value: object, *, name: str) -> float | None:
    if value is None or value == "":
        return None
    try:
        return float(value)
    except (TypeError, ValueError) as exc:
        raise ValueError(f"{name} must be a number") from exc


def _resolve_open_meteo_location(
    arguments: dict[str, object],
    *,
    json_get: Callable[[str, int], object],
) -> dict[str, object]:
    latitude = _coerce_float_arg(arguments.get("latitude"), name="latitude")
    longitude = _coerce_float_arg(arguments.get("longitude"), name="longitude")
    location_text = str(arguments.get("location_text") or "").strip()
    if latitude is not None and longitude is not None:
        return {
            "latitude": latitude,
            "longitude": longitude,
            "name": location_text or f"{latitude:.4f},{longitude:.4f}",
            "source_url": None,
        }
    if not location_text:
        raise ValueError("Provide location_text or latitude and longitude.")
    cached_location = _cached_weather_location(location_text, json_get=json_get)
    if cached_location is not None:
        return cached_location
    source_url = ""
    results: object = None
    searched_locations: list[str] = []
    candidate_locations = _weather_location_text_candidates(location_text)
    for candidate_location in dict.fromkeys(candidate_locations):
        searched_locations.append(candidate_location)
        query = urlencode({"name": candidate_location, "count": 1, "language": "en", "format": "json"})
        source_url = f"https://geocoding-api.open-meteo.com/v1/search?{query}"
        payload = json_get(source_url, _weather_http_timeout_seconds())
        results = payload.get("results") if isinstance(payload, dict) else None
        if isinstance(results, list) and results:
            break
    if not isinstance(results, list) or not results:
        # Open-Meteo's geocoder is city-oriented; retry progressively broader
        # locality candidates through Nominatim before declaring failure.
        for candidate_location in dict.fromkeys(candidate_locations):
            nominatim_result = _resolve_nominatim_weather_location(
                candidate_location,
                json_get=json_get,
            )
            if nominatim_result is not None:
                return _remember_weather_location(location_text, nominatim_result, json_get=json_get)
    if not isinstance(results, list) or not results:
        searched = ", ".join(searched_locations)
        raise ValueError(f"Could not resolve location: {location_text}" + (f" (tried: {searched})" if searched else ""))
    first = results[0]
    if not isinstance(first, dict):
        raise ValueError(f"Could not resolve location: {location_text}")
    resolved_latitude = _coerce_float_arg(first.get("latitude"), name="latitude")
    resolved_longitude = _coerce_float_arg(first.get("longitude"), name="longitude")
    if resolved_latitude is None or resolved_longitude is None:
        raise ValueError(f"Resolved location is missing coordinates: {location_text}")
    label_parts = [
        str(first.get("name") or "").strip(),
        str(first.get("admin1") or "").strip(),
        str(first.get("country") or "").strip(),
    ]
    label = ", ".join(part for part in label_parts if part) or location_text
    return _remember_weather_location(location_text, {
        "latitude": resolved_latitude,
        "longitude": resolved_longitude,
        "name": label,
        "country": first.get("country"),
        "timezone": first.get("timezone"),
        "source_url": source_url,
    }, json_get=json_get)


def _weather_location_text_candidates(location_text: str) -> list[str]:
    raw = str(location_text or "").strip()
    if not raw:
        return []
    candidates = [raw]
    zip_match = re.search(r"(?<!\d)(\d{5})(?:-\d{4})?(?!\d)", raw)
    if zip_match:
        zip_code = zip_match.group(1)
        if zip_code != raw:
            candidates.insert(0, zip_code)
    parts = [" ".join(part.split()).strip() for part in raw.split(",") if part.strip()]
    if len(parts) >= 2:
        first = parts[0]
        last = parts[-1]
        if len(parts) >= 3 and first and last:
            # Degrade from precise private/subaddress text to the nearest
            # resolvable public locality before blocking the weather plan.
            first_last = f"{first}, {last}"
            if first_last not in candidates:
                candidates.append(first_last)
        if first and first not in candidates:
            candidates.append(first)
        if last and len(last) > 2 and last not in candidates:
            candidates.append(last)
    return candidates


def _resolve_nominatim_weather_location(
    location_text: str,
    *,
    json_get: Callable[[str, int], object],
) -> dict[str, object] | None:
    query = urlencode({"q": location_text, "format": "jsonv2", "limit": 1, "addressdetails": 1})
    source_url = f"https://nominatim.openstreetmap.org/search?{query}"
    payload = json_get(source_url, _weather_http_timeout_seconds())
    if not isinstance(payload, list) or not payload:
        return None
    first = payload[0]
    if not isinstance(first, dict):
        return None
    latitude = _coerce_float_arg(first.get("lat"), name="latitude")
    longitude = _coerce_float_arg(first.get("lon"), name="longitude")
    if latitude is None or longitude is None:
        return None
    label = str(first.get("display_name") or location_text).strip() or location_text
    address = first.get("address")
    return {
        "latitude": latitude,
        "longitude": longitude,
        "name": label,
        "country": address.get("country") if isinstance(address, dict) else None,
        "timezone": None,
        "source_url": source_url,
    }


def _weather_code_label(code: object) -> str:
    try:
        value = int(code)
    except (TypeError, ValueError):
        return "unknown"
    labels = {
        0: "clear sky",
        1: "mainly clear",
        2: "partly cloudy",
        3: "overcast",
        45: "fog",
        48: "depositing rime fog",
        51: "light drizzle",
        53: "moderate drizzle",
        55: "dense drizzle",
        56: "light freezing drizzle",
        57: "dense freezing drizzle",
        61: "slight rain",
        63: "moderate rain",
        65: "heavy rain",
        66: "light freezing rain",
        67: "heavy freezing rain",
        71: "slight snow",
        73: "moderate snow",
        75: "heavy snow",
        77: "snow grains",
        80: "slight rain showers",
        81: "moderate rain showers",
        82: "violent rain showers",
        85: "slight snow showers",
        86: "heavy snow showers",
        95: "thunderstorm",
        96: "thunderstorm with slight hail",
        99: "thunderstorm with heavy hail",
    }
    return labels.get(value, f"weather code {value}")


def _daily_open_meteo_forecast(payload: dict[str, object]) -> list[dict[str, object]]:
    daily = payload.get("daily")
    if not isinstance(daily, dict):
        return []
    dates = daily.get("time")
    if not isinstance(dates, list):
        return []
    rows: list[dict[str, object]] = []
    for index, date_value in enumerate(dates):
        row: dict[str, object] = {"date": date_value}
        for source_key, target_key in (
            ("weather_code", "weather_code"),
            ("temperature_2m_max", "temperature_max_f"),
            ("temperature_2m_min", "temperature_min_f"),
            ("precipitation_probability_max", "precipitation_probability_max"),
            ("precipitation_sum", "precipitation_sum_in"),
            ("wind_speed_10m_max", "wind_speed_max_mph"),
            ("sunrise", "sunrise"),
            ("sunset", "sunset"),
        ):
            values = daily.get(source_key)
            if isinstance(values, list) and index < len(values):
                row[target_key] = values[index]
        row["summary"] = _weather_code_label(row.get("weather_code"))
        rows.append(row)
    return rows


def _hourly_open_meteo_forecast(payload: dict[str, object], *, limit_hours: int | None) -> list[dict[str, object]]:
    hourly = payload.get("hourly")
    if not isinstance(hourly, dict):
        return []
    times = hourly.get("time")
    if not isinstance(times, list):
        return []
    rows: list[dict[str, object]] = []
    for index, time_value in enumerate(times):
        if limit_hours is not None and index >= limit_hours:
            break
        row: dict[str, object] = {"time": time_value}
        for source_key, target_key in (
            ("weather_code", "weather_code"),
            ("temperature_2m", "temperature_f"),
            ("apparent_temperature", "apparent_temperature_f"),
            ("precipitation_probability", "precipitation_probability"),
            ("precipitation", "precipitation_in"),
            ("wind_speed_10m", "wind_speed_mph"),
            ("wind_gusts_10m", "wind_gusts_mph"),
        ):
            values = hourly.get(source_key)
            if isinstance(values, list) and index < len(values):
                row[target_key] = values[index]
        row["summary"] = _weather_code_label(row.get("weather_code"))
        rows.append(row)
    return rows


def _coerce_bool_arg(value: object, *, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in {"1", "true", "yes", "on"}:
            return True
        if normalized in {"0", "false", "no", "off"}:
            return False
    return default


def _build_weather_forecast_handler(
    json_get: Callable[[str, int], object],
) -> ToolHandler:
    def _handler(invocation: ToolInvocation) -> ToolResult:
        try:
            arguments = dict(invocation.arguments)
            has_coordinates = arguments.get("latitude") is not None and arguments.get("longitude") is not None
            if not str(arguments.get("location_text") or "").strip() and not has_coordinates:
                try:
                    from nullion.preferences import load_profile

                    profile = load_profile()
                except Exception:
                    profile = {}
                profile_location = (
                    str(profile.get("address") or "").strip()
                    if isinstance(profile, Mapping)
                    else ""
                )
                if profile_location:
                    arguments["location_text"] = profile_location
                    arguments["location_source"] = "saved_profile"
            location = _resolve_open_meteo_location(arguments, json_get=json_get)
            forecast_days = arguments.get("forecast_days", 3)
            try:
                days = min(7, max(1, int(forecast_days)))
            except (TypeError, ValueError):
                days = 3
            timezone = str(arguments.get("timezone") or location.get("timezone") or "auto").strip() or "auto"
            include_hourly = _coerce_bool_arg(arguments.get("include_hourly"), default=False)
            hourly_limit: int | None = None
            if include_hourly:
                requested_hours = arguments.get("hourly_hours")
                try:
                    hourly_limit = min(168, max(1, int(requested_hours))) if requested_hours is not None else min(168, days * 24)
                except (TypeError, ValueError):
                    hourly_limit = min(168, days * 24)
            query_args = {
                "latitude": location["latitude"],
                "longitude": location["longitude"],
                "current": "temperature_2m,relative_humidity_2m,apparent_temperature,precipitation,weather_code,wind_speed_10m",
                "daily": "weather_code,temperature_2m_max,temperature_2m_min,precipitation_probability_max,precipitation_sum,wind_speed_10m_max,sunrise,sunset",
                "temperature_unit": "fahrenheit",
                "wind_speed_unit": "mph",
                "precipitation_unit": "inch",
                "timezone": timezone,
                "forecast_days": days,
            }
            if include_hourly:
                query_args["hourly"] = (
                    "temperature_2m,apparent_temperature,precipitation_probability,"
                    "precipitation,weather_code,wind_speed_10m,wind_gusts_10m"
                )
            query = urlencode(query_args)
            forecast_url = f"https://api.open-meteo.com/v1/forecast?{query}"
            payload = json_get(forecast_url, _weather_http_timeout_seconds())
            if not isinstance(payload, dict):
                raise ValueError("Weather forecast response was not a JSON object.")
            hourly_rows = _hourly_open_meteo_forecast(payload, limit_hours=hourly_limit) if include_hourly else []
            output = {
                "provider": "open-meteo",
                "source_url": forecast_url,
                "geocoding_source_url": location.get("source_url"),
                "location": {
                    "name": location.get("name"),
                    "country": location.get("country"),
                    "latitude": location.get("latitude"),
                    "longitude": location.get("longitude"),
                    "timezone": payload.get("timezone") or timezone,
                },
                "location_source": arguments.get("location_source") or "tool_arguments",
                "current": payload.get("current") if isinstance(payload.get("current"), dict) else {},
                "daily": _daily_open_meteo_forecast(payload),
                "hourly": hourly_rows,
                "units": {
                    "temperature": "fahrenheit",
                    "wind_speed": "mph",
                    "precipitation": "inch",
                },
            }
            return ToolResult(invocation.invocation_id, invocation.tool_name, "completed", output)
        except Exception as exc:
            return ToolResult(
                invocation.invocation_id,
                invocation.tool_name,
                "failed",
                {"reason": "weather_forecast_failed"},
                str(exc),
            )

    return _handler


def register_weather_forecast_tool(
    registry: ToolRegistry,
    *,
    json_get: Callable[[str, int], object] | None = None,
) -> ToolRegistry:
    try:
        registry.get_spec("weather_forecast")
        return registry
    except KeyError:
        pass
    registry.register(
        ToolSpec(
            name="weather_forecast",
            description=(
                "Fetch current, hourly, and multi-day public forecast data from Open-Meteo using structured "
                "coordinates or a resolvable location. Use for read-only weather forecast questions "
                "before browser navigation or general web search when this tool is available."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.READ,
            requires_approval=False,
            timeout_seconds=20,
            capability_tags=("public_web", "weather", "forecast"),
        ),
        _build_weather_forecast_handler(json_get or _weather_json_get),
    )
    return registry


def register_market_quote_tool(
    registry: ToolRegistry,
    *,
    quote_fetcher: Callable[[str, int], dict[str, object]] | None = None,
) -> ToolRegistry:
    try:
        registry.get_spec("market_quote")
        return registry
    except KeyError:
        pass
    registry.register(
        ToolSpec(
            name="market_quote",
            description=(
                "Fetch structured public market quote rows for explicit ticker symbols. "
                "Use before browser navigation or general web search for current equity/ETF quote checks "
                "when concrete symbols are available."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.READ,
            requires_approval=False,
            timeout_seconds=20,
            input_schema=_default_input_schema_for_tool("market_quote"),
            capability_tags=("public_web", "market_data", "quote"),
        ),
        _build_market_quote_handler(quote_fetcher or _default_market_quote_fetcher),
    )
    return registry



def _build_kernel_tool_registry(
    *,
    plugin_registration_allowed: bool,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
    direct_web_fetch_enabled: bool | None = None,
) -> ToolRegistry:
    filesystem_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else (() if workspace_root is None else (Path(workspace_root).resolve(),))
    registry = ToolRegistry(
        plugin_registration_allowed=plugin_registration_allowed,
        filesystem_allowed_roots=filesystem_allowed_roots,
    )
    if allowed_roots is not None:
        workspace_root = None
    if _env_flag("NULLION_FILE_ACCESS_ENABLED"):
        registry.register(
            ToolSpec(
                name="file_read",
                description="Read a local file inside the workspace.",
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
            ),
            _build_file_read_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="file_write",
                description="Write a local file inside the workspace.",
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=20,
            ),
            _build_file_write_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="file_download",
                description="Download an HTTP/HTTPS URL directly to a workspace file without shell execution.",
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=60,
            ),
            _build_file_download_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="archive_extract",
                description=(
                    "Extract supported archive/compression files (.zip, .tar, .tar.gz/.tgz, .tar.bz2/.tbz2, "
                    ".tar.xz/.txz, .gz, .bz2, .xz) into a workspace directory. Skips macOS metadata sidecar "
                    "entries unless include_metadata is true. Also writes full CSV, JSON, and XLSX entry "
                    "manifests so downstream spreadsheet/report tasks can use the complete inventory even when "
                    "the tool result is compacted."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_archive_extract_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="archive_create",
                description=(
                    "Create supported archive/compression files (.zip, .tar, .tar.gz/.tgz, .tar.bz2/.tbz2, "
                    ".tar.xz/.txz, .gz, .bz2, .xz) from workspace files or directories. Skips macOS metadata "
                    "sidecar entries unless include_metadata is true."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_archive_create_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="document_create",
                description=(
                    "Create a real .docx document artifact from structured paragraphs, sections, and existing image files. "
                    "The built-in generator applies a report-quality layout profile; provide structured content instead of raw dumps. "
                    "Use this as the first local document-delivery rung; if it cannot complete, request local shell "
                    "execution as the last-resort local fallback when no external account auth is required."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_document_create_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="spreadsheet_create",
                description=(
                    "Create a real .xlsx spreadsheet artifact from structured rows, links, and existing image files. "
                    "Use the structured sheets argument when the requested workbook needs multiple worksheets, "
                    "including a separate summary sheet. "
                    "If the request or a structured plan gives an explicit total data-row count, pass expected_rows "
                    "and make the provided rows match it; only treat the count as per-source/per-category when that "
                    "is explicitly requested. "
                    "Use this as the first local spreadsheet-delivery rung; if it cannot complete, request local shell "
                    "execution as the last-resort local fallback when no external account auth is required."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_spreadsheet_create_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="presentation_create",
                description=(
                    "Create a real .pptx slide deck artifact from structured slides and existing image files. "
                    "The built-in generator applies a report-quality layout profile and preserves media aspect ratios. "
                    "Use this as the first local presentation-delivery rung; if it cannot complete, request local shell "
                    "execution as the last-resort local fallback when no external account auth is required."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_presentation_create_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="pdf_create",
                description=(
                    "Create a real PDF artifact locally from existing image files and/or report text pages. "
                    "Use text_pages for readable report content; image-only PDFs are only appropriate when the requested deliverable is image pages. "
                    "Try this before terminal_exec; local shell remains the last-resort fallback when this cannot complete."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_pdf_create_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
        registry.register(
            ToolSpec(
                name="pdf_edit",
                description=(
                    "Edit a PDF locally into a new PDF artifact: keep/reorder pages, rotate pages, "
                    "append other PDFs, append image pages, or append text pages. "
                    "Try this before terminal_exec; local shell remains the last-resort fallback when this cannot complete."
                ),
                risk_level=ToolRiskLevel.MEDIUM,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=30,
            ),
            _build_pdf_edit_handler(
                None
                if workspace_root is None
                else Path(workspace_root),
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
            ),
        )
    if _env_flag("NULLION_TERMINAL_ENABLED"):
        registry.register(
            ToolSpec(
                name="terminal_exec",
                description=(
                    "Execute a local shell command with approval. Use this for local system inspection, installed software "
                    "or package-manager changes, service/process control, environment/path checks, and other computer-level "
                    "work when no safer dedicated structured tool can satisfy the task. Use dedicated structured tools and "
                    "reusable Python-backed local tools first when they can satisfy the task."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.DANGEROUS_EXEC,
                requires_approval=True,
                timeout_seconds=20,
            ),
            _build_terminal_exec_handler(
                workspace_root,
                allowed_roots=[Path(root) for root in allowed_roots] if allowed_roots is not None else None,
                terminal_executor_backend=terminal_executor_backend,
                terminal_attestation_verifier=terminal_attestation_verifier,
            ),
        )
    if _env_flag("NULLION_WEB_ACCESS_ENABLED"):
        if (
            direct_web_fetch_enabled
            if direct_web_fetch_enabled is not None
            else _env_flag("NULLION_DIRECT_WEB_FETCH_ENABLED", default=False)
        ):
            register_web_fetch_tool(registry)
        if _env_flag("NULLION_FILE_ACCESS_ENABLED"):
            register_browser_image_collect_tool(registry)
        register_weather_forecast_tool(registry)
        register_market_quote_tool(registry)
    if _connector_access_enabled():
        register_connector_plugin(registry)
    if _env_flag("NULLION_SKILL_PACK_ACCESS_ENABLED", default=False):
        registry.register(
            ToolSpec(
                name="skill_pack_read",
                description=(
                    "Read an installed skill pack reference file by pack id and relative path. "
                    "Use this when an enabled skill pack lists service-specific API reference docs."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                capability_tags=("skill_pack",),
            ),
            _build_skill_pack_read_handler(),
        )
    return registry



def create_core_tool_registry(
    *,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
    direct_web_fetch_enabled: bool | None = None,
) -> ToolRegistry:
    return _build_kernel_tool_registry(
        plugin_registration_allowed=False,
        workspace_root=workspace_root,
        allowed_roots=allowed_roots,
        terminal_executor_backend=terminal_executor_backend,
        terminal_attestation_verifier=terminal_attestation_verifier,
        direct_web_fetch_enabled=direct_web_fetch_enabled,
    )



def create_plugin_tool_registry(
    *,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
    direct_web_fetch_enabled: bool | None = None,
) -> ToolRegistry:
    return _build_kernel_tool_registry(
        plugin_registration_allowed=True,
        workspace_root=workspace_root,
        allowed_roots=allowed_roots,
        terminal_executor_backend=terminal_executor_backend,
        terminal_attestation_verifier=terminal_attestation_verifier,
        direct_web_fetch_enabled=direct_web_fetch_enabled,
    )



def create_extension_tool_registry(
    *,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
) -> ToolRegistry:
    return create_plugin_tool_registry(
        workspace_root=workspace_root,
        allowed_roots=allowed_roots,
        terminal_executor_backend=terminal_executor_backend,
        terminal_attestation_verifier=terminal_attestation_verifier,
    )



def create_default_tool_registry(
    *,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
) -> ToolRegistry:
    return create_core_tool_registry(
        workspace_root=workspace_root,
        allowed_roots=allowed_roots,
        terminal_executor_backend=terminal_executor_backend,
        terminal_attestation_verifier=terminal_attestation_verifier,
    )



def build_default_tool_registry(
    *,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    terminal_executor_backend: TerminalExecutorBackend | None = None,
    terminal_attestation_verifier: TerminalAttestationVerifier | None = None,
) -> ToolRegistry:
    return create_core_tool_registry(
        workspace_root=workspace_root,
        allowed_roots=allowed_roots,
        terminal_executor_backend=terminal_executor_backend,
        terminal_attestation_verifier=terminal_attestation_verifier,
    )


def _build_web_fetch_handler(
    web_fetcher: Callable[[str, int], dict[str, object]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_url = invocation.arguments.get("url")
        if not isinstance(raw_url, str) or not raw_url:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: url",
            )

        try:
            response = web_fetcher(raw_url, 20)
        except Exception as exc:  # pragma: no cover - caller-provided fetcher guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=str(exc),
            )

        body_bytes = response.pop("_body_bytes", None)
        if isinstance(body_bytes, bytes):
            try:
                from nullion.artifacts import artifact_path_for_generated_workspace_file

                suffix = str(response.get("suggested_extension") or ".bin")
                path = artifact_path_for_generated_workspace_file(
                    principal_id=invocation.principal_id,
                    suffix=suffix,
                    stem="web-fetch",
                )
                path.parent.mkdir(parents=True, exist_ok=True)
                path.write_bytes(body_bytes)
                response["artifact_path"] = str(path)
                response["artifact_paths"] = [str(path)]
            except Exception as exc:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output=response,
                    error=f"Could not materialize binary web_fetch response: {exc}",
                )

        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output=response,
            error=None,
        )

    return handler


def _connector_provider_env_prefix(provider_id: str) -> str:
    return re.sub(r"[^A-Z0-9]+", "_", provider_id.upper()).strip("_")


def _connector_credential_candidate_names(connection: object | None, provider_id: str) -> tuple[str, ...]:
    candidates: list[str] = []
    credential_ref = getattr(connection, "credential_ref", None)
    if isinstance(credential_ref, str) and credential_ref.strip():
        candidates.append(credential_ref.strip().removeprefix("env:"))
    prefix = _connector_provider_env_prefix(provider_id)
    if prefix:
        candidates.extend((f"{prefix}_API_KEY", f"{prefix}_TOKEN", f"{prefix}_SECRET_KEY"))
        for suffix in ("_CONNECTOR_PROVIDER", "_CONNECTOR"):
            if prefix.endswith(suffix):
                gateway_prefix = prefix.removesuffix(suffix).strip("_")
                if gateway_prefix:
                    candidates.extend(
                        (
                            f"{gateway_prefix}_API_KEY",
                            f"{gateway_prefix}_TOKEN",
                            f"{gateway_prefix}_SECRET_KEY",
                        )
                    )
        if prefix.startswith("SKILL_PACK_CONNECTOR_"):
            gateway_prefix = prefix.removeprefix("SKILL_PACK_CONNECTOR_").strip("_")
            if gateway_prefix:
                candidates.extend(
                    (
                        f"{gateway_prefix}_API_KEY",
                        f"{gateway_prefix}_TOKEN",
                        f"{gateway_prefix}_SECRET_KEY",
                    )
                )
    return tuple(dict.fromkeys(candidates))


def _account_read_tool_timeout_seconds() -> int:
    return min(
        90,
        _env_int("NULLION_ACCOUNT_READ_TOOL_TIMEOUT_SECONDS", 45, minimum=20),
    )


def _email_search_detail_fetch_limit() -> int:
    return min(
        5,
        _env_int("NULLION_EMAIL_SEARCH_DETAIL_FETCH_LIMIT", 3, minimum=0),
    )


def register_connector_plugin(registry: ToolRegistry) -> None:
    email_tool_timeout_seconds = _account_read_tool_timeout_seconds()
    try:
        registry.get_spec("connector_request")
    except KeyError:
        registry.mark_plugin_installed("connector_plugin")
        registry.register(
            ToolSpec(
                name="connector_request",
                description=(
                    "Make an HTTP request to an installed connector/API gateway using the current workspace's "
                    "configured connection credential. GET/HEAD are always read-only; POST, PUT, PATCH, and "
                    "DELETE require that connection's permission mode to be read_write. Use enabled connector "
                    "skill instructions, keep requests under that provider's configured base URL, and never "
                    "reveal the credential value. If connector skill docs show path-only routes such as "
                    "/app/native/path, pass that connector-relative path directly or combine it with the "
                    "configured connector gateway base; do not prepend the native third-party host. Use public "
                    "web/browser tools for generic public URLs instead."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
                requires_approval=False,
                timeout_seconds=20,
                capability_tags=("connector",),
            ),
            _build_connector_request_handler(),
        )
    try:
        registry.get_spec("email_search")
    except KeyError:
        registry.mark_plugin_installed("connector_plugin")
        registry.register(
            ToolSpec(
                name="email_search",
                description=(
                    "Search messages through an active Google Mail connector. Use this for inbox checks, "
                    "triage, and finding message ids. Search results are metadata-only but include enough "
                    "headers/snippets for broad inbox triage and prioritization. For broad multi-message "
                    "summaries, start from these previews and read only the few selected messages whose full "
                    "body is needed. Call email_read with a returned id before summarizing full body content "
                    "or claiming a specific message body was read."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=email_tool_timeout_seconds,
                input_schema=_email_search_input_schema(),
                capability_tags=("email", "connector", "account_read"),
                continuation_tools=("email_read",),
            ),
            _build_connector_email_search_handler(),
        )
    try:
        registry.get_spec("email_read")
    except KeyError:
        registry.mark_plugin_installed("connector_plugin")
        registry.register(
            ToolSpec(
                name="email_read",
                description=(
                    "Read one full message through an active Google Mail connector. The id must be exactly one "
                    "email_search.results[].id value; do not infer ids or pass placeholders. Use this whenever "
                    "the user asks to read, summarize, inspect, or act on message body content. HTML links are "
                    "returned as structured message.links entries when present, including hidden href targets. "
                    "Attachments are returned as structured message.attachments entries with attachmentId values."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=email_tool_timeout_seconds,
                input_schema=_email_read_input_schema(),
                capability_tags=("email", "connector", "account_read"),
            ),
            _build_connector_email_read_handler(),
        )
    try:
        registry.get_spec("email_attachment_read")
    except KeyError:
        registry.mark_plugin_installed("connector_plugin")
        registry.register(
            ToolSpec(
                name="email_attachment_read",
                description=(
                    "Fetch one Gmail attachment through an active Google Mail connector. The message_id and "
                    "attachment_id must be exact values returned from email_read.message.attachments[]. Use this "
                    "when the user asks for an attachment, checklist, document, file, invoice, receipt, or other "
                    "message attachment content."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=email_tool_timeout_seconds,
                input_schema=_email_attachment_read_input_schema(),
                capability_tags=("email", "connector", "account_read"),
                continuation_tools=("file_read", "terminal_exec"),
            ),
            _build_connector_email_attachment_read_handler(),
        )
    try:
        registry.get_spec("calendar_list")
    except KeyError:
        registry.mark_plugin_installed("connector_plugin")
        registry.register(
            ToolSpec(
                name="calendar_list",
                description=(
                    "List calendar events through an active Google Calendar connector for a specific time window. "
                    "Use this for agenda, schedule, availability, and calendar checks. If the user asks for a "
                    "calendar check without a time window, choose a near-term window such as today through the "
                    "next 7 days instead of asking for options first."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                input_schema=_calendar_list_input_schema(),
                capability_tags=("calendar", "connector", "account_read"),
            ),
            _build_connector_calendar_list_handler(),
        )
    for spec, handler in (
        (
            ToolSpec(
                name="calendar_create",
                description=(
                    "Create a Google Calendar event through an active write-capable Google Calendar connector. "
                    "Use calendar_list first when the event should be checked against existing calendar state. "
                    "Requires explicit approval before mutating the calendar."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
                requires_approval=True,
                timeout_seconds=20,
                input_schema=_calendar_create_input_schema(),
                capability_tags=("calendar", "connector", "account_write"),
                continuation_tools=("calendar_list", "calendar_update", "calendar_delete"),
            ),
            _build_connector_calendar_create_handler(),
        ),
        (
            ToolSpec(
                name="calendar_update",
                description=(
                    "Update or reschedule an existing Google Calendar event through an active write-capable "
                    "Google Calendar connector. The event_id must be an exact id from calendar_list or a prior "
                    "calendar write result. Requires explicit approval before mutating the calendar."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
                requires_approval=True,
                timeout_seconds=20,
                input_schema=_calendar_update_input_schema(),
                capability_tags=("calendar", "connector", "account_write"),
                continuation_tools=("calendar_list", "calendar_delete"),
            ),
            _build_connector_calendar_update_handler(),
        ),
        (
            ToolSpec(
                name="calendar_respond",
                description=(
                    "Accept, decline, or tentatively respond to a Google Calendar event through an active "
                    "write-capable Google Calendar connector. The event_id must be exact. Requires explicit "
                    "approval before mutating the calendar."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
                requires_approval=True,
                timeout_seconds=20,
                input_schema=_calendar_respond_input_schema(),
                capability_tags=("calendar", "connector", "account_write"),
                continuation_tools=("calendar_list",),
            ),
            _build_connector_calendar_respond_handler(),
        ),
        (
            ToolSpec(
                name="calendar_delete",
                description=(
                    "Delete/cancel an existing Google Calendar event through an active write-capable Google "
                    "Calendar connector. The event_id must be exact. Requires explicit approval before mutating "
                    "the calendar."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
                requires_approval=True,
                timeout_seconds=20,
                input_schema=_calendar_delete_input_schema(),
                capability_tags=("calendar", "connector", "account_write"),
                continuation_tools=("calendar_list",),
            ),
            _build_connector_calendar_delete_handler(),
        ),
    ):
        try:
            registry.get_spec(spec.name)
        except KeyError:
            registry.mark_plugin_installed("connector_plugin")
            registry.register(spec, handler)
    try:
        registry.get_spec("email_send")
        return
    except KeyError:
        pass
    registry.mark_plugin_installed("connector_plugin")
    registry.register(
        ToolSpec(
            name="email_send",
            description=(
                "Send a plain-text email, optionally with local artifact/media attachments, through an active "
                "write-capable Google Mail connector. Use this for actual email delivery; use connector_request "
                "only for lower-level connector APIs."
            ),
            risk_level=ToolRiskLevel.HIGH,
            side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
            requires_approval=True,
            timeout_seconds=20,
            capability_tags=("email", "connector"),
        ),
        _build_connector_email_send_handler(),
    )


def _connector_connection_for_invocation(invocation: ToolInvocation, provider_id: str):
    from .connections import require_workspace_connection_for_principal

    return require_workspace_connection_for_principal(invocation.principal_id, provider_id)


def _connector_credential_value(connection: object | None, provider_id: str) -> str:
    candidates = _connector_credential_candidate_names(connection, provider_id)
    for name in candidates:
        value = os.environ.get(name, "").strip()
        if value:
            return value
    labels = " or ".join(candidates) or "a workspace credential_ref"
    raise RuntimeError(f"{provider_id} requires {labels}")


def _connector_request_headers(connection: object | None, provider_id: str) -> dict[str, str]:
    headers = {
        "Accept": "application/json",
        "Authorization": f"Bearer {_connector_credential_value(connection, provider_id)}",
        "Connection": "close",
        "User-Agent": "Nullion connector_request/0.1",
    }
    return headers


_CONNECTOR_READ_METHODS = {"GET", "HEAD"}
_CONNECTOR_WRITE_METHODS = {"POST", "PUT", "PATCH", "DELETE"}
_CONNECTOR_ALLOWED_METHODS = _CONNECTOR_READ_METHODS | _CONNECTOR_WRITE_METHODS


def _connector_request_method(raw_method: object) -> str:
    method = str(raw_method or "GET").strip().upper() or "GET"
    if method not in _CONNECTOR_ALLOWED_METHODS:
        allowed = ", ".join(sorted(_CONNECTOR_ALLOWED_METHODS))
        raise ValueError(f"Unsupported connector_request method: {method}. Allowed methods: {allowed}")
    return method


def _connector_connection_allows_write(connection: object | None) -> bool:
    mode = str(getattr(connection, "permission_mode", "read") or "read").strip().lower().replace("-", "_")
    return mode in {"write", "read_write", "readwrite", "rw", "read_and_write"}


_CALENDAR_ACCOUNT_WRITE_TOOLS = frozenset({"calendar_create", "calendar_update", "calendar_respond", "calendar_delete"})


def _native_email_send_provider_has_write_connection(provider_id: str = "imap_smtp_provider") -> bool:
    try:
        from nullion.connections import load_connection_registry
    except Exception:
        return False
    try:
        connections = load_connection_registry().connections
    except Exception:
        return False
    for connection in connections:
        if str(getattr(connection, "provider_id", "") or "").strip() != provider_id:
            continue
        if not getattr(connection, "active", True):
            continue
        if _connector_connection_allows_write(connection):
            return True
    return False


def _string_list_argument(raw_value: object) -> list[str]:
    if raw_value is None:
        return []
    if isinstance(raw_value, str):
        value = raw_value.strip()
        return [value] if value else []
    if isinstance(raw_value, (list, tuple)):
        values: list[str] = []
        for item in raw_value:
            value = str(item or "").strip()
            if value:
                values.append(value)
        return values
    value = str(raw_value or "").strip()
    return [value] if value else []


def _default_email_connector_provider_id(principal_id: str | None, *, require_write: bool = False) -> str:
    try:
        from nullion.connections import default_email_connector_provider_id

        provider_id = default_email_connector_provider_id(principal_id, require_write=require_write)
        if provider_id:
            return provider_id
    except TypeError:
        try:
            from nullion.connections import default_email_connector_provider_id

            provider_id = default_email_connector_provider_id(principal_id)
            if provider_id:
                return provider_id
        except Exception:
            pass
    except Exception:
        pass
    raise RuntimeError("No active Google Mail connector is available for this workspace/principal.")


def _default_calendar_connector_provider_id(principal_id: str | None, *, require_write: bool = False) -> str:
    try:
        from nullion.connections import (
            connection_has_runtime_credentials,
            connection_for_principal,
            load_connection_registry,
        )
    except Exception:
        raise RuntimeError("No active Google Calendar connector is available for this workspace/principal.")
    try:
        connections = load_connection_registry().connections
    except Exception:
        connections = ()
    for connection in connections:
        provider_id = str(getattr(connection, "provider_id", "") or "").strip()
        if not provider_id or not getattr(connection, "active", True):
            continue
        normalized_provider = provider_id.lower()
        if not (
            normalized_provider.startswith("skill_pack_connector_")
            or normalized_provider.endswith("_connector_provider")
        ):
            continue
        scoped_connection = connection_for_principal(principal_id, provider_id)
        if scoped_connection is not None and require_write and not _connector_connection_allows_write(scoped_connection):
            continue
        if scoped_connection is not None and connection_has_runtime_credentials(scoped_connection):
            return provider_id
    requirement = "write-capable " if require_write else ""
    raise RuntimeError(f"No active {requirement}Google Calendar connector is available for this workspace/principal.")


def _email_messages_endpoint_for_provider(connection: object | None, provider_id: str) -> str:
    for base_url in _connector_allowed_base_urls(connection, provider_id):
        parsed = urlparse(base_url)
        if parsed.scheme and parsed.netloc:
            return f"{parsed.scheme}://{parsed.netloc}/google-mail/gmail/v1/users/me/messages"
    raise RuntimeError(f"{provider_id} does not have a configured connector base URL for Google Mail.")


def _email_send_endpoint_for_provider(connection: object | None, provider_id: str) -> str:
    for base_url in _connector_allowed_base_urls(connection, provider_id):
        parsed = urlparse(base_url)
        if parsed.scheme and parsed.netloc:
            return f"{parsed.scheme}://{parsed.netloc}/google-mail/gmail/v1/users/me/messages/send"
    raise RuntimeError(f"{provider_id} does not have a configured connector base URL for Google Mail.")


def _calendar_events_endpoint_for_provider(
    connection: object | None,
    provider_id: str,
    *,
    calendar_id: str = "primary",
) -> str:
    resolved_calendar_id = quote(str(calendar_id or "primary").strip() or "primary", safe="")
    for base_url in _connector_allowed_base_urls(connection, provider_id):
        parsed = urlparse(base_url)
        if parsed.scheme and parsed.netloc:
            return f"{parsed.scheme}://{parsed.netloc}/google-calendar/calendar/v3/calendars/{resolved_calendar_id}/events"
    raise RuntimeError(f"{provider_id} does not have a configured connector base URL for Google Calendar.")


def _calendar_event_endpoint_for_provider(
    connection: object | None,
    provider_id: str,
    *,
    event_id: str,
    calendar_id: str = "primary",
) -> str:
    return f"{_calendar_events_endpoint_for_provider(connection, provider_id, calendar_id=calendar_id)}/{quote(event_id, safe='')}"


def _gmail_header_map(message: dict[str, object]) -> dict[str, str]:
    payload = message.get("payload")
    if not isinstance(payload, dict):
        return {}
    headers = payload.get("headers")
    if not isinstance(headers, list):
        return {}
    wanted = {"from", "to", "cc", "subject", "date"}
    mapped: dict[str, str] = {}
    for header in headers:
        if not isinstance(header, dict):
            continue
        name = str(header.get("name") or "").strip().lower()
        value = str(header.get("value") or "").strip()
        if name in wanted and value:
            mapped[name] = value
    return mapped


def _gmail_decode_body_data(data: object) -> str:
    if not isinstance(data, str) or not data.strip():
        return ""
    padded = data + ("=" * ((4 - len(data) % 4) % 4))
    try:
        return base64.urlsafe_b64decode(padded.encode("ascii")).decode("utf-8", "replace")
    except Exception:
        return ""


def _gmail_decode_attachment_bytes(data: object) -> bytes:
    if not isinstance(data, str) or not data.strip():
        return b""
    padded = data + ("=" * ((4 - len(data) % 4) % 4))
    try:
        return base64.urlsafe_b64decode(padded.encode("ascii"))
    except Exception:
        return b""


_EMAIL_HTML_DROP_BLOCK_RE = re.compile(
    r"(?is)<\s*(?:script|style|head|noscript|svg|template)\b[^>]*>.*?<\s*/\s*(?:script|style|head|noscript|svg|template)\s*>"
)
_EMAIL_HTML_BLOCK_BREAK_RE = re.compile(
    r"(?i)<\s*/?\s*(?:br|p|div|tr|table|tbody|thead|section|article|header|footer|li|ul|ol|h[1-6])\b[^>]*>"
)
_EMAIL_LINK_MAX_COUNT = 24
_EMAIL_LINK_TEXT_MAX_CHARS = 180
_EMAIL_LINK_URL_MAX_CHARS = 2048
_EMAIL_VISIBLE_URL_RE = re.compile(r"(?i)\bhttps?://[^\s<>()\"']+")
_EMAIL_LINK_FOLLOWUP_TOOLS = ("browser_navigate", "browser_extract_text", "web_fetch")


class _EmailHtmlLinkParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.links: list[dict[str, str]] = []
        self._active_links: list[dict[str, object]] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() != "a":
            return
        href = ""
        for name, value in attrs:
            if name.lower() == "href" and value:
                href = value
                break
        if not href:
            return
        self._active_links.append({"url": href, "text_parts": []})

    def handle_data(self, data: str) -> None:
        if not data or not self._active_links:
            return
        for link in self._active_links:
            parts = link.get("text_parts")
            if isinstance(parts, list):
                parts.append(data)

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() != "a" or not self._active_links:
            return
        link = self._active_links.pop()
        url = _normalize_email_link_url(str(link.get("url") or ""))
        if not url:
            return
        parts = link.get("text_parts")
        text = (
            _clean_email_link_text(" ".join(str(part) for part in parts))
            if isinstance(parts, list)
            else ""
        )
        entry = {"url": url}
        if text:
            entry["text"] = text
        self.links.append(entry)


def _strip_email_invisible_text(text: str) -> str:
    chars: list[str] = []
    for char in str(text or ""):
        category = unicodedata.category(char)
        if category == "Cf" or char in {"\u034f", "\ufeff"}:
            continue
        if category.startswith("C") and char not in {"\n", "\r", "\t"}:
            continue
        chars.append(char)
    return "".join(chars)


def _clean_email_body_text(text: str) -> str:
    cleaned = _strip_email_invisible_text(unescape(str(text or "")))
    cleaned = re.sub(r"(?m)^[ \t]*[.#]?[A-Za-z0-9_-]{1,80}\s*\{[^{}]{0,1000}\}[ \t]*$", " ", cleaned)
    cleaned = re.sub(r"(?m)^[ \t]*(?:[A-Za-z-]+\s*:\s*[^;{}]{0,180};\s*){2,}[ \t]*$", " ", cleaned)
    cleaned = re.sub(r"[ \t\r\f\v]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return " ".join(cleaned.split())


def _plain_text_from_email_html(html_body: str) -> str:
    text = _EMAIL_HTML_DROP_BLOCK_RE.sub(" ", str(html_body or ""))
    text = _EMAIL_HTML_BLOCK_BREAK_RE.sub("\n", text)
    text = re.sub(r"(?is)<[^>]+>", " ", text)
    return _clean_email_body_text(text)


def _clean_email_link_text(text: str) -> str:
    cleaned = _clean_email_body_text(text)
    return cleaned[:_EMAIL_LINK_TEXT_MAX_CHARS].rstrip()


def _normalize_email_link_url(url: str) -> str:
    cleaned = _strip_email_invisible_text(unescape(str(url or ""))).strip()
    cleaned = cleaned.replace("\r", "").replace("\n", "")
    if not cleaned or cleaned.startswith("#"):
        return ""
    parsed = urlparse(cleaned)
    scheme = parsed.scheme.lower()
    if scheme in {"javascript", "data", "cid"}:
        return ""
    return cleaned[:_EMAIL_LINK_URL_MAX_CHARS].rstrip()


def _email_link_entries_from_visible_text(text: str) -> list[dict[str, str]]:
    entries: list[dict[str, str]] = []
    for match in _EMAIL_VISIBLE_URL_RE.finditer(str(text or "")):
        url = _normalize_email_link_url(match.group(0).rstrip(".,;:!?)]}"))
        if url:
            entries.append({"url": url})
    return entries


def _email_link_entries_from_html(html_body: str) -> list[dict[str, str]]:
    parser = _EmailHtmlLinkParser()
    try:
        parser.feed(str(html_body or ""))
        parser.close()
    except Exception:
        return []
    return parser.links


def _gmail_message_attachments(message: dict[str, object], *, limit: int = 25) -> list[dict[str, object]]:
    payload = message.get("payload")
    if not isinstance(payload, dict):
        return []
    message_id = str(message.get("id") or "").strip()
    attachments: list[dict[str, object]] = []
    seen: set[tuple[str, str]] = set()

    def visit(part: object, fallback_part_id: str = "") -> None:
        if not isinstance(part, dict):
            return
        part_id = str(part.get("partId") or fallback_part_id).strip()
        filename = str(part.get("filename") or "").strip()
        mime_type = str(part.get("mimeType") or "").strip()
        body = part.get("body")
        attachment_id = ""
        size = None
        if isinstance(body, dict):
            attachment_id = str(body.get("attachmentId") or "").strip()
            raw_size = body.get("size")
            if isinstance(raw_size, int):
                size = raw_size
        if filename or attachment_id:
            identity = (attachment_id, filename)
            if identity not in seen:
                seen.add(identity)
                entry: dict[str, object] = {}
                if filename:
                    entry["filename"] = filename[:240]
                if mime_type:
                    entry["mimeType"] = mime_type[:160]
                if attachment_id:
                    entry["attachmentId"] = attachment_id
                if message_id:
                    entry["messageId"] = message_id
                if part_id:
                    entry["partId"] = part_id
                if isinstance(size, int):
                    entry["size"] = size
                attachments.append(entry)
        children = part.get("parts")
        if isinstance(children, list):
            for index, child in enumerate(children):
                visit(child, f"{part_id}.{index}" if part_id else str(index))

    visit(payload)
    return attachments[:limit]


def _gmail_message_body_parts(message: dict[str, object]) -> list[tuple[str, str]]:
    payload = message.get("payload")
    if not isinstance(payload, dict):
        return []
    candidates: list[tuple[str, str]] = []

    def visit(part: object) -> None:
        if not isinstance(part, dict):
            return
        mime_type = str(part.get("mimeType") or "").strip().lower()
        body = part.get("body")
        if isinstance(body, dict):
            text = _gmail_decode_body_data(body.get("data"))
            if text.strip():
                candidates.append((mime_type, text))
        children = part.get("parts")
        if isinstance(children, list):
            for child in children:
                visit(child)

    visit(payload)
    return candidates


def _email_text_signal_score(text: str) -> int:
    cleaned = _clean_email_body_text(text)
    if not cleaned:
        return 0
    words = re.findall(r"[A-Za-z0-9][A-Za-z0-9'@./:-]{1,}", cleaned)
    if words:
        normalized_words = [word.casefold() for word in words]
        unique_words = set(normalized_words)
        if len(words) >= 20 and len(unique_words) <= max(3, len(words) // 12):
            return len(unique_words)
    css_markers = len(
        re.findall(
            r"(?:!important|line-height|font-weight|text-decoration|background-color|webkit)",
            cleaned,
            flags=re.I,
        )
    )
    return max(0, len(words) - (css_markers * 8))


def _gmail_message_body_text(message: dict[str, object], *, limit: int = 6000) -> str:
    candidates = _gmail_message_body_parts(message)
    cleaned_candidates: list[tuple[str, str]] = []
    for mime, value in candidates:
        cleaned = _plain_text_from_email_html(value) if mime == "text/html" else _clean_email_body_text(value)
        if cleaned:
            cleaned_candidates.append((mime, cleaned))
    text = next(
        (value for mime, value in cleaned_candidates if mime == "text/plain" and _email_text_signal_score(value) >= 12),
        "",
    )
    if not text:
        text = next((value for mime, value in cleaned_candidates if mime == "text/html"), "")
    if not text and cleaned_candidates:
        text = max((value for _, value in cleaned_candidates), key=_email_text_signal_score)
    snippet = message.get("snippet")
    snippet_text = _clean_email_body_text(str(snippet or "")) if isinstance(snippet, str) else ""
    if snippet_text and (not text or _email_text_signal_score(snippet_text) > _email_text_signal_score(text) * 2):
        text = snippet_text
    return text[:limit]


def _gmail_message_links(message: dict[str, object], *, limit: int = _EMAIL_LINK_MAX_COUNT) -> list[dict[str, str]]:
    entries: list[dict[str, str]] = []
    for mime, value in _gmail_message_body_parts(message):
        if mime == "text/html":
            entries.extend(_email_link_entries_from_html(value))
            entries.extend(_email_link_entries_from_visible_text(_plain_text_from_email_html(value)))
        else:
            entries.extend(_email_link_entries_from_visible_text(value))
    snippet = message.get("snippet")
    if isinstance(snippet, str):
        entries.extend(_email_link_entries_from_visible_text(snippet))

    deduped: list[dict[str, str]] = []
    seen_urls: set[str] = set()
    for entry in entries:
        url = _normalize_email_link_url(str(entry.get("url") or ""))
        if not url or url in seen_urls:
            continue
        seen_urls.add(url)
        cleaned_entry = {"url": url}
        text = _clean_email_link_text(str(entry.get("text") or ""))
        if text:
            cleaned_entry["text"] = text
        deduped.append(cleaned_entry)
        if len(deduped) >= limit:
            break
    return deduped


def _email_links_include_http_url(links: Iterable[dict[str, str]]) -> bool:
    for entry in links:
        parsed = urlparse(str(entry.get("url") or ""))
        if parsed.scheme.lower() in {"http", "https"} and parsed.netloc:
            return True
    return False


def _compact_gmail_message(message: dict[str, object], *, include_body: bool = False) -> dict[str, object]:
    compact: dict[str, object] = {
        "id": message.get("id"),
        "threadId": message.get("threadId"),
    }
    headers = _gmail_header_map(message)
    if headers:
        compact["headers"] = headers
        for key in ("from", "subject", "date"):
            if key in headers:
                compact[key] = headers[key]
    snippet = message.get("snippet")
    if isinstance(snippet, str) and snippet.strip():
        compact["snippet"] = _clean_email_body_text(snippet)
    label_ids = message.get("labelIds")
    if isinstance(label_ids, list):
        compact["labelIds"] = [str(item) for item in label_ids[:12]]
    attachments = _gmail_message_attachments(message)
    if attachments:
        compact["attachments"] = attachments
        compact["attachment_count"] = len(attachments)
        compact["next_tool_for_attachment"] = "email_attachment_read"
        compact["attachment_requires_tool"] = "email_attachment_read"
        compact["next_tool"] = "email_attachment_read"
    if include_body:
        body = _gmail_message_body_text(message)
        if body:
            compact["body"] = body
        links = _gmail_message_links(message)
        if links:
            compact["links"] = links
            if _email_links_include_http_url(links):
                compact["next_tools_for_links"] = list(_EMAIL_LINK_FOLLOWUP_TOOLS)
    return {key: value for key, value in compact.items() if value not in (None, "", [], {})}


def _compact_google_calendar_event(event: dict[str, object]) -> dict[str, object]:
    compact: dict[str, object] = {
        "id": event.get("id"),
        "summary": event.get("summary"),
        "status": event.get("status"),
        "start": event.get("start"),
        "end": event.get("end"),
        "location": event.get("location"),
        "description": event.get("description"),
        "htmlLink": event.get("htmlLink"),
    }
    attendees = event.get("attendees")
    if isinstance(attendees, list):
        compact["attendees"] = [
            {
                key: attendee.get(key)
                for key in ("email", "displayName", "responseStatus")
                if isinstance(attendee, dict) and attendee.get(key) not in (None, "", [], {})
            }
            for attendee in attendees[:20]
            if isinstance(attendee, dict)
        ]
    return {key: value for key, value in compact.items() if value not in (None, "", [], {})}


def _calendar_selected_results_payload(
    results: list[dict[str, object]],
    *,
    query: str,
) -> dict[str, object]:
    normalized_query = str(query or "").strip().casefold()
    if not normalized_query:
        return {}
    selected: list[dict[str, object]] = []
    selected_ids: list[str] = []
    for item in results:
        values = (
            str(item.get("summary") or ""),
            str(item.get("description") or ""),
            str(item.get("location") or ""),
        )
        if not any(normalized_query in value.casefold() for value in values if value):
            continue
        selected.append(item)
        event_id = str(item.get("id") or item.get("event_id") or item.get("uid") or "").strip()
        if event_id:
            selected_ids.append(event_id)
    if not selected:
        return {}
    return {
        "selected_results": selected[:3],
        **({"selected_result_ids": selected_ids[:3]} if selected_ids else {}),
    }


def _connector_json_payload_from_response(
    response,
    *,
    max_body_bytes: int = _CONNECTOR_JSON_MAX_BODY_BYTES,
) -> dict[str, object]:
    body_bytes = response.read(max_body_bytes + 1)
    if len(body_bytes) > max_body_bytes:
        raise RuntimeError(f"connector response exceeded the {max_body_bytes}-byte limit")
    body = body_bytes.decode("utf-8", "ignore")
    try:
        payload = json.loads(body or "{}")
    except json.JSONDecodeError as exc:
        raise RuntimeError("connector returned invalid JSON") from exc
    if not isinstance(payload, dict):
        raise RuntimeError("connector returned non-object JSON")
    return payload


def _connector_json_request(
    invocation: ToolInvocation,
    *,
    provider_id: str,
    url: str,
    params: dict[str, object] | None = None,
    method: str = "GET",
    json_payload: object | None = None,
    max_response_bytes: int = _CONNECTOR_JSON_MAX_BODY_BYTES,
) -> dict[str, object]:
    connection = _connector_connection_for_invocation(invocation, provider_id)
    normalized_method = _connector_request_method(method)
    if normalized_method not in _CONNECTOR_READ_METHODS and not _connector_connection_allows_write(connection):
        raise RuntimeError(f"{provider_id} is configured as read-only for connector_request.")
    request_url = _connector_request_url(url, params or {}, connection, provider_id)
    resolution = _resolve_web_fetch_resolution(request_url)
    headers = _connector_request_headers(connection, provider_id)
    data = None
    if json_payload is not None:
        headers["Content-Type"] = "application/json"
        data = json.dumps(json_payload).encode("utf-8")
    request = urllib.request.Request(request_url, data=data, headers=headers, method=normalized_method)
    opener = _build_web_fetch_opener()
    try:
        with _pinned_web_fetch_resolution(resolution):
            with opener.open(request, timeout=_connector_request_timeout_seconds()) as response:
                return _connector_json_payload_from_response(response, max_body_bytes=max_response_bytes)
    except urllib.error.HTTPError as exc:
        _output, error = _connector_http_error_output(
            exc,
            provider_id=provider_id,
            method=normalized_method,
            url=request_url,
        )
        raise RuntimeError(error) from exc


def _connector_request_timeout_seconds() -> int:
    try:
        raw = int(os.environ.get("NULLION_CONNECTOR_REQUEST_TIMEOUT_SECONDS", "12"))
    except ValueError:
        raw = 12
    return max(3, min(raw, 30))


def _email_message_for_invocation(invocation: ToolInvocation) -> tuple[EmailMessage, list[str]]:
    recipients = _string_list_argument(invocation.arguments.get("to"))
    if not recipients:
        raise ValueError("Missing required argument: to")
    subject = str(invocation.arguments.get("subject") or "").strip()
    body = str(invocation.arguments.get("body") or "")
    html_body = _email_html_body_from_invocation(invocation)
    if html_body and not body.strip():
        body = _plain_text_from_html_body(html_body)
    msg = EmailMessage()
    msg["To"] = ", ".join(recipients)
    cc = _string_list_argument(invocation.arguments.get("cc"))
    bcc = _string_list_argument(invocation.arguments.get("bcc"))
    if cc:
        msg["Cc"] = ", ".join(cc)
    if bcc:
        msg["Bcc"] = ", ".join(bcc)
    msg["Subject"] = subject
    msg.set_content(body)
    if html_body:
        from nullion.artifacts import normalize_html_document

        msg.add_alternative(normalize_html_document(html_body, title=subject or "Email"), subtype="html")

    attached_paths: list[str] = []
    effective_roots = _principal_workspace_file_roots(invocation.principal_id)
    for raw_path in _string_list_argument(invocation.arguments.get("attachment_paths")):
        resolved = _resolve_local_workspace_file_input(
            raw_path,
            principal_id=invocation.principal_id,
            effective_roots=effective_roots,
            trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
        )
        if resolved is None:
            resolved = Path(_resolve_virtual_workspace_path(raw_path, principal_id=invocation.principal_id)).expanduser()
        if not resolved.exists() or not resolved.is_file():
            raise FileNotFoundError(f"Attachment path does not exist or is not a file: {raw_path}")
        content_type, _encoding = mimetypes.guess_type(str(resolved))
        maintype, subtype = (content_type or "application/octet-stream").split("/", 1)
        msg.add_attachment(
            resolved.read_bytes(),
            maintype=maintype,
            subtype=subtype,
            filename=resolved.name,
        )
        attached_paths.append(str(resolved))
    return msg, attached_paths


def _connector_request_payload(invocation: ToolInvocation, method: str) -> tuple[bytes | None, dict[str, str]]:
    if method in _CONNECTOR_READ_METHODS:
        return None, {}
    headers: dict[str, str] = {}
    if "json" in invocation.arguments and invocation.arguments.get("json") is not None:
        headers["Content-Type"] = "application/json"
        return json.dumps(invocation.arguments.get("json")).encode("utf-8"), headers
    raw_body = invocation.arguments.get("body")
    if raw_body is None:
        return None, {}
    headers["Content-Type"] = "text/plain; charset=utf-8"
    return str(raw_body).encode("utf-8"), headers


def _connector_http_error_output(
    exc: urllib.error.HTTPError,
    *,
    provider_id: str,
    method: str | None = None,
    url: str | None = None,
) -> tuple[dict[str, object], str]:
    body = exc.read(20000).decode("utf-8", "ignore")
    output: dict[str, object] = {
        "provider_id": provider_id,
        "status_code": exc.code,
        "content_type": exc.headers.get("content-type"),
    }
    if method:
        output["method"] = method
    if url:
        output["url"] = url
    parsed_json: object | None = None
    try:
        parsed_json = json.loads(body)
    except Exception:
        parsed_json = None
    if isinstance(parsed_json, dict):
        output["json"] = parsed_json
        message = str(parsed_json.get("message") or "").strip()
        app_match = re.search(r"connections for `([^`]+)` are either PENDING or FAILED", message)
        if app_match:
            output["connection_state"] = "pending_or_failed"
            output["connector_app_id"] = app_match.group(1)
            return output, (
                f"{app_match.group(1)} connection is pending or failed. "
                "Reconnect the account before using this connector."
            )
        if message:
            return output, message
    elif body.strip():
        output["text"] = body[:20000]
        return output, body.strip()[:500]
    return output, f"HTTP Error {exc.code}: {exc.reason}"


def _connector_base_url_candidate_names(connection: object | None, provider_id: str) -> tuple[str, ...]:
    names: list[str] = []
    credential_ref = getattr(connection, "credential_ref", None)
    if isinstance(credential_ref, str) and credential_ref.strip():
        ref = credential_ref.strip().removeprefix("env:")
        if ref.endswith(("_API_KEY", "_TOKEN", "_SECRET_KEY")):
            names.append(re.sub(r"_(API_KEY|TOKEN|SECRET_KEY)$", "_BASE_URL", ref))
    prefix = _connector_provider_env_prefix(provider_id)
    if prefix:
        names.append(f"{prefix}_BASE_URL")
        for suffix in ("_CONNECTOR_PROVIDER", "_CONNECTOR"):
            if prefix.endswith(suffix):
                gateway_prefix = prefix.removesuffix(suffix).strip("_")
                if gateway_prefix:
                    names.append(f"{gateway_prefix}_BASE_URL")
        if prefix.startswith("SKILL_PACK_CONNECTOR_"):
            gateway_prefix = prefix.removeprefix("SKILL_PACK_CONNECTOR_").strip("_")
            if gateway_prefix:
                names.append(f"{gateway_prefix}_BASE_URL")
    return tuple(dict.fromkeys(names))


def _connector_allowed_base_urls(connection: object | None, provider_id: str) -> tuple[str, ...]:
    urls: list[str] = []
    provider_profile = getattr(connection, "provider_profile", None)
    if isinstance(provider_profile, str) and provider_profile.strip().lower().startswith(("http://", "https://")):
        urls.append(provider_profile.strip())
    for name in _connector_base_url_candidate_names(connection, provider_id):
        raw = os.environ.get(name, "").strip()
        if raw.lower().startswith(("http://", "https://")):
            urls.append(raw)
    urls.extend(_connector_skill_pack_base_urls(provider_id))
    normalized: list[str] = []
    for url in urls:
        parsed = urlparse(url)
        if parsed.scheme in {"http", "https"} and parsed.hostname:
            normalized.append(url.rstrip("/") + "/")
    return tuple(dict.fromkeys(normalized))


def _connector_skill_pack_base_urls(provider_id: str) -> tuple[str, ...]:
    pack_id = _installed_connector_skill_pack_id_for_provider(provider_id)
    if not pack_id:
        return ()
    try:
        from nullion.skill_pack_installer import get_installed_skill_pack

        pack = get_installed_skill_pack(pack_id)
    except Exception:
        pack = None
    pack_path = Path(str(getattr(pack, "path", "") or ""))
    if not pack_path.exists():
        return ()
    urls: list[str] = []
    for skill_file in sorted(pack_path.rglob("SKILL.md")):
        try:
            text = skill_file.read_text(encoding="utf-8", errors="replace")
        except Exception:
            continue
        trusted_hosts: set[str] = set()
        base_context_lines = 0
        for line in text.splitlines():
            normalized_line = line.lower()
            starts_base_context = (
                "base url" in normalized_line
                or "base_url" in normalized_line
                or "endpoint base" in normalized_line
            )
            if starts_base_context:
                base_context_lines = 8
            is_base_line = starts_base_context or base_context_lines > 0
            for match in re.finditer(r"https?://[^\s`'\"<>)]+", line):
                raw_url = match.group(0).rstrip(".,;:")
                parsed = urlparse(raw_url)
                if parsed.scheme not in {"http", "https"} or not parsed.hostname:
                    continue
                if any(marker in parsed.hostname for marker in ("{", "}", "<", ">")):
                    continue
                path = parsed.path or "/"
                path_has_template = "{" in path or "<" in path
                host_key = parsed.netloc.lower()
                if not is_base_line and not (path_has_template and host_key in trusted_hosts):
                    continue
                cut_positions = [pos for marker in ("{", "<") if (pos := path.find(marker)) >= 0]
                if cut_positions:
                    path = path[: min(cut_positions)]
                urls.append(f"{parsed.scheme}://{parsed.netloc}{path.rstrip('/')}/")
                trusted_hosts.add(host_key)
            if base_context_lines > 0:
                base_context_lines -= 1
    return tuple(dict.fromkeys(urls))


def _url_is_under_base(url: str, base_url: str) -> bool:
    parsed = urlparse(url)
    base = urlparse(base_url)
    if parsed.scheme != base.scheme or parsed.hostname != base.hostname:
        return False
    parsed_port = parsed.port or (443 if parsed.scheme == "https" else 80)
    base_port = base.port or (443 if base.scheme == "https" else 80)
    if parsed_port != base_port:
        return False
    base_path = (base.path or "/").rstrip("/") + "/"
    path = (parsed.path or "/").rstrip("/") + "/"
    return path.startswith(base_path)


def _connector_request_url(raw_url: str, raw_params: object, connection: object | None = None, provider_id: str = "") -> str:
    url = raw_url.strip()
    allowed_bases = _connector_allowed_base_urls(connection, provider_id)
    parsed = urlparse(url)
    host = parsed.hostname
    if (
        parsed.scheme not in {"http", "https"}
        or not isinstance(host, str)
        or not host
    ):
        if not allowed_bases:
            raise ValueError(f"Blocked URL for connector_request: {url}")
        if parsed.scheme or parsed.netloc or not url or url.startswith(("#", "?")):
            raise ValueError(f"Blocked URL for connector_request: {url}")
        url = urljoin(allowed_bases[0], url.lstrip("/"))
        parsed = urlparse(url)
        host = parsed.hostname
        if (
            parsed.scheme not in {"http", "https"}
            or not isinstance(host, str)
            or not host
        ):
            raise ValueError(f"Blocked URL for connector_request: {url}")
    params: dict[str, str] = {}
    if isinstance(raw_params, dict):
        for key, value in raw_params.items():
            if value is None:
                continue
            params[str(key)] = str(value)
    if params:
        separator = "&" if parsed.query else "?"
        url += separator + urlencode(params)
    is_under_configured_base = bool(allowed_bases) and any(_url_is_under_base(url, base_url) for base_url in allowed_bases)
    if allowed_bases and not is_under_configured_base:
        labels = ", ".join(allowed_bases)
        raise ValueError(f"Blocked URL for connector_request: {url} is not under configured connector base URL(s): {labels}")
    resolution = _resolve_web_fetch_resolution(url)
    if resolution is None and not _is_global_literal_ip(host) and not is_under_configured_base:
        raise ValueError(f"Blocked URL for connector_request: {url}")
    return url


def _build_connector_request_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        raw_url = invocation.arguments.get("url")
        if not provider_id:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: provider_id",
            )
        if not isinstance(raw_url, str) or not raw_url.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: url",
            )
        try:
            connection = _connector_connection_for_invocation(invocation, provider_id)
            method = _connector_request_method(invocation.arguments.get("method"))
            if method not in _CONNECTOR_READ_METHODS and not _connector_connection_allows_write(connection):
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"provider_id": provider_id, "method": method, "permission_mode": "read"},
                    error=(
                        f"{provider_id} is configured as read-only for connector_request. "
                        "Change this connection's permission mode to read_write in Settings > Users > Connections "
                        "before using POST, PUT, PATCH, or DELETE."
                    ),
                )
            url = _connector_request_url(raw_url, invocation.arguments.get("params"), connection, provider_id)
            resolution = _resolve_web_fetch_resolution(url)
            payload, payload_headers = _connector_request_payload(invocation, method)
            headers = _connector_request_headers(connection, provider_id)
            headers.update(payload_headers)
            request = urllib.request.Request(
                url,
                data=payload,
                headers=headers,
                method=method,
            )
            opener = _build_web_fetch_opener()
            with _pinned_web_fetch_resolution(resolution):
                with opener.open(request, timeout=20) as response:
                    status_code = getattr(response, "status", 200)
                    content_type = response.headers.get_content_type()
                    body = response.read(1_000_000).decode("utf-8", "ignore")
            output: dict[str, object] = {
                "url": url,
                "provider_id": provider_id,
                "method": method,
                "status_code": status_code,
                "content_type": content_type,
            }
            try:
                output["json"] = json.loads(body)
            except Exception:
                output["text"] = body[:20000]
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output=output,
                error=None,
            )
        except urllib.error.HTTPError as exc:
            output, error = _connector_http_error_output(
                exc,
                provider_id=provider_id,
                method=method if "method" in locals() else None,
                url=url if "url" in locals() else None,
            )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=output,
                error=error,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"provider_id": provider_id},
                error=str(exc),
            )

    return handler


def _build_connector_email_search_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_query = invocation.arguments.get("query")
        raw_limit = invocation.arguments.get("limit", 10)
        if not isinstance(raw_query, str) or not raw_query.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: query",
            )
        if not isinstance(raw_limit, int) or raw_limit <= 0:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="limit must be a positive integer",
            )
        limit = min(raw_limit, 10)
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            if not provider_id:
                provider_id = _default_email_connector_provider_id(invocation.principal_id)
            connection = _connector_connection_for_invocation(invocation, provider_id)
            endpoint = _email_messages_endpoint_for_provider(connection, provider_id)
            listing = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
                params={"q": raw_query.strip(), "maxResults": limit},
            )
            raw_messages = listing.get("messages")
            messages = raw_messages if isinstance(raw_messages, list) else []
            results: list[dict[str, object]] = []
            detail_fetch_limit = _email_search_detail_fetch_limit()
            for index, item in enumerate(messages[:limit]):
                if not isinstance(item, dict):
                    continue
                message_id = str(item.get("id") or "").strip()
                if not message_id:
                    continue
                if index < detail_fetch_limit:
                    try:
                        detail = _connector_json_request(
                            invocation,
                            provider_id=provider_id,
                            url=f"{endpoint}/{quote(message_id, safe='')}",
                            params={
                                "format": "metadata",
                                "fields": "id,threadId,labelIds,snippet,payload(headers)",
                            },
                        )
                    except Exception:
                        detail = item
                    results.append(_compact_gmail_message(detail, include_body=False))
                else:
                    results.append(_compact_gmail_message(item, include_body=False))
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "query": raw_query.strip(),
                    "provider_id": provider_id,
                    "resultSizeEstimate": listing.get("resultSizeEstimate"),
                    "results": results,
                    "body_included": False,
                    "next_tool_for_body": "email_read",
                    "body_requires_tool": "email_read",
                },
                error=None,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=_account_tool_failure_output(invocation.principal_id, query=raw_query),
                error=str(exc),
            )

    return handler


def _email_attachment_safe_filename(value: object, *, mime_type: str = "") -> tuple[str, str, str]:
    raw_name = Path(str(value or "").strip()).name
    suffix = Path(raw_name).suffix.lower()
    if not suffix:
        guessed = mimetypes.guess_extension(mime_type.split(";", 1)[0].strip()) if mime_type else None
        suffix = (guessed or ".bin").lower()
    if len(suffix) > 16 or not re.fullmatch(r"\.[a-z0-9][a-z0-9._-]{0,14}", suffix):
        suffix = ".bin"
    stem = _safe_pdf_stem(Path(raw_name).stem or "email-attachment")
    filename = f"{stem}{suffix}"
    return filename, stem, suffix


def _extract_pdf_text(data: bytes, *, limit: int = 12000) -> tuple[str, int, bool]:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    page_count = len(reader.pages)
    parts: list[str] = []
    total_chars = 0
    truncated = False
    for page in reader.pages:
        page_text = (page.extract_text() or "").strip()
        if not page_text:
            continue
        parts.append(page_text)
        total_chars += len(page_text)
        if total_chars > limit:
            truncated = True
            break
    text = "\n\n".join(parts).replace("\r\n", "\n").replace("\r", "\n").strip()
    text = re.sub(r"\n{3,}", "\n\n", text)
    if len(text) > limit:
        text = text[:limit].rstrip()
        truncated = True
    return text, page_count, truncated


def _email_attachment_text_preview(data: bytes, *, mime_type: str, filename: str, limit: int = 6000) -> str:
    normalized_type = mime_type.split(";", 1)[0].strip().lower()
    suffix = Path(filename).suffix.lower()
    if normalized_type == "application/pdf" or suffix == ".pdf":
        try:
            text, _page_count, _truncated = _extract_pdf_text(data, limit=limit)
        except Exception:
            return ""
        return text
    if normalized_type in {"text/plain", "text/csv", "text/html", "text/calendar", "application/json"} or suffix in {
        ".txt",
        ".csv",
        ".html",
        ".htm",
        ".ics",
        ".json",
        ".md",
        ".xml",
    }:
        text = data.decode("utf-8", "replace")
        if normalized_type == "text/html" or suffix in {".html", ".htm"}:
            text = _plain_text_from_email_html(text)
        else:
            text = _clean_email_body_text(text)
        return text[:limit]
    return ""


def _build_connector_email_attachment_read_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_message_id = invocation.arguments.get("message_id")
        raw_attachment_id = invocation.arguments.get("attachment_id")
        message_id = raw_message_id.strip() if isinstance(raw_message_id, str) else ""
        attachment_id = raw_attachment_id.strip() if isinstance(raw_attachment_id, str) else ""
        if not message_id or not attachment_id:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "reason": "missing_email_attachment_id",
                    "required_source_tool": "email_read",
                    "next_step": "Read the email first and pass one returned message.attachments[].attachmentId to email_attachment_read.",
                },
                error="Missing required message_id or attachment_id.",
            )
        if _invalid_concrete_resource_id_reason(message_id) or _invalid_concrete_resource_id_reason(attachment_id):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "reason": "invalid_email_attachment_id",
                    "required_source_tool": "email_read",
                    "next_step": "Use exact message_id and attachment_id values returned by email_read.message.attachments.",
                },
                error="Invalid email attachment id: use ids returned by email_read.",
            )
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        mime_type = str(invocation.arguments.get("mime_type") or "").strip()
        filename, stem, suffix = _email_attachment_safe_filename(
            invocation.arguments.get("filename"),
            mime_type=mime_type,
        )
        try:
            if not provider_id:
                provider_id = _default_email_connector_provider_id(invocation.principal_id)
            connection = _connector_connection_for_invocation(invocation, provider_id)
            endpoint = _email_messages_endpoint_for_provider(connection, provider_id)
            payload = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=f"{endpoint}/{quote(message_id, safe='')}/attachments/{quote(attachment_id, safe='')}",
                max_response_bytes=_EMAIL_ATTACHMENT_CONNECTOR_MAX_BODY_BYTES,
            )
            data = _gmail_decode_attachment_bytes(payload.get("data"))
            if not data:
                raise RuntimeError("connector returned an empty attachment body")
            from nullion.artifacts import artifact_path_for_generated_workspace_file

            path = artifact_path_for_generated_workspace_file(
                principal_id=invocation.principal_id,
                suffix=suffix,
                stem=stem or "email-attachment",
            )
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(data)
            if not mime_type:
                mime_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
            attachment_role = (
                ARTIFACT_ROLE_DELIVERABLE
                if _flow_context_requires_artifact_delivery_for_extension(invocation.flow_context, suffix)
                else ARTIFACT_ROLE_SOURCE
            )
            output: dict[str, object] = {
                "provider_id": provider_id,
                "message_id": message_id,
                "attachment_id": attachment_id,
                "filename": filename,
                "mime_type": mime_type,
                "size_bytes": len(data),
                "artifact_path": str(path),
                "artifact_paths": [str(path)],
                "artifact_descriptors": [
                    artifact_output_descriptor(
                        path,
                        role=attachment_role,
                        kind="email_attachment",
                        label=filename,
                    )
                ],
            }
            text_preview = _email_attachment_text_preview(data, mime_type=mime_type, filename=filename)
            if text_preview:
                output["text_preview"] = text_preview
                normalized_type = mime_type.split(";", 1)[0].strip().lower()
                if normalized_type == "application/pdf" or suffix == ".pdf":
                    try:
                        _pdf_text, page_count, preview_truncated = _extract_pdf_text(data, limit=6000)
                    except Exception:
                        page_count = 0
                        preview_truncated = False
                    if page_count:
                        output["page_count"] = page_count
                        output["text_extraction_method"] = "pypdf"
                        output["text_preview_truncated"] = preview_truncated
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output=output,
                error=None,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "provider_id": provider_id,
                    "message_id": message_id,
                    "attachment_id": attachment_id,
                    "required_source_tool": "email_read",
                },
                error=str(exc),
            )

    return handler


def _build_connector_email_read_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_id = invocation.arguments.get("id")
        if not isinstance(raw_id, str) or not raw_id.strip():
            return _invalid_email_message_id_result(invocation, raw_id)
        if _invalid_concrete_resource_id_reason(raw_id):
            return _invalid_email_message_id_result(invocation, raw_id)
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            if not provider_id:
                provider_id = _default_email_connector_provider_id(invocation.principal_id)
            connection = _connector_connection_for_invocation(invocation, provider_id)
            endpoint = _email_messages_endpoint_for_provider(connection, provider_id)
            message = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=f"{endpoint}/{quote(raw_id.strip(), safe='')}",
                params={"format": "full"},
            )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "id": raw_id.strip(),
                    "provider_id": provider_id,
                    "message": _compact_gmail_message(message, include_body=True),
                },
                error=None,
            )
        except Exception as exc:
            return _email_read_source_required_failure_result(invocation, raw_id, exc)

    return handler


def _build_connector_calendar_list_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_start = invocation.arguments.get("start")
        raw_end = invocation.arguments.get("end")
        raw_max = invocation.arguments.get("max", 10)
        query = str(invocation.arguments.get("query") or "").strip()
        if not isinstance(raw_start, str) or not raw_start.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: start",
            )
        if not isinstance(raw_end, str) or not raw_end.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: end",
            )
        if not isinstance(raw_max, int) or raw_max <= 0:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="max must be a positive integer",
            )
        limit = min(raw_max, 50)
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            if not provider_id:
                provider_id = _default_calendar_connector_provider_id(invocation.principal_id)
            connection = _connector_connection_for_invocation(invocation, provider_id)
            endpoint = _calendar_events_endpoint_for_provider(connection, provider_id)
            listing = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
                params={
                    "timeMin": raw_start.strip(),
                    "timeMax": raw_end.strip(),
                    "maxResults": limit,
                    "singleEvents": "true",
                    "orderBy": "startTime",
                },
            )
            raw_items = listing.get("items")
            items = raw_items if isinstance(raw_items, list) else []
            results = [
                _compact_google_calendar_event(item)
                for item in items[:limit]
                if isinstance(item, dict)
            ]
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "provider_id": provider_id,
                    "start": raw_start.strip(),
                    "end": raw_end.strip(),
                    "max": limit,
                    **({"query": query} if query else {}),
                    "result_count": len(results),
                    "results": results,
                    **_calendar_selected_results_payload(results, query=query),
                },
                error=None,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=_account_tool_failure_output(invocation.principal_id),
                error=str(exc),
            )

    return handler


def _calendar_write_read_only_result(invocation: ToolInvocation, provider_id: str) -> ToolResult:
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={"provider_id": provider_id, "permission_mode": "read"},
        error=(
            f"{provider_id} is configured as read-only for {invocation.tool_name}. "
            "Change this connection's permission mode to read_write in Settings > Users > Connections "
            "before modifying calendar events."
        ),
    )


def _invalid_calendar_event_id_result(invocation: ToolInvocation, raw_id: object) -> ToolResult:
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={
            "reason": "invalid_event_id",
            "id": raw_id if isinstance(raw_id, str) else "",
            "required_source_tool": "calendar_list",
            "next_step": "List calendar events and pass one returned results[].id value.",
        },
        error="Invalid calendar event id: use an exact id returned by calendar_list.",
    )


def _calendar_write_provider_and_connection(invocation: ToolInvocation) -> tuple[str, object]:
    provider_id = str(invocation.arguments.get("provider_id") or "").strip()
    if not provider_id:
        provider_id = _default_calendar_connector_provider_id(invocation.principal_id, require_write=True)
    connection = _connector_connection_for_invocation(invocation, provider_id)
    if not _connector_connection_allows_write(connection):
        raise PermissionError(provider_id)
    return provider_id, connection


def _calendar_id_from_invocation(invocation: ToolInvocation) -> str:
    raw = str(invocation.arguments.get("calendar_id") or "primary").strip()
    if not raw:
        return "primary"
    if _invalid_concrete_resource_id_reason(raw):
        raise ValueError("Invalid calendar_id.")
    return raw


def _calendar_send_updates_params(invocation: ToolInvocation) -> dict[str, object]:
    raw = str(invocation.arguments.get("send_updates") or "all").strip()
    if raw not in {"all", "externalOnly", "none"}:
        raise ValueError("send_updates must be one of all, externalOnly, or none.")
    return {"sendUpdates": raw}


def _calendar_time_payload(value: object, *, time_zone: str) -> dict[str, object]:
    text = str(value or "").strip()
    if not text:
        raise ValueError("Missing calendar event time value.")
    if re.fullmatch(r"\d{4}-\d{2}-\d{2}", text):
        return {"date": text}
    payload: dict[str, object] = {"dateTime": text}
    if time_zone:
        payload["timeZone"] = time_zone
    return payload


def _calendar_event_payload_from_invocation(
    invocation: ToolInvocation,
    *,
    require_summary: bool = False,
    require_times: bool = False,
) -> dict[str, object]:
    payload: dict[str, object] = {}
    summary = str(invocation.arguments.get("summary") or "").strip()
    if summary:
        payload["summary"] = summary
    elif require_summary:
        raise ValueError("Missing required argument: summary")
    time_zone = str(invocation.arguments.get("time_zone") or "").strip()
    raw_start = invocation.arguments.get("start")
    raw_end = invocation.arguments.get("end")
    if raw_start is not None and str(raw_start).strip():
        payload["start"] = _calendar_time_payload(raw_start, time_zone=time_zone)
    elif require_times:
        raise ValueError("Missing required argument: start")
    if raw_end is not None and str(raw_end).strip():
        payload["end"] = _calendar_time_payload(raw_end, time_zone=time_zone)
    elif require_times:
        raise ValueError("Missing required argument: end")
    for argument_name, event_key in (("location", "location"), ("description", "description")):
        if argument_name in invocation.arguments:
            payload[event_key] = str(invocation.arguments.get(argument_name) or "").strip()
    attendees = _string_list_argument(invocation.arguments.get("attendees"))
    if attendees:
        payload["attendees"] = [{"email": email} for email in attendees]
    return payload


def _calendar_write_result(
    invocation: ToolInvocation,
    *,
    provider_id: str,
    action: str,
    event: dict[str, object] | None = None,
    event_id: str | None = None,
) -> ToolResult:
    compact_event = _compact_google_calendar_event(event or {}) if isinstance(event, dict) else {}
    resolved_id = str(event_id or compact_event.get("id") or "").strip()
    summary = str(compact_event.get("summary") or invocation.arguments.get("summary") or "").strip()
    output: dict[str, object] = {
        "action": action,
        "provider_id": provider_id,
        "id": resolved_id,
        "event_id": resolved_id,
        "summary": summary,
        "message": f"Calendar event {action.removeprefix('calendar_').replace('_', ' ')} completed.",
    }
    if compact_event:
        output["event"] = compact_event
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="completed",
        output={key: value for key, value in output.items() if value not in (None, "", [], {})},
        error=None,
    )


def _build_connector_calendar_create_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            provider_id, connection = _calendar_write_provider_and_connection(invocation)
            endpoint = _calendar_events_endpoint_for_provider(
                connection,
                provider_id,
                calendar_id=_calendar_id_from_invocation(invocation),
            )
            event = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
                params=_calendar_send_updates_params(invocation),
                method="POST",
                json_payload=_calendar_event_payload_from_invocation(
                    invocation,
                    require_summary=True,
                    require_times=True,
                ),
            )
            return _calendar_write_result(invocation, provider_id=provider_id, action="calendar_create", event=event)
        except PermissionError as exc:
            return _calendar_write_read_only_result(invocation, str(exc))
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"provider_id": provider_id},
                error=str(exc),
            )

    return handler


def _build_connector_calendar_update_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_event_id = invocation.arguments.get("event_id")
        if not isinstance(raw_event_id, str) or _invalid_concrete_resource_id_reason(raw_event_id):
            return _invalid_calendar_event_id_result(invocation, raw_event_id)
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            provider_id, connection = _calendar_write_provider_and_connection(invocation)
            payload = _calendar_event_payload_from_invocation(invocation)
            if not payload:
                raise ValueError("Provide at least one event field to update.")
            endpoint = _calendar_event_endpoint_for_provider(
                connection,
                provider_id,
                calendar_id=_calendar_id_from_invocation(invocation),
                event_id=raw_event_id.strip(),
            )
            event = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
                params=_calendar_send_updates_params(invocation),
                method="PATCH",
                json_payload=payload,
            )
            return _calendar_write_result(
                invocation,
                provider_id=provider_id,
                action="calendar_update",
                event=event,
                event_id=raw_event_id.strip(),
            )
        except PermissionError as exc:
            return _calendar_write_read_only_result(invocation, str(exc))
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"provider_id": provider_id},
                error=str(exc),
            )

    return handler


def _build_connector_calendar_respond_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_event_id = invocation.arguments.get("event_id")
        if not isinstance(raw_event_id, str) or _invalid_concrete_resource_id_reason(raw_event_id):
            return _invalid_calendar_event_id_result(invocation, raw_event_id)
        response_status = str(invocation.arguments.get("response_status") or "").strip()
        if response_status not in {"accepted", "declined", "tentative", "needsAction"}:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"event_id": raw_event_id},
                error="response_status must be one of accepted, declined, tentative, or needsAction.",
            )
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            provider_id, connection = _calendar_write_provider_and_connection(invocation)
            attendee_email = str(invocation.arguments.get("attendee_email") or "").strip()
            endpoint = _calendar_event_endpoint_for_provider(
                connection,
                provider_id,
                calendar_id=_calendar_id_from_invocation(invocation),
                event_id=raw_event_id.strip(),
            )
            existing_event = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
            )
            existing_attendees = existing_event.get("attendees") if isinstance(existing_event, dict) else None
            attendees = [dict(item) for item in existing_attendees if isinstance(item, dict)] if isinstance(existing_attendees, list) else []
            if attendee_email:
                matched = False
                for attendee in attendees:
                    if str(attendee.get("email") or "").strip().casefold() == attendee_email.casefold():
                        attendee["responseStatus"] = response_status
                        matched = True
                if not matched:
                    attendees.append({"email": attendee_email, "responseStatus": response_status})
            else:
                matched = False
                for attendee in attendees:
                    if attendee.get("self") is True:
                        attendee["responseStatus"] = response_status
                        matched = True
                if not matched:
                    raise ValueError("Could not find the current user's attendee record; provide attendee_email.")
            event = _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
                params=_calendar_send_updates_params(invocation),
                method="PATCH",
                json_payload={"attendees": attendees},
            )
            return _calendar_write_result(
                invocation,
                provider_id=provider_id,
                action="calendar_respond",
                event=event,
                event_id=raw_event_id.strip(),
            )
        except PermissionError as exc:
            return _calendar_write_read_only_result(invocation, str(exc))
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"provider_id": provider_id, "event_id": raw_event_id},
                error=str(exc),
            )

    return handler


def _build_connector_calendar_delete_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_event_id = invocation.arguments.get("event_id")
        if not isinstance(raw_event_id, str) or _invalid_concrete_resource_id_reason(raw_event_id):
            return _invalid_calendar_event_id_result(invocation, raw_event_id)
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            provider_id, connection = _calendar_write_provider_and_connection(invocation)
            endpoint = _calendar_event_endpoint_for_provider(
                connection,
                provider_id,
                calendar_id=_calendar_id_from_invocation(invocation),
                event_id=raw_event_id.strip(),
            )
            _connector_json_request(
                invocation,
                provider_id=provider_id,
                url=endpoint,
                params=_calendar_send_updates_params(invocation),
                method="DELETE",
            )
            return _calendar_write_result(
                invocation,
                provider_id=provider_id,
                action="calendar_delete",
                event_id=raw_event_id.strip(),
            )
        except PermissionError as exc:
            return _calendar_write_read_only_result(invocation, str(exc))
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"provider_id": provider_id, "event_id": raw_event_id},
                error=str(exc),
            )

    return handler


def _build_connector_email_send_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        provider_id = str(invocation.arguments.get("provider_id") or "").strip()
        try:
            if not provider_id:
                provider_id = _default_email_connector_provider_id(invocation.principal_id, require_write=True)
            connection = _connector_connection_for_invocation(invocation, provider_id)
            if not _connector_connection_allows_write(connection):
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"provider_id": provider_id, "permission_mode": "read"},
                    error=(
                        f"{provider_id} is configured as read-only for email_send. "
                        "Change this connection's permission mode to read_write in Settings > Users > Connections "
                        "before sending email."
                    ),
                )
            message, attached_paths = _email_message_for_invocation(invocation)
            raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode("ascii").rstrip("=")
            endpoint = _email_send_endpoint_for_provider(connection, provider_id)
            url = _connector_request_url(endpoint, None, connection, provider_id)
            resolution = _resolve_web_fetch_resolution(url)
            payload = json.dumps({"raw": raw_message}).encode("utf-8")
            headers = _connector_request_headers(connection, provider_id)
            headers["Content-Type"] = "application/json"
            request = urllib.request.Request(url, data=payload, headers=headers, method="POST")
            opener = _build_web_fetch_opener()
            with _pinned_web_fetch_resolution(resolution):
                with opener.open(request, timeout=20) as response:
                    status_code = getattr(response, "status", 200)
                    content_type = response.headers.get_content_type()
                    body = response.read(1_000_000).decode("utf-8", "ignore")
            output: dict[str, object] = {
                "url": url,
                "provider_id": provider_id,
                "method": "POST",
                "status_code": status_code,
                "content_type": content_type,
                "to": _string_list_argument(invocation.arguments.get("to")),
                "subject": str(invocation.arguments.get("subject") or "").strip(),
                "attachment_count": len(attached_paths),
                "attachment_paths": attached_paths,
            }
            try:
                output["json"] = json.loads(body)
            except Exception:
                output["text"] = body[:20000]
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output=output,
                error=None,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"provider_id": provider_id},
                error=str(exc),
            )

    return handler


def _build_skill_pack_read_handler() -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        pack_id = str(invocation.arguments.get("pack_id") or "").strip()
        path = str(invocation.arguments.get("path") or "").strip()
        if not pack_id:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: pack_id",
            )
        if not path:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"pack_id": pack_id},
                error="Missing required argument: path",
            )
        try:
            from .skill_pack_installer import read_skill_pack_reference

            text = read_skill_pack_reference(pack_id, path)
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"pack_id": pack_id, "path": path},
                error=str(exc),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"pack_id": pack_id, "path": path, "text": text},
            error=None,
        )

    return handler


def _build_web_search_handler(
    web_searcher: Callable[[str, int], list[dict[str, object]]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_query = invocation.arguments.get("query")
        raw_limit = invocation.arguments.get("limit", 5)
        if not isinstance(raw_query, str) or not raw_query:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: query",
            )

        try:
            limit = int(raw_limit)
        except (TypeError, ValueError):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Invalid argument: limit",
            )
        if limit < 1:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"limit must be at least 1, got {limit}",
            )

        try:
            results = web_searcher(raw_query, limit)
        except Exception as exc:  # pragma: no cover - caller-provided searcher guard
            status_code = getattr(exc, "code", None)
            systemic_http_failure = isinstance(status_code, int) and (
                status_code in {401, 402, 403, 408, 429} or status_code >= 500
            )
            failure_scope = "tool" if systemic_http_failure or isinstance(exc, urllib.error.URLError) else "invocation"
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "query": raw_query,
                    "reason": "web_search_failed",
                    "failure_scope": failure_scope,
                    **({"status_code": status_code} if isinstance(status_code, int) else {}),
                },
                error=str(exc),
            )

        usable_results = [
            result
            for result in results
            if _web_search_result_has_usable_evidence(result)
        ]
        if not usable_results:
            candidates = [
                _web_search_candidate_for_fallback(result)
                for result in results[:limit]
            ]
            candidates = [candidate for candidate in candidates if candidate is not None]
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "query": raw_query,
                    "reason": "no_usable_search_results",
                    "result_count": len(results),
                    "candidates": candidates,
                },
                error="web_search returned result rows without usable source evidence",
            )

        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"query": raw_query, "results": usable_results},
            error=None,
        )

    return handler


def _web_search_result_has_usable_evidence(result: dict[str, object]) -> bool:
    url = result.get("url")
    if not isinstance(url, str) or not url.strip():
        return False
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        return False
    snippet = str(
        result.get("snippet")
        or result.get("summary")
        or result.get("description")
        or ""
    ).strip()
    return bool(snippet)


def _web_search_candidate_for_fallback(result: dict[str, object]) -> dict[str, str] | None:
    url = result.get("url")
    if not isinstance(url, str) or not url.strip():
        return None
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        return None
    candidate: dict[str, str] = {"url": url}
    title = result.get("title")
    if isinstance(title, str) and title.strip():
        candidate["title"] = title.strip()
    return candidate



def _build_file_search_handler(
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
    *,
    include_principal_workspace: bool = True,
) -> ToolHandler:
    resolved_root = Path(workspace_root).resolve() if workspace_root is not None else None
    resolved_allowed_roots = tuple(Path(root).resolve() for root in allowed_roots) if allowed_roots is not None else None

    def _resolve_requested_root(raw_root: object, search_roots: tuple[Path, ...]) -> tuple[Path, ...] | ToolResult | None:
        root_text = str(raw_root or "").strip()
        if not root_text:
            return None
        candidates: list[Path] = []
        requested = Path(root_text).expanduser()
        if requested.is_absolute():
            candidates.append(requested)
        else:
            for root in search_roots:
                candidates.append(root / requested)
        for candidate in candidates:
            try:
                resolved_candidate = candidate.resolve()
            except OSError:
                continue
            if not _path_within_any_root(resolved_candidate, search_roots):
                continue
            if not resolved_candidate.exists():
                continue
            if not resolved_candidate.is_dir():
                return ToolResult(
                    invocation_id="",
                    tool_name="file_search",
                    status="failed",
                    output={
                        "reason": "search_root_not_directory",
                        "root": str(resolved_candidate),
                    },
                    error=f"Search root is not a directory: {resolved_candidate}",
                )
            return (resolved_candidate,)
        return ToolResult(
            invocation_id="",
            tool_name="file_search",
            status="failed",
            output={
                "reason": "search_root_outside_allowed_roots",
                "root": root_text,
                "allowed_roots": [str(root) for root in search_roots],
            },
            error="Search root must be inside the configured workspace or allowed roots.",
        )

    def _file_contains_pattern(path: Path, pattern_text: str, *, max_bytes: int) -> bool:
        try:
            with path.open("rb") as handle:
                data = handle.read(max_bytes + 1)
        except OSError:
            return False
        if not data:
            return False
        if b"\x00" in data:
            return False
        try:
            text = data[:max_bytes].decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = data[:max_bytes].decode("latin-1")
            except UnicodeDecodeError:
                return False
        return pattern_text in text.lower()

    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_pattern = invocation.arguments.get("pattern")
        if not isinstance(raw_pattern, str) or not raw_pattern:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: pattern",
            )

        search_roots = _effective_filesystem_roots(
            invocation=invocation,
            resolved_root=resolved_root,
            resolved_allowed_roots=resolved_allowed_roots,
            include_principal_workspace=include_principal_workspace,
        )
        if not search_roots:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="file_search requires workspace_root or allowed_roots",
            )
        requested_roots = _resolve_requested_root(invocation.arguments.get("root"), search_roots)
        if isinstance(requested_roots, ToolResult):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status=requested_roots.status,
                output=requested_roots.output,
                error=requested_roots.error,
            )
        if requested_roots:
            search_roots = requested_roots
        else:
            search_roots = tuple(sorted(search_roots, key=lambda root: len(root.parts), reverse=True))

        raw_limit = invocation.arguments.get("limit")
        limit = 100
        if isinstance(raw_limit, int) and raw_limit > 0:
            limit = min(raw_limit, 500)

        pattern = raw_pattern.lower()
        search_contents = bool(invocation.arguments.get("search_contents"))
        max_entries = _env_int("NULLION_FILE_SEARCH_MAX_ENTRIES", 25_000, minimum=1)
        max_content_bytes = _env_int("NULLION_FILE_SEARCH_CONTENT_MAX_BYTES", 65_536, minimum=1)
        matches: list[str] = []
        match_details: list[dict[str, object]] = []
        seen_matches: set[str] = set()
        scanned_entries = 0
        scanned_files = 0
        truncated = False
        for root in search_roots:
            for current_dir, dirnames, filenames in os.walk(root, topdown=True, followlinks=False):
                scanned_entries += 1
                if scanned_entries > max_entries:
                    truncated = True
                    dirnames[:] = []
                    break
                current_path = Path(current_dir)
                scoped_dirnames: list[str] = []
                for dirname in sorted(dirnames):
                    scanned_entries += 1
                    if scanned_entries > max_entries:
                        truncated = True
                        break
                    candidate_dir = current_path / dirname
                    if candidate_dir.is_symlink() or _should_prune_filesystem_walk_dir(candidate_dir):
                        continue
                    try:
                        resolved_dir = candidate_dir.resolve()
                    except OSError:
                        continue
                    if not _path_within_any_root(resolved_dir, search_roots):
                        continue
                    scoped_dirnames.append(dirname)
                dirnames[:] = scoped_dirnames

                for filename in sorted(filenames):
                    scanned_entries += 1
                    if scanned_entries > max_entries:
                        truncated = True
                        break
                    candidate_file = current_path / filename
                    if candidate_file.is_symlink():
                        continue
                    try:
                        if not candidate_file.is_file():
                            continue
                        resolved_path = candidate_file.resolve()
                    except OSError:
                        continue
                    if not _path_within_any_root(resolved_path, search_roots):
                        continue
                    scanned_files += 1
                    filename_match = pattern in filename.lower()
                    content_match = False
                    if not filename_match and search_contents:
                        content_match = _file_contains_pattern(
                            resolved_path,
                            pattern,
                            max_bytes=max_content_bytes,
                        )
                    if not filename_match and not content_match:
                        continue
                    resolved_text = str(resolved_path)
                    if resolved_text in seen_matches:
                        continue
                    seen_matches.add(resolved_text)
                    matches.append(resolved_text)
                    match_details.append({
                        "path": resolved_text,
                        "match_type": "filename" if filename_match else "content",
                    })
                    if len(matches) >= limit:
                        break
                if len(matches) >= limit:
                    break
                if truncated:
                    break
            if len(matches) >= limit:
                break
            if truncated:
                break

        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "matches": matches,
                "match_details": match_details,
                "searched_roots": [str(root) for root in search_roots],
                "searched_file_count": scanned_files,
                "truncated": truncated,
                "search_contents": search_contents,
            },
            error=None,
        )

    return handler


def register_search_plugin(
    registry: ToolRegistry,
    *,
    web_fetcher: Callable[[str, int], dict[str, object]] | None = None,
    web_searcher: Callable[[str, int], list[dict[str, object]]] | None = None,
) -> ToolRegistry:
    del web_fetcher
    registry.require_plugin_registration_allowed()
    registry.mark_plugin_installed("search_plugin")
    if web_searcher is not None:
        try:
            registry.get_spec("web_search")
        except KeyError:
            registry.register(
                ToolSpec(
                    name="web_search",
                    description=(
                        "Search the web and return ranked source links with evidence snippets. "
                        "Use this when no specific URL is already available. If the result has "
                        "no usable evidence snippets, treat it as insufficient for a final answer "
                        "and continue from structured candidate URLs or another structured source."
                    ),
                    risk_level=ToolRiskLevel.LOW,
                    side_effect_class=ToolSideEffectClass.READ,
                    requires_approval=False,
                    timeout_seconds=20,
                    continuation_tools=(
                        "web_fetch",
                        "browser_navigate",
                    ),
                ),
                _build_web_search_handler(web_searcher),
            )
    return registry



def register_web_extension(
    registry: ToolRegistry,
    *,
    web_fetcher: Callable[[str, int], dict[str, object]],
    web_searcher: Callable[[str, int], list[dict[str, object]]] | None = None,
) -> ToolRegistry:
    return register_search_plugin(
        registry,
        web_fetcher=web_fetcher,
        web_searcher=web_searcher,
    )



def register_email_plugin(
    registry: ToolRegistry,
    *,
    email_searcher: Callable[[str, int], list[dict[str, object]]] | None = None,
    email_reader: Callable[[str], dict[str, object]] | None = None,
    email_sender: Callable[[EmailMessage, list[str]], dict[str, object]] | None = None,
) -> ToolRegistry:
    registry.require_plugin_registration_allowed()
    registry.mark_plugin_installed("email_plugin")
    if email_searcher is None:
        raise ValueError("email_plugin requires email_searcher")
    registry.unregister("email_search")
    registry.register(
        ToolSpec(
            name="email_search",
                description=(
                    "Search email messages via the configured provider. Search results are metadata-only; "
                    "use their headers/snippets for broad inbox triage, then call email_read with a returned "
                    "id only when a specific message's full body is needed."
                ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.READ,
            requires_approval=False,
            timeout_seconds=20,
            input_schema=_email_search_input_schema(),
            capability_tags=("email", "connector", "account_read"),
            continuation_tools=("email_read",),
        ),
        _build_email_search_handler(email_searcher),
    )
    if email_reader is not None:
        registry.unregister("email_read")
        registry.register(
            ToolSpec(
                name="email_read",
                description=(
                    "Read a single email message via the configured provider. The id must be exactly one "
                    "email_search.results[].id value; do not infer ids or pass placeholders."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                input_schema=_email_read_input_schema(),
                capability_tags=("email", "connector", "account_read"),
            ),
            _build_email_read_handler(email_reader),
        )
    if email_sender is not None:
        try:
            registry.get_spec("email_send")
            existing_email_send = True
        except KeyError:
            existing_email_send = False
        if existing_email_send and not _native_email_send_provider_has_write_connection():
            return registry
        registry.unregister("email_send")
        registry.register(
            ToolSpec(
                name="email_send",
                description=(
                    "Send a plain-text email, optionally with local artifact/media attachments, "
                    "through the configured email provider."
                ),
                risk_level=ToolRiskLevel.HIGH,
                side_effect_class=ToolSideEffectClass.ACCOUNT_WRITE,
                requires_approval=True,
                timeout_seconds=20,
                capability_tags=("email", "account_write"),
            ),
            _build_email_send_handler(email_sender),
        )
    return registry



def register_calendar_plugin(
    registry: ToolRegistry,
    *,
    calendar_lister: Callable[[str, str, int], list[dict[str, object]]] | None = None,
) -> ToolRegistry:
    registry.require_plugin_registration_allowed()
    registry.mark_plugin_installed("calendar_plugin")
    if calendar_lister is None:
        raise ValueError("calendar_plugin requires calendar_lister")
    try:
        registry.get_spec("calendar_list")
    except KeyError:
        registry.register(
            ToolSpec(
                name="calendar_list",
                description="List calendar events via the configured provider.",
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                input_schema=_calendar_list_input_schema(include_provider_id=False),
                capability_tags=("calendar", "connector", "account_read"),
            ),
            _build_calendar_list_handler(calendar_lister),
        )
    return registry



def _build_email_search_handler(
    email_searcher: Callable[[str, int], list[dict[str, object]]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_query = invocation.arguments.get("query")
        raw_limit = invocation.arguments.get("limit", 10)
        if not isinstance(raw_query, str) or not raw_query:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: query",
            )
        if not isinstance(raw_limit, int) or raw_limit <= 0:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="limit must be a positive integer",
            )
        try:
            results = _call_provider_with_principal(email_searcher, raw_query, raw_limit, principal_id=invocation.principal_id)
        except Exception as exc:  # pragma: no cover - provider guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=_account_tool_failure_output(invocation.principal_id, query=raw_query),
                error=str(exc),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "query": raw_query,
                "results": results,
                "body_included": False,
                "next_tool_for_body": "email_read",
                "body_requires_tool": "email_read",
            },
            error=None,
        )

    return handler



def _build_email_read_handler(
    email_reader: Callable[[str], dict[str, object]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_id = invocation.arguments.get("id")
        if not isinstance(raw_id, str) or not raw_id:
            return _invalid_email_message_id_result(invocation, raw_id)
        if _invalid_concrete_resource_id_reason(raw_id):
            return _invalid_email_message_id_result(invocation, raw_id)
        try:
            message = _call_provider_with_principal(email_reader, raw_id, principal_id=invocation.principal_id)
        except Exception as exc:  # pragma: no cover - provider guard
            return _email_read_source_required_failure_result(invocation, raw_id, exc)
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"id": raw_id, "message": message},
            error=None,
        )

    return handler


def _build_email_send_handler(
    email_sender: Callable[[EmailMessage, list[str]], dict[str, object]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        try:
            message, attached_paths = _email_message_for_invocation(invocation)
            output = _call_provider_with_principal(
                email_sender,
                message,
                attached_paths,
                principal_id=invocation.principal_id,
            )
            if not isinstance(output, dict):
                output = {"result": output}
            output.setdefault("to", _string_list_argument(invocation.arguments.get("to")))
            output.setdefault("subject", str(invocation.arguments.get("subject") or "").strip())
            output.setdefault("attachment_count", len(attached_paths))
            output.setdefault("attachment_paths", attached_paths)
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output=output,
                error=None,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=_account_tool_failure_output(invocation.principal_id),
                error=str(exc),
            )

    return handler


def _call_provider_with_principal(provider: Callable, *args, principal_id: str):
    try:
        parameters = inspect.signature(provider).parameters
    except (TypeError, ValueError):
        parameters = {}
    if "principal_id" in parameters:
        return provider(*args, principal_id=principal_id)
    return provider(*args)



def _build_calendar_list_handler(
    calendar_lister: Callable[[str, str, int], list[dict[str, object]]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_start = invocation.arguments.get("start")
        raw_end = invocation.arguments.get("end")
        raw_max = invocation.arguments.get("max", 10)
        query = str(invocation.arguments.get("query") or "").strip()
        if not isinstance(raw_start, str) or not raw_start:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: start",
            )
        if not isinstance(raw_end, str) or not raw_end:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: end",
            )
        if not isinstance(raw_max, int) or raw_max <= 0:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="max must be a positive integer",
            )
        try:
            results = calendar_lister(raw_start, raw_end, raw_max)
        except Exception as exc:  # pragma: no cover - provider guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=_account_tool_failure_output(invocation.principal_id),
                error=str(exc),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "start": raw_start,
                "end": raw_end,
                "max": raw_max,
                **({"query": query} if query else {}),
                "results": results,
                **_calendar_selected_results_payload(results, query=query),
            },
            error=None,
        )

    return handler


def _missing_media_provider_result(invocation: ToolInvocation, capability: str) -> ToolResult:
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={
            "reason": "provider_not_configured",
            "capability": capability,
            "setup": format_setup_tip(MEDIA_PROVIDER_SETUP_TIP),
        },
        error=f"{capability} provider is not configured",
    )


def _media_tool_enabled(enabled_env: str) -> bool:
    enabled = os.environ.get(enabled_env)
    return enabled is None or enabled.strip().lower() not in {"0", "false", "no", "off"}


def _build_audio_transcribe_handler(
    audio_transcriber: Callable[[str, str | None], dict[str, object]] | None,
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_path = invocation.arguments.get("path")
        raw_language = invocation.arguments.get("language")
        if not isinstance(raw_path, str) or not raw_path:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: path",
            )
        if raw_language is not None and not isinstance(raw_language, str):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="language must be a string when provided",
            )
        if audio_transcriber is None:
            return _missing_media_provider_result(invocation, "audio_transcribe")
        try:
            payload = audio_transcriber(raw_path, raw_language)
        except Exception as exc:  # pragma: no cover - provider guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=str(exc),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"path": raw_path, **payload},
            error=None,
        )

    return handler


def _build_image_extract_text_handler(
    image_text_extractor: Callable[[str], dict[str, object]] | None,
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_path = invocation.arguments.get("path")
        if not isinstance(raw_path, str) or not raw_path:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: path",
            )
        if image_text_extractor is None:
            return _missing_media_provider_result(invocation, "image_extract_text")
        try:
            payload = image_text_extractor(raw_path)
        except Exception as exc:  # pragma: no cover - provider guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=str(exc),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"path": raw_path, **payload},
            error=None,
        )

    return handler


def _build_image_generate_handler(
    image_generator: Callable[..., dict[str, object]] | None,
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_prompt = invocation.arguments.get("prompt")
        raw_output_path = invocation.arguments.get("output_path")
        raw_size = invocation.arguments.get("size")
        raw_source_path = invocation.arguments.get("source_path")
        if not isinstance(raw_prompt, str) or not raw_prompt.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: prompt",
            )
        if not isinstance(raw_output_path, str) or not raw_output_path.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: output_path",
            )
        if raw_size is not None and not isinstance(raw_size, str):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="size must be a string when provided",
            )
        if raw_source_path is not None and not isinstance(raw_source_path, str):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="source_path must be a string when provided",
            )
        if image_generator is None:
            try:
                from nullion.providers import _fallback_svg_image_generate

                # Keep the core image tool visible even without a raster model:
                # callers get a real local artifact plus setup guidance.
                payload = _fallback_svg_image_generate(
                    raw_prompt,
                    raw_output_path,
                    raw_size,
                    fallback_error="No image-generation provider was registered for this runtime.",
                )
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="completed",
                    output={"output_path": raw_output_path, **payload},
                    error=None,
                )
            except Exception as exc:  # pragma: no cover - fallback filesystem guard
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "reason": "provider_not_configured",
                        "capability": "image_generate",
                        "setup": format_setup_tip(MEDIA_PROVIDER_SETUP_TIP),
                    },
                    error=str(exc),
                )
        try:
            if raw_source_path:
                try:
                    accepts_source = len(inspect.signature(image_generator).parameters) >= 4
                except (TypeError, ValueError):
                    accepts_source = False
                if accepts_source:
                    payload = image_generator(raw_prompt, raw_output_path, raw_size, raw_source_path)
                else:
                    payload = image_generator(raw_prompt, raw_output_path, raw_size)
            else:
                payload = image_generator(raw_prompt, raw_output_path, raw_size)
        except Exception as exc:  # pragma: no cover - provider guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=str(exc),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"output_path": raw_output_path, **payload},
            error=None,
        )

    return handler


def register_media_plugin(
    registry: ToolRegistry,
    *,
    audio_transcriber: Callable[[str, str | None], dict[str, object]] | None = None,
    image_text_extractor: Callable[[str], dict[str, object]] | None = None,
    image_generator: Callable[[str, str, str | None], dict[str, object]] | None = None,
) -> ToolRegistry:
    registry.require_plugin_registration_allowed()
    registry.mark_plugin_installed("media_plugin")
    try:
        registry.get_spec("audio_transcribe")
    except KeyError:
        registry.register(
            ToolSpec(
                name="audio_transcribe",
                description="Transcribe a local audio file to text using the configured local media provider.",
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=60,
                filesystem_boundary_policy=_FILESYSTEM_BOUNDARY_TRUSTED_ROOTS_ONLY,
            ),
            _build_audio_transcribe_handler(audio_transcriber),
        )
    try:
        registry.get_spec("image_extract_text")
    except KeyError:
        registry.register(
            ToolSpec(
                name="image_extract_text",
                description="Extract visible text from a local image using the configured local OCR provider.",
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=60,
                filesystem_boundary_policy=_FILESYSTEM_BOUNDARY_TRUSTED_ROOTS_ONLY,
            ),
            _build_image_extract_text_handler(image_text_extractor),
        )
    if _media_tool_enabled("NULLION_IMAGE_GENERATE_ENABLED"):
        try:
            registry.get_spec("image_generate")
        except KeyError:
            registry.register(
                ToolSpec(
                    name="image_generate",
                    description=(
                        "Generate an image file using the configured image model when available, "
                        "or a safe local SVG fallback when no image model is configured."
                    ),
                    risk_level=ToolRiskLevel.MEDIUM,
                    side_effect_class=ToolSideEffectClass.WRITE,
                    requires_approval=False,
                    timeout_seconds=120,
                    capability_tags=("media", "image_generation"),
                ),
                _build_image_generate_handler(image_generator),
            )
    return registry



def _build_browser_navigate_handler(
    browser_navigator: Callable[[str], dict[str, object]],
) -> ToolHandler:
    def handler(invocation: ToolInvocation) -> ToolResult:
        raw_url = invocation.arguments.get("url")
        if not isinstance(raw_url, str) or not raw_url:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: url",
            )

        try:
            snapshot = browser_navigator(raw_url)
        except Exception as exc:  # pragma: no cover - caller-provided navigator guard
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=str(exc),
            )

        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output=snapshot,
            error=None,
        )

    return handler


def register_browser_plugin(
    registry: ToolRegistry,
    *,
    browser_navigator: Callable[[str], dict[str, object]],
) -> ToolRegistry:
    registry.require_plugin_registration_allowed()
    registry.mark_plugin_installed("browser_plugin")
    try:
        registry.get_spec("browser_navigate")
    except KeyError:
        registry.register(
            ToolSpec(
                name="browser_navigate",
                description=(
                    "Drive the configured browser automation backend to a SPECIFIC URL you already have, "
                    "and return page metadata. Use this only with a URL from explicit "
                    "runtime evidence, structured tool output, or a model-produced "
                    "structured recovery plan."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                continuation_tools=(
                    "browser_snapshot",
                    "browser_extract_text",
                    "browser_extract_detail",
                    "browser_extract_items",
                    "browser_run_js",
                ),
            ),
            _build_browser_navigate_handler(browser_navigator),
        )
    return registry



def register_workspace_plugin(
    registry: ToolRegistry,
    *,
    workspace_root: str | Path | None = None,
    allowed_roots: list[Path] | tuple[Path, ...] | None = None,
) -> ToolRegistry:
    if workspace_root is None and allowed_roots is None:
        raise ValueError("workspace_plugin requires workspace_root or allowed_roots")
    registry.require_plugin_registration_allowed()
    registry.mark_plugin_installed("workspace_plugin")
    try:
        registry.get_spec("file_search")
    except KeyError:
        registry.register(
            ToolSpec(
                name="file_search",
                description=(
                    "Search for local files inside the configured workspace or allowed roots. "
                    "Use root to target a specific allowed folder and search_contents for text inside files."
                ),
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                continuation_tools=("file_read",),
            ),
            _build_file_search_handler(
                workspace_root=workspace_root,
                allowed_roots=allowed_roots,
            ),
        )
    try:
        registry.get_spec("file_patch")
    except KeyError:
        registry.register(
            ToolSpec(
                name="file_patch",
                description="Replace text inside a local file within the workspace.",
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.WRITE,
                requires_approval=False,
                timeout_seconds=20,
            ),
            _build_file_patch_handler(
                workspace_root=workspace_root,
                allowed_roots=allowed_roots,
            ),
        )
    try:
        registry.get_spec("workspace_summary")
    except KeyError:
        registry.register(
            ToolSpec(
                name="workspace_summary",
                description="Summarize local workspace contents inside the workspace scope.",
                risk_level=ToolRiskLevel.LOW,
                side_effect_class=ToolSideEffectClass.READ,
                requires_approval=False,
                timeout_seconds=20,
                continuation_tools=("file_search", "file_read"),
            ),
            _build_workspace_summary_handler(
                workspace_root=workspace_root,
                allowed_roots=allowed_roots,
            ),
        )
    return registry



def register_browser_extension(
    registry: ToolRegistry,
    *,
    browser_navigator: Callable[[str], dict[str, object]],
) -> ToolRegistry:
    return register_browser_plugin(
        registry,
        browser_navigator=browser_navigator,
    )


def _build_set_reminder_handler(runtime, *, default_chat_id: str | None) -> ToolHandler:
    """Return a handler that schedules a reminder via the runtime store."""

    def handler(invocation: ToolInvocation) -> ToolResult:
        from nullion.reminders import (
            current_reminder_chat_id,
            due_at_from_relative_seconds,
            normalize_reminder_due_at,
            reminder_due_at_output,
        )

        text = invocation.arguments.get("text")
        due_at_str = invocation.arguments.get("due_at")
        due_in_seconds = invocation.arguments.get("due_in_seconds")
        chat_id = invocation.arguments.get("chat_id") or current_reminder_chat_id() or default_chat_id

        if not isinstance(text, str) or not text.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: text",
            )
        has_due_at = isinstance(due_at_str, str) and bool(due_at_str.strip())
        relative_delay_seconds: float | None = None
        if isinstance(due_in_seconds, (int, float)) and not isinstance(due_in_seconds, bool):
            relative_delay_seconds = float(due_in_seconds)
        elif isinstance(due_in_seconds, str) and due_in_seconds.strip():
            try:
                relative_delay_seconds = float(due_in_seconds)
            except ValueError:
                relative_delay_seconds = None
        has_relative_delay = relative_delay_seconds is not None
        if not has_due_at and not has_relative_delay:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=(
                    "Missing required argument: due_at or due_in_seconds "
                    "(ISO 8601 datetime with timezone, or relative delay in seconds)"
                ),
            )
        if not chat_id:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="No chat_id available for reminder delivery.",
            )

        try:
            if has_relative_delay:
                due_at = due_at_from_relative_seconds(relative_delay_seconds)
            else:
                due_at = datetime.fromisoformat(str(due_at_str).replace("Z", "+00:00"))
                due_at = normalize_reminder_due_at(due_at)
        except (ValueError, TypeError, OverflowError) as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Invalid due_at datetime: {exc}",
            )

        try:
            task = runtime.schedule_reminder(
                chat_id=str(chat_id),
                text=text.strip(),
                due_at=due_at,
            )
            try:
                runtime.checkpoint(force=True)
            except TypeError:
                runtime.checkpoint()
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "task_id": task.task_id,
                    "text": text.strip(),
                    "chat_id": str(chat_id),
                    **reminder_due_at_output(due_at),
                    "message": f"Reminder scheduled: {text.strip()}",
                    "action_receipt": _action_receipt(
                        action="scheduled",
                        object_type="reminder",
                        object_id=str(task.task_id),
                        object_name=text.strip(),
                        summary=f"Reminder scheduled: {text.strip()}",
                        details=[
                            f"Reminder text: {_receipt_value_summary(text.strip())}.",
                            f"Delivery chat: {str(chat_id)}.",
                            f"Due at: {reminder_due_at_output(due_at).get('due_at_display') or reminder_due_at_output(due_at).get('due_at')}.",
                        ],
                    ),
                },
            )
        except Exception as exc:  # noqa: BLE001
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Failed to schedule reminder: {exc}",
            )

    return handler


def _build_list_reminders_handler(runtime) -> ToolHandler:
    """Return a handler that lists all pending (undelivered) reminders."""

    def handler(invocation: ToolInvocation) -> ToolResult:
        from nullion.reminders import reminder_due_at_output

        try:
            all_reminders = runtime.store.list_reminders()
            pending = [
                {
                    "task_id": r.task_id,
                    "text": r.text,
                    "due_at": r.due_at.isoformat(),
                    **reminder_due_at_output(r.due_at),
                    "chat_id": r.chat_id,
                }
                for r in all_reminders
                if r.delivered_at is None
            ]
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={"reminders": pending, "count": len(pending)},
            )
        except Exception as exc:  # noqa: BLE001
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Failed to list reminders: {exc}",
            )

    return handler


def _build_delete_reminder_handler(runtime) -> ToolHandler:
    """Return a handler that deletes a one-off reminder by task id."""

    def handler(invocation: ToolInvocation) -> ToolResult:
        from nullion.reminders import reminder_due_at_output

        task_id = invocation.arguments.get("task_id")
        if not isinstance(task_id, str) or not task_id.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: task_id",
            )
        normalized_task_id = task_id.strip()
        if _invalid_concrete_resource_id_reason(normalized_task_id):
            return _invalid_structured_identifier_result(
                invocation,
                field_name="task_id",
                raw_value=normalized_task_id,
                reason="invalid_reminder_task_id",
                required_source_tool="list_reminders",
            )
        try:
            reminder = runtime.store.get_reminder(normalized_task_id)
            if reminder is None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={},
                    error=f"Reminder not found: {normalized_task_id}",
                )
            removed = runtime.store.remove_reminder(normalized_task_id)
            runtime.store.scheduled_tasks.pop(normalized_task_id, None)
            try:
                runtime.checkpoint(force=True)
            except TypeError:
                runtime.checkpoint()
            if not removed:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={},
                    error=f"Reminder not found: {normalized_task_id}",
                )
            due_at_details = reminder_due_at_output(reminder.due_at)
            due_at_display = due_at_details.get("due_at_display") or due_at_details.get("due_at") or reminder.due_at.isoformat()
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "task_id": normalized_task_id,
                    "text": reminder.text,
                    "chat_id": reminder.chat_id,
                    **due_at_details,
                    "message": f"Reminder deleted: {reminder.text}",
                    "action_receipt": _action_receipt(
                        action="deleted",
                        object_type="reminder",
                        object_id=normalized_task_id,
                        object_name=reminder.text,
                        summary=f"Reminder deleted: {reminder.text}",
                        details=[
                            f"Reminder text: {_receipt_value_summary(reminder.text)}.",
                            f"Due at: {due_at_display}.",
                        ],
                    ),
                },
            )
        except Exception as exc:  # noqa: BLE001
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Failed to delete reminder: {exc}",
            )

    return handler


def _build_update_reminder_handler(runtime) -> ToolHandler:
    """Return a handler that updates a pending one-off reminder."""

    def handler(invocation: ToolInvocation) -> ToolResult:
        from nullion.reminders import (
            due_at_from_relative_seconds,
            normalize_reminder_due_at,
            reminder_due_at_output,
        )

        task_id = invocation.arguments.get("task_id")
        if not isinstance(task_id, str) or not task_id.strip():
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Missing required argument: task_id",
            )
        normalized_task_id = task_id.strip()
        if _invalid_concrete_resource_id_reason(normalized_task_id):
            return _invalid_structured_identifier_result(
                invocation,
                field_name="task_id",
                raw_value=normalized_task_id,
                reason="invalid_reminder_task_id",
                required_source_tool="list_reminders",
            )
        reminder = runtime.store.get_reminder(normalized_task_id)
        if reminder is None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Reminder not found: {normalized_task_id}",
            )

        text_arg = invocation.arguments.get("text")
        due_at_str = invocation.arguments.get("due_at")
        due_in_seconds = invocation.arguments.get("due_in_seconds")
        new_text = reminder.text
        if isinstance(text_arg, str) and text_arg.strip():
            new_text = text_arg.strip()
        elif text_arg is not None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="Invalid argument: text must be a non-empty string when provided",
            )

        has_due_at = isinstance(due_at_str, str) and bool(due_at_str.strip())
        relative_delay_seconds: float | None = None
        if isinstance(due_in_seconds, (int, float)) and not isinstance(due_in_seconds, bool):
            relative_delay_seconds = float(due_in_seconds)
        elif isinstance(due_in_seconds, str) and due_in_seconds.strip():
            try:
                relative_delay_seconds = float(due_in_seconds)
            except ValueError:
                relative_delay_seconds = None
        has_relative_delay = relative_delay_seconds is not None
        if not has_due_at and not has_relative_delay and new_text == reminder.text:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="No reminder updates provided: text, due_at, or due_in_seconds is required",
            )

        try:
            if has_relative_delay:
                new_due_at = due_at_from_relative_seconds(relative_delay_seconds)
            elif has_due_at:
                new_due_at = datetime.fromisoformat(str(due_at_str).replace("Z", "+00:00"))
                new_due_at = normalize_reminder_due_at(new_due_at)
            else:
                new_due_at = normalize_reminder_due_at(reminder.due_at)
        except (ValueError, TypeError, OverflowError) as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Invalid due_at datetime: {exc}",
            )

        updated = replace(reminder, text=new_text, due_at=new_due_at, delivered_at=None)
        runtime.store.add_reminder(updated)
        task = runtime.store.get_scheduled_task(normalized_task_id)
        if task is not None:
            runtime.store.add_scheduled_task(replace(task, enabled=True))
        try:
            runtime.checkpoint(force=True)
        except TypeError:
            runtime.checkpoint()
        due_at_details = reminder_due_at_output(new_due_at)
        due_at_display = due_at_details.get("due_at_display") or due_at_details.get("due_at") or new_due_at.isoformat()
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "task_id": normalized_task_id,
                "text": new_text,
                "chat_id": updated.chat_id,
                **due_at_details,
                "message": f"Reminder updated: {new_text}",
                "action_receipt": _action_receipt(
                    action="updated",
                    object_type="reminder",
                    object_id=normalized_task_id,
                    object_name=new_text,
                    summary=f"Reminder updated: {new_text}",
                    details=[
                        f"Reminder text: {_receipt_value_summary(new_text)}.",
                        f"Due at: {due_at_display}.",
                    ],
                ),
            },
        )

    return handler


def register_reminder_tools(
    registry: ToolRegistry,
    runtime,
    *,
    default_chat_id: str | None = None,
) -> None:
    """Register reminder tools into an existing ToolRegistry.

    These tools are runtime-store-aware: they require a live ``runtime`` reference
    because they write/read directly from ``runtime.store``.  Wire them in after
    the registry is created but before it is attached to the runtime.
    """
    registry.register(
        ToolSpec(
            name="set_reminder",
            description=(
                "Schedule a reminder message to be delivered at a specific time. "
                "Requires: text and either due_in_seconds for relative requests like "
                "'in 2 minutes', "
                "or due_at as an ISO 8601 datetime with timezone offset for absolute requests. "
                "If due_at has no offset, it is interpreted in the user's configured timezone. "
                "Optional: chat_id (defaults to operator chat)."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "reminder"),
            continuation_tools=("delete_reminder", "update_reminder"),
            input_schema={
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Reminder message to deliver."},
                    "due_at": {
                        "type": "string",
                        "description": "Absolute ISO 8601 due time. Include timezone offset when known.",
                    },
                    "due_in_seconds": {
                        "type": "number",
                        "description": (
                            "Relative delay from the current moment, in seconds. "
                            "Prefer this for requests like 'in 2 minutes'."
                        ),
                    },
                    "chat_id": {"type": "string", "description": "Optional delivery chat ID."},
                },
                "required": ["text"],
                "additionalProperties": False,
            },
        ),
        _build_set_reminder_handler(runtime, default_chat_id=default_chat_id),
    )
    registry.register(
        ToolSpec(
            name="list_reminders",
            description=(
                "List all pending one-off reminders. This is not for scheduled cron jobs; "
                "use list_crons/run_cron for cron jobs."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.READ,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "reminder"),
            continuation_tools=("delete_reminder", "update_reminder"),
            input_schema={
                "type": "object",
                "properties": {},
                "additionalProperties": False,
            },
        ),
        _build_list_reminders_handler(runtime),
    )
    registry.register(
        ToolSpec(
            name="delete_reminder",
            description=(
                "Delete or cancel a pending one-off reminder by task_id. "
                "Use list_reminders first when the task_id is not already known."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "reminder"),
            input_schema={
                "type": "object",
                "properties": {
                    "task_id": {"type": "string", "description": "Reminder task id to delete."},
                },
                "required": ["task_id"],
                "additionalProperties": False,
            },
        ),
        _build_delete_reminder_handler(runtime),
    )
    registry.register(
        ToolSpec(
            name="update_reminder",
            description=(
                "Update a pending one-off reminder by task_id. "
                "Use list_reminders first when the task_id is not already known. "
                "Provide text to change the message and either due_in_seconds or due_at to change the time."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "reminder"),
            input_schema={
                "type": "object",
                "properties": {
                    "task_id": {"type": "string", "description": "Reminder task id to update."},
                    "text": {"type": "string", "description": "New reminder message. Omit to keep the current text."},
                    "due_at": {
                        "type": "string",
                        "description": "New absolute ISO 8601 due time. Include timezone offset when known.",
                    },
                    "due_in_seconds": {
                        "type": "number",
                        "description": "New relative delay from the current moment, in seconds.",
                    },
                },
                "required": ["task_id"],
                "additionalProperties": False,
            },
        ),
        _build_update_reminder_handler(runtime),
    )


# ── Cron tools ────────────────────────────────────────────────────────────────

def _scheduler_creation_capsule_key(invocation: ToolInvocation) -> str:
    capsule_id = str(getattr(invocation, "capsule_id", "") or "").strip()
    if capsule_id:
        return capsule_id
    flow_context = invocation.flow_context if isinstance(invocation.flow_context, dict) else {}
    for key in ("turn_id", "conversation_turn_id", "request_id"):
        value = str(flow_context.get(key) or "").strip()
        if value:
            return value
    return ""


def _claim_scheduler_creation_capsule(invocation: ToolInvocation) -> bool:
    capsule_key = _scheduler_creation_capsule_key(invocation)
    if not capsule_key:
        return True
    now = datetime.now(UTC).timestamp()
    with _SCHEDULER_CREATION_CAPSULE_LOCK:
        expired_before = now - _SCHEDULER_CREATION_CAPSULE_TTL_SECONDS
        expired_keys = [
            key for key, claimed_at in _SCHEDULER_CREATION_CAPSULES.items()
            if claimed_at < expired_before
        ]
        for key in expired_keys:
            _SCHEDULER_CREATION_CAPSULES.pop(key, None)
        while len(_SCHEDULER_CREATION_CAPSULES) >= _SCHEDULER_CREATION_CAPSULE_LIMIT:
            oldest_key = min(_SCHEDULER_CREATION_CAPSULES, key=_SCHEDULER_CREATION_CAPSULES.get)
            _SCHEDULER_CREATION_CAPSULES.pop(oldest_key, None)
        if capsule_key in _SCHEDULER_CREATION_CAPSULES:
            return False
        _SCHEDULER_CREATION_CAPSULES[capsule_key] = now
    return True


def _build_create_cron_handler(*, default_delivery_channel: str = "", default_delivery_target: str = ""):
    def _matching_existing_cron(
        *,
        name: str,
        workspace_id: str,
        delivery_channel: str,
        delivery_target: str,
    ):
        from nullion.crons import list_crons

        normalized_name = str(name or "").strip().casefold()
        normalized_workspace = str(workspace_id or "").strip()
        normalized_channel = str(delivery_channel or "").strip()
        normalized_target = str(delivery_target or "").strip()
        for job in list_crons(workspace_id=normalized_workspace):
            if str(getattr(job, "name", "") or "").strip().casefold() != normalized_name:
                continue
            if str(getattr(job, "delivery_channel", "") or "").strip() != normalized_channel:
                continue
            if str(getattr(job, "delivery_target", "") or "").strip() != normalized_target:
                continue
            return job
        return None

    def _delivery_context_from_identifier(identifier: object, fallback_channel: str = "") -> tuple[str, str]:
        from nullion.cron_delivery import normalize_cron_delivery_channel

        value = str(identifier or "").strip()
        if not value:
            return "", ""
        channel, separator, target = value.partition(":")
        normalized_channel = normalize_cron_delivery_channel(channel)
        if separator and normalized_channel and target.strip():
            if normalized_channel == "web":
                return normalized_channel, value
            return normalized_channel, target.strip()
        fallback = normalize_cron_delivery_channel(fallback_channel)
        if fallback and value:
            return fallback, value
        return "", ""

    def _current_delivery_context_defaults() -> tuple[str, str]:
        try:
            from nullion.reminders import current_reminder_chat_id

            return _delivery_context_from_identifier(current_reminder_chat_id(), default_delivery_channel)
        except Exception:
            pass
        return "", ""

    def _invocation_delivery_context_defaults(invocation: ToolInvocation) -> tuple[str, str]:
        from nullion.cron_delivery import normalize_cron_delivery_channel

        configured_default_channel = normalize_cron_delivery_channel(default_delivery_channel)
        channel, target = _delivery_context_from_identifier(invocation.principal_id)
        if channel and target:
            return channel, target
        flow_context = invocation.flow_context if isinstance(invocation.flow_context, dict) else {}
        for key in ("conversation_id", "chat_id", "delivery_target", "delivery_target_id"):
            channel, target = _delivery_context_from_identifier(flow_context.get(key), configured_default_channel)
            if channel and target:
                return channel, target
        return "", ""

    def _workspace_id_from_invocation(invocation: ToolInvocation, args: dict[str, object]) -> str:
        explicit = str(args.get("workspace_id") or "").strip()
        if explicit:
            return explicit
        try:
            from nullion.connections import workspace_id_for_principal

            return workspace_id_for_principal(invocation.principal_id)
        except Exception:
            return "workspace_admin"

    def handle(invocation: ToolInvocation) -> ToolResult:
        from nullion.cron_delivery import normalize_cron_delivery_channel
        from nullion.crons import add_cron, cron_display_fields, update_cron
        args = invocation.arguments or {}
        name     = str(args.get("name", "")).strip()
        schedule = str(args.get("schedule", "")).strip()
        task     = str(args.get("task", "")).strip()
        enabled  = bool(args.get("enabled", True))
        html_image_delivery_mode = str(args.get("html_image_delivery_mode") or "").strip()
        workspace_id = _workspace_id_from_invocation(invocation, args)
        context_channel, context_target = _current_delivery_context_defaults()
        if not context_channel or not context_target:
            context_channel, context_target = _invocation_delivery_context_defaults(invocation)
        configured_default_channel = normalize_cron_delivery_channel(default_delivery_channel)
        default_channel = context_channel or configured_default_channel
        default_target = context_target if context_channel else str(default_delivery_target or "").strip()
        delivery_channel = normalize_cron_delivery_channel(args.get("delivery_channel")) or default_channel
        explicit_target = str(args.get("delivery_target") or "").strip()
        delivery_target = explicit_target
        if not delivery_target and (not delivery_channel or delivery_channel == default_channel):
            delivery_target = default_target
        if not name or not schedule or not task:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="name, schedule and task are all required",
            )
        existing = None
        flow_context = invocation.flow_context if isinstance(invocation.flow_context, dict) else {}
        if (
            flow_context.get("allow_multiple_scheduler_creations") is not True
            and flow_context.get("allow_duplicate_cron_creation") is not True
            and not _claim_scheduler_creation_capsule(invocation)
        ):
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "reason": "multiple_scheduler_creation_tools_in_turn",
                    "requested_scheduler_creation_tool": invocation.tool_name,
                },
                error=(
                    "This turn already created a scheduler object. Do not create another "
                    "cron unless structured flow state explicitly allows multiple scheduler creations."
                ),
            )
        if flow_context.get("allow_duplicate_cron_creation") is not True:
            existing = _matching_existing_cron(
                name=name,
                workspace_id=workspace_id,
                delivery_channel=delivery_channel,
                delivery_target=delivery_target,
            )
        if existing is not None:
            updates = {
                "name": name,
                "schedule": schedule,
                "task": task,
                "enabled": enabled,
                "delivery_channel": delivery_channel,
                "delivery_target": delivery_target,
                "workspace_id": workspace_id,
                "html_image_delivery_mode": html_image_delivery_mode,
            }
            try:
                job = update_cron(existing.id, **updates)
            except Exception as exc:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "id": getattr(existing, "id", ""),
                        "name": getattr(existing, "name", name),
                        "reason": "duplicate_cron_create_update_failed",
                    },
                    error=f"Failed to update existing cron from duplicate create request: {exc}",
                )
            if job is None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "id": getattr(existing, "id", ""),
                        "name": getattr(existing, "name", name),
                        "reason": "duplicate_cron_create_existing_missing",
                    },
                    error=f"Existing cron {getattr(existing, 'id', '')!r} was not found during duplicate create handling.",
                )
            display = cron_display_fields(job)
            schedule_description = display["schedule_description"]
            next_description = display["next_run_description"]
            changes = _cron_update_changes(existing, job, updates)
            changed_labels = [str(change.get("label")) for change in changes if change.get("changed") is not False]
            if changed_labels:
                update_message = (
                    f"Existing cron updated: '{job.name}' (id={job.id}). "
                    f"Changed: {', '.join(changed_labels)}."
                )
            else:
                update_message = (
                    f"Existing cron matched: '{job.name}' (id={job.id}). "
                    "No values changed; the requested fields were already set."
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "id": job.id,
                    "existing_cron_id": job.id,
                    "created_new": False,
                    "name": job.name,
                    "schedule": job.schedule,
                    "task": job.task,
                    "workspace_id": job.workspace_id,
                    "delivery_channel": job.delivery_channel,
                    "delivery_target": job.delivery_target,
                    "html_image_delivery_mode": getattr(job, "html_image_delivery_mode", ""),
                    "enabled": job.enabled,
                    "next_run": job.next_run,
                    "schedule_description": schedule_description,
                    "next_run_description": next_description,
                    "changed_fields": [change["field"] for change in changes if change.get("changed") is not False],
                    "changes": changes,
                    "reason": "duplicate_cron_create_resolved_as_update",
                    "message": update_message,
                    "action_receipt": _action_receipt(
                        action="updated",
                        object_type="cron",
                        object_id=job.id,
                        object_name=job.name,
                        summary=update_message,
                        details=[
                            *_cron_change_details(changes),
                            f"Schedule: {schedule_description}.",
                            *( [f"Next run: {next_description}."] if next_description else [] ),
                            f"Workspace: {job.workspace_id}.",
                            f"Delivery: {job.delivery_channel or '(default)'}{(':' + job.delivery_target) if job.delivery_target else ''}.",
                        ],
                        changes=changes,
                    ),
                },
                error=None,
            )
        try:
            job = add_cron(
                name=name,
                schedule=schedule,
                task=task,
                enabled=enabled,
                delivery_channel=delivery_channel,
                delivery_target=delivery_target,
                html_image_delivery_mode=html_image_delivery_mode,
                workspace_id=workspace_id,
            )
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error=f"Failed to create cron: {exc}",
            )
        display = cron_display_fields(job)
        schedule_description = display["schedule_description"]
        next_description = display["next_run_description"]
        next_info = f" Next run: {next_description}." if next_description else ""
        create_message = (
            f"Cron created: '{job.name}' in workspace {job.workspace_id}. "
            f"Schedule: {schedule_description}.{next_info}"
        )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "id": job.id,
                "name": job.name,
                "schedule": job.schedule,
                "task": job.task,
                "workspace_id": job.workspace_id,
                "delivery_channel": job.delivery_channel,
                "delivery_target": job.delivery_target,
                "html_image_delivery_mode": getattr(job, "html_image_delivery_mode", ""),
                "enabled": job.enabled,
                "next_run": job.next_run,
                "schedule_description": schedule_description,
                "next_run_description": next_description,
                "message": create_message,
                "action_receipt": _action_receipt(
                    action="created",
                    object_type="cron",
                    object_id=job.id,
                    object_name=job.name,
                    summary=create_message,
                    details=[
                        f"Schedule: {schedule_description}.",
                        *( [f"Next run: {next_description}."] if next_description else [] ),
                        f"Workspace: {job.workspace_id}.",
                        f"Delivery: {job.delivery_channel or '(default)'}{(':' + job.delivery_target) if job.delivery_target else ''}.",
                        f"Status: {'enabled' if job.enabled else 'disabled'}.",
                    ],
                ),
            },
            error=None,
        )
    return handle


def _build_list_crons_handler():
    def _cron_display_line(index: int, job: object, display: dict[str, str]) -> str:
        name = str(getattr(job, "name", "") or "Untitled scheduled task").strip()
        schedule = str(getattr(job, "schedule", "") or "").strip()
        enabled = bool(getattr(job, "enabled", False))
        next_run = str(getattr(job, "next_run", "") or "").strip()
        status = "enabled" if enabled else "disabled"
        parts = [f"{index}. {name}", f"   Status: {status}"]
        if schedule:
            parts.append(f"   Schedule: {display['schedule_description']}")
        if next_run:
            next_description = display["next_run_description"] or next_run
            parts.append(f"   Next run: {next_description}")
        return "\n".join(parts)

    def handle(invocation: ToolInvocation) -> ToolResult:
        from nullion.connections import workspace_id_for_principal
        from nullion.crons import cron_display_fields, cron_display_timezone, list_crons

        args = invocation.arguments or {}
        include_all = bool(args.get("include_all_workspaces", False))
        workspace_id = str(args.get("workspace_id") or "").strip() or workspace_id_for_principal(invocation.principal_id)
        jobs = list_crons(workspace_id=None if include_all else workspace_id)
        if not jobs:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={"crons": [], "message": "No crons scheduled."},
                error=None,
            )
        lines = []
        crons = []
        cron_tz = cron_display_timezone()
        for index, j in enumerate(jobs, start=1):
            display = cron_display_fields(j, tz=cron_tz)
            lines.append(_cron_display_line(index, j, display))
            run_by_name = f'run_cron name="{j.name}"'
            crons.append(
                {
                    "selection_index": index,
                    "id": j.id,
                    "name": j.name,
                    "display_name": j.name,
                    "workspace_id": j.workspace_id,
                    "delivery_channel": j.delivery_channel,
                    "delivery_target": j.delivery_target,
                    "html_image_delivery_mode": getattr(j, "html_image_delivery_mode", ""),
                    "enabled": j.enabled,
                    "schedule_description": display["schedule_description"],
                    "next_run_description": display["next_run_description"],
                    "last_run": j.last_run,
                    "run_by_name": run_by_name,
                    "presentation_hint": (
                        "Show schedule_description and next_run_description for timing. "
                        "Do not show cron expressions, raw ids, ISO timestamps, or UTC conversions unless the user asks for technical details. "
                        "When asking the user to choose, show exactly one numbered options list and accept the number."
                    ),
                    "task": j.task,
                    "has_task": bool(str(j.task or "").strip()),
                    "has_last_result": bool(str(j.last_result or "").strip()),
                }
            )
        header = f"Here are your {len(crons)} crons:"
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={"crons": crons, "message": f"{header}\n\n" + "\n\n".join(lines)},
            error=None,
        )
    return handle


def _cron_admin_workspace_allowed(workspace_id: str) -> bool:
    return str(workspace_id or "").strip() == "workspace_admin"


def _cron_workspace_denial(
    *,
    invocation: ToolInvocation,
    cron_id: str,
    owner_workspace_id: str,
) -> ToolResult:
    return ToolResult(
        invocation_id=invocation.invocation_id,
        tool_name=invocation.tool_name,
        status="failed",
        output={"id": cron_id, "workspace_id": owner_workspace_id},
        error=f"Cron {cron_id!r} belongs to workspace {owner_workspace_id}.",
    )


_CRON_RECEIPT_FIELD_LABELS = {
    "name": "Name",
    "schedule": "Schedule",
    "task": "Task instructions",
    "enabled": "Status",
    "delivery_channel": "Delivery channel",
    "delivery_target": "Delivery target",
    "workspace_id": "Workspace",
    "html_image_delivery_mode": "HTML image delivery mode",
}


def _cron_update_changes(existing: object | None, updated: object, updates: dict[str, object]) -> list[dict[str, object]]:
    changes: list[dict[str, object]] = []
    for field in updates:
        label = _CRON_RECEIPT_FIELD_LABELS.get(field, field.replace("_", " ").title())
        before = getattr(existing, field, None) if existing is not None else None
        after = getattr(updated, field, None)
        before_text = _receipt_value_summary(before)
        after_text = _receipt_value_summary(after)
        changes.append(
            {
                "field": field,
                "label": label,
                "before": before_text or "(blank)",
                "after": after_text or "(blank)",
                "changed": before_text != after_text,
            }
        )
    return changes


def _cron_change_details(changes: list[dict[str, object]]) -> list[str]:
    details: list[str] = []
    for change in changes:
        label = str(change.get("label") or "Field")
        before = str(change.get("before") or "(blank)")
        after = str(change.get("after") or "(blank)")
        if change.get("changed") is False:
            details.append(f"{label} was already set to {after}.")
        elif before and before != "(blank)":
            details.append(f"{label}: {before} -> {after}.")
        else:
            details.append(f"{label}: {after}.")
    return details


def _build_delete_cron_handler():
    def handle(invocation: ToolInvocation) -> ToolResult:
        from nullion.connections import workspace_id_for_principal
        from nullion.crons import get_cron, remove_cron
        args = invocation.arguments or {}
        cron_id = str(args.get("id", "")).strip()
        if not cron_id:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="id is required",
            )
        if _invalid_concrete_resource_id_reason(cron_id):
            return _invalid_structured_identifier_result(
                invocation,
                field_name="id",
                raw_value=cron_id,
                reason="invalid_cron_id",
                required_source_tool="list_crons",
            )
        job = get_cron(cron_id)
        workspace_id = workspace_id_for_principal(invocation.principal_id)
        admin_cross_workspace = job is not None and job.workspace_id != workspace_id and _cron_admin_workspace_allowed(workspace_id)
        if job is not None and job.workspace_id != workspace_id and not admin_cross_workspace:
            return _cron_workspace_denial(invocation=invocation, cron_id=cron_id, owner_workspace_id=job.workspace_id)
        removed = remove_cron(cron_id)
        if not removed:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"id": cron_id},
                error=f"No cron found with id={cron_id!r}",
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "id": cron_id,
                "name": getattr(job, "name", ""),
                "workspace_id": getattr(job, "workspace_id", ""),
                "admin_cross_workspace": admin_cross_workspace,
                "message": f"Cron deleted: '{getattr(job, 'name', cron_id)}' (id={cron_id}).",
                "action_receipt": _action_receipt(
                    action="deleted",
                    object_type="cron",
                    object_id=cron_id,
                    object_name=str(getattr(job, "name", "") or ""),
                    summary=f"Cron deleted: '{getattr(job, 'name', cron_id)}' (id={cron_id}).",
                    details=[
                        f"Workspace: {getattr(job, 'workspace_id', '') or '(unknown)'}.",
                    ],
                ),
            },
            error=None,
        )
    return handle


def _build_update_cron_handler():
    def handle(invocation: ToolInvocation) -> ToolResult:
        from nullion.connections import workspace_id_for_principal
        from nullion.crons import get_cron, update_cron

        args = invocation.arguments or {}
        cron_id = str(args.get("id", "")).strip()
        if not cron_id:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="id is required",
            )
        if _invalid_concrete_resource_id_reason(cron_id):
            return _invalid_structured_identifier_result(
                invocation,
                field_name="id",
                raw_value=cron_id,
                reason="invalid_cron_id",
                required_source_tool="list_crons",
            )
        existing = get_cron(cron_id)
        workspace_id = workspace_id_for_principal(invocation.principal_id)
        admin_cross_workspace = (
            existing is not None
            and existing.workspace_id != workspace_id
            and _cron_admin_workspace_allowed(workspace_id)
        )
        if existing is not None and existing.workspace_id != workspace_id and not admin_cross_workspace:
            return _cron_workspace_denial(invocation=invocation, cron_id=cron_id, owner_workspace_id=existing.workspace_id)
        mutable_fields = (
            "name",
            "schedule",
            "task",
            "enabled",
            "delivery_channel",
            "delivery_target",
            "workspace_id",
            "html_image_delivery_mode",
        )
        updates: dict[str, object] = {}
        for field in mutable_fields:
            if field not in args:
                continue
            value = args[field]
            if field in {"name", "schedule", "task", "delivery_channel", "delivery_target", "workspace_id", "html_image_delivery_mode"}:
                value = str(value or "").strip()
            updates[field] = value
        if not updates:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"id": cron_id},
                error="at least one cron field is required",
            )
        try:
            job = update_cron(cron_id, **updates)
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"id": cron_id},
                error=f"Failed to update cron: {exc}",
            )
        if job is None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"id": cron_id},
                error=f"No cron found with id={cron_id!r}",
            )
        changes = _cron_update_changes(existing, job, updates)
        changed_labels = [str(change.get("label")) for change in changes if change.get("changed") is not False]
        if changed_labels:
            update_message = (
                f"Cron updated: '{job.name}' (id={job.id}). "
                f"Changed: {', '.join(changed_labels)}."
            )
        else:
            update_message = (
                f"Cron update checked: '{job.name}' (id={job.id}). "
                "No values changed; the requested fields were already set."
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "id": job.id,
                "name": job.name,
                "schedule": job.schedule,
                "workspace_id": job.workspace_id,
                "delivery_channel": job.delivery_channel,
                "delivery_target": job.delivery_target,
                "html_image_delivery_mode": getattr(job, "html_image_delivery_mode", ""),
                "enabled": job.enabled,
                "admin_cross_workspace": admin_cross_workspace,
                "next_run": job.next_run,
                "has_task": bool(str(job.task or "").strip()),
                "has_last_result": bool(str(job.last_result or "").strip()),
                "changed_fields": [change["field"] for change in changes if change.get("changed") is not False],
                "changes": changes,
                "message": update_message,
                "action_receipt": _action_receipt(
                    action="updated",
                    object_type="cron",
                    object_id=job.id,
                    object_name=job.name,
                    summary=update_message,
                    details=_cron_change_details(changes),
                    changes=changes,
                ),
            },
            error=None,
        )
    return handle


def _build_toggle_cron_handler():
    def handle(invocation: ToolInvocation) -> ToolResult:
        from nullion.connections import workspace_id_for_principal
        from nullion.crons import get_cron, list_crons, toggle_cron
        args = invocation.arguments or {}
        cron_id = str(args.get("id", "")).strip()
        raw_ids = args.get("ids")
        ids = tuple(
            dict.fromkeys(
                str(item or "").strip()
                for item in (raw_ids if isinstance(raw_ids, list) else ())
                if str(item or "").strip()
            )
        )
        all_current_workspace = bool(args.get("all_current_workspace") is True)
        enabled = bool(args.get("enabled", True))
        selector_count = int(bool(cron_id)) + int(bool(ids)) + int(all_current_workspace)
        if selector_count > 1:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "id": cron_id or None,
                    "ids": list(ids),
                    "all_current_workspace": all_current_workspace,
                },
                error="Use exactly one of id, ids, or all_current_workspace.",
            )
        if not cron_id and not ids and not all_current_workspace:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="id, ids, or all_current_workspace is required",
            )
        workspace_id = workspace_id_for_principal(invocation.principal_id)
        if all_current_workspace:
            target_ids = tuple(
                str(job.id)
                for job in list_crons(workspace_id=workspace_id)
                if str(getattr(job, "id", "") or "").strip()
                and bool(getattr(job, "enabled", False)) != enabled
            )
        else:
            target_ids = (cron_id,) if cron_id else ids
        invalid_ids = [
            target_id
            for target_id in target_ids
            if _invalid_concrete_resource_id_reason(target_id)
        ]
        if invalid_ids:
            if len(invalid_ids) == 1:
                return _invalid_structured_identifier_result(
                    invocation,
                    field_name="id",
                    raw_value=invalid_ids[0],
                    reason="invalid_cron_id",
                    required_source_tool="list_crons",
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"ids": list(target_ids), "invalid_ids": invalid_ids},
                error="One or more cron ids are invalid. Use exact ids returned by list_crons.",
            )
        if not all_current_workspace and len(target_ids) == 1:
            cron_id = target_ids[0]
            existing = get_cron(cron_id)
            admin_cross_workspace = (
                existing is not None
                and existing.workspace_id != workspace_id
                and _cron_admin_workspace_allowed(workspace_id)
            )
            if existing is not None and existing.workspace_id != workspace_id and not admin_cross_workspace:
                return _cron_workspace_denial(
                    invocation=invocation,
                    cron_id=cron_id,
                    owner_workspace_id=existing.workspace_id,
                )
            job = toggle_cron(cron_id, enabled)
            if job is None:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"id": cron_id},
                    error=f"No cron found with id={cron_id!r}",
                )
            state = "enabled" if enabled else "disabled"
            toggle_message = f"Cron '{job.name}' ({cron_id}) is now {state}."
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="completed",
                output={
                    "id": job.id,
                    "name": job.name,
                    "workspace_id": job.workspace_id,
                    "enabled": job.enabled,
                    "admin_cross_workspace": admin_cross_workspace,
                    "message": toggle_message,
                    "action_receipt": _action_receipt(
                        action="updated",
                        object_type="cron",
                        object_id=job.id,
                        object_name=job.name,
                        summary=toggle_message,
                        details=[f"Status: {state}.", f"Workspace: {job.workspace_id}."],
                        changes=[
                            {
                                "field": "enabled",
                                "label": "Status",
                                "before": "enabled" if bool(getattr(existing, "enabled", False)) else "disabled",
                                "after": state,
                                "changed": bool(getattr(existing, "enabled", None)) != bool(job.enabled),
                            }
                        ],
                    ),
                },
                error=None,
            )

        results: list[dict[str, object]] = []
        details: list[str] = []
        missing_ids: list[str] = []
        denied: list[ToolResult] = []
        for target_id in target_ids:
            existing = get_cron(target_id)
            admin_cross_workspace = (
                existing is not None
                and existing.workspace_id != workspace_id
                and _cron_admin_workspace_allowed(workspace_id)
            )
            if existing is not None and existing.workspace_id != workspace_id and not admin_cross_workspace:
                denied.append(
                    _cron_workspace_denial(
                        invocation=invocation,
                        cron_id=target_id,
                        owner_workspace_id=existing.workspace_id,
                    )
                )
                continue
            job = toggle_cron(target_id, enabled)
            if job is None:
                missing_ids.append(target_id)
                continue
            changed = bool(getattr(existing, "enabled", None)) != bool(job.enabled)
            item = {
                "id": job.id,
                "name": job.name,
                "workspace_id": job.workspace_id,
                "enabled": job.enabled,
                "changed": changed,
                "admin_cross_workspace": admin_cross_workspace,
            }
            results.append(item)
            state = "enabled" if enabled else "disabled"
            details.append(f"{state.title()}: {job.name} ({job.id}).")
        if denied:
            first = denied[0]
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "ids": list(target_ids),
                    "results": results,
                    "missing_ids": missing_ids,
                    "denied_id": first.output.get("id") if isinstance(first.output, dict) else None,
                },
                error=first.error,
            )
        if missing_ids:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"ids": list(target_ids), "results": results, "missing_ids": missing_ids},
                error="One or more cron jobs were not found.",
            )
        state = "enabled" if enabled else "disabled"
        verb = "Enabled" if enabled else "Disabled"
        if not results:
            summary = f"No cron jobs needed to be {state}."
        else:
            summary = f"{verb} {len(results)} cron job{'s' if len(results) != 1 else ''}."
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output={
                "ids": [str(item.get("id")) for item in results],
                "results": results,
                "enabled": enabled,
                "updated_count": len(results),
                "all_current_workspace": all_current_workspace,
                "message": summary,
                "action_receipt": _action_receipt(
                    action="updated",
                    object_type="cron",
                    summary=summary,
                    details=details,
                ),
            },
            error=None,
        )
    return handle


def _build_run_cron_handler(cron_runner: Callable[..., str | dict[str, object] | None] | None):
    def _foreground_cron_no_output_text() -> str:
        from nullion.cron_delivery import DEFAULT_CRON_NO_OUTPUT_MESSAGE

        return DEFAULT_CRON_NO_OUTPUT_MESSAGE

    def _cron_lookup_parts(value: object) -> tuple[str, tuple[str, ...]]:
        from nullion.text_match import ascii_match_text

        text = ascii_match_text(value).casefold()
        tokens: list[str] = []
        current: list[str] = []
        for char in text:
            if char.isalnum():
                current.append(char)
            elif current:
                tokens.append("".join(current))
                current = []
        if current:
            tokens.append("".join(current))
        expanded_tokens: list[str] = []
        for token in tokens:
            expanded_tokens.append(token)
            for suffix in ("ing", "ers", "er", "ed", "s"):
                if token.endswith(suffix) and len(token) - len(suffix) >= 4:
                    expanded_tokens.append(token[: -len(suffix)])
        return "".join(tokens), tuple(dict.fromkeys(expanded_tokens))

    def _cron_name_match_rank(query: str, candidate_name: str) -> int | None:
        query_text = str(query or "").strip()
        candidate_text = str(candidate_name or "").strip()
        if not query_text or not candidate_text:
            return None
        if query_text.casefold() == candidate_text.casefold():
            return 0

        query_compact, query_tokens = _cron_lookup_parts(query_text)
        candidate_compact, candidate_tokens = _cron_lookup_parts(candidate_text)
        if not query_compact or not candidate_compact:
            return None
        if query_compact == candidate_compact:
            return 1
        if query_tokens and query_tokens == candidate_tokens:
            return 1
        if query_tokens and all(token in candidate_tokens for token in query_tokens):
            return 2
        if len(query_compact) >= 6 and query_compact in candidate_compact:
            return 3
        return None

    def _cron_reference_match_rank(query: str, candidate_text: str) -> int | None:
        rank = _cron_name_match_rank(query, candidate_text)
        if rank is not None:
            return rank
        _query_compact, query_tokens = _cron_lookup_parts(query)
        _candidate_compact, candidate_tokens = _cron_lookup_parts(candidate_text)
        overlap = set(query_tokens).intersection(candidate_tokens)
        if len(overlap) >= 2:
            return 5
        fuzzy_overlap = _fuzzy_cron_token_overlap(query_tokens, candidate_tokens)
        if len(fuzzy_overlap) >= 1 and len(query_tokens) <= 4:
            return 6
        return None

    def _fuzzy_cron_token_overlap(query_tokens: tuple[str, ...], candidate_tokens: tuple[str, ...]) -> set[str]:
        from difflib import SequenceMatcher

        matches: set[str] = set()
        candidate_set = {token for token in candidate_tokens if len(token) >= 6}
        for query_token in query_tokens:
            if len(query_token) < 6:
                continue
            for candidate_token in candidate_set:
                length_delta = abs(len(query_token) - len(candidate_token))
                if length_delta > max(2, len(candidate_token) // 4):
                    continue
                if SequenceMatcher(None, query_token, candidate_token).ratio() >= 0.86:
                    matches.add(query_token)
                    break
        return matches

    def _significant_cron_tokens(value: object) -> set[str]:
        _compact, tokens = _cron_lookup_parts(value)
        return {token for token in tokens if len(token) >= 5}

    def _cron_search_text(candidate: object) -> str:
        return " ".join(
            (
                str(getattr(candidate, "name", "") or ""),
                str(getattr(candidate, "task", "") or ""),
            )
        )

    def _descriptive_cron_matches(query: str, jobs: list[object]) -> list[object]:
        query_tokens = _significant_cron_tokens(query)
        if not query_tokens:
            return []
        scored: list[tuple[int, object]] = []
        for candidate in jobs:
            if not bool(getattr(candidate, "enabled", True)):
                continue
            candidate_tokens = _significant_cron_tokens(_cron_search_text(candidate))
            overlap = query_tokens.intersection(candidate_tokens)
            fuzzy_overlap = _fuzzy_cron_token_overlap(tuple(query_tokens), tuple(candidate_tokens))
            score = len(overlap) * 2 + len(fuzzy_overlap)
            if score:
                scored.append((score, candidate))
        if not scored:
            return []
        best_score = max(score for score, _candidate in scored)
        return [candidate for score, candidate in scored if score == best_score]

    def _unique_cron_name_match(query: str, jobs: list[object]) -> tuple[object | None, list[object]]:
        ranked: list[tuple[int, object]] = []
        for candidate in jobs:
            rank = _cron_name_match_rank(query, str(getattr(candidate, "name", "") or ""))
            if rank is None and bool(getattr(candidate, "enabled", True)):
                search_rank = _cron_reference_match_rank(query, _cron_search_text(candidate))
                if search_rank is not None:
                    rank = search_rank + 2
            if rank is not None:
                ranked.append((rank, candidate))
        if not ranked:
            descriptive_matches = _descriptive_cron_matches(query, jobs)
            if len(descriptive_matches) == 1:
                return descriptive_matches[0], descriptive_matches
            if descriptive_matches:
                return None, descriptive_matches
            return None, []
        best_rank = min(rank for rank, _candidate in ranked)
        matches = [candidate for rank, candidate in ranked if rank == best_rank]
        if len(matches) == 1:
            return matches[0], matches
        descriptive_matches = _descriptive_cron_matches(query, matches)
        if len(descriptive_matches) == 1:
            return descriptive_matches[0], descriptive_matches
        if descriptive_matches:
            return None, descriptive_matches
        return None, matches

    def _numbered_cron_matches(matches: list[object]) -> list[dict[str, object]]:
        return [
            {
                "selection_index": index,
                "id": item.id,
                "name": item.name,
                "schedule": item.schedule,
                "workspace_id": item.workspace_id,
                "enabled": item.enabled,
                "next_run": item.next_run,
                "reply_with": str(index),
            }
            for index, item in enumerate(matches, start=1)
        ]

    def _foreground_cron_result_view(text: str) -> tuple[str, int]:
        from nullion.artifacts import parse_media_directive_line

        removed_media_count = 0
        kept_lines: list[str] = []
        for raw_line in str(text or "").splitlines():
            directive = parse_media_directive_line(raw_line)
            if directive is not None:
                removed_media_count += 1
                continue
            kept_lines.append(raw_line)
        if removed_media_count:
            return "", removed_media_count
        return "\n".join(kept_lines).strip(), 0

    def _foreground_cron_result_text(text: str) -> str:
        return _foreground_cron_result_view(text)[0]

    def _foreground_cron_status_text(job: object, status: object) -> str:
        normalized = str(status or "").strip()
        name = str(getattr(job, "name", "") or "cron").strip()
        if normalized in {"sent", "saved"}:
            return f"Manual scheduled task run started: {name}. The result was delivered to the configured destination."
        if normalized == "silent":
            return "Cron ran successfully; no output was produced."
        if normalized == "deferred":
            return (
                f"Manual scheduled task run started: {name}. "
                "The result will be delivered to this chat when ready."
            )
        return ""

    def _manual_cron_progress_status(job: object, invocation_id: object) -> tuple[str, str]:
        from nullion.cron_delivery import manual_cron_running_status_text, manual_cron_status_group_id

        return (
            manual_cron_status_group_id(job, run_id=invocation_id),
            manual_cron_running_status_text(job),
        )

    def _call_cron_runner(runner: Callable[..., str | dict[str, object] | None], job: object, invocation: ToolInvocation):
        try:
            parameters = inspect.signature(runner).parameters
            accepts_invocation = (
                any(parameter.kind is inspect.Parameter.VAR_POSITIONAL for parameter in parameters.values())
                or "invocation" in parameters
                or len(
                    [
                        parameter
                        for parameter in parameters.values()
                        if parameter.kind
                        in {
                            inspect.Parameter.POSITIONAL_ONLY,
                            inspect.Parameter.POSITIONAL_OR_KEYWORD,
                        }
                    ]
                )
                >= 2
            )
        except (TypeError, ValueError):
            accepts_invocation = False
        if accepts_invocation:
            return runner(job, invocation)
        return runner(job)

    def _foreground_cron_failure_text(reason: str) -> str:
        if reason == "cron_run_raw_tool_payload":
            return (
                "Cron run was blocked because it produced raw structured tool output "
                "instead of a deliverable report."
            )
        if reason == "cron_run_internal_tool_output_leaked":
            return "Cron run was blocked because it tried to deliver internal tool reference content."
        if reason == "cron_run_reached_iteration_limit":
            return "Cron run stopped before producing a deliverable result."
        if reason == "cron_run_waiting_for_approval":
            return "Cron run is waiting for approval."
        if reason:
            return "Cron run did not deliver its result to the configured platform."
        return ""

    def _foreground_cron_runner_output(runner_output: str | dict[str, object] | None) -> str | dict[str, object] | None:
        if not isinstance(runner_output, dict):
            return runner_output
        delivery_status = str(runner_output.get("cron_delivery_status") or "").strip()
        if delivery_status or runner_output.get("cron_delivery_failed") or runner_output.get("cron_run_failed"):
            allowed_keys = {
                "cron_delivery_status",
                "cron_delivery_failed",
                "cron_run_failed",
                "reached_iteration_limit",
                "raw_tool_payload_blocked",
                "suspended_for_approval",
                "approval_id",
                "reason",
            }
            if delivery_status == "deferred":
                allowed_keys.update({
                    "text",
                    "final_text",
                    "message",
                    "result_text",
                    "planner_status_text",
                    "status_delivered",
                    "task_group_id",
                    "progress_status_text",
                    "planner_status_text",
                })
            return {
                key: value
                for key, value in runner_output.items()
                if key in allowed_keys
            }
        sanitized = dict(runner_output)
        sanitized.pop("artifact_paths", None)
        sanitized.pop("artifacts", None)
        for key in ("text", "result_text", "message", "final_text"):
            value = sanitized.get(key)
            if isinstance(value, str):
                sanitized_value = _foreground_cron_result_text(value)
                if sanitized_value:
                    sanitized[key] = sanitized_value
                else:
                    sanitized.pop(key, None)
        return sanitized

    def handle(invocation: ToolInvocation) -> ToolResult:
        from datetime import timezone

        from nullion.connections import workspace_id_for_principal
        from nullion.crons import get_cron, list_crons, load_crons, save_crons

        args = invocation.arguments or {}
        raw_id = str(args.get("id", "")).strip()
        raw_name = str(args.get("name", "")).strip()
        raw_ids = args.get("ids")
        raw_all_enabled = args.get("all_enabled")
        workspace_id = workspace_id_for_principal(invocation.principal_id)
        if raw_id and _invalid_concrete_resource_id_reason(raw_id):
            return _invalid_structured_identifier_result(
                invocation,
                field_name="id",
                raw_value=raw_id,
                reason="invalid_cron_id",
                required_source_tool="list_crons",
            )
        if raw_name and _invalid_concrete_resource_name_reason(raw_name):
            return _invalid_structured_identifier_result(
                invocation,
                field_name="name",
                raw_value=raw_name,
                reason="invalid_cron_name",
                required_source_tool="list_crons",
            )
        if raw_all_enabled is True:
            if raw_id or raw_name or raw_ids:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"all_enabled": True},
                    error="all_enabled cannot be combined with id, ids, or name",
                )
            enabled_jobs = [job for job in list_crons(workspace_id=workspace_id) if bool(getattr(job, "enabled", True))]
            raw_ids = [str(job.id) for job in enabled_jobs if str(job.id or "").strip()]
            if not raw_ids:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"all_enabled": True, "started_count": 0},
                    error="No enabled cron jobs found in this workspace.",
                )
        if isinstance(raw_ids, (list, tuple)):
            ids = tuple(
                dict.fromkeys(
                    str(item or "").strip()
                    for item in raw_ids
                    if str(item or "").strip()
                )
            )
            if not ids:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={"ids": []},
                    error="ids must include at least one cron id",
                )
            invalid_ids = [cron_id for cron_id in ids if _invalid_concrete_resource_id_reason(cron_id)]
            if invalid_ids:
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "reason": "invalid_cron_id",
                        "ids": list(ids),
                        "invalid_ids": invalid_ids,
                        "required_source_tool": "list_crons",
                    },
                    error="Invalid id in ids: use exact cron ids returned by list_crons.",
                )
            results: list[dict[str, object]] = []
            completed_count = 0
            failed_count = 0
            for index, cron_id in enumerate(ids, start=1):
                child = handle(
                    ToolInvocation(
                        invocation_id=f"{invocation.invocation_id}:{index}",
                        tool_name=invocation.tool_name,
                        principal_id=invocation.principal_id,
                        arguments={"id": cron_id},
                        capsule_id=invocation.capsule_id,
                        trusted_filesystem_selectors=invocation.trusted_filesystem_selectors,
                        flow_context=invocation.flow_context,
                    )
                )
                item: dict[str, object] = {
                    "id": cron_id,
                    "status": child.status,
                }
                if child.output.get("name"):
                    item["name"] = child.output["name"]
                if child.output.get("delivery_status"):
                    item["delivery_status"] = child.output["delivery_status"]
                if child.output.get("cron_delivery_status"):
                    item["cron_delivery_status"] = child.output["cron_delivery_status"]
                if child.output.get("task_group_id"):
                    item["task_group_id"] = child.output["task_group_id"]
                if child.output.get("status_delivered") is True:
                    item["status_delivered"] = True
                if child.error:
                    item["error"] = child.error
                results.append(item)
                if child.status == "completed":
                    completed_count += 1
                else:
                    failed_count += 1
            status = "completed" if completed_count else "failed"
            message = (
                f"Started {completed_count} cron run"
                f"{'' if completed_count == 1 else 's'}."
            )
            if failed_count:
                message = f"{message} {failed_count} failed to start."
            background_count = sum(
                1
                for item in results
                if str(item.get("delivery_status") or item.get("cron_delivery_status") or "").strip() == "deferred"
            )
            if background_count:
                message = (
                    f"{message} {background_count} background run"
                    f"{'' if background_count == 1 else 's'} will deliver to the configured destination"
                    f"{'' if background_count == 1 else 's'} when ready."
                )
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status=status,
                output={
                    "ids": list(ids),
                    "results": results,
                    "started_count": completed_count,
                    "failed_count": failed_count,
                    "message": message,
                    "action_receipt": _action_receipt(
                        action="started",
                        object_type="cron",
                        summary=message,
                        details=[
                            *[
                                f"Started: {item.get('name') or item.get('id')} ({item.get('id')})."
                                for item in results
                                if item.get("status") == "completed"
                            ],
                            *[
                                f"Failed: {item.get('name') or item.get('id')} ({item.get('id')}): {item.get('error') or 'could not start'}."
                                for item in results
                                if item.get("status") != "completed"
                            ],
                        ],
                    ),
                    **({"all_enabled": True} if raw_all_enabled is True else {}),
                },
                error=None if completed_count else "No cron runs started.",
            )
        if not raw_id and not raw_name:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={},
                error="id or name is required",
            )
        job = get_cron(raw_id) if raw_id else None
        if job is not None and job.workspace_id != workspace_id:
            job = None
        if job is None and raw_name:
            job, matches = _unique_cron_name_match(raw_name, list_crons(workspace_id=workspace_id))
            if job is None and matches:
                numbered_matches = _numbered_cron_matches(matches)
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="failed",
                    output={
                        "name": raw_name,
                        "matches": numbered_matches,
                        "message": "\n".join(
                            f"{item['selection_index']}. {item['name']}" for item in numbered_matches
                        ),
                        "presentation_hint": "Ask the user to choose by number.",
                    },
                    error="Multiple crons matched; ask the user to choose by number.",
                )
        if job is None:
            lookup = raw_id or raw_name
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"id": raw_id, "name": raw_name},
                error=f"No cron found for {lookup!r}",
            )
        if cron_runner is None:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={
                    "id": job.id,
                    "name": job.name,
                    "has_task": bool(str(job.task or "").strip()),
                    "reason": "cron_runner_not_configured",
                },
                error="This runtime can list crons but cannot run them on demand.",
            )
        try:
            runner_output = _call_cron_runner(cron_runner, job, invocation)
        except Exception as exc:
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output={"id": job.id, "name": job.name},
                error=f"Failed to run cron: {exc}",
            )
        runner_failed = False
        runner_failure_reason = ""
        if isinstance(runner_output, dict):
            if runner_output.get("reached_iteration_limit"):
                runner_failed = True
                runner_failure_reason = "cron_run_reached_iteration_limit"
            elif runner_output.get("cron_delivery_failed"):
                runner_failed = True
                runner_failure_reason = str(runner_output.get("reason") or "cron_delivery_failed")
            elif runner_output.get("suspended_for_approval"):
                runner_failed = True
                runner_failure_reason = "cron_run_waiting_for_approval"
        jobs = load_crons()
        now = datetime.now(timezone.utc).isoformat(timespec="seconds")
        from nullion.response_fulfillment_contract import guaranteed_user_visible_text

        delivery_status = runner_output.get("cron_delivery_status") if isinstance(runner_output, dict) else ""
        delivery_status_text = str(delivery_status or "").strip()
        deferred_status_owned_by_background = (
            delivery_status_text == "deferred"
            and isinstance(runner_output, dict)
            and (
                runner_output.get("status_delivered") is True
                or runner_output.get("planner_status_owned_by_background") is True
                or runner_output.get("foreground_reply_suppressed") is True
            )
        )
        foreground_reply_suppressed = delivery_status_text in {
            "saved",
            "sent",
            "partial_success",
        } or deferred_status_owned_by_background
        result_text = (
            ""
            if foreground_reply_suppressed
            else _foreground_cron_failure_text(runner_failure_reason)
            or _foreground_cron_status_text(job, delivery_status)
        )
        removed_media_count = 0
        if not result_text and not foreground_reply_suppressed:
            raw_result_text = guaranteed_user_visible_text(subject=job, output=runner_output, kind="cron")
            result_text, removed_media_count = _foreground_cron_result_view(raw_result_text)
            if not result_text:
                result_text = _foreground_cron_no_output_text()
        updated_stored_job = False
        for stored in jobs:
            if stored.id == job.id:
                stored.last_run = now
                stored.last_result = (
                    f"manual run failed: {runner_failure_reason}"
                    if runner_failed
                    else result_text
                )
                updated_stored_job = True
                break
        if updated_stored_job:
            save_crons(jobs)
        run_message = f"Ran cron '{job.name}' ({job.id}) now."
        if delivery_status_text == "deferred":
            run_message = _foreground_cron_status_text(job, delivery_status) or run_message
        progress_group_id = ""
        progress_status_text = ""
        if delivery_status_text == "deferred":
            if isinstance(runner_output, dict):
                progress_group_id = str(runner_output.get("task_group_id") or "").strip()
                progress_status_text = str(
                    runner_output.get("progress_status_text")
                    or runner_output.get("planner_status_text")
                    or ""
                ).strip()
            if not progress_group_id or not progress_status_text:
                progress_group_id, progress_status_text = _manual_cron_progress_status(job, invocation.invocation_id)
        output: dict[str, object] = {
            "id": job.id,
            "name": job.name,
            "has_task": bool(str(job.task or "").strip()),
            "workspace_id": job.workspace_id,
            "last_run": now,
            "message": run_message,
            "action_receipt": _action_receipt(
                action="started",
                object_type="cron",
                object_id=job.id,
                object_name=job.name,
                summary=run_message,
                details=[
                    f"Cron: {job.name} ({job.id}).",
                    f"Workspace: {job.workspace_id}.",
                    *( [f"Delivery status: {delivery_status}."] if delivery_status else [] ),
                    *( [f"Result: {_receipt_value_summary(result_text)}."] if result_text else [] ),
                ],
            ),
            "foreground_auto_attach_created_artifacts": False,
        }
        if result_text:
            output["result_text"] = result_text
        if foreground_reply_suppressed:
            output["foreground_reply_suppressed"] = True
        if delivery_status:
            output["delivery_status"] = str(delivery_status)
            output["cron_delivery_status"] = str(delivery_status)
            if delivery_status_text == "deferred":
                output["status_delivered"] = False
        if (
            isinstance(runner_output, dict)
            and delivery_status_text == "deferred"
            and runner_output.get("status_delivered") is True
        ):
            output["status_delivered"] = True
        if (
            isinstance(runner_output, dict)
            and delivery_status_text == "deferred"
            and runner_output.get("planner_status_owned_by_background") is True
        ):
            output["planner_status_owned_by_background"] = True
        if delivery_status_text == "deferred" and progress_group_id and progress_status_text:
            output["task_group_id"] = progress_group_id
            output["progress_status_text"] = progress_status_text
            # Keep the legacy key populated so existing platform status helpers
            # can render the same card without parsing the user-visible receipt.
            output["planner_status_text"] = progress_status_text
        if removed_media_count:
            output["foreground_media_directive_count"] = removed_media_count
        if isinstance(runner_output, dict):
            output["result"] = _foreground_cron_runner_output(runner_output)
        if runner_failed:
            approval_id = runner_output.get("approval_id") if isinstance(runner_output, dict) else None
            if runner_failure_reason == "cron_run_waiting_for_approval":
                output["reason"] = "approval_required"
                output["requires_approval"] = True
                if isinstance(approval_id, str) and approval_id:
                    output["approval_id"] = approval_id
                return ToolResult(
                    invocation_id=invocation.invocation_id,
                    tool_name=invocation.tool_name,
                    status="denied",
                    output=output,
                    error="Approval required before the cron can continue.",
                )
            output["reason"] = runner_failure_reason
            return ToolResult(
                invocation_id=invocation.invocation_id,
                tool_name=invocation.tool_name,
                status="failed",
                output=output,
                error=(
                    "Cron run did not finish cleanly; check the Doctor action or retry after cleanup."
                    if runner_failure_reason == "cron_run_reached_iteration_limit"
                    else "Cron run produced raw structured tool output instead of a deliverable report."
                    if runner_failure_reason == "cron_run_raw_tool_payload"
                    else "Cron run tried to deliver internal tool reference content."
                    if runner_failure_reason == "cron_run_internal_tool_output_leaked"
                    else "Cron run did not deliver its result to the configured platform."
                    if runner_failure_reason.startswith("cron_delivery") or runner_failure_reason.startswith("cron_run")
                    else "Cron run is waiting for approval."
                ),
            )
        return ToolResult(
            invocation_id=invocation.invocation_id,
            tool_name=invocation.tool_name,
            status="completed",
            output=output,
            error=None,
        )
    return handle


def register_cron_tools(
    registry,
    *,
    cron_runner: Callable[..., str | dict[str, object] | None] | None = None,
    default_delivery_channel: str = "",
    default_delivery_target: str = "",
    skip_existing: bool = False,
) -> None:
    """Register cron management tools into an existing ToolRegistry.

    Call this after building the registry so the agent can create, list,
    update, toggle, run, and delete scheduled cron jobs.
    """
    def _register(spec: ToolSpec, handler: Callable[..., object]) -> None:
        if skip_existing:
            try:
                registry.get_spec(spec.name)
            except KeyError:
                pass
            else:
                return
        registry.register(spec, handler)

    _register(
        ToolSpec(
            name="create_cron",
            description=(
                "Create a new scheduled cron job. "
                "Required args: name (human-readable label), schedule (5-field cron expression, "
                "e.g. '0 9 * * 1-5' for weekdays at 9 AM), task (the natural-language instruction "
                "Nullion will execute when the cron fires). "
                "Optional: enabled (bool, default true), workspace_id, delivery_channel, delivery_target, "
                "html_image_delivery_mode (linked|auto|self_contained). "
                "If workspace and delivery are omitted, Nullion uses the current workspace and chat adapter defaults."
            ),
            risk_level=ToolRiskLevel.MEDIUM,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "cron"),
        ),
        _build_create_cron_handler(
            default_delivery_channel=default_delivery_channel,
            default_delivery_target=default_delivery_target,
        ),
    )
    _register(
        ToolSpec(
            name="list_crons",
            description=(
                "List scheduled cron jobs for the current workspace with their id, workspace, schedule, "
                "enabled state, next run time, and stored task instructions in structured output. "
                "Optional args: workspace_id, include_all_workspaces."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.READ,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "cron"),
            continuation_tools=("run_cron", "delete_cron", "update_cron", "toggle_cron"),
        ),
        _build_list_crons_handler(),
    )
    _register(
        ToolSpec(
            name="delete_cron",
            description="Delete a scheduled cron job by id. Required args: id (the cron job id).",
            risk_level=ToolRiskLevel.MEDIUM,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=True,
            timeout_seconds=10,
            capability_tags=("scheduler", "cron"),
        ),
        _build_delete_cron_handler(),
    )
    _register(
        ToolSpec(
            name="update_cron",
            description=(
                "Update a scheduled cron job by id. Required args: id. "
                "Optional mutable fields: name, schedule, task, enabled, workspace_id, delivery_channel, "
                "delivery_target, html_image_delivery_mode."
            ),
            risk_level=ToolRiskLevel.MEDIUM,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "cron"),
        ),
        _build_update_cron_handler(),
    )
    _register(
        ToolSpec(
            name="toggle_cron",
            description=(
                "Enable or disable scheduled cron jobs. Required args: enabled "
                "(bool — true to enable, false to disable), plus exactly one selector: id for one cron, "
                "ids for exact ids returned by list_crons, or all_current_workspace=true for every cron in "
                "the current workspace."
            ),
            risk_level=ToolRiskLevel.LOW,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=10,
            capability_tags=("scheduler", "cron"),
        ),
        _build_toggle_cron_handler(),
    )
    _register(
        ToolSpec(
            name="run_cron",
            description=(
                "Run one or more existing scheduled cron jobs immediately. For bulk runs, call list_crons first "
                "and pass exact cron ids as ids; this starts each selected cron as one structured operation. "
                "When the structured request selects every enabled scheduled task in the current workspace, "
                "pass all_enabled=true. "
                "For one cron, use the exact visible cron name from list_crons when it is known, or pass the "
                "user's partial/descriptive reference as name so the scheduler can resolve it against the "
                "structured cron records. Use id only when names are ambiguous or the user explicitly provides "
                "an id. Required args: ids, id, or name. Matching is conservative and punctuation-insensitive; "
                "ambiguous references return numbered candidate options instead of running a job."
            ),
            risk_level=ToolRiskLevel.MEDIUM,
            side_effect_class=ToolSideEffectClass.WRITE,
            requires_approval=False,
            timeout_seconds=120,
            capability_tags=("scheduler", "cron"),
            input_schema={
                "type": "object",
                "properties": {
                    "id": {"type": "string", "description": "Exact cron id when known."},
                    "ids": {
                        "type": "array",
                        "description": "Exact cron ids to run immediately as one bulk operation after list_crons.",
                        "items": {"type": "string"},
                    },
                    "all_enabled": {
                        "type": "boolean",
                        "description": "Run every enabled cron in the current workspace as one bulk operation.",
                    },
                    "name": {
                        "type": "string",
                        "description": "Exact, partial, or descriptive cron name/reference to resolve conservatively.",
                    },
                },
                "additionalProperties": False,
            },
        ),
        _build_run_cron_handler(cron_runner),
    )


__all__ = [
    "TerminalBackendDescriptor",
    "ToolExecutor",
    "ToolInvocation",
    "ToolRegistry",
    "ToolResult",
    "ToolRiskLevel",
    "ToolSideEffectClass",
    "ToolSpec",
    "build_default_tool_registry",
    "create_core_tool_registry",
    "create_default_tool_registry",
    "create_plugin_tool_registry",
    "create_extension_tool_registry",
    "normalize_tool_result",
    "normalize_tool_status",
    "register_browser_plugin",
    "register_browser_extension",
    "register_calendar_plugin",
    "register_connector_plugin",
    "register_email_plugin",
    "register_media_plugin",
    "register_search_plugin",
    "register_reminder_tools",
    "register_market_quote_tool",
    "register_weather_forecast_tool",
    "register_web_extension",
    "register_workspace_plugin",
]
